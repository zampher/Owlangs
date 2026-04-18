# SPDX-FileCopyrightText: 2026 Zampherss
# SPDX-License-Identifier: MPL-2.0
import asyncio
import os
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Self

from converter.base import ConverterConfig
from converter.converter_identity import ConverterIdentity
from exporter.base import ExporterConfig
from glossary.glossary import Glossary
from ir.document import Document
from logger.logger import LogModule
from translator.ai_translator.pptx_translator import PptxTranslatorConfig, PptxTranslator
from workflow.base import Workflow, WorkflowConfig
from workflow.interfaces import HTMLExportable


@dataclass(kw_only=True)
class Pptx2HTMLExporterConfig:
    """Configuration for PPTX to HTML exporter"""
    pass


@dataclass(kw_only=True)
class PptxWorkflowConfig(WorkflowConfig):
    translator_config: PptxTranslatorConfig
    html_exporter_config: Pptx2HTMLExporterConfig = None  # Will be created if None
    translate_notes: bool = False
    translate_master: bool = False


class PptxWorkflow(Workflow[PptxWorkflowConfig, Document, Document], HTMLExportable[Pptx2HTMLExporterConfig]):

    def __init__(self, config: PptxWorkflowConfig):
        super().__init__(config=config)
        if config.logger:
            for sub_config in [self.config.translator_config]:
                if sub_config:
                    sub_config.logger = config.logger
        # Initialize html_exporter_config if not provided
        if self.config.html_exporter_config is None:
            self.config.html_exporter_config = Pptx2HTMLExporterConfig()
        self.translator = None  # Store translator instance for token_counter access

    def _pre_translate(self, document_original: Document):
        document = document_original.copy()
        translate_config = self.config.translator_config
        # Set translate_notes and translate_master from workflow config
        translate_config.translate_notes = self.config.translate_notes
        translate_config.translate_master = self.config.translate_master
        translator = PptxTranslator(translate_config)
        return document, translator

    def translate(self) -> Self:
        document, translator = self._pre_translate(self.document_original)
        
        # Store translator instance in workflow for token_counter access
        self.translator = translator
        
        translator.translate(document)
        if translator.glossary_dict_gen:
            self.attachment.add_document("glossary", Glossary.glossary_dict2csv(translator.glossary_dict_gen))
        self.document_translated = document
        return self

    async def translate_async(self, progress_callback=None, task_id: str = None, 
                              original_filename: str = None, workflow_type: str = None, temp_dir: str = None) -> Self:
        document, translator = self._pre_translate(self.document_original)
        
        # Store translator instance in workflow for token_counter access
        self.translator = translator
        
        # Set task_id and other attributes if provided (for segment recording)
        if task_id:
            translator._task_id = task_id
        if original_filename:
            translator._original_filename = original_filename
        if workflow_type:
            translator._workflow_type = workflow_type
        
        # Pass temp_dir to translator for PPTX file copying
        await translator.translate_async(document, progress_callback=progress_callback, temp_dir=temp_dir)
        if translator.glossary_dict_gen:
            self.attachment.add_document("glossary", Glossary.glossary_dict2csv(translator.glossary_dict_gen))
        self.document_translated = document
        return self

    def export_to_html(self, config: Pptx2HTMLExporterConfig = None) -> str:
        """
        Export PPTX to HTML for preview.
        For now, returns a simple HTML representation.
        TODO: Implement full HTML export with slide rendering.
        """
        config = config or self.config.html_exporter_config
        # Simple HTML export - can be enhanced later
        try:
            from pptx import Presentation
            from io import BytesIO
            
            prs = Presentation(BytesIO(self.document_translated.content))
            html_parts = ['<html><head><meta charset="UTF-8"><title>PPTX Preview</title></head><body>']
            
            for slide_idx, slide in enumerate(prs.slides):
                html_parts.append(f'<div class="slide" style="page-break-after: always; padding: 20px;">')
                html_parts.append(f'<h2>Slide {slide_idx + 1}</h2>')
                
                # Add title
                if slide.shapes.title and slide.shapes.title.text:
                    html_parts.append(f'<h3>{slide.shapes.title.text}</h3>')
                
                # Add text content
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                html_parts.append(f'<p>{text}</p>')
                    elif shape.has_table:
                        html_parts.append('<table border="1" style="border-collapse: collapse; width: 100%;">')
                        for row in shape.table.rows:
                            html_parts.append('<tr>')
                            for cell in row.cells:
                                html_parts.append(f'<td>{cell.text}</td>')
                            html_parts.append('</tr>')
                        html_parts.append('</table>')
                
                html_parts.append('</div>')
            
            html_parts.append('</body></html>')
            return ''.join(html_parts)
        except Exception as e:
            self.logger.error(LogModule.WORKFLOW, f"Failed to export PPTX to HTML: {e}")
            return f"<html><body><p>Error exporting PPTX: {e}</p></body></html>"

    def save_as_html(self, name: str = None, output_dir: Path | str = "./output",
                     config: Pptx2HTMLExporterConfig | None = None) -> Self:
        config = config or self.config.html_exporter_config
        html_content = self.export_to_html(config)
        output_path = Path(output_dir) / (name or "presentation.html")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(html_content, encoding='utf-8')
        return self

    def export_to_pptx(self) -> bytes:
        """
        Export translated PPTX file as bytes.
        Prefers temporary file if available (preserves file structure better).
        
        Returns:
            Translated PPTX file content as bytes
        """
        # Priority 1: Use temporary file if available (preferred for preserving file structure)
        if self.translator and hasattr(self.translator, 'temp_pptx_path') and self.translator.temp_pptx_path:
            temp_path = self.translator.temp_pptx_path
            if os.path.exists(temp_path):
                self.logger.debug(LogModule.WORKFLOW, f"[PPTX] Exporting from temporary file: {temp_path}")
                with open(temp_path, 'rb') as f:
                    return f.read()
        
        # Priority 2: Fallback to document_translated.content
        if not self.document_translated or not self.document_translated.content:
            raise ValueError("No translated document available. Please run translate() or translate_async() first.")
        
        # document_translated.content is already bytes from PptxTranslator
        if isinstance(self.document_translated.content, bytes):
            return self.document_translated.content
        else:
            # Fallback: convert to bytes if needed
            from io import BytesIO
            if hasattr(self.document_translated.content, 'read'):
                return self.document_translated.content.read()
            else:
                return bytes(self.document_translated.content)

    def save_as_pptx(self, name: str = None, output_dir: Path | str = "./output") -> Self:
        """
        Save translated PPTX file to disk.
        Prefers copying from temporary file if available (preserves file structure better).
        
        Args:
            name: Output filename (without extension)
            output_dir: Output directory path
            
        Returns:
            Self for method chaining
        """
        # Priority 1: Copy from temporary file if available (preserves file structure)
        if self.translator and hasattr(self.translator, 'temp_pptx_path') and self.translator.temp_pptx_path:
            temp_path = self.translator.temp_pptx_path
            if os.path.exists(temp_path):
                output_path = Path(output_dir) / (f"{name}.pptx" if name else "presentation_translated.pptx")
                output_path.parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(temp_path, output_path)
                self.logger.debug(LogModule.WORKFLOW, f"[PPTX] Copied temporary file to output: {temp_path} -> {output_path}")
                return self
        
        # Priority 2: Fallback to export_to_pptx (reads from memory)
        pptx_content = self.export_to_pptx()
        output_path = Path(output_dir) / (f"{name}.pptx" if name else "presentation_translated.pptx")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(pptx_content)
        return self

