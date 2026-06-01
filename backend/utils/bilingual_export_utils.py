# SPDX-FileCopyrightText: 2026 Zampher
# SPDX-License-Identifier: MPL-2.0

"""Utilities for bilingual (source + target) document export.

This module provides shared helpers to rebuild documents that interleave
original and translated text at the segment level. It is used by
DownloadService when the user enables bilingual export.
"""

import html as html_module
import os
import re
import zipfile
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

# Preset colors for bilingual export (shared by markdown HTML spans and DOCX rebuild)
BILINGUAL_EXPORT_COLOR_NAMES = frozenset({"gray", "blue", "red", "green", "orange", "black"})
BILINGUAL_EXPORT_COLOR_HEX: Dict[str, str] = {
    "gray": "#808080",
    "red": "#FF0000",
    "blue": "#0000FF",
    "green": "#008000",
    "orange": "#FFA500",
    "black": "#000000",
}

from logger import unified_logger as logger
from logger.logger import LogModule

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

try:
    import openpyxl
    from openpyxl.cell.rich_text import CellRichText, TextBlock
    from openpyxl.cell.text import InlineFont
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    openpyxl = None
    CellRichText = None
    TextBlock = None
    InlineFont = None

try:
    from pptx.dml.color import RGBColor as PptxRGBColor
    PPTX_COLOR_AVAILABLE = PPTX_AVAILABLE
except ImportError:
    PPTX_COLOR_AVAILABLE = False
    PptxRGBColor = None


def get_bilingual_config(task_state: Optional[Dict[str, Any]]) -> Tuple[bool, bool]:
    """Return (enabled, target_first) from task_state format settings.

    Args:
        task_state: Task state dictionary (may be None).

    Returns:
        Tuple of (bilingual_export_enabled, target_first).
        target_first=True means target text comes before source text.
    """
    if not task_state:
        return False, False

    bilingual_export = task_state.get("bilingual_export")
    if isinstance(bilingual_export, str):
        enabled = bilingual_export.lower() in ("true", "1", "yes", "on")
    else:
        enabled = bool(bilingual_export)

    bilingual_order = task_state.get("bilingual_order")
    # "target_before_source" => target_first=True
    # "target_after_source"  => target_first=False (default)
    target_first = bool(bilingual_order) and str(bilingual_order).lower() == "target_before_source"

    return enabled, target_first


def _normalize_bilingual_color(color: Optional[str]) -> Optional[str]:
    if not color:
        return None
    normalized = str(color).strip().lower()
    if normalized not in BILINGUAL_EXPORT_COLOR_NAMES:
        return None
    return normalized


def get_bilingual_style_config(
    task_state: Optional[Dict[str, Any]],
) -> Tuple[bool, Optional[str], bool, Optional[str]]:
    """Return (source_italic, source_color, target_italic, target_color) from task_state."""
    if not task_state:
        return False, None, True, "gray"
    source_italic = bool(task_state.get("source_text_italic", False))
    source_color = _normalize_bilingual_color(task_state.get("source_text_color"))
    target_italic = bool(task_state.get("target_text_italic", True))
    target_color = _normalize_bilingual_color(task_state.get("target_text_color")) or "gray"
    return source_italic, source_color, target_italic, target_color


def wrap_bilingual_html_part(
    text: str,
    italic: bool = False,
    color: Optional[str] = None,
) -> str:
    """Wrap text in an inline HTML span for styled bilingual markdown/DOCX export."""
    if not text:
        return ""
    escaped = html_module.escape(text).replace("\n", "<br/>")
    styles: List[str] = []
    if italic:
        styles.append("font-style:italic")
    normalized_color = _normalize_bilingual_color(color)
    if normalized_color:
        styles.append(f"color:{BILINGUAL_EXPORT_COLOR_HEX[normalized_color]}")
    if not styles:
        return escaped
    return f'<span style="{";".join(styles)}">{escaped}</span>'


def build_bilingual_segment_text(
    source_text: str,
    target_text: str,
    target_first: bool,
    is_excluded: bool = False,
    is_cleared: bool = False,
    inner_separator: str = "\n\n",
    *,
    source_text_italic: bool = False,
    source_text_color: Optional[str] = None,
    target_text_italic: bool = False,
    target_text_color: Optional[str] = None,
    use_html_styles: bool = False,
) -> str:
    """Build bilingual text for a single segment.

    Rules:
      - Excluded or cleared segments: emit only source_text.
      - target_text == source_text (untranslated/failed): emit once to avoid duplication.
      - Otherwise: emit both parts in the requested order, separated by inner_separator.

    Args:
        source_text: Original text.
        target_text: Translated text.
        target_first: If True, place target before source.
        is_excluded: Whether the segment was excluded from translation.
        is_cleared: Whether the segment was cleared (empty target).
        inner_separator: Separator between source and target within one segment.

    Returns:
        Combined text for this segment.
    """
    source = source_text or ""
    target = target_text or ""

    if is_excluded or is_cleared:
        return source

    # Avoid duplication when translation failed or was identical
    if target.strip() and target.strip() != source.strip():
        if use_html_styles:
            if target_first:
                first_part = wrap_bilingual_html_part(
                    target, italic=target_text_italic, color=target_text_color
                )
                second_part = wrap_bilingual_html_part(
                    source, italic=source_text_italic, color=source_text_color
                )
            else:
                first_part = wrap_bilingual_html_part(
                    source, italic=source_text_italic, color=source_text_color
                )
                second_part = wrap_bilingual_html_part(
                    target, italic=target_text_italic, color=target_text_color
                )
            return f"{first_part}{inner_separator}{second_part}"
        first = target if target_first else source
        second = source if target_first else target
        return f"{first}{inner_separator}{second}"

    # Untranslated / identical / empty target
    return source


def _bilingual_rgb_triplet(color: Optional[str]) -> Optional[Tuple[int, int, int]]:
    normalized = _normalize_bilingual_color(color)
    if not normalized:
        return None
    hex_str = BILINGUAL_EXPORT_COLOR_HEX[normalized].lstrip("#")
    return int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)


def get_bilingual_styled_run_parts(
    source_text: str,
    target_text: str,
    target_first: bool,
    is_excluded: bool = False,
    is_cleared: bool = False,
    *,
    source_text_italic: bool = False,
    source_text_color: Optional[str] = None,
    target_text_italic: bool = False,
    target_text_color: Optional[str] = None,
) -> Optional[List[Tuple[str, bool, Optional[str]]]]:
    """Return ordered (text, italic, color) parts for rich bilingual export, or None if plain text only."""
    source = source_text or ""
    target = target_text or ""

    if is_excluded or is_cleared:
        return None
    if not target.strip() or target.strip() == source.strip():
        return None

    if target_first:
        return [
            (target, target_text_italic, target_text_color),
            (source, source_text_italic, source_text_color),
        ]
    return [
        (source, source_text_italic, source_text_color),
        (target, target_text_italic, target_text_color),
    ]


def bilingual_segment_to_html(
    source_text: str,
    target_text: str,
    target_first: bool,
    is_excluded: bool = False,
    is_cleared: bool = False,
    inner_separator: str = "<br/>",
    *,
    source_text_italic: bool = False,
    source_text_color: Optional[str] = None,
    target_text_italic: bool = False,
    target_text_color: Optional[str] = None,
) -> str:
    """Build styled bilingual HTML for one spreadsheet cell."""
    source = source_text or ""
    target = target_text or ""

    if is_excluded or is_cleared:
        return html_module.escape(source).replace("\n", "<br/>")
    if not target.strip() or target.strip() == source.strip():
        return html_module.escape(source).replace("\n", "<br/>")

    styled_parts = get_bilingual_styled_run_parts(
        source_text,
        target_text,
        target_first,
        is_excluded=is_excluded,
        is_cleared=is_cleared,
        source_text_italic=source_text_italic,
        source_text_color=source_text_color,
        target_text_italic=target_text_italic,
        target_text_color=target_text_color,
    )
    if not styled_parts:
        return html_module.escape(source).replace("\n", "<br/>")

    chunks: List[str] = []
    for idx, (text, italic, color) in enumerate(styled_parts):
        if idx > 0 and inner_separator:
            chunks.append(inner_separator)
        chunks.append(
            wrap_bilingual_html_part(text, italic=italic, color=color)
        )
    return "".join(chunks)


def _render_bilingual_xlsx_sheet_html(sheet, cell_html_rows: List[List[str]]) -> str:
    if not cell_html_rows:
        return ""
    max_cols = max(len(row) for row in cell_html_rows)
    lines = [f'<h2>{html_module.escape(sheet.title)}</h2>', "<table>"]
    for row in cell_html_rows:
        lines.append("<tr>")
        for col_idx in range(max_cols):
            cell_html = row[col_idx] if col_idx < len(row) else ""
            lines.append(f"<td>{cell_html}</td>")
        lines.append("</tr>")
    lines.append("</table>")
    return "\n".join(lines)


def rebuild_bilingual_xlsx_html_from_segments(
    task_state: Dict[str, Any],
    target_first: bool = False,
) -> Optional[str]:
    """Rebuild bilingual XLSX content as an HTML table with styled source/target spans."""
    if not XLSX_AVAILABLE:
        logger.warning(LogModule.EXPORT, "[BILINGUAL] openpyxl not available for XLSX bilingual HTML rebuild")
        return None

    segments_data = task_state.get("translation_segments")
    if not isinstance(segments_data, dict):
        logger.warning(LogModule.EXPORT, "[BILINGUAL] No translation_segments dict found for XLSX HTML rebuild")
        return None

    segments = segments_data.get("segments", [])
    if not segments:
        logger.warning(LogModule.EXPORT, "[BILINGUAL] Empty segments list for XLSX HTML rebuild")
        return None

    segments_by_index: Dict[int, Dict[str, Any]] = {}
    for seg in segments:
        idx = seg.get("segment_index")
        if idx is not None:
            segments_by_index[int(idx)] = seg

    (
        source_text_italic,
        source_text_color,
        target_text_italic,
        target_text_color,
    ) = get_bilingual_style_config(task_state)

    content = _get_source_document_bytes(task_state, "xlsx")
    if not content:
        return None

    try:
        workbook = openpyxl.load_workbook(BytesIO(content))
    except Exception as e:
        logger.error(LogModule.EXPORT, f"[BILINGUAL] Failed to load XLSX for HTML rebuild: {e}", exc_info=True)
        return None

    segment_index = 0
    sheet_sections: List[str] = []

    try:
        for sheet in workbook.worksheets:
            cell_html_rows: List[List[str]] = []
            for row in sheet.iter_rows():
                row_html: List[str] = []
                for cell in row:
                    if isinstance(cell.value, str) and cell.data_type == "s":
                        seg = segments_by_index.get(segment_index)
                        source = seg.get("source_text", "") if seg else ""
                        target = seg.get("modified_text") or seg.get("target_text", "") if seg else ""
                        is_excluded = bool(seg.get("is_excluded", False)) if seg else False
                        is_cleared = bool(seg.get("status") == "cleared") if seg else False
                        row_html.append(
                            bilingual_segment_to_html(
                                source_text=source,
                                target_text=target,
                                target_first=target_first,
                                is_excluded=is_excluded,
                                is_cleared=is_cleared,
                                inner_separator="<br/>",
                                source_text_italic=source_text_italic,
                                source_text_color=source_text_color,
                                target_text_italic=target_text_italic,
                                target_text_color=target_text_color,
                            )
                        )
                        segment_index += 1
                    elif cell.value is not None:
                        row_html.append(html_module.escape(str(cell.value)))
                    else:
                        row_html.append("")
                if any(cell.strip() for cell in row_html):
                    cell_html_rows.append(row_html)

            section = _render_bilingual_xlsx_sheet_html(sheet, cell_html_rows)
            if section:
                sheet_sections.append(section)
    finally:
        workbook.close()

    if not sheet_sections:
        return None

    body = "\n<br/>\n".join(sheet_sections)
    html_doc = (
        "<!DOCTYPE html>\n"
        '<html lang="en">\n'
        "<head>\n"
        '<meta charset="utf-8"/>\n'
        "<title>Translated Spreadsheet</title>\n"
        "<style>\n"
        "table { border-collapse: collapse; margin-bottom: 1.5em; }\n"
        "td, th { border: 1px solid #ccc; padding: 6px 8px; vertical-align: top; }\n"
        "</style>\n"
        "</head>\n"
        f"<body>\n{body}\n</body>\n"
        "</html>\n"
    )
    logger.info(
        LogModule.EXPORT,
        f"[BILINGUAL] Rebuilt XLSX HTML: {segment_index} text cells processed, "
        f"target_first={target_first}",
    )
    return html_doc


def _clear_pptx_paragraph_runs(paragraph) -> None:
    from pptx.oxml.ns import qn

    p_elem = paragraph._p
    for run_elem in list(p_elem.findall(qn("a:r"))):
        p_elem.remove(run_elem)


def _apply_pptx_run_style(run, italic: bool, color: Optional[str]) -> None:
    if not run:
        return
    try:
        if italic:
            run.font.italic = True
        rgb = _bilingual_rgb_triplet(color)
        if rgb and PptxRGBColor is not None:
            run.font.color.rgb = PptxRGBColor(*rgb)
    except Exception:
        pass


def apply_bilingual_styled_segment_to_pptx_paragraph(
    paragraph,
    source_text: str,
    target_text: str,
    target_first: bool,
    is_excluded: bool = False,
    is_cleared: bool = False,
    inner_separator: str = "\n",
    *,
    source_text_italic: bool = False,
    source_text_color: Optional[str] = None,
    target_text_italic: bool = False,
    target_text_color: Optional[str] = None,
) -> None:
    """Write bilingual segment text into a PPTX paragraph with per-run styling."""
    styled_parts = get_bilingual_styled_run_parts(
        source_text,
        target_text,
        target_first,
        is_excluded=is_excluded,
        is_cleared=is_cleared,
        source_text_italic=source_text_italic,
        source_text_color=source_text_color,
        target_text_italic=target_text_italic,
        target_text_color=target_text_color,
    )

    if styled_parts is None:
        plain = build_bilingual_segment_text(
            source_text=source_text,
            target_text=target_text,
            target_first=target_first,
            is_excluded=is_excluded,
            is_cleared=is_cleared,
            inner_separator=inner_separator,
        )
        _clear_pptx_paragraph_runs(paragraph)
        run = paragraph.add_run()
        run.text = plain
        return

    _clear_pptx_paragraph_runs(paragraph)
    for idx, (text, italic, color) in enumerate(styled_parts):
        if idx > 0 and inner_separator:
            sep_run = paragraph.add_run()
            sep_run.text = inner_separator
        run = paragraph.add_run()
        run.text = text
        _apply_pptx_run_style(run, italic, color)


def _xlsx_escape_cell_text(text: str) -> str:
    """Escape control characters for Excel OOXML text nodes.

    Literal newlines inside ``<t>`` without ``xml:space=\"preserve\"`` cause Excel
    to repair worksheet string properties. Use ST_Xstring escapes instead.
    """
    if not text:
        return text
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    return normalized.replace("\n", "_x000A_")


def _sanitize_xlsx_worksheet_xml(xml_text: str) -> str:
    """Replace literal newlines inside ``<t>`` nodes with Excel ST_Xstring escapes."""
    def _fix_t_element(match: re.Match[str]) -> str:
        attrs = match.group(1) or ""
        content = match.group(2)
        if "\n" not in content and "\r" not in content:
            return match.group(0)
        fixed = (
            content.replace("\r\n", "\n")
            .replace("\r", "\n")
            .replace("\n", "_x000A_")
        )
        return f"<t{attrs}>{fixed}</t>"

    return re.sub(r"<t([^>]*)>(.*?)</t>", _fix_t_element, xml_text, flags=re.DOTALL)


def _sanitize_xlsx_bytes(xlsx_bytes: bytes) -> bytes:
    """Post-process saved XLSX bytes so worksheet inline strings comply with Excel OOXML."""
    in_buf = BytesIO(xlsx_bytes)
    out_buf = BytesIO()
    sanitized_parts = 0

    with zipfile.ZipFile(in_buf, "r") as zin:
        with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if (
                    item.filename.startswith("xl/worksheets/sheet")
                    and item.filename.endswith(".xml")
                ):
                    xml_text = data.decode("utf-8")
                    sanitized = _sanitize_xlsx_worksheet_xml(xml_text)
                    if sanitized != xml_text:
                        sanitized_parts += 1
                        data = sanitized.encode("utf-8")
                zout.writestr(item, data)

    if sanitized_parts:
        logger.info(
            LogModule.EXPORT,
            f"[BILINGUAL] Sanitized inline string XML in {sanitized_parts} worksheet part(s)",
        )
    return out_buf.getvalue()


def _xlsx_cell_display_text(value: Any) -> str:
    """Flatten a cell value to display text for line-count heuristics."""
    if value is None:
        return ""
    if isinstance(value, str):
        return value.replace("_x000A_", "\n")
    if CellRichText is not None and isinstance(value, CellRichText):
        parts: List[str] = []
        for block in value:
            if isinstance(block, str):
                parts.append(block)
            else:
                parts.append(getattr(block, "text", str(block)))
        return "".join(parts).replace("_x000A_", "\n")
    return str(value)


def _xlsx_base_inline_font_kwargs(cell: Any) -> Dict[str, Any]:
    """Return base InlineFont fields required by Excel OOXML (rFont + sz at minimum)."""
    base: Dict[str, Any] = {"rFont": "Calibri", "sz": 11}
    try:
        cell_font = getattr(cell, "font", None)
        if cell_font is not None:
            if cell_font.name:
                base["rFont"] = cell_font.name
            if cell_font.size:
                base["sz"] = int(cell_font.size)
    except Exception:
        pass
    return base


def _xlsx_inline_font(
    italic: bool,
    color: Optional[str],
    *,
    base_font: Optional[Dict[str, Any]] = None,
) -> Optional[Any]:
    """Build an InlineFont for XLSX rich-text runs with Excel-required base fields."""
    if InlineFont is None:
        return None

    kwargs = dict(base_font or {"rFont": "Calibri", "sz": 11})
    if italic:
        kwargs["i"] = True
    rgb = _bilingual_rgb_triplet(color)
    if rgb:
        kwargs["color"] = f"FF{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    if not italic and not rgb:
        return None
    return InlineFont(**kwargs)


def apply_bilingual_styled_segment_to_xlsx_cell(
    cell,
    source_text: str,
    target_text: str,
    target_first: bool,
    is_excluded: bool = False,
    is_cleared: bool = False,
    inner_separator: str = "\n",
    *,
    source_text_italic: bool = False,
    source_text_color: Optional[str] = None,
    target_text_italic: bool = False,
    target_text_color: Optional[str] = None,
) -> None:
    """Write bilingual segment text into an XLSX cell with rich-text styling when supported."""
    styled_parts = get_bilingual_styled_run_parts(
        source_text,
        target_text,
        target_first,
        is_excluded=is_excluded,
        is_cleared=is_cleared,
        source_text_italic=source_text_italic,
        source_text_color=source_text_color,
        target_text_italic=target_text_italic,
        target_text_color=target_text_color,
    )

    if styled_parts is None:
        cell.value = _xlsx_escape_cell_text(
            build_bilingual_segment_text(
                source_text=source_text,
                target_text=target_text,
                target_first=target_first,
                is_excluded=is_excluded,
                is_cleared=is_cleared,
                inner_separator=inner_separator,
            )
        )
        return

    if CellRichText is None or TextBlock is None:
        cell.value = _xlsx_escape_cell_text(
            build_bilingual_segment_text(
                source_text=source_text,
                target_text=target_text,
                target_first=target_first,
                is_excluded=is_excluded,
                is_cleared=is_cleared,
                inner_separator=inner_separator,
            )
        )
        return

    base_font = _xlsx_base_inline_font_kwargs(cell)
    rich_blocks: List[Any] = []
    escaped_separator = _xlsx_escape_cell_text(inner_separator) if inner_separator else ""
    last_index = len(styled_parts) - 1
    for idx, (text, italic, color) in enumerate(styled_parts):
        block_text = _xlsx_escape_cell_text(text)
        # Append bilingual separator to the styled run to avoid a bare unstyled newline run.
        if idx < last_index and escaped_separator:
            block_text += escaped_separator
        font = _xlsx_inline_font(italic, color, base_font=base_font)
        if font is not None:
            rich_blocks.append(TextBlock(font, block_text))
        else:
            rich_blocks.append(block_text)

    cell.value = CellRichText(*rich_blocks)


def should_skip_bilingual_for_image_render(
    segment: Dict[str, Any],
    block_types: List[str],
    *,
    table_body_format: Optional[str] = None,
    equation_format: Optional[str] = None,
    is_table_body: bool = False,
) -> bool:
    """Return True when a segment is rendered as an image and cannot be bilingual.

    Image/table captions share layout block indices with image/table blocks but contain
    real text; they must NOT be skipped. Only placeholder-only image segments and
    segments explicitly exported as image (table body / equation) are skipped.
    """
    from utils.translation_segments import _is_image_segment

    source_text = segment.get("source_text") or segment.get("text") or ""
    if _is_image_segment(source_text):
        return True

    tbl_fmt = (table_body_format or "html").strip().lower()
    if tbl_fmt == "image" and is_table_body:
        return True

    eq_fmt = (equation_format or "text").strip().lower()
    if eq_fmt == "image" and "interline_equation" in block_types:
        return True

    return False


def rebuild_bilingual_plain_text_from_segments(
    task_state: Dict[str, Any],
    target_first: bool = False,
    paragraph_separator: str = "\n\n",
) -> str:
    """Rebuild a plain-text document from segments in bilingual mode.

    Used for TXT workflow (and any other workflow that needs simple text rebuild).

    Args:
        task_state: Task state containing translation_segments.
        target_first: If True, place target before source for each segment.
        paragraph_separator: Separator between segments.

    Returns:
        Rebuilt bilingual text, or empty string if no segments.
    """
    segments_data = task_state.get("translation_segments")
    if not isinstance(segments_data, dict):
        logger.warning(LogModule.EXPORT, "[BILINGUAL] No translation_segments dict found for plain-text rebuild")
        return ""

    segments = segments_data.get("segments", [])
    if not segments:
        logger.warning(LogModule.EXPORT, "[BILINGUAL] Empty segments list for plain-text rebuild")
        return ""

    # Sort by segment_index for consistent ordering
    segments = sorted(segments, key=lambda s: s.get("segment_index", 0))

    parts: list[str] = []
    for seg in segments:
        source = seg.get("source_text", "")
        target = seg.get("modified_text") or seg.get("target_text", "")
        is_excluded = bool(seg.get("is_excluded", False))
        is_cleared = bool(
            seg.get("status") == "cleared"
            or (not target and seg.get("modified", False) and seg.get("target_length", -1) == 0)
        )

        segment_text = build_bilingual_segment_text(
            source_text=source,
            target_text=target,
            target_first=target_first,
            is_excluded=is_excluded,
            is_cleared=is_cleared,
            inner_separator="\n\n",
        )
        if segment_text:
            parts.append(segment_text)

    result = paragraph_separator.join(parts)
    logger.info(
        LogModule.EXPORT,
        f"[BILINGUAL] Rebuilt plain text: {len(segments)} segments -> {len(parts)} paragraphs, "
        f"target_first={target_first}",
    )
    return result


def rebuild_bilingual_srt_from_segments(
    task_state: Dict[str, Any],
    target_first: bool = False,
) -> str:
    """Rebuild an SRT subtitle file from segments in bilingual mode.

    Each subtitle entry keeps its original index and timecode.
    Source and target texts are placed on separate lines inside the same entry.

    Args:
        task_state: Task state containing translation_segments.
        target_first: If True, place target before source within each entry.

    Returns:
        Rebuilt bilingual SRT content, or empty string if no segments.
    """
    segments_data = task_state.get("translation_segments")
    if not isinstance(segments_data, dict):
        logger.warning(LogModule.EXPORT, "[BILINGUAL] No translation_segments dict found for SRT rebuild")
        return ""

    segments = segments_data.get("segments", [])
    if not segments:
        logger.warning(LogModule.EXPORT, "[BILINGUAL] Empty segments list for SRT rebuild")
        return ""

    # Sort by segment_index for consistent ordering
    segments = sorted(segments, key=lambda s: s.get("segment_index", 0))

    # Try to get SRT metadata (timecodes) from task_state if available
    srt_entries = task_state.get("srt_entries")
    if isinstance(srt_entries, list) and len(srt_entries) == len(segments):
        # Use stored timecodes
        return _rebuild_srt_with_entries(segments, srt_entries, target_first)

    # Fallback: generate simple sequential entries without timecodes
    # This preserves structure but loses timing; caller should ensure srt_entries is stored.
    logger.warning(
        LogModule.EXPORT,
        "[BILINGUAL] srt_entries not available in task_state; generating sequential entries without timecodes",
    )
    return _rebuild_srt_sequential(segments, target_first)


def _rebuild_srt_with_entries(
    segments: list,
    srt_entries: list,
    target_first: bool,
) -> str:
    """Rebuild SRT using stored timecode entries."""
    lines: list[str] = []
    for idx, (seg, entry) in enumerate(zip(segments, srt_entries), start=1):
        source = seg.get("source_text", "")
        target = seg.get("modified_text") or seg.get("target_text", "")
        is_excluded = bool(seg.get("is_excluded", False))
        is_cleared = bool(
            seg.get("status") == "cleared"
            or (not target and seg.get("modified", False) and seg.get("target_length", -1) == 0)
        )

        segment_text = build_bilingual_segment_text(
            source_text=source,
            target_text=target,
            target_first=target_first,
            is_excluded=is_excluded,
            is_cleared=is_cleared,
            inner_separator="\n",  # Within one SRT entry, use single newline
        )

        timecode = entry.get("timecode", "00:00:00,000 --> 00:00:00,000") if isinstance(entry, dict) else "00:00:00,000 --> 00:00:00,000"
        lines.append(str(idx))
        lines.append(timecode)
        lines.append(segment_text)
        lines.append("")  # Blank line between entries

    return "\n".join(lines)


def _rebuild_srt_sequential(
    segments: list,
    target_first: bool,
) -> str:
    """Fallback: rebuild SRT with sequential indices and dummy timecodes."""
    lines: list[str] = []
    for idx, seg in enumerate(segments, start=1):
        source = seg.get("source_text", "")
        target = seg.get("modified_text") or seg.get("target_text", "")
        is_excluded = bool(seg.get("is_excluded", False))
        is_cleared = bool(
            seg.get("status") == "cleared"
            or (not target and seg.get("modified", False) and seg.get("target_length", -1) == 0)
        )

        segment_text = build_bilingual_segment_text(
            source_text=source,
            target_text=target,
            target_first=target_first,
            is_excluded=is_excluded,
            is_cleared=is_cleared,
            inner_separator="\n",
        )

        lines.append(str(idx))
        lines.append("00:00:00,000 --> 00:00:00,000")
        lines.append(segment_text)
        lines.append("")

    return "\n".join(lines)


def _get_source_document_bytes(task_state: Dict[str, Any], suffix_hint: str) -> Optional[bytes]:
    """Get document bytes for bilingual rebuild from multiple fallback sources.

    Priority:
    1. workflow_instance.document_original.content (in-memory original)
    2. task_state temp_dir + original_filename (file on disk)
    3. task_state temp_dir + output/<stem>_translated.<suffix> (translated output file)

    Args:
        task_state: Task state dictionary.
        suffix_hint: File extension hint (e.g. "pptx", "xlsx").

    Returns:
        Document bytes, or None if all sources fail.
    """
    # Priority 1: workflow_instance.document_original.content
    workflow_instance = task_state.get("workflow_instance")
    if workflow_instance and hasattr(workflow_instance, "document_original"):
        doc = workflow_instance.document_original
        if doc and doc.content:
            logger.debug(LogModule.EXPORT, f"[BILINGUAL] Got bytes from workflow_instance.document_original for .{suffix_hint}")
            return doc.content

    # Priority 2: Read original file from temp_dir
    temp_dir = task_state.get("temp_dir")
    original_filename = task_state.get("original_filename")
    if temp_dir and original_filename:
        candidate = os.path.join(str(temp_dir), original_filename)
        if os.path.isfile(candidate):
            logger.debug(LogModule.EXPORT, f"[BILINGUAL] Got bytes from temp_dir/{original_filename} for .{suffix_hint}")
            with open(candidate, 'rb') as f:
                return f.read()

    # Priority 3: Read translated output file (same structure, text will be overwritten)
    if temp_dir:
        output_dir = os.path.join(str(temp_dir), "output")
        if os.path.isdir(output_dir):
            export_filename = task_state.get("original_filename_stem", "rebuilt")
            export_path = os.path.join(output_dir, f"{export_filename}_translated.{suffix_hint}")
            if os.path.isfile(export_path):
                logger.debug(LogModule.EXPORT, f"[BILINGUAL] Got bytes from output file for .{suffix_hint}")
                with open(export_path, 'rb') as f:
                    return f.read()

    logger.warning(LogModule.EXPORT, f"[BILINGUAL] No source document bytes found for .{suffix_hint} rebuild")
    return None


def rebuild_bilingual_pptx_from_segments(
    task_state: Dict[str, Any],
    target_first: bool = False,
) -> Optional[bytes]:
    """Rebuild a PPTX document from segments in bilingual mode.

    Opens the original PPTX from the workflow instance, iterates text elements
    in the same order as PptxTranslator._pre_translate(), and writes source+target
    interleaved text into each element.

    Args:
        task_state: Task state containing translation_segments and workflow_instance.
        target_first: If True, place target before source for each segment.

    Returns:
        Rebuilt PPTX file bytes, or None on failure.
    """
    if not PPTX_AVAILABLE:
        logger.warning(LogModule.EXPORT, "[BILINGUAL] python-pptx not available for PPTX bilingual rebuild")
        return None

    segments_data = task_state.get("translation_segments")
    if not isinstance(segments_data, dict):
        logger.warning(LogModule.EXPORT, "[BILINGUAL] No translation_segments dict found for PPTX rebuild")
        return None

    segments = segments_data.get("segments", [])
    if not segments:
        logger.warning(LogModule.EXPORT, "[BILINGUAL] Empty segments list for PPTX rebuild")
        return None

    # Build segment lookup by segment_index
    segments_by_index: Dict[int, Dict[str, Any]] = {}
    for seg in segments:
        idx = seg.get("segment_index")
        if idx is not None:
            segments_by_index[int(idx)] = seg

    (
        source_text_italic,
        source_text_color,
        target_text_italic,
        target_text_color,
    ) = get_bilingual_style_config(task_state)

    # Get source PPTX content from multiple fallback sources
    content = _get_source_document_bytes(task_state, "pptx")
    if not content:
        return None

    try:
        prs = Presentation(BytesIO(content))
    except Exception as e:
        logger.error(LogModule.EXPORT, f"[BILINGUAL] Failed to load PPTX from content: {e}", exc_info=True)
        return None

    segment_index = 0

    for slide in prs.slides:
        # --- Slide title ---
        if slide.shapes.title and slide.shapes.title.text.strip():
            seg = segments_by_index.get(segment_index)
            source = seg.get("source_text", "") if seg else ""
            target = seg.get("modified_text") or seg.get("target_text", "") if seg else ""
            is_excluded = bool(seg.get("is_excluded", False)) if seg else False
            is_cleared = bool(seg.get("status") == "cleared") if seg else False

            # Write back to title
            text_frame = slide.shapes.title.text_frame
            if text_frame.paragraphs:
                para = text_frame.paragraphs[0]
                apply_bilingual_styled_segment_to_pptx_paragraph(
                    para,
                    source_text=source,
                    target_text=target,
                    target_first=target_first,
                    is_excluded=is_excluded,
                    is_cleared=is_cleared,
                    inner_separator="\n",
                    source_text_italic=source_text_italic,
                    source_text_color=source_text_color,
                    target_text_italic=target_text_italic,
                    target_text_color=target_text_color,
                )

            segment_index += 1

        # --- Shapes (text_frame + table_cell) ---
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue

            # Text frames
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    has_text = any(run.text.strip() for run in paragraph.runs)
                    if not has_text:
                        continue

                    seg = segments_by_index.get(segment_index)
                    source = seg.get("source_text", "") if seg else ""
                    target = seg.get("modified_text") or seg.get("target_text", "") if seg else ""
                    is_excluded = bool(seg.get("is_excluded", False)) if seg else False
                    is_cleared = bool(seg.get("status") == "cleared") if seg else False

                    apply_bilingual_styled_segment_to_pptx_paragraph(
                        paragraph,
                        source_text=source,
                        target_text=target,
                        target_first=target_first,
                        is_excluded=is_excluded,
                        is_cleared=is_cleared,
                        inner_separator="\n",
                        source_text_italic=source_text_italic,
                        source_text_color=source_text_color,
                        target_text_italic=target_text_italic,
                        target_text_color=target_text_color,
                    )

                    segment_index += 1

            # Table cells
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if not cell_text:
                            continue

                        seg = segments_by_index.get(segment_index)
                        source = seg.get("source_text", "") if seg else ""
                        target = seg.get("modified_text") or seg.get("target_text", "") if seg else ""
                        is_excluded = bool(seg.get("is_excluded", False)) if seg else False
                        is_cleared = bool(seg.get("status") == "cleared") if seg else False

                        if cell.text_frame.paragraphs:
                            apply_bilingual_styled_segment_to_pptx_paragraph(
                                cell.text_frame.paragraphs[0],
                                source_text=source,
                                target_text=target,
                                target_first=target_first,
                                is_excluded=is_excluded,
                                is_cleared=is_cleared,
                                inner_separator="\n",
                                source_text_italic=source_text_italic,
                                source_text_color=source_text_color,
                                target_text_italic=target_text_italic,
                                target_text_color=target_text_color,
                            )
                        else:
                            cell.text = build_bilingual_segment_text(
                                source_text=source,
                                target_text=target,
                                target_first=target_first,
                                is_excluded=is_excluded,
                                is_cleared=is_cleared,
                                inner_separator="\n",
                            )
                        segment_index += 1

    try:
        bio = BytesIO()
        prs.save(bio)
        result = bio.getvalue()
    except Exception as e:
        logger.error(LogModule.EXPORT, f"[BILINGUAL] Failed to save bilingual PPTX: {e}", exc_info=True)
        return None

    logger.info(
        LogModule.EXPORT,
        f"[BILINGUAL] Rebuilt PPTX: {segment_index} elements processed, "
        f"target_first={target_first}",
    )
    return result


def rebuild_bilingual_xlsx_from_segments(
    task_state: Dict[str, Any],
    target_first: bool = False,
) -> Optional[bytes]:
    """Rebuild an XLSX document from segments in bilingual mode.

    Opens the original XLSX from the workflow instance, iterates text cells
    in the same order as XlsxTranslator._pre_translate(), and writes source+target
    interleaved text into each cell.

    Args:
        task_state: Task state containing translation_segments and workflow_instance.
        target_first: If True, place target before source for each segment.

    Returns:
        Rebuilt XLSX file bytes, or None on failure.
    """
    if not XLSX_AVAILABLE:
        logger.warning(LogModule.EXPORT, "[BILINGUAL] openpyxl not available for XLSX bilingual rebuild")
        return None

    segments_data = task_state.get("translation_segments")
    if not isinstance(segments_data, dict):
        logger.warning(LogModule.EXPORT, "[BILINGUAL] No translation_segments dict found for XLSX rebuild")
        return None

    segments = segments_data.get("segments", [])
    if not segments:
        logger.warning(LogModule.EXPORT, "[BILINGUAL] Empty segments list for XLSX rebuild")
        return None

    # Build segment lookup by segment_index
    segments_by_index: Dict[int, Dict[str, Any]] = {}
    for seg in segments:
        idx = seg.get("segment_index")
        if idx is not None:
            segments_by_index[int(idx)] = seg

    (
        source_text_italic,
        source_text_color,
        target_text_italic,
        target_text_color,
    ) = get_bilingual_style_config(task_state)

    # Get source XLSX content from multiple fallback sources
    content = _get_source_document_bytes(task_state, "xlsx")
    if not content:
        return None

    try:
        workbook = openpyxl.load_workbook(BytesIO(content))
    except Exception as e:
        logger.error(LogModule.EXPORT, f"[BILINGUAL] Failed to load XLSX from content: {e}", exc_info=True)
        return None

    segment_index = 0
    modified_rows_by_sheet: Dict[str, set] = {}  # Track rows with modified cells per sheet

    for sheet in workbook.worksheets:
        modified_rows_by_sheet[sheet.title] = set()
        for row in sheet.iter_rows():
            for cell in row:
                if not (isinstance(cell.value, str) and cell.data_type == "s"):
                    continue

                seg = segments_by_index.get(segment_index)
                source = seg.get("source_text", "") if seg else ""
                target = seg.get("modified_text") or seg.get("target_text", "") if seg else ""
                is_excluded = bool(seg.get("is_excluded", False)) if seg else False
                is_cleared = bool(seg.get("status") == "cleared") if seg else False

                apply_bilingual_styled_segment_to_xlsx_cell(
                    cell,
                    source_text=source,
                    target_text=target,
                    target_first=target_first,
                    is_excluded=is_excluded,
                    is_cleared=is_cleared,
                    inner_separator="\n",
                    source_text_italic=source_text_italic,
                    source_text_color=source_text_color,
                    target_text_italic=target_text_italic,
                    target_text_color=target_text_color,
                )

                # Enable text wrapping so bilingual newlines display correctly
                from openpyxl.styles import Alignment as OpenpyxlAlignment
                cell.alignment = OpenpyxlAlignment(wrap_text=True)

                modified_rows_by_sheet[sheet.title].add(cell.row)
                segment_index += 1

    # Auto-fit row heights based on line count
    # Default Excel row height is 15 points; each line of text adds ~15 points
    _DEFAULT_LINE_HEIGHT = 15.0
    for sheet in workbook.worksheets:
        for row_number in modified_rows_by_sheet.get(sheet.title, set()):
            max_lines = 1
            for cell in sheet[row_number]:
                if cell.value is None:
                    continue
                line_count = _xlsx_cell_display_text(cell.value).count("\n") + 1
                max_lines = max(max_lines, line_count)
            sheet.row_dimensions[row_number].height = max_lines * _DEFAULT_LINE_HEIGHT

    try:
        bio = BytesIO()
        workbook.save(bio)
        result = _sanitize_xlsx_bytes(bio.getvalue())
    except Exception as e:
        logger.error(LogModule.EXPORT, f"[BILINGUAL] Failed to save bilingual XLSX: {e}", exc_info=True)
        return None
    finally:
        workbook.close()

    logger.info(
        LogModule.EXPORT,
        f"[BILINGUAL] Rebuilt XLSX: {segment_index} cells processed, "
        f"target_first={target_first}",
    )
    return result
