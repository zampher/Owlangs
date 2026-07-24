# SPDX-FileCopyrightText: 2026 Zampher
# SPDX-License-Identifier: MPL-2.0

"""Lightweight Office → HTML helpers for compare-reading previews (no translation)."""

from __future__ import annotations

import html
from io import BytesIO
from typing import List

from logger import unified_logger as logger
from logger.logger import LogModule

_BODY_STYLE = (
    "body{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Roboto,"
    "'Helvetica Neue',Arial;padding:12px;line-height:1.5}"
)
_TABLE_STYLE = (
    ".table{border-collapse:collapse;width:100%;margin:8px 0}"
    ".table td,.table th{border:1px solid #ddd;padding:8px;text-align:left}"
    ".table-striped tr:nth-child(even){background-color:#f2f2f2}"
)
_SLIDE_STYLE = (
    ".slide{border:1px solid #ddd;border-radius:8px;padding:16px;"
    "margin:0 0 16px 0;background:#fff}"
    ".slide h2{margin:0 0 12px 0;font-size:1.1rem}"
)


def wrap_preview_html(body_inner: str, extra_css: str = "") -> str:
    """Wrap HTML body fragments in a minimal document shell."""
    return (
        "<html><head><meta charset='utf-8'>"
        f"<style>{_BODY_STYLE}{_TABLE_STYLE}{_SLIDE_STYLE}{extra_css}</style>"
        "</head><body>"
        f"{body_inner}"
        "</body></html>"
    )


def docx_bytes_to_html(content: bytes) -> str:
    """Convert DOCX bytes to HTML via mammoth."""
    try:
        import mammoth  # type: ignore
    except Exception as exc:
        logger.error(
            LogModule.EXPORT,
            f"[OFFICE_PREVIEW] mammoth unavailable for DOCX preview: {exc}",
        )
        raise RuntimeError(
            "DOCX preview requires 'mammoth' package. Please install it."
        ) from exc

    result = mammoth.convert_to_html(BytesIO(content))
    body = result.value or ""
    messages = getattr(result, "messages", None) or []
    if messages:
        logger.info(
            LogModule.EXPORT,
            f"[OFFICE_PREVIEW] mammoth messages count={len(messages)}",
        )
    return wrap_preview_html(body)


def xlsx_bytes_to_html(content: bytes, max_rows: int = 200) -> str:
    """Convert XLSX bytes to HTML tables (all sheets, capped rows per sheet)."""
    sections: List[str] = []

    try:
        import openpyxl  # type: ignore

        workbook = openpyxl.load_workbook(
            BytesIO(content), data_only=True, read_only=True
        )
        sheet_names = list(workbook.sheetnames)
        logger.info(
            LogModule.EXPORT,
            f"[OFFICE_PREVIEW] XLSX sheets={len(sheet_names)} max_rows={max_rows}",
        )
        for sheet_name in sheet_names:
            worksheet = workbook[sheet_name]
            rows_html: List[str] = []
            for row_idx, row in enumerate(
                worksheet.iter_rows(values_only=True), start=1
            ):
                if row_idx > max_rows:
                    rows_html.append(
                        "<tr><td colspan='99'>"
                        f"<em>… truncated after {max_rows} rows</em>"
                        "</td></tr>"
                    )
                    break
                cells = "".join(
                    f"<td>{html.escape('' if cell is None else str(cell))}</td>"
                    for cell in row
                )
                rows_html.append(f"<tr>{cells}</tr>")
            sections.append(
                f"<h2>{html.escape(str(sheet_name))}</h2>"
                f"<table class='table table-striped'>"
                f"{''.join(rows_html)}"
                f"</table>"
            )
        workbook.close()
        if not sections:
            sections.append("<p><em>Empty workbook</em></p>")
        return wrap_preview_html("".join(sections))
    except Exception as openpyxl_exc:
        logger.warning(
            LogModule.EXPORT,
            f"[OFFICE_PREVIEW] openpyxl XLSX preview failed: {openpyxl_exc}; "
            "trying pandas",
        )

    try:
        import pandas as pd  # type: ignore

        sheets = pd.read_excel(BytesIO(content), sheet_name=None, nrows=max_rows)
        logger.info(
            LogModule.EXPORT,
            f"[OFFICE_PREVIEW] pandas XLSX sheets={len(sheets)} max_rows={max_rows}",
        )
        for sheet_name, frame in sheets.items():
            table_html = frame.to_html(
                classes="table table-striped",
                table_id=None,
                border=0,
                escape=True,
            )
            sections.append(f"<h2>{html.escape(str(sheet_name))}</h2>{table_html}")
        if not sections:
            sections.append("<p><em>Empty workbook</em></p>")
        return wrap_preview_html("".join(sections))
    except Exception as pandas_exc:
        logger.error(
            LogModule.EXPORT,
            f"[OFFICE_PREVIEW] XLSX preview failed (openpyxl+pandas): {pandas_exc}",
            exc_info=True,
        )
        raise RuntimeError(f"XLSX preview failed: {pandas_exc}") from pandas_exc


def _shape_paragraphs_html(shape) -> List[str]:
    parts: List[str] = []
    if not getattr(shape, "has_text_frame", False):
        return parts
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text and paragraph.text:
            text = paragraph.text.strip()
        if text:
            parts.append(f"<p>{html.escape(text)}</p>")
    return parts


def _table_html(shape) -> str:
    rows_html: List[str] = []
    table = shape.table
    for row in table.rows:
        cells = "".join(
            f"<td>{html.escape((cell.text or '').strip())}</td>"
            for cell in row.cells
        )
        rows_html.append(f"<tr>{cells}</tr>")
    return (
        "<table class='table table-striped'>"
        f"{''.join(rows_html)}"
        "</table>"
    )


def pptx_bytes_to_html(content: bytes) -> str:
    """Convert PPTX bytes to HTML slides via python-pptx."""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception as exc:
        logger.error(
            LogModule.EXPORT,
            f"[OFFICE_PREVIEW] python-pptx unavailable for PPTX preview: {exc}",
        )
        raise RuntimeError(
            "PPTX preview requires 'python-pptx' package. Please install it."
        ) from exc

    try:
        prs = Presentation(BytesIO(content))
    except Exception as exc:
        logger.error(
            LogModule.EXPORT,
            f"[OFFICE_PREVIEW] Failed to open PPTX: {exc}",
            exc_info=True,
        )
        raise RuntimeError(f"PPTX preview failed: {exc}") from exc

    slide_sections: List[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = [f"<div class='slide'><h2>Slide {slide_idx}</h2>"]
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            try:
                if title_shape is not None and shape == title_shape:
                    title_text = (title_shape.text or "").strip()
                    if title_text:
                        parts.append(f"<h3>{html.escape(title_text)}</h3>")
                    continue
                if getattr(shape, "has_table", False):
                    parts.append(_table_html(shape))
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for child in shape.shapes:
                        parts.extend(_shape_paragraphs_html(child))
                    continue
                parts.extend(_shape_paragraphs_html(shape))
            except Exception as shape_exc:
                logger.warning(
                    LogModule.EXPORT,
                    f"[OFFICE_PREVIEW] Skip shape on slide {slide_idx}: {shape_exc}",
                )
        has_notes = getattr(slide, "has_notes_slide", False)
        if has_notes:
            notes_slide = slide.notes_slide
            if notes_slide is not None and notes_slide.notes_text_frame is not None:
                notes = (notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(f"<p><em>Notes:</em> {html.escape(notes)}</p>")
        parts.append("</div>")
        slide_sections.append("".join(parts))

    logger.info(
        LogModule.EXPORT,
        f"[OFFICE_PREVIEW] PPTX slides={len(slide_sections)}",
    )
    if not slide_sections:
        slide_sections.append("<p><em>Empty presentation</em></p>")
    return wrap_preview_html("".join(slide_sections))
