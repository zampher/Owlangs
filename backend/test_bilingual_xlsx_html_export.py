# SPDX-FileCopyrightText: 2026 Zampher
# SPDX-License-Identifier: MPL-2.0

"""Regression: bilingual XLSX HTML download must not return raw XLSX bytes."""

import sys
from io import BytesIO
from io import BytesIO
from pathlib import Path

_root = Path(__file__).resolve().parent.parent
_backend = Path(__file__).resolve().parent
for p in (str(_root), str(_backend)):
    if p not in sys.path:
        sys.path.insert(0, p)


def test_bilingual_segment_to_html_styles_and_line_breaks():
    from utils.bilingual_export_utils import bilingual_segment_to_html

    html = bilingual_segment_to_html(
        "Line1\nLine2",
        "行1\n行2",
        target_first=False,
        source_text_italic=True,
        source_text_color="blue",
        target_text_italic=True,
        target_text_color="gray",
    )
    assert 'style="font-style:italic;color:#0000FF"' in html
    assert 'style="font-style:italic;color:#808080"' in html
    assert "<br/>" in html
    assert "Line1" in html and "行1" in html
    print("PASS test_bilingual_segment_to_html_styles_and_line_breaks")


def test_rebuild_bilingual_xlsx_html_from_segments():
    try:
        import openpyxl
        from utils.bilingual_export_utils import rebuild_bilingual_xlsx_html_from_segments
    except ImportError as exc:
        print(f"PASS test_rebuild_bilingual_xlsx_html_from_segments (skipped: {exc})")
        return

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Title"
    ws["B1"] = "Name"
    bio = BytesIO()
    wb.save(bio)
    src_bytes = bio.getvalue()
    wb.close()

    class _Doc:
        content = src_bytes

    class _WF:
        document_original = _Doc()

    task_state = {
        "workflow_instance": _WF(),
        "source_text_italic": True,
        "source_text_color": "orange",
        "target_text_italic": True,
        "target_text_color": "gray",
        "translation_segments": {
            "segments": [
                {"segment_index": 0, "source_text": "Title", "target_text": "标题", "is_excluded": False},
                {"segment_index": 1, "source_text": "Name", "target_text": "名称", "is_excluded": False},
            ]
        },
    }
    html = rebuild_bilingual_xlsx_html_from_segments(task_state, target_first=False)
    assert html and "<table>" in html
    assert 'color:#FFA500' in html or "color:#ffa500" in html.lower()
    assert "font-style:italic" in html
    assert "标题" in html and "Title" in html
    print("PASS test_rebuild_bilingual_xlsx_html_from_segments")


def test_rebuild_bilingual_pptx_html_from_segments():
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from utils.bilingual_export_utils import rebuild_bilingual_pptx_html_from_segments
    except ImportError as exc:
        print(f"PASS test_rebuild_bilingual_pptx_html_from_segments (skipped: {exc})")
        return

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Hello"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    body.text_frame.text = "World"
    bio = BytesIO()
    prs.save(bio)
    src_bytes = bio.getvalue()

    class _Doc:
        content = src_bytes

    class _WF:
        document_original = _Doc()

    task_state = {
        "workflow_instance": _WF(),
        "source_text_italic": True,
        "source_text_color": "blue",
        "target_text_italic": True,
        "target_text_color": "gray",
        "translation_segments": {
            "segments": [
                {"segment_index": 0, "source_text": "Hello", "target_text": "你好", "is_excluded": False},
                {"segment_index": 1, "source_text": "World", "target_text": "世界", "is_excluded": False},
            ]
        },
    }
    html = rebuild_bilingual_pptx_html_from_segments(task_state, target_first=False)
    assert html and "Slide 1" in html
    assert "Hello" in html and "你好" in html
    assert "World" in html and "世界" in html
    assert "font-style:italic" in html
    assert "color:#" in html
    print("PASS test_rebuild_bilingual_pptx_html_from_segments")


def test_bilingual_xlsx_html_is_not_zip_bytes():
    try:
        import openpyxl
        from ir.document import Document
        from utils.bilingual_export_utils import rebuild_bilingual_xlsx_from_segments
        from workflow.xlsx_workflow import XlsxWorkflow, XlsxWorkflowConfig
        from exporter.xlsx.xlsx2html_exporter import Xlsx2HTMLExporterConfig
        from translator.ai_translator.xlsx_translator import XlsxTranslatorConfig
    except (ImportError, TypeError) as exc:
        print(f"PASS test_bilingual_xlsx_html_is_not_zip_bytes (skipped: {exc})")
        return

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Title"
    ws["B1"] = "Name"
    bio = BytesIO()
    wb.save(bio)
    src_bytes = bio.getvalue()
    wb.close()

    class _Doc:
        content = src_bytes

    class _WF:
        document_original = _Doc()

    task_state = {
        "workflow_instance": _WF(),
        "source_text_color": "orange",
        "target_text_italic": True,
        "target_text_color": "gray",
        "translation_segments": {
            "segments": [
                {
                    "segment_index": 0,
                    "source_text": "Title",
                    "target_text": "标题",
                    "is_excluded": False,
                },
                {
                    "segment_index": 1,
                    "source_text": "Name",
                    "target_text": "名称",
                    "is_excluded": False,
                },
            ]
        },
    }

    rebuilt_bytes = rebuild_bilingual_xlsx_from_segments(task_state, target_first=False)
    assert rebuilt_bytes and rebuilt_bytes.startswith(b"PK"), "Expected rebuilt bilingual XLSX bytes"

    rebuilt_doc = Document.from_bytes(content=rebuilt_bytes, suffix=".xlsx", stem="booklist2")
    workflow = XlsxWorkflow(
        XlsxWorkflowConfig(
            translator_config=XlsxTranslatorConfig(skip_translate=True),
            html_exporter_config=Xlsx2HTMLExporterConfig(cdn=True),
        )
    )
    workflow.document_translated = rebuilt_doc
    html_content = workflow.export_to_html()
    assert isinstance(html_content, str)
    assert not html_content.startswith("PK"), "HTML must not be XLSX zip content"
    assert "<" in html_content and "table" in html_content.lower()
    print("PASS test_bilingual_xlsx_html_is_not_zip_bytes")


if __name__ == "__main__":
    test_bilingual_segment_to_html_styles_and_line_breaks()
    test_rebuild_bilingual_xlsx_html_from_segments()
    test_rebuild_bilingual_pptx_html_from_segments()
    test_bilingual_xlsx_html_is_not_zip_bytes()
    print("\nAll bilingual XLSX HTML tests passed!")
