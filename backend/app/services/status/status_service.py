# SPDX-FileCopyrightText: 2026 Zampher
# SPDX-License-Identifier: MPL-2.0

"""
Status Service

Handles task status queries, logs, and preview endpoints.
"""

import asyncio
import logging
import multiprocessing
import os
import pickle
import tempfile
import threading
import time
import zipfile
import io
import base64
import mimetypes
import re
import json
import sys
from pathlib import Path
from queue import Empty
from typing import Dict, Any, Optional, List, Tuple, Callable
from concurrent.futures import ThreadPoolExecutor, as_completed
from fastapi import HTTPException
from fastapi.responses import JSONResponse

from backend.app.services.task import task_manager, MSG_LEVEL_WARNING
from logger import unified_logger as logger
from logger.logger import LogModule
from utils.pagination import parse_pagination_params, PaginatedResponse

# Track task_ids running background language detection so we don't start duplicate threads
_language_detection_tasks: set = set()
_language_detection_lock = threading.Lock()

# Multiprocessing: task_id -> {progress_queue, result_queue, process, segments_path} for draining progress in get_status
_detection_progress_queues: Dict[str, Dict[str, Any]] = {}
_detection_progress_queues_lock = threading.Lock()

# Message prefixes that indicate translation phase; do not overwrite with language-detection progress/message
# CRITICAL: Includes "Retranslating" for batch retry progress during failed segment retranslation
_TRANSLATION_PHASE_PREFIXES = (
    "Translating",
    "Retranslating",
    "Sending translation",
    "Generating output",
    "Translation completed",
    "Retranslation completed",
    "Translation initialized",
    "Preparing retranslation",
    "Batch retry",
)

# Minimal lang code normalization for frozen fallback (no anonymize/spacy/torch)
_LANG_NORMALIZE = {"zh-cn": "zh", "zh-tw": "zh", "no": "nb"}


def _strip_lang_detect_status_downgrade(
    task_manager_ref: Any,
    task_id: str,
    updates: Dict[str, Any],
) -> Dict[str, Any]:
    """
    Language-detection workers emit progress/status='processing'. Draining those updates into
    task_manager must not downgrade a task that already reached a terminal status (e.g. translation
    completed), otherwise the frontend poll loop never observes status=='completed'.
    """
    if not updates:
        return updates
    ts = task_manager_ref.get_task(task_id) or {}
    cur = str(ts.get("status") or "").lower()
    if cur not in ("completed", "failed", "cancelled"):
        return updates
    new_status = str(updates.get("status") or "").lower()
    if new_status != "processing":
        return updates
    out = {k: v for k, v in updates.items() if k != "status"}
    logger.debug(
        LogModule.WORKFLOW,
        f"[STATUS] Task {task_id}: preserving terminal status={cur}; omitted lang-detect "
        f"status downgrade (had progress/message keys={list(out.keys())})",
    )
    return out


def _language_detection_fallback_langdetect_only(
    segments: List[str],
    total_segments: int,
    progress_queue: multiprocessing.Queue,
    progress_interval: int = 500,
) -> Dict[str, Any]:
    """
    Fallback language detection using only langdetect (no anonymize/spacy/torch).
    Used when frozen build fails to load torch/spacy dylibs.
    """
    import langdetect
    from langdetect import detect, DetectorFactory

    DetectorFactory.seed = 0
    language_counts: Dict[str, int] = {}
    min_length = 10

    for idx, segment in enumerate(segments):
        if idx > 0 and idx % progress_interval == 0:
            progress_queue.put({
                "message": f"Detect Language: {idx}/{total_segments} segments ({int(100 * idx / total_segments)}%)",
                "progress": min(100, int(100 * idx / total_segments)),
                "status": "processing",
                "message_level": MSG_LEVEL_WARNING,
            })
        s = (segment or "").strip()
        if len(s) < min_length:
            continue
        try:
            code = detect(s)
            code = _LANG_NORMALIZE.get(code, code)
            language_counts[code] = language_counts.get(code, 0) + 1
        except Exception:
            continue

    total = sum(language_counts.values())
    language_distribution = {lang: count / total for lang, count in language_counts.items()} if total > 0 else {}
    recommended_language = "en"
    is_multilingual = False
    if language_distribution:
        recommended_language, max_pct = max(language_distribution.items(), key=lambda x: x[1])
        is_multilingual = max_pct < 0.9

    return {
        "language_distribution": language_distribution,
        "is_multilingual": is_multilingual,
        "recommended_language": recommended_language,
        "progress": 100,
        "message": f"Detect Language: {total_segments}/{total_segments} segments (100%)",
        "status": "processing",
        "message_level": MSG_LEVEL_WARNING,
    }


def _language_detection_worker(
    segments_path: str,
    task_id: str,
    total_segments: int,
    progress_queue: multiprocessing.Queue,
    result_queue: multiprocessing.Queue,
) -> None:
    """
    Run language detection in a subprocess. Put progress updates to progress_queue
    and final result to result_queue. Callable from multiprocessing.Process (module-level).
    On frozen runtime, falls back to langdetect-only if torch/spacy dylib load fails.
    """
    segments: List[str] = []
    try:
        with open(segments_path, "rb") as f:
            segments = pickle.load(f)
    except Exception as e:
        result_queue.put({"error": str(e), "status": "failed"})
        return

    def _progress_callback(completed: int, total: int, **kwargs: Any) -> None:
        if total == 0:
            return
        phase = kwargs.get("phase")
        if phase == "aggregated_short":
            pct = int((completed / total) * 100) if total else 0
            progress_queue.put({
                "message": f"Detect Language: short segments {completed}/{total} ({pct}%)",
                "progress": min(99, 90 + pct // 10),
                "status": "processing",
                "message_level": MSG_LEVEL_WARNING,
            })
        else:
            pct = int((completed / total) * 100)
            progress_queue.put({
                "message": f"Detect Language: {completed}/{total} segments ({pct}%)",
                "progress": min(100, pct),
                "status": "processing",
                "message_level": MSG_LEVEL_WARNING,
            })

    try:
        from utils.language_detection_utils import detect_language_per_segment_with_distribution

        lang_info = detect_language_per_segment_with_distribution(
            segments=segments,
            log_context=f"task_{task_id}",
            progress_callback=_progress_callback,
            progress_interval=500,
        )
        language_distribution = lang_info.get("language_distribution", {})
        is_multilingual = lang_info.get("is_multilingual", False)
        recommended_language = lang_info.get("recommended_language")
        result_queue.put({
            "language_distribution": language_distribution,
            "is_multilingual": is_multilingual,
            "recommended_language": recommended_language,
            "progress": 100,
            "message": f"Detect Language: {total_segments}/{total_segments} segments (100%)",
            "status": "processing",
            "message_level": MSG_LEVEL_WARNING,
        })
    except Exception as e:
        err_str = str(e).lower()
        if any(x in err_str for x in ("frozen", "dynlib", "dll", "torch", "was not found when the application was frozen")):
            logger.info(
                LogModule.WORKFLOW,
                "[STATUS] Language detection falling back to langdetect-only (torch/spacy unavailable in frozen build): " + str(e)[:200],
            )
            try:
                lang_info = _language_detection_fallback_langdetect_only(
                    segments=segments,
                    total_segments=total_segments,
                    progress_queue=progress_queue,
                    progress_interval=500,
                )
                result_queue.put(lang_info)
            except Exception as e2:
                result_queue.put({"error": str(e2), "status": "failed"})
        else:
            result_queue.put({"error": str(e), "status": "failed"})


class StatusService:
    """Service for querying task status, logs, and previews."""
    
    def __init__(self, task_manager_instance=None):
        """
        Initialize status service.
        
        Args:
            task_manager_instance: Optional TaskManager instance (defaults to global instance)
        """
        self.task_manager = task_manager_instance or task_manager
        self.logger = logging.getLogger(__name__)
    
    def _get_optimal_thread_count(self) -> int:
        """
        Calculate optimal thread count: 2/3 of CPU cores.
        
        Returns:
            Thread count (at least 1)
        """
        import multiprocessing
        cpu_count = multiprocessing.cpu_count()
        thread_count = max(1, int(cpu_count * 2 / 3))
        return thread_count

    def _sanitize_task_state_for_json(self, task_state: Dict[str, Any], task_id: str) -> None:
        """
        Remove or convert non-JSON-serializable keys from task_state in-place.
        Call this before returning task_state in JSONResponse to avoid TypeError
        (e.g. Object of type EpubWorkflow is not JSON serializable).
        """
        # Remove complex Python objects that cannot be serialized
        for key in (
            "workflow_instance",
            "workflow",
            "layout_document",
            "payload",
            "layout_source_zip",
            "mobi_book",
            "mobi_items_to_translate",
        ):
            if key in task_state:
                del task_state[key]
        # Convert current_task_ref (Task/future) to a string status
        if "current_task_ref" in task_state and task_state["current_task_ref"] is not None:
            ref = task_state["current_task_ref"]
            if not isinstance(ref, str):
                try:
                    task_state["current_task_ref"] = "running" if not ref.done() else "completed"
                except Exception:
                    task_state["current_task_ref"] = "unknown"
        # Ensure language_distribution (and other numeric dicts from worker) use native Python types for JSON
        if "language_distribution" in task_state and isinstance(task_state["language_distribution"], dict):
            task_state["language_distribution"] = {
                str(k): float(v) for k, v in task_state["language_distribution"].items()
            }
        # Recursively convert numpy/other non-standard numeric types to native float/int (e.g. after multiprocessing pickle)
        self._convert_to_native_json_types(task_state)

    def _slim_source_chunks_cache(self, cache: Dict[str, Any]) -> Dict[str, Any]:
        """Build a small dict for source_chunks_cache for status response (omit segments to avoid MemoryError)."""
        if not cache:
            return cache
        out = {k: v for k, v in cache.items() if k != "segments"}
        if "total_segments" not in out and "segments" in cache:
            segs = cache["segments"]
            if isinstance(segs, list):
                out["total_segments"] = len(segs)
        return out

    def _slim_source_preview(self, preview: Dict[str, Any]) -> Dict[str, Any]:
        """Build a small dict for source_preview for status response (omit segments to avoid MemoryError)."""
        if not preview:
            return preview
        out = {k: v for k, v in preview.items() if k != "segments"}
        if "total_segments" not in out and "segments" in preview:
            segs = preview["segments"]
            if isinstance(segs, list):
                out["total_segments"] = len(segs)
        return out

    def _slim_translation_segments(self, data: Dict[str, Any]) -> Dict[str, Any]:
        """Build a small dict for translation_segments for status response (omit segments list to avoid MemoryError)."""
        if not data:
            return data
        out = {k: v for k, v in data.items() if k != "segments"}
        if "segments" in data:
            segs = data["segments"]
            if isinstance(segs, list):
                out["segment_count"] = len(segs)
        return out

    def _build_slim_status_response(self, task_state: Dict[str, Any], task_id: str) -> Dict[str, Any]:
        """
        Build a response dict with large fields omitted/summarized to avoid MemoryError on json.dumps.
        Does not mutate task_state; only the returned dict is used for JSONResponse.
        """
        response: Dict[str, Any] = {}
        for key, value in task_state.items():
            if key == "source_chunks_cache" and isinstance(value, dict):
                response[key] = self._slim_source_chunks_cache(value)
            elif key == "source_preview" and isinstance(value, dict):
                response[key] = self._slim_source_preview(value)
            elif key == "translation_segments" and isinstance(value, dict):
                response[key] = self._slim_translation_segments(value)
            else:
                response[key] = value
        return response

    def _convert_to_native_json_types(self, obj: Any) -> None:
        """Convert numpy and other non-JSON-serializable types in-place (dict/list/values)."""
        if obj is None:
            return
        if isinstance(obj, dict):
            for k, v in list(obj.items()):
                if isinstance(v, (dict, list)):
                    self._convert_to_native_json_types(v)
                elif isinstance(v, (str, int, float, bool)) or v is None:
                    pass
                elif hasattr(v, "item") and callable(getattr(v, "item", None)):
                    try:
                        obj[k] = v.item()
                    except Exception:
                        try:
                            obj[k] = float(v)
                        except Exception:
                            obj[k] = int(v) if hasattr(v, "__int__") else str(v)
                else:
                    try:
                        obj[k] = float(v)
                    except (TypeError, ValueError):
                        try:
                            obj[k] = int(v)
                        except (TypeError, ValueError):
                            obj[k] = str(v)
        elif isinstance(obj, list):
            for i, v in enumerate(obj):
                if isinstance(v, (dict, list)):
                    self._convert_to_native_json_types(v)
                elif isinstance(v, (str, int, float, bool)) or v is None:
                    pass
                elif hasattr(v, "item") and callable(getattr(v, "item", None)):
                    try:
                        obj[i] = v.item()
                    except Exception:
                        try:
                            obj[i] = float(v)
                        except Exception:
                            obj[i] = int(v) if hasattr(v, "__int__") else str(v)
                else:
                    try:
                        obj[i] = float(v)
                    except (TypeError, ValueError):
                        try:
                            obj[i] = int(v)
                        except (TypeError, ValueError):
                            obj[i] = str(v)

    def _update_progress(
        self,
        task_id: str,
        operation: str,
        completed: int,
        total: int,
        base_progress: int = 0,
        progress_range: int = 100
    ) -> None:
        """
        Update progress for an operation.
        
        Args:
            task_id: Task ID
            operation: Operation name (e.g., 'identifying', 'language_matching')
            completed: Number of completed items
            total: Total number of items
            base_progress: Base progress percentage (default: 0)
            progress_range: Progress range percentage (default: 100)
        """
        if total == 0:
            return
        
        percent = int((completed / total) * progress_range)
        progress = base_progress + percent
        
        # Map operation names to friendly display names
        operation_display_map = {
            'identifying': 'Detect Identifier',
            'language_matching': 'Detect Language',
            'detecting_exclusions': 'Detect Exclusions',
        }
        operation_display = operation_display_map.get(operation, operation.replace('_', ' ').title())
        message = f"{operation_display}: {completed}/{total} segments ({percent}%)"
        
        try:
            task_state = self.task_manager.get_task(task_id)
            if task_state:
                task_state["progress"] = min(100, progress)
                task_state["message"] = message
                self.task_manager.update_last_logged_status(
                    task_id,
                    {'status': task_state.get('status', 'processing'), 'progress': progress, 'message': message}
                )
        except Exception as e:
            logger.warning(LogModule.EXCLUSION, f"Failed to update progress: {e}")
    
    def _process_segments_concurrently(
        self,
        segments: List[Any],
        process_func: Callable[[int, Any], Tuple[int, Optional[Any]]],
        task_id: str,
        operation: str,
        base_progress: int = 0,
        progress_range: int = 100,
        progress_interval: int = 500
    ) -> Dict[int, Any]:
        """
        Process segments concurrently with progress updates.
        
        Args:
            segments: List of segments to process
            process_func: Function that processes a single segment, returns (index, result)
            task_id: Task ID for progress updates
            operation: Operation name for progress messages
            base_progress: Base progress percentage
            progress_range: Progress range percentage
            progress_interval: Update progress every N segments
            
        Returns:
            Dictionary mapping segment index to result
        """
        total = len(segments)
        if total == 0:
            return {}
        
        results = {}
        thread_count = self._get_optimal_thread_count()
        
        # Initial progress update
        self._update_progress(task_id, operation, 0, total, base_progress, progress_range)
        
        with ThreadPoolExecutor(max_workers=thread_count) as executor:
            # Submit all tasks
            future_to_index = {
                executor.submit(process_func, idx, seg): idx
                for idx, seg in enumerate(segments)
            }
            
            completed = 0
            last_progress_update = 0
            
            # Process completed tasks
            for future in as_completed(future_to_index):
                idx = future_to_index[future]
                try:
                    result_idx, result = future.result()
                    if result is not None:
                        results[result_idx] = result
                except Exception as e:
                    logger.warning(
                        LogModule.EXCLUSION,
                        f"Error processing segment {idx} in {operation}: {e}"
                    )
                
                completed += 1
                
                # Update progress every progress_interval segments
                if completed - last_progress_update >= progress_interval or completed == total:
                    self._update_progress(
                        task_id, operation, completed, total, base_progress, progress_range
                    )
                    last_progress_update = completed
        
        # Final progress update
        self._update_progress(task_id, operation, total, total, base_progress, progress_range)
        
        return results

    def _stash_only_status_payload(self, task_id: str) -> Optional[Dict[str, Any]]:
        """
        Build status JSON when the task is no longer in memory but stashed outputs exist.
        Matches download_service stash fallback so GET /service/status and GET /service/download agree.
        """
        from backend.app.services.translation.translation_result_stash import (
            get_stashed_file_path,
            load_meta,
        )

        meta = load_meta(task_id)
        if not meta:
            return None
        files_map = meta.get("files") or {}
        if not isinstance(files_map, dict) or not files_map:
            return None
        downloads: Dict[str, str] = {}
        for file_type in list(files_map.keys()):
            path = get_stashed_file_path(task_id, file_type)
            if path and os.path.isfile(path):
                downloads[file_type] = f"/service/download/{task_id}/{file_type}"
        if not downloads:
            return None
        return {
            "task_id": task_id,
            "status": "completed",
            "message": "Translated outputs available for download.",
            "progress": 100,
            "original_filename": meta.get("original_filename"),
            "download_ready": True,
            "downloads": downloads,
            "attachments": {},
            "in_memory": False,
            "results_stashed": True,
        }
    
    async def get_status(self, task_id: str):
        """
        Get task status.
        
        Args:
            task_id: Unique task identifier
            
        Returns:
            JSONResponse with task status
            
        Raises:
            HTTPException: If task not found
        """
        task_state = self.task_manager.get_task(task_id)
        if task_state is None:
            stash_payload = self._stash_only_status_payload(task_id)
            if stash_payload is not None:
                return JSONResponse(content=stash_payload)
            raise HTTPException(status_code=404, detail=f"Task ID '{task_id}' not found.")

        # Drain progress from language detection worker (multiprocessing) so response has latest progress.
        # Do not overwrite progress/message if task has already entered translation phase (MOBI/EPUB etc.),
        # so frontend sees translation chunk progress instead of stale "Detect Language: 100%".
        with _detection_progress_queues_lock:
            entry = _detection_progress_queues.get(task_id)
        if entry:
            progress_queue = entry.get("progress_queue")
            if progress_queue is not None:
                while True:
                    try:
                        msg = progress_queue.get_nowait()
                    except Empty:
                        break
                    current = self.task_manager.get_task(task_id)
                    cur_msg = (current.get("message") or "") if current else ""
                    if cur_msg.startswith(_TRANSLATION_PHASE_PREFIXES):
                        msg_filtered = {k: v for k, v in msg.items() if k not in ("progress", "message")}
                        if "progress" in msg or "message" in msg:
                            # Use system log module for status polling debug logs
                            logger.debug(
                                LogModule.SYSTEM,
                                f"[STATUS] Task {task_id}: skipping language-detection progress/message (already in translation phase: {cur_msg[:50]}...)"
                            )
                        if msg_filtered:
                            self.task_manager.update_task(task_id, msg_filtered)
                    else:
                        merged = _strip_lang_detect_status_downgrade(
                            self.task_manager, task_id, msg
                        )
                        self.task_manager.update_task(task_id, merged)
            task_state = self.task_manager.get_task(task_id)

        if task_state is None:
            stash_payload = self._stash_only_status_payload(task_id)
            if stash_payload is not None:
                return JSONResponse(content=stash_payload)
            raise HTTPException(status_code=404, detail=f"Task ID '{task_id}' not found.")
        # Create a copy for response (to avoid modifying original)
        task_state = task_state.copy()
        task_state["task_id"] = task_id

        # Log status query only when status/progress/message changes (reduce log verbosity)
        status = task_state.get('status', 'N/A')
        error_msg = task_state.get('error', '')
        message = task_state.get('message', '')
        progress = task_state.get('progress', 'N/A')
        status_lower = str(status).lower()

        # Normalize status for completed translations so frontend can display results consistently.
        # Some flows (e.g. markdown_based with language-detection or layout post-processing) may leave
        # status as "processing" even after translation has fully completed. If we detect:
        # - status == "processing"
        # - progress >= 100
        # - message starts with "Translation completed"
        # then we treat this as a completed translation and update status to "completed".
        try:
            is_processing = status_lower == "processing"
            is_full_progress = isinstance(progress, (int, float)) and int(progress) >= 100
            msg_str = str(message or "")
        except Exception:
            is_processing = False
            is_full_progress = False
            msg_str = str(message or "")

        if is_processing and is_full_progress and msg_str.startswith("Translation completed"):
            status = "completed"
            status_lower = "completed"
            task_state["status"] = status
            # Persist normalized status so subsequent calls also see "completed"
            try:
                self.task_manager.update_task(task_id, {"status": status})
            except Exception:
                # Best-effort; do not break status endpoint if update fails
                pass
        
        # Check if status/progress/message has changed since last log
        last_logged = self.task_manager.get_last_logged_status(task_id) or {}
        status_changed = (
            last_logged.get('status') != status or
            last_logged.get('progress') != progress or
            last_logged.get('message') != message
        )
        
        # Always log errors and failures, or when status changes
        if status_lower == "failed" or error_msg:
            status_line = f"[STATUS] task={task_id} progress={progress} status={status} message={message}"
            logger.error(LogModule.WORKFLOW, status_line)
            if error_msg:
                logger.error(LogModule.WORKFLOW, f"[STATUS-ERROR] task={task_id} error={error_msg}")
            # Update last logged status
            self.task_manager.update_last_logged_status(task_id, {'status': status, 'progress': progress, 'message': message})
        elif status_changed:
            # Only log when status/progress/message changes
            status_line = f"[STATUS] task={task_id} progress={progress} status={status} message={message}"
            logger.info(LogModule.WORKFLOW, status_line)
            # Update last logged status
            self.task_manager.update_last_logged_status(task_id, {'status': status, 'progress': progress, 'message': message})
        # Otherwise, skip logging (same status as before)
        
        # Remove non-serializable objects (current_task_ref may already be "running"/"completed" after re-fetch in already_running branch)
        if "current_task_ref" in task_state and task_state["current_task_ref"] is not None:
            ref = task_state["current_task_ref"]
            if not isinstance(ref, str):
                try:
                    task_state["current_task_ref"] = "running" if not ref.done() else "completed"
                except Exception:
                    task_state["current_task_ref"] = "unknown"
        
        # Extract serializable task params from payload before removing it,
        # so the frontend can restore original task settings in reedit mode.
        task_params: Dict[str, Any] = {}
        if "payload" in task_state:
            raw_payload = task_state["payload"]
            if isinstance(raw_payload, dict):
                payload_dict = raw_payload
            else:
                # Pydantic model or SimpleNamespace
                try:
                    payload_dict = dict(raw_payload)
                except Exception:
                    try:
                        payload_dict = raw_payload.model_dump() if hasattr(raw_payload, "model_dump") else {}
                    except Exception:
                        payload_dict = {}
            # Extract fields relevant to frontend quick settings
            _relevant_keys = [
                "to_lang", "workflow_type", "deep_split", "temperature",
                "prompt_mode", "prompt_style", "custom_note", "skip_translate",
                "convert_engine", "formula_ocr", "table_ocr", "model_version",
                "ocr_language", "insert_mode", "separator", "chunk_size",
                "base_url", "model_id", "concurrent", "timeout", "retry",
                "segment_auto_retry_rounds", "thinking", "copy_source_only",
                "glossary_dict", "glossary_generate_enable",
            ]
            for k in _relevant_keys:
                if k in payload_dict and payload_dict[k] is not None:
                    val = payload_dict[k]
                    # Skip complex non-serializable types
                    if isinstance(val, (str, int, float, bool, list, dict)) or val is None:
                        task_params[k] = val
            del task_state["payload"]
        task_state["task_params"] = task_params
        
        # Check if layout_document is available BEFORE removing it (for PDF files)
        original_filename = task_state.get("original_filename", "")
        is_pdf_file = original_filename.lower().endswith('.pdf')
        has_layout_for_pdf = False
        has_tables = False
        has_interline_equations = False
        has_charts = False
        page_count = 0  # Initialize page count
        
        # Extract page count from layout_document for PDF files
        if is_pdf_file:
            layout_doc = task_state.get("layout_document")
            if layout_doc is not None:
                try:
                    from layout.base import LayoutDocument as _LD
                    if isinstance(layout_doc, _LD):
                        has_layout_for_pdf = True
                        page_count = layout_doc.page_count
                        logger.trace(LogModule.WORKFLOW, f"PDF file with layout_document available for task {task_id}: {page_count} pages")
                        # Check if document has tables, interline_equations, and charts
                        for page in layout_doc.pages:
                            for block in page.blocks:
                                if block.type == "table":
                                    has_tables = True
                                elif block.type == "interline_equation":
                                    has_interline_equations = True
                                elif block.type == "chart":
                                    has_charts = True
                                if has_tables and has_interline_equations and has_charts:
                                    break
                            if has_tables and has_interline_equations and has_charts:
                                break
                        if has_tables:
                            logger.trace(LogModule.WORKFLOW, f"PDF file has tables for task {task_id}")
                        if has_interline_equations:
                            logger.trace(LogModule.WORKFLOW, f"PDF file has interline_equations for task {task_id}")
                        if has_charts:
                            logger.trace(LogModule.WORKFLOW, f"PDF file has charts for task {task_id}")
                except Exception as e:
                    logger.warning(LogModule.WORKFLOW, f"[STATUS] Failed to extract page count from layout_document: {e}")
        
        # If layout_document is not available, check if page_count was stored
        # early in task_state (e.g. by translation_service on task creation).
        if page_count == 0:
            early_page_count = task_state.get("page_count")
            if early_page_count:
                page_count = int(early_page_count)
                logger.trace(LogModule.WORKFLOW, f"Using early page_count={page_count} from task_state for task {task_id}")
        
        # Extract page count for other document formats
        if page_count == 0:
            page_count = self._extract_page_count_from_document(
                task_state, original_filename
            )
        
        # Store page_count in task_state before removing layout_document
        if page_count > 0:
            task_state["page_count"] = page_count
            logger.trace(LogModule.WORKFLOW, f"Extracted page_count={page_count} for task {task_id}, filename={original_filename}")
        
        # Remove layout_document and workflow_instance (complex Python objects, not JSON serializable)
        if "layout_document" in task_state:
            del task_state["layout_document"]
        if "workflow_instance" in task_state:
            del task_state["workflow_instance"]
        if "layout_source_zip" in task_state:
            del task_state["layout_source_zip"]
        
        # Remove MOBI-related objects that are not JSON serializable (for delayed DOM generation)
        if "mobi_book" in task_state:
            del task_state["mobi_book"]
        if "mobi_items_to_translate" in task_state:
            del task_state["mobi_items_to_translate"]
        
        # Ensure downloads is always a dict (not a string)
        if "downloads" in task_state and not isinstance(task_state["downloads"], dict):
            task_state["downloads"] = {}
        
        # Generate download links if task is completed and has downloadable files
        # For PDF files, include PDF link if layout_document is available (even if PDF not yet generated)
        # Also allow downloads if downloadable_files exist (e.g., for preview-generated files)
        downloads = {}
        has_downloadable_files = bool(task_state.get("downloadable_files"))
        is_format_conversion = task_state.get("is_format_conversion", False) or task_state.get("convert_only", False)
        # Resolve workflow_type for MOBI (used to hide EPUB when ebook-convert is not available)
        workflow_type = None
        _segs_data = task_state.get("translation_segments")
        if isinstance(_segs_data, dict):
            _meta = _segs_data.get("metadata", {})
            if isinstance(_meta, dict):
                workflow_type = _meta.get("workflow_type")
        if not workflow_type and original_filename:
            _ext = (original_filename or "").lower().split(".")[-1] if "." in (original_filename or "") else ""
            if _ext == "mobi":
                workflow_type = "mobi"
            elif _ext == "epub":
                workflow_type = "epub"
        if task_state.get("download_ready") or has_downloadable_files:
            is_md_file = original_filename.lower().endswith('.md')
            is_png_file = original_filename.lower().endswith(('.png', '.jpg', '.jpeg'))
            
            # Add existing downloadable files
            if task_state.get("downloadable_files"):
                for file_type, file_path in task_state["downloadable_files"].items():
                    # Skip PDF download for non-PDF files
                    if file_type == "pdf" and not is_pdf_file:
                        logger.trace(LogModule.WORKFLOW, f"[STATUS] Skipping PDF download link for task {task_id}: original file is not PDF")
                        continue
                    # Skip PDF download for md and png files
                    if file_type == "pdf" and (is_md_file or is_png_file):
                        logger.trace(LogModule.WORKFLOW, f"Skipping PDF download link for task {task_id}: original file is MD or PNG")
                        continue
                    # Check if file exists (handle both string and dict formats)
                    file_path_str = file_path.get("path", "") if isinstance(file_path, dict) else str(file_path)
                    if not file_path_str or not os.path.exists(file_path_str):
                        logger.trace(LogModule.WORKFLOW, f"Skipping {file_type} download link for task {task_id}: file does not exist at {file_path_str}")
                        continue
                    logger.trace(LogModule.WORKFLOW, f"Adding {file_type} download link for task {task_id}: file exists at {file_path_str}")
                    downloads[file_type] = f"/service/download/{task_id}/{file_type}"
            
            # For PDF files, add PDF download link if layout is available (even if not yet generated)
            # Skip PDF download for format conversion tasks (Convert) - not supported yet
            if is_pdf_file and has_layout_for_pdf and "pdf" not in downloads and not is_format_conversion:
                downloads["pdf"] = f"/service/download/{task_id}/pdf"
                logger.info(LogModule.WORKFLOW, f"[STATUS] Added PDF download link for task {task_id}: layout-based PDF available (will be generated on-demand)")
        
        # Completed tasks: merge on-demand export URLs for all supported workflows (see
        # download_service._build_stash_export_plan) so Translation queue shows the same buttons as
        # the in-app download bar.
        status_lower = (task_state.get("status") or "").lower()
        if status_lower == "completed":
            from backend.app.services.download.download_service import (
                completed_task_download_urls,
                resolve_task_export_workflow_type,
            )

            palette = completed_task_download_urls(task_id, task_state)
            before_ct = len(downloads)
            for key, url in palette.items():
                downloads.setdefault(key, url)
            # Local `downloads` is rebuilt on every status query; merge always grows from 0 when only
            # the palette applies — INFO here would log once per poll. Keep diagnosis at DEBUG.
            if len(downloads) > before_ct:
                logger.debug(
                    LogModule.WORKFLOW,
                    f"[STATUS] Task {task_id}: merged completed-task download palette: {list(downloads.keys())}",
                )
            wt_resolved = resolve_task_export_workflow_type(task_state)
            if wt_resolved:
                workflow_type = wt_resolved
        
        attachments = {}
        if task_state.get("download_ready") and task_state.get("attachment_files"):
            for identifier in task_state["attachment_files"].keys():
                attachments[identifier] = f"/service/attachment/{task_id}/{identifier}"
        
        # MOBI workflow: do not offer EPUB when Calibre (ebook-convert) is not available
        if workflow_type == "mobi" and "epub" in downloads:
            try:
                from app.services.download.output_generator import _calibre_cmd_path
                if not _calibre_cmd_path():
                    downloads.pop("epub", None)
            except Exception:
                pass

        # Add download links to response (always ensure it's a dict)
        task_state["downloads"] = downloads
        task_state["attachments"] = attachments
        
        # Log downloads only when they change (not on every status query)
        if downloads:
            # Only log if this is the first time we're seeing downloads for this task
            if not task_state.get("_downloads_logged", False):
                logger.info(LogModule.WORKFLOW, f"[STATUS] Task {task_id} downloads: {list(downloads.keys())}")
                self.task_manager.update_task(task_id, {"_downloads_logged": True})
        
        # Calculate translation statistics from segments if available
        segments_data = task_state.get("translation_segments")
        if segments_data:
            segments = segments_data.get("segments", [])
            if segments:
                total_segments = len(segments)
                # Count success: segments that are not failed OR are image segments
                # Image segments are always considered successful (they don't need translation)
                success_count = sum(1 for seg in segments if not seg.get("is_failed", False) or seg.get("is_image", False))
                # Count failures: segments that are failed AND not image segments
                # Image segments should never be counted as failures
                fail_count = sum(1 for seg in segments if seg.get("is_failed", False) and not seg.get("is_image", False))
                task_state["translation_stats"] = {
                    "total_segments": total_segments,
                    "success_count": success_count,
                    "fail_count": fail_count,
                }
        
        # Ensure token_usage is included in response if it exists
        # token_usage is already in task_state (from _process_translation_task), so it will be included automatically
        token_usage = task_state.get("token_usage")
        if token_usage:
            # Verify token_usage is a dict and has required fields
            if isinstance(token_usage, dict) and "total_tokens" in token_usage:
                # Ensure all required fields are present
                if "input_tokens" not in token_usage:
                    token_usage["input_tokens"] = 0
                if "cached_tokens" not in token_usage:
                    token_usage["cached_tokens"] = 0
                if "output_tokens" not in token_usage:
                    token_usage["output_tokens"] = 0
                if "reasoning_tokens" not in token_usage:
                    token_usage["reasoning_tokens"] = 0
            else:
                # Log error if token_usage format is invalid (for troubleshooting)
                logger.warning(LogModule.WORKFLOW, f"[TokenStats] Status API: task {task_id} has token_usage but invalid format: {type(token_usage)}")
        
        # Detect language distribution for multilingual detection (if segments are available)
        # Use cached result if available to avoid repeated detection and logging
        language_distribution = task_state.get("language_distribution")
        is_multilingual = task_state.get("is_multilingual", False)
        detected_language = task_state.get("detected_language")
        
        if language_distribution is None:
            try:
                # Try to get segments from cache first
                cache_info = task_state.get("source_chunks_cache", {})
                segments_for_lang = cache_info.get("segments", [])
                
                # Fallback to preview segments if cache not available
                if not segments_for_lang:
                    source_preview = task_state.get("source_preview", {})
                    if isinstance(source_preview, dict):
                        segments_for_lang = source_preview.get("segments", [])
                
                # If we have segments, run language detection (in background for large files so status returns quickly)
                if segments_for_lang and len(segments_for_lang) > 0:
                    total_segments = len(segments_for_lang)
                    # Do not start this worker while translation is active: it overwrites progress/message with
                    # "Detect Language: …%". Translation tasks reuse segment caches and would otherwise hit 100%
                    # before chunk callbacks, breaking the progress bar (see workflow_executor translation phase).
                    fresh_live = self.task_manager.get_task(task_id) or {}
                    skip_detect_start = str(fresh_live.get("message") or "").startswith(
                        _TRANSLATION_PHASE_PREFIXES
                    )
                    with _language_detection_lock:
                        already_running = task_id in _language_detection_tasks
                    if skip_detect_start and not already_running:
                        logger.debug(
                            LogModule.WORKFLOW,
                            f"[STATUS] Task {task_id}: skip language-detection worker start "
                            f"(translation phase active)",
                        )
                        for key in ("language_distribution", "is_multilingual", "detected_language"):
                            if key in fresh_live:
                                task_state[key] = fresh_live[key]
                        language_distribution = task_state.get("language_distribution")
                        is_multilingual = task_state.get("is_multilingual", False)
                        detected_language = task_state.get("detected_language")
                    elif already_running:
                        # Detection already in progress; copy only updated fields from fresh task_state
                        # into our response copy so we do not mutate the original (e.g. payload must stay)
                        fresh = self.task_manager.get_task(task_id)
                        if fresh:
                            for key in ("language_distribution", "is_multilingual", "detected_language"):
                                if key in fresh:
                                    task_state[key] = fresh[key]
                        language_distribution = task_state.get("language_distribution")
                        is_multilingual = task_state.get("is_multilingual", False)
                        detected_language = task_state.get("detected_language")
                    else:
                        # Start language detection in a subprocess so main process event loop is not blocked (no GIL contention)
                        with _language_detection_lock:
                            _language_detection_tasks.add(task_id)
                        self.task_manager.update_task(task_id, {
                            "message": f"Detect Language: 0/{total_segments} segments (0%)",
                            "message_level": MSG_LEVEL_WARNING,
                            "progress": 0,
                            "status": task_state.get("status", "processing")
                        })
                        task_manager_ref = self.task_manager

                        # Write segments to temp file for worker process (avoid large pickle over Process args)
                        fd, segments_path = tempfile.mkstemp(suffix=".pkl", prefix="langdetect_")
                        try:
                            with os.fdopen(fd, "wb") as f:
                                pickle.dump(segments_for_lang, f)
                        except Exception:
                            os.close(fd)
                            with _language_detection_lock:
                                _language_detection_tasks.discard(task_id)
                            raise

                        progress_queue: multiprocessing.Queue = multiprocessing.Queue()
                        result_queue: multiprocessing.Queue = multiprocessing.Queue()
                        
                        # In frozen (PyInstaller) desktop builds, multiprocessing.Process can be fragile.
                        # Use a thread-based worker in that environment to avoid spawn/exec issues.
                        use_subprocess = not getattr(sys, "frozen", False)
                        if use_subprocess:
                            process = multiprocessing.Process(
                                target=_language_detection_worker,
                                args=(segments_path, task_id, total_segments, progress_queue, result_queue),
                            )
                        else:
                            logger.info(
                                LogModule.WORKFLOW,
                                f"[STATUS] Using thread-based language detection worker for task {task_id} (frozen runtime detected)",
                            )
                            process = threading.Thread(
                                target=_language_detection_worker,
                                args=(segments_path, task_id, total_segments, progress_queue, result_queue),
                                daemon=True,
                            )
                        process.start()
                        with _detection_progress_queues_lock:
                            _detection_progress_queues[task_id] = {
                                "progress_queue": progress_queue,
                                "result_queue": result_queue,
                                "process": process,
                                "segments_path": segments_path,
                            }

                        def _join_and_finish() -> None:
                            process.join()
                            try:
                                result = result_queue.get(timeout=1)
                            except Exception:
                                # Worker process did not put a result; treat as language detection failure
                                result = {"error": "Worker did not return result", "status": "failed"}
                            try:
                                ts = task_manager_ref.get_task(task_id) or {}
                                cur_status = ts.get("status")
                                cur_msg = ts.get("message") or ""

                                if result.get("error"):
                                    # Log detailed error for troubleshooting
                                    logger.error(
                                        LogModule.WORKFLOW,
                                        f"[STATUS] Language detection worker error for task {task_id}: {result.get('error')}",
                                    )
                                    # If translation is already in progress/completed, do NOT fail the whole task.
                                    # Store error in task state and keep existing status/progress.
                                    if cur_status in ("processing", "completed") and cur_msg.startswith(
                                        _TRANSLATION_PHASE_PREFIXES
                                    ):
                                        update_fields = {
                                            "language_detection_error": result.get("error", "Language detection failed")
                                        }
                                        task_manager_ref.update_task(task_id, update_fields)
                                    else:
                                        # Before translation starts (pure format-conversion or early detect failure),
                                        # it is still reasonable to mark task as failed.
                                        task_manager_ref.update_task(task_id, {
                                            "status": "failed",
                                            "error": result.get("error", "Language detection failed"),
                                            "message": f"Language detection error: {result.get('error', 'Language detection failed')}",
                                            "message_level": MSG_LEVEL_WARNING,
                                        })
                                else:
                                    # Do not overwrite progress/message if task already in translation phase
                                    ts = task_manager_ref.get_task(task_id)
                                    cur_msg = (ts or {}).get("message") or ""
                                    if cur_msg.startswith(_TRANSLATION_PHASE_PREFIXES):
                                        # Do not overwrite status: language detection result has "processing",
                                        # which would overwrite "completed" if translation already finished
                                        filtered = {
                                            k: v
                                            for k, v in result.items()
                                            if k not in ("progress", "message", "status")
                                        }
                                        if filtered:
                                            task_manager_ref.update_task(task_id, filtered)
                                    else:
                                        merged = _strip_lang_detect_status_downgrade(
                                            task_manager_ref, task_id, result
                                        )
                                        task_manager_ref.update_task(task_id, merged)
                                    if result.get("recommended_language"):
                                        task_manager_ref.update_task(
                                            task_id, {"detected_language": result["recommended_language"]}
                                        )
                                    # For format-conversion-only tasks, set status to completed so frontend finishes progress
                                    ts = task_manager_ref.get_task(task_id)
                                    if ts and (ts.get("convert_only") or ts.get("is_format_conversion")):
                                        task_manager_ref.update_task(task_id, {
                                            "status": "completed",
                                            "progress": 100,
                                            "message": result.get("message", "Format conversion completed successfully"),
                                        })
                                        logger.info(
                                            LogModule.WORKFLOW,
                                            f"[STATUS] Task {task_id}: Language detection done, set status=completed (format-conversion task)",
                                        )
                            except Exception as e:
                                logger.trace(
                                    LogModule.WORKFLOW,
                                    f"Language distribution detection result apply failed for task {task_id}: {e}"
                                )
                            finally:
                                with _detection_progress_queues_lock:
                                    _detection_progress_queues.pop(task_id, None)
                                with _language_detection_lock:
                                    _language_detection_tasks.discard(task_id)
                                try:
                                    os.remove(segments_path)
                                except Exception:
                                    pass

                        threading.Thread(target=_join_and_finish, daemon=True).start()
                        # Return immediately with progress 0 so frontend shows Extract tab and can poll for progress
                        language_distribution = None
                        task_state["message"] = f"Detect Language: 0/{total_segments} segments (0%)"
                        task_state["message_level"] = MSG_LEVEL_WARNING
                        task_state["progress"] = 0
            except Exception as e:
                # If detection fails, log but don't fail the request
                logger.trace(LogModule.WORKFLOW, f"Language distribution detection failed for task {task_id}: {e}")
        
        # Extract detected_language from language_distribution if not already set
        if not detected_language and language_distribution:
            # Get the language with highest percentage
            if language_distribution:
                detected_language = max(language_distribution.items(), key=lambda x: x[1])[0]
                self.task_manager.update_task(task_id, {"detected_language": detected_language})
        
        # Store in task_state for response (use cached or newly detected values)
        task_state["language_distribution"] = language_distribution
        task_state["is_multilingual"] = is_multilingual
        if detected_language:
            task_state["detected_language"] = detected_language
        
        # Add has_tables, has_interline_equations, and has_charts fields for PDF files (for format selection dialog)
        if is_pdf_file:
            task_state["has_tables"] = has_tables
            task_state["has_interline_equations"] = has_interline_equations
            task_state["has_charts"] = has_charts

        # Ensure response is JSON-serializable (e.g. after re-fetch in already_running branch)
        self._sanitize_task_state_for_json(task_state, task_id)

        # Use slim response to avoid MemoryError when serializing tasks with huge segment lists
        response_payload = self._build_slim_status_response(task_state, task_id)
        return JSONResponse(content=response_payload)

    def get_logs(self, task_id: str) -> Dict[str, Any]:
        """
        Get task logs.
        
        Args:
            task_id: Unique task identifier
            
        Returns:
            Dictionary with task_id and logs list
            
        Raises:
            HTTPException: If task not found
        """
        task_state = self.task_manager.get_task(task_id)
        if task_state is None:
            raise HTTPException(status_code=404, detail=f"Task ID '{task_id}' not found.")
        
        # Get logs from task manager
        logs = self.task_manager.get_logs(task_id)
        
        return {
            "task_id": task_id,
            "logs": logs
        }
    
    async def get_source_preview(
        self,
        task_id: str,
        offset: int = 0,
        limit: Optional[int] = None,
        target_lang: Optional[str] = None
    ):
        """
        Get source text preview segments.
        
        Args:
            task_id: Unique task identifier
            offset: Number of segments to skip
            limit: Maximum number of segments to return
            
        Returns:
            JSONResponse with paginated segments and metadata
            
        Raises:
            HTTPException: If task not found
        """
        logger.info(
            LogModule.EXTRACT,
            f"[PREVIEW-API] get_source_preview called: task_id={task_id}, "
            f"offset={offset}, limit={limit}, target_lang={target_lang}"
        )

        task_state = self.task_manager.get_task(task_id)
        if task_state is None:
            raise HTTPException(status_code=404, detail=f"Task ID '{task_id}' not found.")
        
        st = task_state
        
        # Check if task has failed
        # CRITICAL: Even if task failed, we should still allow source preview access
        # if source_preview data is available (extraction may have succeeded before translation failed)
        task_status = st.get("status")
        cache_info = st.get("source_chunks_cache", {})
        preview = st.get("source_preview") or {}
        has_source_data = bool(cache_info.get("segments")) or bool(preview.get("segments"))
        
        if task_status == "failed" and not has_source_data:
            # Only block if task failed AND no source data is available
            error_message = st.get("error") or st.get("message") or "Task processing failed"
            raise HTTPException(
                status_code=400,
                detail=f"Task '{task_id}' has failed and cannot provide source preview. Error: {error_message}"
            )
        
        # Parse pagination parameters
        # Import pagination config for consistent limits
        from app.config.pagination_config import MAX_PAGINATION_LIMIT, DEFAULT_PAGINATION_LIMIT, DEFAULT_SEGMENT_PREVIEW_LIMIT
        pagination = parse_pagination_params(offset=offset, limit=limit, max_limit=MAX_PAGINATION_LIMIT, default_limit=DEFAULT_PAGINATION_LIMIT)
        
        # Try to get full segments from cache first (more reliable)
        cache_info = st.get("source_chunks_cache", {})
        all_segments = cache_info.get("segments", [])
        
        # CRITICAL: Get total_segments from cache if available (more accurate than len(all_segments))
        # This ensures we return the correct total even if cache segments are truncated
        cache_total_segments = cache_info.get("total_segments")
        if cache_total_segments is not None:
            total = cache_total_segments
        else:
            # Fallback to len(all_segments) if total_segments not in cache
            total = len(all_segments) if all_segments else 0
        
        # Removed frequent "Getting segments" logging to reduce log verbosity
        # Only log when there are issues (e.g., cache mismatch)
        
        # CRITICAL: If cache has segments but count doesn't match total_segments, log warning
        # This indicates cache segments might be truncated or incomplete
        if all_segments and cache_total_segments is not None:
            if len(all_segments) != cache_total_segments:
                logger.warning(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: WARNING - Cache segments count ({len(all_segments)}) "
                    f"does not match cache total_segments ({cache_total_segments}). "
                    f"This may cause pagination issues. Using total_segments={cache_total_segments} for total, "
                    f"but only {len(all_segments)} segments available for processing."
                    )
        
        # Fallback to preview segments if cache not available
        if not all_segments:
            preview = st.get("source_preview") or {"segments": [], "total_segments": 0, "ready": False}
            all_segments = preview.get("segments", [])
            # Use preview total_segments if cache total not available
            if cache_total_segments is None:
                preview_total = preview.get("total_segments", 0)
                if preview_total > 0:
                    total = preview_total
                else:
                    total = len(all_segments) if all_segments else 0
            logger.debug(
                LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Using source_preview fallback - "
                f"preview segments count: {len(all_segments) if all_segments else 0}, "
                f"preview total_segments: {preview.get('total_segments', 0)}, "
                f"calculated total: {total}, "
                f"preview segments type: {type(all_segments).__name__}"
                )
        
        # Check both source_preview.ready and cache_info existence
        source_preview = st.get("source_preview", {})
        preview_ready = source_preview.get("ready", False) if isinstance(source_preview, dict) else False
        has_cache = bool(cache_info)
        ready = bool(preview_ready or has_cache)
        
        # Log segment count information for debugging (reduced verbosity - only log on first request or when cache changes)
        # Removed frequent logging to reduce log noise
        # logger.info(...)  # Removed: too frequent
        
        # Import unified exclusion detection utility
        from utils.translation_segments import _is_image_segment
        
        # Get target_lang from function parameter or fallback to payload
        if not target_lang:
            payload = st.get("payload")
            target_lang = self._get_payload_target_lang(payload)
        
        # Removed target_lang logging to reduce log verbosity
        
        # CRITICAL: Priority order for is_excluded status:
        # 1. translation_segments (most up-to-date, updated by update_excluded_segments_for_language)
        # 2. segments_metadata.excluded_segment_indices (pre-computed during extraction)
        # 3. Real-time detection (fallback)
        
        # First, try to get is_excluded and target_text from translation_segments (most reliable)
        translation_segments_data = st.get("translation_segments")
        translation_segments_map = {}  # Map segment_index -> {is_excluded, is_image, target_text, is_failed, etc.}
        failed_segments_in_translation = []  # Track failed segments for debugging
        if translation_segments_data and isinstance(translation_segments_data, dict):
            segments_list = translation_segments_data.get("segments", [])
            excluded_segments_in_translation = []
            for segment in segments_list:
                if isinstance(segment, dict):
                    segment_index = segment.get("segment_index")
                    if segment_index is not None:
                        is_excluded_seg = segment.get("is_excluded", False)
                        exclusion_reason_seg = segment.get("exclusion_reason")
                        is_failed_seg = segment.get("is_failed", False)
                        failure_reason_seg = segment.get("failure_reason")
                        translation_segments_map[segment_index] = {
                            "is_excluded": is_excluded_seg,
                            "is_image": segment.get("is_image", False),
                            "exclusion_reason": exclusion_reason_seg,  # CRITICAL: Get exclusion_reason for frontend display
                            "target_text": segment.get("target_text"),  # CRITICAL: Get target_text for frontend display
                            "is_failed": is_failed_seg,  # CRITICAL: Get is_failed for failure detection
                            "failure_reason": failure_reason_seg,  # CRITICAL: Get failure_reason for debugging
                            "status": segment.get("status"),  # CRITICAL: Get status (translated, failed, etc.)
                            "modified": segment.get("modified"),  # CRITICAL: Get modified for user-translated detection
                        }
                        # Debug: Track excluded segments without exclusion_reason
                        if is_excluded_seg and not exclusion_reason_seg and segment_index < 30:
                            excluded_segments_in_translation.append(segment_index)
                        # Debug: Track ALL failed segments for logging (index + failure_reason)
                        if is_failed_seg:
                            failed_segments_in_translation.append((segment_index, failure_reason_seg))
            if excluded_segments_in_translation:
                logger.warning(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Found {len(excluded_segments_in_translation)} excluded segments "
                    f"in translation_segments without exclusion_reason: {excluded_segments_in_translation}",
                )
            # Count total failed segments and log all (or first 20) for troubleshooting UI vs API mismatch
            total_failed_count = sum(1 for seg in segments_list if isinstance(seg, dict) and seg.get("is_failed", False))
            if total_failed_count > 0:
                logger.info(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Total failed segments in translation_segments: {total_failed_count}",
                )
                # Log failed segment indices and reasons so user can compare with API debug file
                to_log = failed_segments_in_translation[:20] if len(failed_segments_in_translation) > 20 else failed_segments_in_translation
                logger.info(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Failed segments (index, reason): {to_log}"
                    + (f" ... and {len(failed_segments_in_translation) - 20} more" if len(failed_segments_in_translation) > 20 else ""),
                )
        
        # Get pre-computed excluded segment indices from segments_metadata (if available)
        # This allows exclusion to be determined during extraction phase
        segments_metadata_dict = st.get("segments_metadata", {})
        precomputed_excluded_indices = set(segments_metadata_dict.get("excluded_segment_indices", []))
        
        # CRITICAL: Get exclusion reasons from segments_metadata.excluded_segments
        # This is needed for frontend to display exclusion types (identifier, language_match, etc.)
        excluded_segments_dict = segments_metadata_dict.get("excluded_segments", {})
        excluded_reasons_map = {}  # {segment_index: exclusion_reason_string}
        
        # CRITICAL: Also get all detected exclusion reasons (including non-excluded ones)
        # This allows frontend to display identifier, language_match, etc. even if not excluded
        detected_exclusion_reasons_dict = segments_metadata_dict.get("detected_exclusion_reasons", {})
        detected_reasons_map = {}  # {segment_index: exclusion_reason_string}
        
        if detected_exclusion_reasons_dict and isinstance(detected_exclusion_reasons_dict, dict):
            for seg_idx_str, detection_info in detected_exclusion_reasons_dict.items():
                try:
                    seg_idx = int(seg_idx_str)
                    if isinstance(detection_info, dict):
                        reason = detection_info.get("reason", "unknown")
                        # CRITICAL: For language_match, check if target_lang matches
                        # If target_lang is provided and doesn't match stored target_lang, skip this detection
                        # This prevents using incorrect language_match detections from previous target_lang
                        if reason == "language_match" and target_lang:
                            stored_metadata = detection_info.get("metadata", {})
                            stored_target_lang = stored_metadata.get("target_lang")
                            if stored_target_lang and stored_target_lang != target_lang:
                                logger.debug(
                                    LogModule.EXTRACT,
                                    f"[PREVIEW-API] Task {task_id}: Skipping language_match detection for segment {seg_idx} "
                                    f"(stored target_lang={stored_target_lang} != current target_lang={target_lang})"
                                    )
                                continue  # Skip this language_match detection, will be re-detected below
                    elif isinstance(detection_info, str):
                        reason = detection_info
                        # For string format, we can't check target_lang, but if current target_lang is set,
                        # we should re-detect language_match to ensure correctness
                        if reason == "language_match" and target_lang:
                            # Don't skip, but we'll re-detect below if not in excluded_reasons_map
                            # NOTE: Removed verbose debug log for each segment to reduce log noise
                            pass
                    else:
                        reason = "unknown"
                    detected_reasons_map[seg_idx] = reason
                except (ValueError, TypeError):
                    continue
            
            # Log detected reasons for debugging
            if detected_reasons_map:
                reason_counts = {}
                for reason in detected_reasons_map.values():
                    reason_counts[reason] = reason_counts.get(reason, 0) + 1
                reason_summary = ', '.join(f'{count} {reason}' for reason, count in sorted(reason_counts.items()))
                logger.info(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Found {len(detected_reasons_map)} detected exclusion reasons "
                    f"(including non-excluded): {reason_summary}"
                    )
                # Log identifier and language_match counts separately
                identifier_count = reason_counts.get("identifier", 0)
                language_match_count = reason_counts.get("language_match", 0)
                if identifier_count > 0 or language_match_count > 0:
                    logger.info(
                        LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Detected reasons breakdown - Identifier: {identifier_count}, "
                        f"Language Match: {language_match_count}"
                        )
        
        # Removed segments_metadata_dict structure logging to reduce verbosity
        
        if excluded_segments_dict and isinstance(excluded_segments_dict, dict):
            # Removed frequent logging to reduce log noise
            # logger.info(...)  # Removed: too frequent
            for seg_idx_str, exclusion_info in excluded_segments_dict.items():
                try:
                    seg_idx = int(seg_idx_str)
                    if isinstance(exclusion_info, dict):
                        reason = exclusion_info.get("reason", "unknown")
                        # NOTE: Removed verbose debug log for each excluded segment to reduce log noise
                    elif isinstance(exclusion_info, str):
                        reason = exclusion_info
                    else:
                        reason = "unknown"
                        logger.warning(
                            LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Unexpected exclusion_info type for segment {seg_idx}: "
                            f"{type(exclusion_info)}, value: {exclusion_info}"
                            )
                    excluded_reasons_map[seg_idx] = reason
                except (ValueError, TypeError) as e:
                    logger.warning(
                        LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Failed to parse excluded_segments entry: "
                        f"seg_idx_str={seg_idx_str}, exclusion_info={exclusion_info}, exclusion_info type={type(exclusion_info)}, error: {e}"
                        )
                    continue
        else:
            logger.debug(
                LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: excluded_segments_dict is empty or not a dict. "
                f"segments_metadata_dict keys: {list(segments_metadata_dict.keys())}"
                )
        
        # Debug: Log excluded indices for troubleshooting (only if present)
        if precomputed_excluded_indices:
            logger.trace(
                LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Found {len(precomputed_excluded_indices)} pre-computed excluded segment indices"
                )
        if excluded_reasons_map:
            reason_counts = {}
            for reason in excluded_reasons_map.values():
                reason_counts[reason] = reason_counts.get(reason, 0) + 1
            # Removed frequent exclusion breakdown logging to reduce log verbosity
            # Only log summary when there are issues or significant changes
            # logger.info(...)  # Removed: too frequent
        else:
            # CRITICAL: Log warning if excluded_reasons_map is empty but we have excluded segments
            if precomputed_excluded_indices:
                logger.warning(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: WARNING - excluded_reasons_map is empty but "
                    f"precomputed_excluded_indices has {len(precomputed_excluded_indices)} segments. "
                    f"This indicates exclusion reasons were not properly stored in segments_metadata.excluded_segments. "
                    f"excluded_segments_dict type: {type(excluded_segments_dict)}, "
                    f"excluded_segments_dict is None: {excluded_segments_dict is None}, "
                    f"excluded_segments_dict is dict: {isinstance(excluded_segments_dict, dict)}, "
                    f"excluded_segments_dict keys: {list(excluded_segments_dict.keys())[:10] if excluded_segments_dict else 'N/A'}"
                    )
        if translation_segments_map:
            excluded_from_translation = sum(1 for v in translation_segments_map.values() if v.get("is_excluded", False))
            logger.trace(
                LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Found {excluded_from_translation} excluded segments from translation_segments"
                )
        
        # CRITICAL: Determine if we're in Translate phase or Extract phase (before processing segments)
        # Translate phase: translation_segments_map exists (translation has started)
        # Extract phase: translation_segments_map does not exist (only extraction completed)
        is_translate_phase = bool(translation_segments_map)
        
        # Get workflow_type early (needed for MOBI image segment detection)
        segments_metadata = st.get("segments_metadata", {})
        workflow_type = segments_metadata.get("workflow_type") or (st.get("payload", {}).get("workflow_type") if isinstance(st.get("payload"), dict) else (getattr(st.get("payload"), 'workflow_type', None) if st.get("payload") else None))
        
        # CRITICAL: For MOBI files, check HTML templates for img tags and create image placeholder segments
        # This ensures images are displayed in Extract preview and Exclusion statistics
        # Try to position images based on their location in HTML relative to text segments
        mobi_image_segments = []  # List of (insert_index, placeholder_id, image_path, context_text) tuples
        if workflow_type == "mobi":
            html_templates = st.get("mobi_html_templates", {})
            image_data_map_raw = st.get("image_data_map", {})
            if html_templates and image_data_map_raw and all_segments:
                from bs4 import BeautifulSoup
                import re
                # Check each HTML template for img tags
                for item_id, html_content in html_templates.items():
                    try:
                        soup = BeautifulSoup(html_content, 'html.parser')
                        img_tags = soup.find_all('img')
                        
                        # Extract text nodes in order (same as MobiExtractor does)
                        # This helps us map images to segments
                        text_nodes_in_order = []
                        for text_node in soup.find_all(string=True):
                            if (
                                text_node.parent.name not in ['style', 'script', 'head', 'title', 'meta', '[document]']
                                and not text_node.isspace()
                            ):
                                text = text_node.get_text(strip=True)
                                if text:
                                    text_nodes_in_order.append({
                                        'text': text,
                                        'node': text_node,
                                    })
                        
                        for img in img_tags:
                            src = img.get('src', '')
                            if not src:
                                continue
                            # Try to match src with image_data_map keys
                            # Handle encoding issues (e.g., "一mages" vs "images")
                            matched_image_path = None
                            for image_path in image_data_map_raw.keys():
                                # Try exact match first
                                if src == image_path:
                                    matched_image_path = image_path
                                    break
                                # Try filename match
                                import os
                                src_filename = os.path.basename(src)
                                image_filename = os.path.basename(image_path)
                                if src_filename == image_filename:
                                    matched_image_path = image_path
                                    break
                                # Try path contains match (handle encoding issues)
                                if src in image_path or image_path in src:
                                    matched_image_path = image_path
                                    break
                                # Try fixing common encoding issues
                                src_fixed = src.replace('一', 'i').replace('mages', 'images')
                                if src_fixed in image_path or image_path in src_fixed:
                                    matched_image_path = image_path
                                    break
                            
                            if matched_image_path:
                                # Create placeholder ID from image path
                                placeholder_id = matched_image_path
                                
                                # Try to find the position of this img tag relative to text segments
                                # Strategy: Find text nodes before/after the img tag at BODY level, then match to segments
                                insert_index = None
                                context_before = None
                                context_after = None
                                
                                # Get all elements in document order to find img position
                                try:
                                    # Find text nodes before and after this img tag
                                    before_text_nodes = []
                                    after_text_nodes = []
                                    
                                    # Strategy 1: Check siblings within parent (original approach)
                                    parent = img.parent
                                    if parent:
                                        siblings = list(parent.children)
                                        img_sibling_index = None
                                        for i, sibling in enumerate(siblings):
                                            if sibling == img:
                                                img_sibling_index = i
                                                break
                                        
                                        if img_sibling_index is not None:
                                            # Get text before img (from previous siblings)
                                            for i in range(img_sibling_index - 1, -1, -1):
                                                sibling = siblings[i]
                                                if isinstance(sibling, str) and sibling.strip():
                                                    before_text_nodes.insert(0, sibling.strip())
                                                elif hasattr(sibling, 'get_text'):
                                                    text = sibling.get_text(strip=True)
                                                    if text:
                                                        before_text_nodes.insert(0, text)
                                            
                                            # Get text after img (from next siblings)
                                            for i in range(img_sibling_index + 1, len(siblings)):
                                                sibling = siblings[i]
                                                if isinstance(sibling, str) and sibling.strip():
                                                    after_text_nodes.append(sibling.strip())
                                                elif hasattr(sibling, 'get_text'):
                                                    text = sibling.get_text(strip=True)
                                                    if text:
                                                        after_text_nodes.append(text)
                                    
                                    # Strategy 2: If no text found in siblings, check body-level adjacent elements
                                    # This handles cases where image is in its own <p> tag
                                    if not before_text_nodes and not after_text_nodes:
                                        body = soup.find('body')
                                        if body and parent:
                                            # Find all direct children of body
                                            body_children = list(body.children)
                                            parent_index = None
                                            for i, child in enumerate(body_children):
                                                if child == parent:
                                                    parent_index = i
                                                    break
                                            
                                            if parent_index is not None:
                                                # Find text in previous sibling elements (at body level)
                                                for i in range(parent_index - 1, -1, -1):
                                                    prev_elem = body_children[i]
                                                    if hasattr(prev_elem, 'get_text'):
                                                        text = prev_elem.get_text(strip=True)
                                                        if text:
                                                            before_text_nodes.insert(0, text)
                                                            # Only take the last (closest) text node
                                                            if len(before_text_nodes) > 0:
                                                                before_text_nodes = [before_text_nodes[-1]]
                                                                break
                                                
                                                # Find text in next sibling elements (at body level)
                                                for i in range(parent_index + 1, len(body_children)):
                                                    next_elem = body_children[i]
                                                    if hasattr(next_elem, 'get_text'):
                                                        text = next_elem.get_text(strip=True)
                                                        if text:
                                                            after_text_nodes.append(text)
                                                            # Only take the first (closest) text node
                                                            if len(after_text_nodes) > 0:
                                                                after_text_nodes = [after_text_nodes[0]]
                                                                break
                                    
                                    # Try to match before_text_nodes to segments to find insertion point
                                    # Strategy: Use the most specific text match (longer text is more unique)
                                    if before_text_nodes:
                                        # Try to match the longest text first (more unique)
                                        sorted_before_texts = sorted(before_text_nodes, key=len, reverse=True)
                                        for before_text in sorted_before_texts:
                                            if not before_text:
                                                continue
                                            context_before = before_text[:100]  # Store context for logging
                                            # Find segment that contains this text (exact match or substring)
                                            for seg_idx, seg_text in enumerate(all_segments):
                                                if isinstance(seg_text, str):
                                                    # Try exact match first (most accurate)
                                                    if before_text == seg_text.strip():
                                                        insert_index = seg_idx + 1
                                                        break
                                                    # Try substring match (segment contains the text)
                                                    elif before_text in seg_text:
                                                        insert_index = seg_idx + 1
                                                        break
                                                    # Try reverse match (text contains segment - segment might be split)
                                                    elif len(before_text) > 20 and seg_text.strip() in before_text:
                                                        insert_index = seg_idx + 1
                                                        break
                                                if insert_index is not None:
                                                    break
                                            if insert_index is not None:
                                                break
                                    
                                    # If not found, try matching after_text_nodes
                                    if insert_index is None and after_text_nodes:
                                        # Try to match the shortest text first (more specific)
                                        sorted_after_texts = sorted(after_text_nodes, key=len)
                                        for after_text in sorted_after_texts:
                                            if not after_text:
                                                continue
                                            context_after = after_text[:100]  # Store context for logging
                                            # Find segment that contains this text
                                            for seg_idx, seg_text in enumerate(all_segments):
                                                if isinstance(seg_text, str):
                                                    # Try exact match first
                                                    if after_text == seg_text.strip():
                                                        insert_index = seg_idx
                                                        break
                                                    # Try substring match
                                                    elif after_text in seg_text:
                                                        insert_index = seg_idx
                                                        break
                                                    # Try reverse match
                                                    elif len(after_text) > 20 and seg_text.strip() in after_text:
                                                        insert_index = seg_idx
                                                        break
                                                if insert_index is not None:
                                                    break
                                            if insert_index is not None:
                                                break
                                    
                                    # Fallback: if we can't determine position, append to end
                                    if insert_index is None:
                                        insert_index = len(all_segments) + len(mobi_image_segments)
                                        
                                        # Write debug file with DOM/HTML information for analysis
                                        try:
                                            import os
                                            import json
                                            from pathlib import Path
                                            
                                            # Get debug directory (use temp_dir if available, otherwise logs_dir)
                                            debug_dir = None
                                            temp_dir = st.get("temp_dir")
                                            if temp_dir and os.path.isdir(temp_dir):
                                                debug_dir = os.path.join(temp_dir, "debug", "image_position")
                                                os.makedirs(debug_dir, exist_ok=True)
                                            else:
                                                # Fallback to logs directory
                                                try:
                                                    from utils.path_utils import get_logs_dir
                                                    logs_dir = get_logs_dir()
                                                    debug_dir = logs_dir / "debug" / "image_position"
                                                    debug_dir.mkdir(parents=True, exist_ok=True)
                                                    debug_dir = str(debug_dir)
                                                except Exception:
                                                    debug_dir = None
                                            
                                            if debug_dir:
                                                debug_file = os.path.join(debug_dir, f"{task_id}_image_{len(mobi_image_segments)}_debug.txt")
                                                
                                                with open(debug_file, 'w', encoding='utf-8') as f:
                                                    f.write("=" * 80 + "\n")
                                                    f.write(f"MOBI Image Position Detection Debug Info\n")
                                                    f.write(f"Task ID: {task_id}\n")
                                                    f.write(f"Image Index: {len(mobi_image_segments)}\n")
                                                    f.write(f"Image Src: {src}\n")
                                                    f.write(f"Matched Path: {matched_image_path}\n")
                                                    f.write(f"Item ID: {item_id}\n")
                                                    f.write("=" * 80 + "\n\n")
                                                    
                                                    # 1. Image tag info
                                                    f.write("1. IMAGE TAG INFORMATION\n")
                                                    f.write("-" * 80 + "\n")
                                                    f.write(f"Tag: {img}\n")
                                                    f.write(f"Parent: {img.parent.name if img.parent else 'None'}\n")
                                                    f.write(f"Parent Tag: {str(img.parent)[:200] if img.parent else 'None'}\n")
                                                    f.write("\n")
                                                    
                                                    # 2. Siblings info
                                                    f.write("2. SIBLINGS INFORMATION\n")
                                                    f.write("-" * 80 + "\n")
                                                    if parent:
                                                        siblings = list(parent.children)
                                                        f.write(f"Total siblings: {len(siblings)}\n")
                                                        f.write(f"Image sibling index: {img_sibling_index}\n")
                                                        f.write("\nSiblings (with index):\n")
                                                        for i, sibling in enumerate(siblings):
                                                            marker = " <-- IMAGE" if sibling == img else ""
                                                            if isinstance(sibling, str):
                                                                f.write(f"  [{i}] STRING: {repr(sibling[:100])}{marker}\n")
                                                            else:
                                                                sibling_str = str(sibling)[:200] if hasattr(sibling, '__str__') else str(type(sibling))
                                                                f.write(f"  [{i}] {type(sibling).__name__}: {sibling_str}{marker}\n")
                                                    else:
                                                        f.write("Parent is None, cannot get siblings\n")
                                                    f.write("\n")
                                                    
                                                    # 3. Text nodes before/after image
                                                    f.write("3. TEXT NODES BEFORE/AFTER IMAGE\n")
                                                    f.write("-" * 80 + "\n")
                                                    f.write(f"Before text nodes count: {len(before_text_nodes)}\n")
                                                    if before_text_nodes:
                                                        for i, text in enumerate(before_text_nodes):
                                                            f.write(f"  [{i}] {repr(text[:200])}\n")
                                                    else:
                                                        f.write("  (None found)\n")
                                                    f.write(f"\nAfter text nodes count: {len(after_text_nodes)}\n")
                                                    if after_text_nodes:
                                                        for i, text in enumerate(after_text_nodes):
                                                            f.write(f"  [{i}] {repr(text[:200])}\n")
                                                    else:
                                                        f.write("  (None found)\n")
                                                    f.write("\n")
                                                    
                                                    # 4. All text nodes in order (from soup)
                                                    f.write("4. ALL TEXT NODES IN HTML (in order)\n")
                                                    f.write("-" * 80 + "\n")
                                                    f.write(f"Total text nodes found: {len(text_nodes_in_order)}\n")
                                                    f.write("First 20 text nodes:\n")
                                                    for i, node_info in enumerate(text_nodes_in_order[:20]):
                                                        text = node_info['text']
                                                        f.write(f"  [{i}] {repr(text[:150])}\n")
                                                    if len(text_nodes_in_order) > 20:
                                                        f.write(f"  ... (and {len(text_nodes_in_order) - 20} more)\n")
                                                    f.write("\n")
                                                    
                                                    # 5. Segments info (for comparison)
                                                    f.write("5. SEGMENTS INFORMATION (for comparison)\n")
                                                    f.write("-" * 80 + "\n")
                                                    f.write(f"Total segments: {len(all_segments)}\n")
                                                    f.write("First 30 segments:\n")
                                                    for i, seg in enumerate(all_segments[:30]):
                                                        if isinstance(seg, str):
                                                            f.write(f"  [{i}] {repr(seg[:150])}\n")
                                                        else:
                                                            f.write(f"  [{i}] {type(seg).__name__}: {str(seg)[:150]}\n")
                                                    if len(all_segments) > 30:
                                                        f.write(f"  ... (and {len(all_segments) - 30} more)\n")
                                                    f.write("\n")
                                                    
                                                    # 6. HTML structure around image
                                                    f.write("6. HTML STRUCTURE AROUND IMAGE\n")
                                                    f.write("-" * 80 + "\n")
                                                    if img.parent:
                                                        # Get parent's parent for more context
                                                        grandparent = img.parent.parent if img.parent.parent else None
                                                        if grandparent:
                                                            f.write(f"Grandparent: {grandparent.name}\n")
                                                            f.write(f"Grandparent HTML (first 500 chars):\n{str(grandparent)[:500]}\n\n")
                                                        f.write(f"Parent: {img.parent.name}\n")
                                                        f.write(f"Parent HTML (first 1000 chars):\n{str(img.parent)[:1000]}\n\n")
                                                    f.write(f"Image tag HTML:\n{str(img)}\n\n")
                                                    
                                                    # 7. Full HTML template (truncated)
                                                    f.write("7. FULL HTML TEMPLATE (first 2000 chars)\n")
                                                    f.write("-" * 80 + "\n")
                                                    f.write(f"{html_content[:2000]}\n")
                                                    if len(html_content) > 2000:
                                                        f.write(f"\n... (truncated, total length: {len(html_content)} chars)\n")
                                                    
                                                logger.info(
                                                    LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Wrote image position debug file: {debug_file}"
                                                    )
                                        except Exception as debug_error:
                                            logger.warning(
                                                LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Failed to write debug file: {debug_error}",
                                                exc_info=True,
                                                )
                                        
                                        logger.debug(
                                            LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Could not determine image position, "
                                            f"appending to end. src={src}, matched_path={matched_image_path}, "
                                            f"before_text_nodes_count={len(before_text_nodes)}, "
                                            f"after_text_nodes_count={len(after_text_nodes)}, "
                                            f"before_texts_sample={[t[:50] for t in before_text_nodes[:2]] if before_text_nodes else []}, "
                                            f"after_texts_sample={[t[:50] for t in after_text_nodes[:2]] if after_text_nodes else []}"
                                            )
                                    else:
                                        logger.debug(
                                            LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Determined image position: "
                                            f"insert_index={insert_index}, src={src}, "
                                            f"context_before={context_before}, context_after={context_after}"
                                            )
                                except Exception as pos_error:
                                    # If position detection fails, append to end
                                    insert_index = len(all_segments) + len(mobi_image_segments)
                                    logger.debug(
                                        LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Position detection failed, appending to end: {pos_error}"
                                        )
                                
                                mobi_image_segments.append((
                                    insert_index,
                                    placeholder_id,
                                    matched_image_path,
                                    {'context_before': context_before, 'context_after': context_after}
                                ))
                                logger.debug(
                                    LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Found MOBI image in HTML template - "
                                    f"item_id={item_id}, src={src}, matched_path={matched_image_path}, "
                                    f"placeholder_id={placeholder_id}, insert_index={insert_index}"
                                    )
                    except Exception as e:
                        logger.warning(
                            LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Failed to parse HTML template for item_id={item_id}: {e}",
                            exc_info=True,
                            )
                        continue
                
                if mobi_image_segments:
                    logger.info(
                        LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Found {len(mobi_image_segments)} MOBI image segments to add to preview. "
                        f"Sample image paths: {[path for _, _, path, _ in mobi_image_segments[:3]]}"
                        )
                    # CRITICAL: Save image segments info to task_state for use in _record_mobi_segments
                    # This ensures image segments are included in translation_segments
                    # IMPORTANT: Always reset to avoid duplicate accumulation when get_source_preview is called multiple times
                    st["mobi_image_segments_info"] = []
                    # Store image segment info: (insert_index, placeholder_id, image_path, placeholder_text)
                    for insert_idx, placeholder_id, image_path, context_info in mobi_image_segments:
                        image_info = image_data_map_raw.get(image_path, {})
                        data_uri = image_info.get("data", "")
                        placeholder_text = f"<ph-{placeholder_id}>"
                        st["mobi_image_segments_info"].append({
                            "insert_index": insert_idx,
                            "placeholder_id": placeholder_id,
                            "image_path": image_path,
                            "placeholder_text": placeholder_text,
                            "image_data": data_uri,
                        })
                    logger.debug(
                        LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Saved {len(mobi_image_segments)} image segments info to task_state for translation_segments"
                        )
                else:
                    logger.debug(
                        LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: No MOBI image segments found. "
                        f"html_templates count: {len(html_templates) if html_templates else 0}, "
                        f"image_data_map count: {len(image_data_map_raw) if image_data_map_raw else 0}"
                        )
        
        # Convert segments to objects with metadata (unified format for all workflows)
        # This ensures consistent handling of is_excluded, is_image, etc. across PDF, DOCX, XLSX, etc.
        segments_with_metadata = []
        for idx, seg in enumerate(all_segments):
            # Handle both string and dict formats
            if isinstance(seg, dict):
                # Already has metadata, use it
                seg_text = seg.get("text", str(seg.get("source_text", "")))
                seg_obj = seg.copy()
            else:
                # Plain string, convert to object with metadata
                seg_text = str(seg)
                seg_obj = {"text": seg_text}
            
            # CRITICAL: Always include segment_index in the response for proper mapping
            # This ensures frontend can correctly match source segments with translation segments
            # even when pagination is used
            if "segment_index" not in seg_obj:
                seg_obj["segment_index"] = idx
            
            if is_translate_phase:
                # ============================================================
                # TRANSLATE PHASE: COMPLETELY READ-ONLY, NO DETECTION
                # ============================================================
                # CRITICAL: Translate phase should NEVER re-detect or re-validate exclusions
                # It should ONLY use exclusion data from Extract phase (segments_metadata.excluded_segments)
                # This ensures "what you see is what you get" - Extract phase results are preserved
                # User can manually modify exclusions via API, but system should NOT auto-detect
                
                # Priority 1: Use translation_segments if available (most up-to-date from Translate phase)
                if idx in translation_segments_map:
                    translation_info = translation_segments_map[idx]
                    is_excluded = translation_info.get("is_excluded", False)
                    is_image = translation_info.get("is_image", _is_image_segment(seg_text))
                    # Get exclusion_reason from translation_segments if available
                    exclusion_reason = translation_info.get("exclusion_reason")
                    # CRITICAL: If segment is excluded but exclusion_reason is missing, try to get from excluded_reasons_map
                    # This handles the case where exclusion_reason was not properly set during record_translation_segments
                    # But we do NOT re-detect - we only use stored data from Extract phase
                    if is_excluded and not exclusion_reason and idx in excluded_reasons_map:
                        exclusion_reason = excluded_reasons_map.get(idx)
                        # NOTE: Removed verbose debug log for each segment to reduce log noise
                    # Debug: Log exclusion_reason from translation_segments for excluded segments
                    # PERFORMANCE: Reduced logging to only first 5 segments to avoid I/O overhead
                    if is_excluded and idx < 5:
                        logger.trace(
                            LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Translate phase - Segment {idx} from translation_segments - "
                            f"is_excluded={is_excluded}, exclusion_reason={exclusion_reason}, "
                            f"translation_info keys: {list(translation_info.keys())}"
                            )
                    
                    # CRITICAL: Get target_text from translation_segments for frontend display
                    # This ensures frontend shows translated text instead of source text
                    # Priority: modified_text > target_text (same as export logic)
                    modified_text = translation_info.get("modified_text")
                    target_text = translation_info.get("target_text")
                    # Use modified_text if available, otherwise use target_text
                    final_target_text = modified_text if modified_text is not None else target_text
                    if final_target_text is not None:
                        seg_obj["target_text"] = final_target_text
                    # Also include modified_text field if it exists (for frontend to use)
                    if modified_text is not None:
                        seg_obj["modified_text"] = modified_text
                    # Also set is_failed and failure_reason if available
                    # CRITICAL: Always set is_failed explicitly (True or False) so frontend can detect it
                    is_failed = translation_info.get("is_failed", False)
                    seg_obj["is_failed"] = is_failed
                    if is_failed:
                        failure_reason = translation_info.get("failure_reason")
                        if failure_reason:
                            seg_obj["failure_reason"] = failure_reason
                        # Debug: Log failed segments for troubleshooting
                        logger.debug(
                            LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Translate phase - Segment {idx} marked as FAILED: "
                            f"failure_reason={failure_reason}"
                            )
                    # Set status if available
                    status = translation_info.get("status")
                    if status:
                        seg_obj["status"] = status
                    # CRITICAL: Copy modified flag so is_user_translated check works
                    modified = translation_info.get("modified")
                    if modified is not None:
                        seg_obj["modified"] = modified
                # Priority 2: Use excluded_reasons_map from Extract phase (if not in translation_segments)
                elif idx in excluded_reasons_map:
                    # Segment was excluded during Extract phase - use stored reason
                    # CRITICAL: Do NOT re-detect, just use stored data
                    is_image = _is_image_segment(seg_text)
                    is_excluded = True
                    exclusion_reason = excluded_reasons_map.get(idx)
                    # CRITICAL: Excluded segments should not be marked as failed
                    seg_obj["is_failed"] = False
                    # PERFORMANCE: Reduced logging to only first 5 segments to avoid I/O overhead
                    if idx < 5:
                        logger.trace(
                            LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Translate phase - Segment {idx} excluded from Extract phase "
                            f"(exclusion_reason={exclusion_reason}, NOT re-detected)"
                            )
                # Priority 3: Use detected_reasons_map for non-excluded but detected segments (for frontend display)
                elif idx in detected_reasons_map:
                    # Segment has detected reason but is not excluded (for frontend display)
                    # CRITICAL: Do NOT re-validate, just use stored detection from Extract phase
                    detected_reason = detected_reasons_map.get(idx)
                    is_image = _is_image_segment(seg_text)
                    is_excluded = False  # Not excluded, but has detected reason
                    exclusion_reason = detected_reason
                    # CRITICAL: If not in translation_segments, assume not failed (translation may not have started yet)
                    seg_obj["is_failed"] = False
                    # PERFORMANCE: Reduced logging to only first 5 segments to avoid I/O overhead
                    if idx < 5:
                        logger.trace(
                            LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Translate phase - Segment {idx} has detected_reason={detected_reason} "
                            f"from Extract phase (NOT re-validated, NOT re-detected)"
                            )
                else:
                    # Segment is not excluded and has no detected reason
                    # CRITICAL: Do NOT re-detect, just mark as not excluded
                    is_image = _is_image_segment(seg_text)
                    is_excluded = False
                    exclusion_reason = None
                    # CRITICAL: If not in translation_segments, assume not failed (translation may not have started yet)
                    seg_obj["is_failed"] = False
                    # PERFORMANCE: Reduced logging - removed verbose logging for non-excluded segments to avoid I/O overhead
            else:
                # ============================================================
                # EXTRACT PHASE: DETECTION ALLOWED
                # ============================================================
                # Extract phase can perform real-time detection as fallback
                # This is the ONLY phase where exclusion detection should happen
                
                # Priority 1: Check if this segment was pre-marked as excluded during extraction
                # OR if it's in excluded_reasons_map (more reliable than precomputed_excluded_indices)
                if idx in precomputed_excluded_indices or idx in excluded_reasons_map:
                    # Use pre-computed exclusion status
                    is_image = _is_image_segment(seg_text)
                    is_excluded = True
                    # Get exclusion_reason from excluded_segments if available
                    # CRITICAL: excluded_reasons_map is the source of truth for exclusion reasons
                    exclusion_reason = excluded_reasons_map.get(idx)
                    # If not in excluded_reasons_map but in precomputed_excluded_indices, 
                    # it might be a legacy exclusion without a reason
                    if not exclusion_reason and idx in precomputed_excluded_indices:
                        logger.debug(
                            LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Extract phase - Segment {idx} is in precomputed_excluded_indices "
                            f"but not in excluded_reasons_map. This may indicate a legacy exclusion without a reason."
                            )
                # Priority 2: Check detected_reasons_map for non-excluded but detected segments
                # This allows frontend to display identifier, language_match, etc. even if not excluded
                # CRITICAL: For language_match, re-validate with current target_lang to ensure correctness
                elif idx in detected_reasons_map:
                    detected_reason = detected_reasons_map.get(idx)
                    # For language_match, always re-detect with current target_lang to ensure correctness
                    # This prevents using incorrect language_match detections from previous target_lang
                    if detected_reason == "language_match" and target_lang:
                        # Re-detect with current target_lang
                        # NOTE: Removed verbose debug log for each segment to reduce log noise
                        from exclusion.core.exclusion_detector import detect_exclusion_reason
                        re_detected_result = detect_exclusion_reason(
                            text=seg_text,
                            block_type=None,
                            target_lang=target_lang,
                            is_image=_is_image_segment(seg_text),
                            is_table=False
                        )
                        if re_detected_result:
                            re_detected_reason, re_detected_metadata = re_detected_result
                            if re_detected_reason.value == "language_match":
                                # Still language_match with current target_lang, use it
                                is_image = _is_image_segment(seg_text)
                                is_excluded = False  # Not excluded, but has detected reason
                                exclusion_reason = re_detected_reason.value
                            else:
                                # No longer language_match with current target_lang, clear detection
                                logger.debug(
                                    LogModule.EXCLUSION,
                    f"[PREVIEW-API] Task {task_id}: Extract phase - Segment {idx} is no longer language_match "
                                    f"with target_lang={target_lang}"
                                    )
                                is_image = _is_image_segment(seg_text)
                                is_excluded = False
                                exclusion_reason = None
                                # Remove from detected_reasons_map to prevent future use
                                detected_reasons_map.pop(idx, None)
                        else:
                            # No longer language_match with current target_lang, skip and clear detection
                            logger.debug(
                                LogModule.EXCLUSION,
                    f"[PREVIEW-API] Task {task_id}: Extract phase - Segment {idx} is no longer language_match "
                                f"with target_lang={target_lang}"
                                )
                            is_image = _is_image_segment(seg_text)
                            is_excluded = False
                            exclusion_reason = None
                            # Remove from detected_reasons_map to prevent future use
                            detected_reasons_map.pop(idx, None)
                    else:
                        # Not language_match or target_lang not set, use stored detection
                        is_image = _is_image_segment(seg_text)
                        is_excluded = False  # Not excluded, but has detected reason
                        exclusion_reason = detected_reason
                else:
                    # Priority 4: Real-time detection (ONLY for Extract phase)
                    # CRITICAL: Use detect_exclusion_reason instead of should_exclude_text
                    # This ensures all excluded segments have a proper exclusion_reason
                    is_image = _is_image_segment(seg_text)
                    from exclusion.core.exclusion_detector import detect_exclusion_reason
                    from exclusion.core.exclusion_reason import ExclusionReason
                    # Determine if strict_table_priority should be used (PDF format only)
                    # For non-PDF formats, Identifier takes priority over Table
                    workflow_type = segments_metadata_dict.get("workflow_type") or (st.get("payload", {}).get("workflow_type") if isinstance(st.get("payload"), dict) else (getattr(st.get("payload"), 'workflow_type', None) if st.get("payload") else None))
                    strict_table_priority = (workflow_type == "pdf" or workflow_type == "markdown_based")
                    detected_result = detect_exclusion_reason(
                        text=seg_text,
                        block_type=None,
                        target_lang=target_lang,
                        is_image=is_image,
                        is_table=False,
                        strict_table_priority=strict_table_priority
                    )
                    if detected_result:
                        detected_reason, _ = detected_result
                        # CRITICAL: Only mark as excluded when reason is in exclusion_defaults (system.json)
                        default_excluded = ExclusionReason.get_default_excluded()
                        if detected_reason in default_excluded:
                            is_excluded = True
                            exclusion_reason = detected_reason.value
                            logger.debug(
                                LogModule.EXTRACT,
                                f"[PREVIEW-API] Task {task_id}: Extract phase - Segment {idx} excluded by real-time detection, "
                                f"detected exclusion_reason={exclusion_reason}"
                            )
                        else:
                            is_excluded = False
                            exclusion_reason = detected_reason.value  # Show tag in UI even when not excluded by default
                            logger.trace(
                                LogModule.EXTRACT,
                                f"[PREVIEW-API] Task {task_id}: Extract phase - Segment {idx} detected as {exclusion_reason} "
                                f"(not in exclusion_defaults, checkbox unchecked)"
                            )
                    else:
                        is_excluded = False
                        exclusion_reason = None
            
            # Update metadata
            seg_obj["is_image"] = is_image
            seg_obj["is_excluded"] = is_excluded
            
            # CRITICAL DEBUG: Log exclusion_reason value before setting seg_obj
            # PERFORMANCE: Reduced logging to only first 5 segments to avoid I/O overhead
            if is_excluded and idx < 5:
                logger.trace(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: BEFORE setting seg_obj - Segment {idx}: "
                    f"is_excluded={is_excluded}, exclusion_reason={exclusion_reason}, "
                    f"in_translation_segments_map={idx in translation_segments_map}, "
                    f"in_excluded_reasons_map={idx in excluded_reasons_map}"
                    )
            
            # CRITICAL: Include exclusion_reason for frontend to display exclusion types
            # Always set exclusion_reason if segment is excluded, even if it's None (for debugging)
            if is_excluded:
                # If exclusion_reason is missing but segment is excluded, try to get from excluded_reasons_map
                if not exclusion_reason and idx in excluded_reasons_map:
                    exclusion_reason = excluded_reasons_map.get(idx)
                    # PERFORMANCE: Reduced logging to only first 5 segments to avoid I/O overhead
                    if idx < 5:
                        logger.trace(
                            LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Segment {idx} is excluded but missing exclusion_reason, "
                            f"using exclusion_reason={exclusion_reason} from excluded_reasons_map"
                            )
                # CRITICAL: Set exclusion_reason to seg_obj BEFORE any other logic
                seg_obj["exclusion_reason"] = exclusion_reason  # Set even if None (for debugging)
                
                # CRITICAL DEBUG: Log after setting to verify
                # PERFORMANCE: Reduced logging to only first 5 segments to avoid I/O overhead
                if idx < 5:
                    logger.trace(
                        LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: AFTER setting seg_obj - Segment {idx}: "
                        f"seg_obj['exclusion_reason']={seg_obj.get('exclusion_reason')}, "
                        f"exclusion_reason variable={exclusion_reason}"
                        )
            elif exclusion_reason:
                # For non-excluded segments, only set exclusion_reason if it exists (for detected_reasons)
                seg_obj["exclusion_reason"] = exclusion_reason
            # CRITICAL: Also include detected_exclusion_reason for all detected types (including non-excluded)
            # This allows frontend to display identifier, language_match, etc. even if not excluded
            # BUT: If user manually retried and successfully translated this segment, don't overlay old reason
            if idx in detected_reasons_map:
                detected_reason = detected_reasons_map[idx]
                is_user_translated = seg_obj.get("status") == "translated" and seg_obj.get("modified") is True
                if not is_user_translated:
                    seg_obj["detected_exclusion_reason"] = detected_reason
                # Also include metadata from detected_exclusion_reasons if available
                detected_info = segments_metadata_dict.get("detected_exclusion_reasons", {}).get(str(idx), {})
                if isinstance(detected_info, dict) and "metadata" in detected_info:
                    seg_obj["exclusion_metadata"] = detected_info["metadata"]
                # CRITICAL: Also include exclusion_metadata if available (for language_match, contains detected_lang and target_lang)
                # This is needed for frontend to properly display and filter language_match segments
                if idx in excluded_reasons_map:
                    # Get exclusion_metadata from excluded_segments if available
                    excluded_segments_dict = segments_metadata_dict.get("excluded_segments", {})
                    if excluded_segments_dict and isinstance(excluded_segments_dict, dict):
                        seg_idx_str = str(idx)
                        if seg_idx_str in excluded_segments_dict:
                            exclusion_info = excluded_segments_dict[seg_idx_str]
                            if isinstance(exclusion_info, dict):
                                exclusion_metadata = exclusion_info.get("metadata", {})
                                if exclusion_metadata:
                                    seg_obj["exclusion_metadata"] = exclusion_metadata
                # Debug: Log exclusion_reason assignment for first 30 segments (especially formula, identifier, and language_match)
                if idx < 30:
                    # Always log formula, identifier, and language_match for troubleshooting
                    if exclusion_reason in ["formula", "identifier", "language_match"]:
                        # For language_match, also log detected language and target language for debugging
                        if exclusion_reason == "language_match":
                            # Try to get detected_lang and target_lang from metadata if available
                            detected_lang_info = "N/A"
                            target_lang_info = target_lang or "N/A"
                            # Check if we can get detected_lang from segment metadata
                            exclusion_metadata = seg_obj.get("exclusion_metadata", {})
                            if isinstance(exclusion_metadata, dict):
                                detected_lang_info = exclusion_metadata.get("detected_lang", "N/A")
                            # NOTE: Removed verbose info log for each segment to reduce log noise
                        else:
                            # NOTE: Removed verbose info log for each segment to reduce log noise
                            pass
                    else:
                        # NOTE: Removed verbose debug log for each segment to reduce log noise
                        pass
            # Check if segment is excluded but has no exclusion_reason (only if not already handled above)
            # This check should only trigger if the segment was excluded but exclusion_reason was not set in the if/elif blocks above
            if is_excluded and not seg_obj.get('exclusion_reason'):
                # Log WARNING: data inconsistency that we repair below (not a request failure)
                logger.warning(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Segment {idx} excluded but missing exclusion_reason, repairing. "
                    f"in_precomputed={idx in precomputed_excluded_indices}, in_excluded_reasons_map={idx in excluded_reasons_map}."
                )
                # Try to fix the inconsistency by setting exclusion_reason from variable or map
                if exclusion_reason:
                    seg_obj["exclusion_reason"] = exclusion_reason
                elif idx in excluded_reasons_map:
                    seg_obj["exclusion_reason"] = excluded_reasons_map.get(idx)
                else:
                    # Last resort: set as unknown
                    seg_obj["exclusion_reason"] = "unknown"
            
            # Attach has_latex flag so frontend knows whether to show PDF-compat check button
            try:
                from utils.latex_repair_payload import has_latex_content
                _preview_text = seg_obj.get("target_text") or seg_obj.get("text") or seg_obj.get("source_text", "")
                seg_obj["has_latex"] = has_latex_content(_preview_text)
            except Exception:
                seg_obj["has_latex"] = False

            # Attach cached PDF compat result if available
            _pdf_compat_results = st.get("pdf_compat_results")
            if isinstance(_pdf_compat_results, dict):
                _seg_compat = _pdf_compat_results.get(str(idx)) or _pdf_compat_results.get(idx)
                if _seg_compat is not None:
                    seg_obj["pdf_compat"] = _seg_compat

            segments_with_metadata.append(seg_obj)

        # CRITICAL: For MOBI files, insert image placeholder segments
        # Insert them at their calculated positions (based on HTML structure) or at the end if position cannot be determined
        if workflow_type == "mobi" and mobi_image_segments:
            image_data_map_raw = st.get("image_data_map", {})
            
            # Sort image segments by insert_index to ensure correct insertion order
            # Images with same insert_index will be inserted in the order they appear in HTML
            sorted_image_segments = sorted(mobi_image_segments, key=lambda x: (x[0], x[1]))
            
            # Track how many images we've inserted to adjust subsequent indices
            images_inserted_count = 0
            
            for insert_idx, placeholder_id, image_path, context_info in sorted_image_segments:
                # Create image placeholder segment
                image_info = image_data_map_raw.get(image_path, {})
                data_uri = image_info.get("data", "")
                mime_type = image_info.get("mime", "image/png")
                import os
                alt_text = os.path.basename(image_path) or image_path
                
                # Create placeholder text (frontend expects <ph-xxx> format)
                # Use image_path as placeholder_id (e.g., "mobi7/Images/image00044.jpeg")
                placeholder_text = f"<ph-{placeholder_id}>"
                
                # Create segment object with image metadata
                image_seg_obj = {
                    "text": placeholder_text,
                    "segment_index": insert_idx + images_inserted_count,  # Adjust index based on already inserted images
                    "is_image": True,
                    "is_excluded": True,  # Images are excluded from translation
                    "exclusion_reason": "image",
                    "block_type": "image",
                    "placeholder_id": placeholder_id,
                    "image_path": image_path,
                    "image_data": data_uri,  # Include image data for frontend preview
                    "has_latex": False,  # Images never contain LaTeX
                }
                
                # Insert at calculated position, or append if position is beyond current segments
                # Adjust insert_idx by images_inserted_count to account for previously inserted images
                adjusted_insert_idx = insert_idx + images_inserted_count
                actual_insert_idx = min(adjusted_insert_idx, len(segments_with_metadata))
                segments_with_metadata.insert(actual_insert_idx, image_seg_obj)
                images_inserted_count += 1
                
                # Safely extract context info for logging
                context_before_str = 'N/A'
                context_after_str = 'N/A'
                if context_info:
                    context_before = context_info.get('context_before')
                    context_after = context_info.get('context_after')
                    if context_before:
                        context_before_str = str(context_before)[:50]
                    if context_after:
                        context_after_str = str(context_after)[:50]
                
                logger.debug(
                    LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Inserted MOBI image placeholder segment at position {actual_insert_idx} "
                    f"(calculated={insert_idx}, adjusted={adjusted_insert_idx}): "
                    f"placeholder_id={placeholder_id}, image_path={image_path}, "
                    f"context_before={context_before_str}, context_after={context_after_str}"
                    )
            
            # Update total segments count to include image segments
            # Re-index all segments after insertion to ensure segment_index is correct
            if mobi_image_segments:
                for idx, seg_obj in enumerate(segments_with_metadata):
                    seg_obj["segment_index"] = idx
                
                total = len(segments_with_metadata)
                positioned_count = sum(1 for x in sorted_image_segments if x[0] < len(all_segments))
                logger.info(
                    LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Inserted {len(mobi_image_segments)} MOBI image segments "
                    f"({positioned_count} positioned based on HTML structure, {len(mobi_image_segments) - positioned_count} appended to end). "
                    f"Total segments: {total} (original: {len(all_segments)}, images: {len(mobi_image_segments)})"
                    )

        # CRITICAL: For DOCX files with textbox/SDT content, append these segments
        # from translation_segments_data so the frontend can display them
        if translation_segments_data and isinstance(translation_segments_data, dict):
            _tb_segments_list = translation_segments_data.get("segments", [])
            _textbox_sdt_segments = []
            for _seg in _tb_segments_list:
                if isinstance(_seg, dict):
                    _seg_type = _seg.get("segment_type", "")
                    if "textbox" in _seg_type:
                        _textbox_sdt_segments.append(_seg)

            if _textbox_sdt_segments:
                _added_count = 0
                for _tb_seg in _textbox_sdt_segments:
                    _tb_obj = {
                        "segment_index": len(segments_with_metadata),
                        "source_text": _tb_seg.get("source_text", ""),
                        "target_text": _tb_seg.get("target_text", ""),
                        "is_excluded": _tb_seg.get("is_excluded", False),
                        "is_image": _tb_seg.get("is_image", False),
                        "is_failed": _tb_seg.get("is_failed", False),
                        "segment_type": _tb_seg.get("segment_type", "textbox_sdt"),
                        "textbox_key": _tb_seg.get("textbox_key", ""),
                    }
                    segments_with_metadata.append(_tb_obj)
                    _added_count += 1

                total += _added_count
                logger.info(
                    LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Appended {_added_count} textbox/SDT segments "
                    f"to segments_with_metadata. New total: {total}"
                )

        # CRITICAL: Handle case where available segments < total_segments
        # This can happen if source_chunks_cache.segments was truncated during storage
        available_segments_count = len(segments_with_metadata)
        
        # Removed frequent segments_with_metadata logging to reduce log verbosity
        # Only log when there's a mismatch (error case)
        
        if available_segments_count < total:
            logger.error(
                LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: CRITICAL ERROR - segments_with_metadata count ({available_segments_count}) "
                f"is less than total_segments ({total}). This indicates source_chunks_cache.segments was truncated "
                f"or not properly saved during extraction. Only {available_segments_count} segments are available, "
                f"but {total} segments were expected. This is likely a bug in prepare_source_preview_for_extractor_based "
                f"or task_state storage. Returning available segments but frontend will see incomplete data."
                )
            # CRITICAL: If offset is beyond available segments, return empty list immediately
            # This prevents slow pagination where each request returns only 1 segment
            if pagination.offset >= available_segments_count:
                logger.warning(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Requested offset ({pagination.offset}) is beyond available segments "
                    f"({available_segments_count}). Returning empty list to prevent slow pagination."
                    )
                paginated_segments = []
                offset_clamped = pagination.offset
                end = pagination.offset
            else:
                # Clamp pagination to available segments
                offset_clamped = pagination.offset
                end = min(offset_clamped + pagination.limit, available_segments_count)
                paginated_segments = segments_with_metadata[offset_clamped:end]
            effective_total_for_pagination = available_segments_count
        else:
            # Normal case: all segments are available
            offset_clamped = min(pagination.offset, max(0, total - 1)) if total > 0 else 0
            end = min(offset_clamped + pagination.limit, total)
            paginated_segments = segments_with_metadata[offset_clamped:end]
            effective_total_for_pagination = total
        
        # Removed frequent pagination logging to reduce log verbosity
        # Only log when there are issues (e.g., mismatch between available and total)
        
        # Debug: Count excluded segments in paginated slice
        excluded_in_slice = sum(1 for seg in paginated_segments if seg.get("is_excluded", False))
        if excluded_in_slice > 0:
            logger.trace(
                LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Paginated slice (offset={offset_clamped}, limit={pagination.limit}) "
                f"contains {excluded_in_slice} excluded segments out of {len(paginated_segments)} total"
                )
        
        # Get workflow-specific metadata
        # NOTE: segments_metadata and workflow_type are already defined earlier (before MOBI image segment detection)
        # Only re-assign if not already defined (shouldn't happen, but safety check)
        if 'segments_metadata' not in locals() or not segments_metadata:
            segments_metadata = st.get("segments_metadata", {})
        if 'workflow_type' not in locals() or workflow_type is None:
            workflow_type = segments_metadata.get("workflow_type") or st.get("payload", {}).get("workflow_type") if isinstance(st.get("payload"), dict) else getattr(st.get("payload"), 'workflow_type', None) if st.get("payload") else None
        # Get chunk_size from centralized service (priority: payload → platform config → fallback)
        chunk_size = None
        payload = st.get("payload")
        if payload:
            from app.services.translation.chunk_size_service import chunk_size_service
            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
        if not chunk_size or chunk_size == 0:
            chunk_size = 3000  # Default fallback
        
        # Generate chunks from segments if chunk_to_segment_map is available
        # This allows frontend to display chunks for DOCX/Excel files (not just PDF)
        all_chunks = []
        chunks_text = []
        chunk_to_segment_map = st.get("chunk_to_segment_map")
        # Track original cached chunk_size (if any) for validation
        # Check both segments_metadata and cache_info to detect chunk_size changes
        cached_chunk_size = None
        segments_metadata_chunk_size = None
        if isinstance(segments_metadata, dict):
            segments_metadata_chunk_size = segments_metadata.get("chunk_size")
            cached_chunk_size = segments_metadata_chunk_size
        if not cached_chunk_size and cache_info:
            cached_chunk_size = cache_info.get("chunk_size")
        
        # If chunk_size changed or existing chunk tokens exceed limit, force rebuild chunk map
        chunk_tokens_info_cached = st.get("chunk_tokens_info") or []
        need_rebuild_chunk_map = False
        if chunk_to_segment_map and chunk_size:
            if cached_chunk_size and cached_chunk_size != chunk_size:
                need_rebuild_chunk_map = True
                logger.info(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: chunk_size changed "
                    f"(cached={cached_chunk_size}, current={chunk_size}), will rebuild chunk_to_segment_map"
                )
            # Case 3: Cached chunk tokens exceed current limit
            elif any((t or 0) > chunk_size for t in chunk_tokens_info_cached):
                need_rebuild_chunk_map = True
                logger.info(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Cached chunk tokens exceed current chunk_size={chunk_size}, "
                    f"will rebuild chunk_to_segment_map"
                )
        if need_rebuild_chunk_map:
            st.pop("chunk_to_segment_map", None)
            st.pop("chunk_tokens_info", None)
            st.pop("total_estimated_input_tokens", None)
            chunk_to_segment_map = None
        # Initialize text_token_limit for logging (will be set if using segments2json_chunks)
        text_token_limit = None
        
        # Calculate text token limit for logging (used when regenerating)
        text_token_limit = None
        if chunk_size:
            from utils.chunk_size_converter import get_text_content_token_limit
            text_token_limit = get_text_content_token_limit(chunk_size)
        
        # Check if chunk_to_segment_map is valid (not None and not empty)
        has_valid_chunk_map = chunk_to_segment_map is not None and len(chunk_to_segment_map) > 0 if chunk_to_segment_map else False

        # Backfill chunk_id into segments_with_metadata for merged paragraph preview
        if has_valid_chunk_map:
            seg_to_chunk = {}
            for c_idx, seg_indices in enumerate(chunk_to_segment_map):
                for seg_idx in seg_indices:
                    seg_to_chunk[seg_idx] = c_idx
            for seg_obj in segments_with_metadata:
                seg_idx = seg_obj.get("segment_index")
                if seg_idx is not None and seg_idx in seg_to_chunk:
                    seg_obj["chunk_id"] = seg_to_chunk[seg_idx]

        if has_valid_chunk_map and segments_with_metadata:
            try:
                # Removed frequent logging to reduce log verbosity
                # logger.info(...)  # Removed: too frequent
                # Build chunks from chunk_to_segment_map using segments_with_metadata
                # Skip excluded segments when building chunks
                for chunk_idx, chunk_segment_indices in enumerate(chunk_to_segment_map):
                    if chunk_segment_indices:
                        # Merge segments in this chunk (extract text from metadata objects)
                        # Skip excluded segments
                        chunk_segments = []
                        for seg_idx in chunk_segment_indices:
                            if seg_idx < len(segments_with_metadata):
                                seg_obj = segments_with_metadata[seg_idx]
                                # Skip excluded segments
                                if isinstance(seg_obj, dict) and seg_obj.get("is_excluded", False):
                                    continue
                                if isinstance(seg_obj, dict):
                                    seg_text = seg_obj.get("text", seg_obj.get("source_text", ""))
                                else:
                                    seg_text = str(seg_obj)
                                if seg_text.strip():  # Only add non-empty segments
                                    chunk_segments.append(seg_text)
                        # Only add chunk if it has at least one non-excluded segment
                        if chunk_segments:
                            chunk_text = "\n\n".join(chunk_segments)
                            # Estimate tokens for this chunk (including system prompt)
                            from utils.token_estimator import estimate_chunk_input_tokens
                            # Convert chunk segments to JSON format for token estimation
                            chunk_dict = {str(i): seg for i, seg in enumerate(chunk_segments)}
                            import json
                            chunk_json = json.dumps(chunk_dict, ensure_ascii=False)
                            estimated_tokens = estimate_chunk_input_tokens(chunk_json)
                            # Calculate text-only tokens (without system prompt and overhead)
                            text_only_tokens = max(0, estimated_tokens - 500)  # system prompt (450) + overhead (50)
                            chunk_bytes = len(chunk_json.encode('utf-8'))
                            all_chunks.append({
                                "text": chunk_text,
                                "chunk_type": "text",
                                "estimated_input_tokens": estimated_tokens,
                            })
                            chunks_text.append(chunk_text)
                            # Log chunk details (trace level for detailed chunk information)
                            # Note: chunk_idx is the original index in chunk_to_segment_map, 
                            # len(all_chunks)-1 is the index in the filtered all_chunks list
                            logger.trace(
                                LogModule.EXTRACT,
                                f"[PREVIEW-API] Task {task_id}: Chunk #{chunk_idx} (mapped to #{len(all_chunks)-1} in filtered list) - "
                                f"estimated_tokens={estimated_tokens} (text_only={text_only_tokens}, system+overhead=500), "
                                f"bytes={chunk_bytes}, segments={len(chunk_segments)}, "
                                f"chunk_size={chunk_size}, text_token_limit={text_token_limit}"
                                )
                
                logger.info(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Generated {len(all_chunks)} chunks from "
                    f"{len(all_segments)} segments using chunk_to_segment_map"
                )
            except Exception as e:
                logger.warning(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Failed to generate chunks from chunk_to_segment_map: {e}",
                    exc_info=True
                )
        elif chunk_size and segments_with_metadata:
            # Check if this is a format conversion task (convert_only mode)
            # Format conversion tasks skip chunk_to_segment_map building, so this is expected
            convert_only = st.get("convert_only", False)
            is_format_conversion = st.get("is_format_conversion", False)
            
            if convert_only or is_format_conversion:
                # Format conversion task - chunks are not needed, skip silently
                logger.debug(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Skipping chunk generation for format conversion task "
                    f"(convert_only={convert_only}, is_format_conversion={is_format_conversion})"
                )
                all_chunks = []
                chunks_text = []
            else:
                # CRITICAL: chunk_to_segment_map is required for translation tasks, do not use fallback
                # If chunk_to_segment_map is not available, it means the extraction phase failed to build it
                # This should not happen if all formats properly call _build_chunk_to_segment_map
                logger.warning(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: chunk_to_segment_map not available for chunk generation. "
                    f"This indicates a bug in the extraction phase - _build_chunk_to_segment_map was not called. "
                    f"chunk_size={chunk_size}, segments={len(segments_with_metadata)}, workflow_type={workflow_type}",
                )
                # Do not generate chunks - return empty chunks list
                # This will cause the API to return chunks as empty, which is better than incorrect fallback behavior
                all_chunks = []
                chunks_text = []
        else:
            # No chunk_size or segments_with_metadata available
            if not chunk_size:
                logger.trace(LogModule.EXTRACT, f"[PREVIEW-API] Task {task_id}: chunk_size is None or 0, cannot generate chunks")
            if not segments_with_metadata:
                logger.trace(LogModule.EXTRACT, f"[PREVIEW-API] Task {task_id}: segments_with_metadata is empty, cannot generate chunks")
        
        # Build paginated response
        paginated_response = PaginatedResponse(
            items=paginated_segments,
            offset=offset_clamped,
            limit=pagination.limit,
            total=total,
        )
        
        # Convert paginated segments to strings for backward compatibility (items field)
        # But also include full objects in segments field for metadata access
        # For frontend to access is_excluded flag, we need to keep objects in items field
        # But also provide strings for backward compatibility
        paginated_segments_strings = []
        paginated_segments_objects = []  # Keep objects for metadata access
        for seg_obj in paginated_segments:
            if isinstance(seg_obj, dict):
                # Extract text from object
                seg_text = seg_obj.get("text", seg_obj.get("source_text", ""))
                paginated_segments_strings.append(seg_text)
                paginated_segments_objects.append(seg_obj)  # Keep object with metadata
            else:
                seg_text = str(seg_obj)
                paginated_segments_strings.append(seg_text)
                # Create object from string for metadata access
                paginated_segments_objects.append({"text": seg_text})
        
        # Build paginated response with objects (for frontend to access is_excluded)
        # Debug: Verify that objects contain is_excluded flag
        excluded_in_response = sum(1 for seg in paginated_segments_objects if isinstance(seg, dict) and seg.get("is_excluded", False))
        if excluded_in_response > 0:
            logger.debug(
                LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: Response items contain {excluded_in_response} excluded segments "
                f"out of {len(paginated_segments_objects)} total in paginated slice"
            )
        
        # Check for missing segment_index (critical for frontend mapping)
        if paginated_segments_objects:
            segments_without_index = []
            for i, seg_obj in enumerate(paginated_segments_objects):
                if isinstance(seg_obj, dict) and seg_obj.get("segment_index") is None:
                    segments_without_index.append(i)
            
            if segments_without_index:
                logger.warning(
                    LogModule.EXTRACT,
                    f"[PREVIEW-API] Task {task_id}: {len(segments_without_index)} segments missing segment_index: "
                    f"{segments_without_index[:10]}"
                )
        
        # CRITICAL: Return original total to frontend so it knows the real count
        # Frontend can then detect the mismatch and handle it appropriately
        # But pagination is clamped to available segments to prevent IndexError
        paginated_response = PaginatedResponse(
            items=paginated_segments_objects,  # Use objects instead of strings
            offset=offset_clamped,
            limit=pagination.limit,
            total=total,  # Return original total (not effective_total) so frontend knows real count
        )
        
        # Log response details before returning
        logger.info(
            LogModule.EXTRACT,
            f"[PREVIEW-API] Task {task_id}: Preparing response - "
            f"total_segments={total}, "
            f"returned_items={len(paginated_segments_objects)}, "
            f"offset={offset_clamped}, "
            f"limit={pagination.limit}, "
            f"has_next={paginated_response.has_next}, "
            f"ready={ready}"
            )
        
        # Convert to dict and add extra fields
        # CRITICAL: Return original total to frontend so it knows the real count
        # Frontend can detect mismatch between total_segments and returned segments
        response = paginated_response.to_dict(
            task_id=task_id,
            ready=ready,
            segments=paginated_segments_objects,  # CRITICAL: Use objects instead of strings so frontend can access exclusion_reason
            total_segments=total,  # Return original total (not effective_total) so frontend knows real count
        )
        
        # Add full segment objects with metadata for paginated slice (for frontend to access is_excluded, etc.)
        # Note: This only includes metadata for the current page, not all segments
        # Frontend should check items/segments directly for full metadata if needed
        response["segments_metadata"] = paginated_segments
        
        # Add chunks if available (for frontend Extract page)
        if all_chunks:
            response["chunks"] = all_chunks
            response["chunks_text"] = chunks_text
            response["total_chunks"] = len(all_chunks)
            # Removed frequent logging to reduce log verbosity
            # logger.info(...)  # Removed: too frequent
            # Calculate total estimated input tokens for all chunks
        else:
            logger.warning(LogModule.WORKFLOW, f"[PREVIEW-API] Task {task_id}: No chunks generated, chunks will not be available in response")
            total_estimated_tokens = sum(
                chunk.get("estimated_input_tokens", 0) 
                for chunk in all_chunks 
                if isinstance(chunk, dict)
            )
            
            # Log detailed token calculation for debugging
            if all_chunks:
                import json
                chunk_tokens_detail = []
                for i, chunk in enumerate(all_chunks):
                    if isinstance(chunk, dict):
                        estimated_tokens = chunk.get("estimated_input_tokens", 0)
                        chunk_text = chunk.get("text", "")
                        # Calculate bytes (UTF-8 encoding)
                        chunk_bytes = len(chunk_text.encode('utf-8')) if chunk_text else 0
                        # Calculate text-only tokens (without system prompt and overhead)
                        # estimated_tokens includes system prompt (450) + overhead (50) = 500
                        text_only_tokens = max(0, estimated_tokens - 500) if estimated_tokens > 0 else 0
                        # Use format_content_for_log to respect content_display config
                        from logger.logger import format_content_for_log
                        chunk_tokens_detail.append({
                            'index': i,
                            'estimated_tokens': estimated_tokens,
                            'text_only_tokens': text_only_tokens,
                            'bytes': chunk_bytes,
                            'text_preview': format_content_for_log(chunk_text, max_length=100) if chunk_text else ''
                        })
                    else:
                        chunk_tokens_detail.append({
                            'index': i,
                            'estimated_tokens': 0,
                            'text_only_tokens': 0,
                            'bytes': 0,
                            'text_preview': ''
                        })
                
                chunks_with_tokens = [ct for ct in chunk_tokens_detail if ct['estimated_tokens'] > 0]
                chunks_without_tokens = [ct for ct in chunk_tokens_detail if ct['estimated_tokens'] == 0]
                
                logger.debug(
                    LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Token calculation summary - "
                    f"total_chunks={len(all_chunks)}, "
                    f"chunks_with_tokens={len(chunks_with_tokens)}, "
                    f"chunks_without_tokens={len(chunks_without_tokens)}, "
                    f"total_estimated_tokens={total_estimated_tokens}, "
                    f"chunk_size_limit={chunk_size} (total), "
                    f"text_token_limit={text_token_limit if 'text_token_limit' in locals() else 'N/A'}"
                )
                
                # Log detailed info for each chunk
                from logger.logger import format_content_for_log
                for ct in chunk_tokens_detail:
                    # Use format_content_for_log to respect content_display config
                    text_preview = format_content_for_log(ct['text_preview'], max_length=50) if ct['text_preview'] else ''
                    logger.debug(
                        LogModule.WORKFLOW,
                        f"[PREVIEW-API] Task {task_id}: Chunk #{ct['index']} - "
                        f"estimated_tokens={ct['estimated_tokens']} "
                        f"(text_only={ct['text_only_tokens']}, system+overhead=500), "
                        f"bytes={ct['bytes']}, "
                        f"text_preview={text_preview}"
                    )
                
                # Check if any chunk exceeds the limit
                if chunk_size:
                    exceeded_chunks = [ct for ct in chunk_tokens_detail if ct['estimated_tokens'] > chunk_size]
                    if exceeded_chunks:
                        logger.warning(
                            LogModule.WORKFLOW,
                            f"[PREVIEW-API] Task {task_id}: {len(exceeded_chunks)} chunks exceed chunk_size limit ({chunk_size}): "
                            f"{[(ct['index'], ct['estimated_tokens']) for ct in exceeded_chunks]}"
                        )
                    else:
                        logger.debug(
                            LogModule.WORKFLOW,
                            f"[PREVIEW-API] Task {task_id}: All chunks within chunk_size limit ({chunk_size})"
                        )
            
            if total_estimated_tokens > 0:
                response["total_estimated_input_tokens"] = total_estimated_tokens
                logger.info(LogModule.WORKFLOW, f"[PREVIEW-API] Task {task_id}: Added total_estimated_input_tokens={total_estimated_tokens} to response")
            else:
                logger.warning(LogModule.WORKFLOW, f"[PREVIEW-API] Task {task_id}: total_estimated_tokens is 0 or negative, NOT adding to response")
            
            logger.trace(LogModule.WORKFLOW, f"[PREVIEW-API] Task {task_id}: Including {len(all_chunks)} chunks in response (total estimated tokens: {total_estimated_tokens})")
        
        # Add workflow metadata if available
        if workflow_type:
            response["workflow_type"] = workflow_type
        # chunk_size removed from response - frontend should use global settings instead
        if cache_info:
            response["cache_info"] = {
                "content_hash": cache_info.get("content_hash"),
                "cached_chunk_size": cache_info.get("chunk_size"),
                "cached_at": cache_info.get("created_at"),
            }
        
        # Add image data map for frontend to display images instead of placeholders
        # CRITICAL: For MOBI files, convert image_data_map to frontend-expected format
        # Frontend expects: {placeholder_id: {"data": data_uri, "alt": alt_text}}
        # Backend stores: {image_path: {"data": data_uri, "mime": mime_type, "size": size}}
        image_data_map_raw = st.get("image_data_map", {})
        image_data_map_for_frontend = {}
        if image_data_map_raw:
            # Check if this is MOBI workflow
            if workflow_type == "mobi":
                # For MOBI, convert image_path keys to placeholder_id format
                # Use image_path as placeholder_id (e.g., "mobi7/Images/image00044.jpeg")
                for image_path, image_info in image_data_map_raw.items():
                    if isinstance(image_info, dict):
                        data_uri = image_info.get("data", "")
                        mime_type = image_info.get("mime", "image/png")
                        # Use image_path as placeholder_id, and extract filename for alt text
                        import os
                        alt_text = os.path.basename(image_path) or image_path
                        image_data_map_for_frontend[image_path] = {
                            "data": data_uri,
                            "alt": alt_text
                        }
                logger.info(
                    LogModule.WORKFLOW,
                    f"[PREVIEW-API] Task {task_id}: Converted {len(image_data_map_for_frontend)} MOBI images "
                    f"to frontend format. Sample keys: {list(image_data_map_for_frontend.keys())[:3]}"
                    )
            else:
                # For other workflows (PDF, DOCX, etc.), use as-is (already in correct format)
                image_data_map_for_frontend = image_data_map_raw
        
        if image_data_map_for_frontend:
            response["image_data_map"] = image_data_map_for_frontend
            logger.trace(LogModule.WORKFLOW, f"[PREVIEW-API] Task {task_id}: Including {len(image_data_map_for_frontend)} images in response")

        # Build pdf_compat_summary from cached results for quick frontend overview
        _pdf_compat_results = st.get("pdf_compat_results")
        if isinstance(_pdf_compat_results, dict) and _pdf_compat_results:
            _checked = 0
            _passed = 0
            _failed = 0
            _failed_indices = []
            for _k, _v in _pdf_compat_results.items():
                if isinstance(_v, dict):
                    _checked += 1
                    if _v.get("passed"):
                        _passed += 1
                    else:
                        _failed += 1
                        try:
                            _failed_indices.append(int(_k))
                        except (ValueError, TypeError):
                            pass
            response["pdf_compat_summary"] = {
                "checked_segments": _checked,
                "passed": _passed,
                "failed": _failed,
                "failed_segment_indices": sorted(_failed_indices),
            }

        return JSONResponse(content=response)
    
    @staticmethod
    def _get_payload_target_lang(payload: Any) -> Optional[str]:
        """Extract target_lang from payload (supports both dict and object)."""
        if not payload:
            return None
        if isinstance(payload, dict):
            return payload.get("to_lang") or payload.get("target_lang")
        return getattr(payload, 'to_lang', None) or getattr(payload, 'target_lang', None)
    
    def _is_re_extract(self, task_id: str, chunk_size: Optional[int] = None) -> bool:
        """
        Determine if this is a re-extract operation.
        
        Conditions for re-extract:
        1. segments_metadata.excluded_segments already exists and is not empty
        2. chunk_size has changed from the last extraction
        
        Args:
            task_id: Unique task identifier
            chunk_size: Current chunk size (optional)
            
        Returns:
            True if this is a re-extract, False otherwise
        """
        task_state = self.task_manager.get_task(task_id)
        if not task_state:
            return False
        
        segments_metadata = task_state.get("segments_metadata", {})
        existing_excluded = segments_metadata.get("excluded_segments", {})
        
        # If exclusion data already exists, this is a re-extract
        if existing_excluded:
            logger.trace(
                LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Detected re-extract (existing excluded_segments has {len(existing_excluded)} entries)"
                )
            return True
        
        # If chunk_size has changed, this is a re-extract
        if chunk_size:
            last_chunk_size = segments_metadata.get("last_chunk_size")
            if last_chunk_size is not None and last_chunk_size != chunk_size:
                logger.trace(
                    LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Detected re-extract (chunk_size changed: {last_chunk_size} -> {chunk_size})"
                    )
                return True
        
        return False
    
    async def get_layout_extract(
        self,
        task_id: str,
        chunk_size: Optional[int] = None,
        excluded_segment_indices: Optional[str] = None,
        target_lang: Optional[str] = None
    ):
        """
        Get layout extraction result for PDF files.
        
        Args:
            task_id: Unique task identifier
            chunk_size: Override chunk size for regenerating chunks (optional)
            excluded_segment_indices: Comma-separated list of segment indices to exclude (optional)
            
        Returns:
            JSONResponse with layout extract data
            
        Raises:
            HTTPException: If task not found or not a PDF file
        """
        task_state = self.task_manager.get_task(task_id)
        if task_state is None:
            logger.info(
                LogModule.EXTRACT,
                f"[LAYOUT-EXTRACT] Task {task_id}: Returning 404 - task not found (task may have been lost after server reload).",
            )
            raise HTTPException(status_code=404, detail=f"Task ID '{task_id}' not found.")
        
        st = task_state
        # CRITICAL: Set task_id early so progress/message updates (segment building + batch detection) are visible to getStatus
        st["task_id"] = task_id
        
        # CRITICAL: Determine if this is a re-extract operation
        # If re-extract, clear existing exclusion data to regenerate it
        is_re_extract = self._is_re_extract(task_id, chunk_size)
        if is_re_extract:
            segments_metadata = st.get("segments_metadata", {})
            # CRITICAL: Preserve user_selected exclusions even during re-extract
            # User manually excluded segments should not be cleared
            existing_excluded = segments_metadata.get("excluded_segments", {})
            preserved_user_selected = {}
            if existing_excluded and isinstance(existing_excluded, dict):
                for seg_idx_str, exclusion_info in existing_excluded.items():
                    # Handle both dict and string formats
                    if isinstance(exclusion_info, dict):
                        reason_str = exclusion_info.get("reason", "")
                    else:
                        reason_str = str(exclusion_info)
                    
                    if reason_str == "user_selected":
                        preserved_user_selected[seg_idx_str] = exclusion_info
            
            # Clear exclusion data for re-extract
            segments_metadata["excluded_segments"] = preserved_user_selected
            segments_metadata["excluded_segment_indices"] = []
            # NOTE: We preserve user_unexcluded_segments to prevent re-detection
            # of segments that user explicitly unexcluded
            logger.trace(
                LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Re-extract detected, cleared excluded_segments and excluded_segment_indices. "
                f"Preserved {len(preserved_user_selected)} user_selected exclusions: {list(preserved_user_selected.keys())[:10]}{'...' if len(preserved_user_selected) > 10 else ''}. "
                f"Preserved user_unexcluded_segments: {segments_metadata.get('user_unexcluded_segments', [])}"
                )
        
        # CRITICAL: If excluded_segment_indices are provided via query parameter,
        # we need to interpret them VERY carefully:
        #
        # - 对于 PDF 布局预览（当前函数场景），这些 indices 主要来源于后端已经计算好的
        #   segments_metadata.excluded_segments / excluded_segment_indices，前端只是把当前
        #   “已排除片段列表”原样传回，用于这一轮 layout-extract 中跳过对应段。
        # - 它们**不是**“新增的用户排除请求”，否则会把所有内容型排除（identifier / reference 等）
        #   都错误地标记为 user_selected，导致在目标语言切换时，Language Match 无法正确重新检测。
        #
        # 因此这里的处理策略：
        # - 空字符串 ""：表示“清理用户排除”（例如前端取消了某些手动排除），我们后面会在
        #   segment 遍历时按 clear_user_exclusions 标志做细粒度过滤；
        # - 非空列表：仅用于本次调用的 excluded_segment_indices 传递（已经通过 query 参数给到我们），
        #   **不再**写回 segments_metadata，也不再把这些 indices 统一标记成 USER_SELECTED。
        #   持久化的排除状态仍然由 ExclusionDetectionBatch / ExclusionManager.update_excluded_segments
        #   等统一路径维护。
        #
        # Empty string means "clear user-selected exclusions" (keep auto-detected exclusions like equations)
        # We'll use a special marker to indicate "clear user-selected exclusions"
        clear_user_exclusions = False
        if excluded_segment_indices is not None:
            try:
                if excluded_segment_indices == "":
                    # Empty string means clear user-selected exclusions (e.g., references)
                    # We'll filter out ref_text blocks from pre_existing_excluded_indices later
                    clear_user_exclusions = True
                    logger.trace(LogModule.EXCLUSION, f"[LAYOUT-EXTRACT] Task {task_id}: Will clear user-selected excluded segment indices (empty string provided)")
                else:
                    # Parse comma-separated list of indices，仅用于本次调用的 excluded_segment_indices 语义，
                    # 不再写回 segments_metadata，也不更改排除 reason（避免把内容型排除错误标记为 user_selected）。
                    excluded_indices_list = [int(x.strip()) for x in excluded_segment_indices.split(',') if x.strip().isdigit()]
                    if excluded_indices_list:
                        logger.trace(
                            LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Received {len(excluded_indices_list)} excluded_segment_indices from query parameter "
                            f"(used only for current layout-extract call, not persisted as user_selected)."
                            )
            except Exception as e:
                logger.trace(LogModule.EXCLUSION, f"[LAYOUT-EXTRACT] Task {task_id}: Failed to parse excluded_segment_indices from query parameter: {e}")
        
        # Check if this is a PDF file with layout document
        original_filename = st.get("original_filename", "")
        is_pdf_file = original_filename.lower().endswith('.pdf')
        if not is_pdf_file:
            raise HTTPException(
                status_code=400,
                detail="Layout extract is only available for PDF files processed with MinerU."
            )
        
        layout_doc = st.get("layout_document")
        
        # If layout_document is not available, try to load from layout_source_zip or attachments or disk
        if layout_doc is None:
            logger.trace(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] layout_document not in task_state, attempting to load from layout_source_zip or attachments")
            # Try layout_source_zip first
            layout_source_zip = st.get("layout_source_zip")
            if layout_source_zip:
                try:
                    from layout.registry import load_layout_from_engine_zip
                    layout_doc = load_layout_from_engine_zip("mineru", layout_source_zip)
                    if layout_doc:
                        # Store in task_state for future use (bbox at extraction so export does not need layout_doc)
                        st["layout_document"] = layout_doc
                        from utils.format_convert_utils import get_layout_block_bbox
                        st["layout_block_bbox"] = get_layout_block_bbox(layout_doc)
                        logger.trace(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Loaded layout_document from layout_source_zip for task {task_id}")
                except Exception as load_error:
                    logger.trace(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Failed to load layout_document from layout_source_zip: {load_error}")
            
            # If still not available, try attachments
            if layout_doc is None:
                attachments = st.get("attachments", {})
                if "mineru" in attachments:
                    mineru_attachment = attachments["mineru"]
                    zip_bytes = None
                    if hasattr(mineru_attachment, "content"):
                        zip_bytes = mineru_attachment.content
                    elif hasattr(mineru_attachment, "document") and hasattr(mineru_attachment.document, "content"):
                        zip_bytes = mineru_attachment.document.content
                    
                    if zip_bytes:
                        try:
                            from layout.registry import load_layout_from_engine_zip
                            layout_doc = load_layout_from_engine_zip("mineru", zip_bytes)
                            if layout_doc:
                                # Store in task_state for future use (bbox at extraction so export does not need layout_doc)
                                st["layout_document"] = layout_doc
                                from utils.format_convert_utils import get_layout_block_bbox
                                st["layout_block_bbox"] = get_layout_block_bbox(layout_doc)
                                # Also store layout_source_zip for future use
                                st["layout_source_zip"] = zip_bytes
                                logger.trace(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Loaded layout_document from MinerU attachment for task {task_id}")
                        except Exception as load_error:
                            logger.trace(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Failed to load layout_document from MinerU attachment: {load_error}")
            
            # If still not available, try loading from disk (mineru_zip_path or temp_dir/mineru_layout.zip)
            if layout_doc is None:
                mineru_zip_path = st.get("mineru_zip_path")
                if not mineru_zip_path:
                    temp_dir = st.get("temp_dir")
                    if temp_dir:
                        mineru_zip_path = os.path.join(temp_dir, "mineru_layout.zip")
                if mineru_zip_path and os.path.isfile(mineru_zip_path):
                    try:
                        with open(mineru_zip_path, "rb") as f:
                            zip_bytes = f.read()
                        from layout.registry import load_layout_from_engine_zip
                        layout_doc = load_layout_from_engine_zip("mineru", zip_bytes)
                        if layout_doc:
                            st["layout_document"] = layout_doc
                            from utils.format_convert_utils import get_layout_block_bbox
                            st["layout_block_bbox"] = get_layout_block_bbox(layout_doc)
                            st["layout_source_zip"] = zip_bytes
                            logger.info(
                                LogModule.EXTRACT,
                                f"[LAYOUT-EXTRACT] Task {task_id}: Loaded layout_document from disk (mineru_zip_path).",
                            )
                    except Exception as load_error:
                        logger.trace(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Failed to load layout from mineru_zip_path: {load_error}")
        
        if layout_doc is None:
            has_zip = bool(st.get("layout_source_zip"))
            has_att = bool(st.get("attachments", {}).get("mineru"))
            _zip_path = st.get("mineru_zip_path") or (os.path.join(st.get("temp_dir", ""), "mineru_layout.zip") if st.get("temp_dir") else "")
            has_path = bool(_zip_path and os.path.isfile(_zip_path))
            logger.info(
                LogModule.EXTRACT,
                f"[LAYOUT-EXTRACT] Task {task_id}: Returning 404 - layout document not available "
                f"(has_layout_source_zip={has_zip}, has_attachments_mineru={has_att}, has_mineru_zip_path={has_path}). "
                "Ensure the file was processed with MinerU and layout is ready.",
            )
            raise HTTPException(
                status_code=404,
                detail="Layout document not available. Please ensure the file was processed with MinerU."
            )
        
        try:
            from layout.base import LayoutDocument as _LD
            if not isinstance(layout_doc, _LD):
                raise HTTPException(
                    status_code=404,
                    detail="Invalid layout document type."
                )
        except Exception as e:
            logger.trace(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Failed to validate layout document: {e}")
            raise HTTPException(
                status_code=500,
                detail=f"Failed to validate layout document: {e}"
            )
        
        # Extract images from ZIP if available
        from layout.pdf_renderer.shared.block_processor import BlockProcessor
        
        zip_bytes = st.get("layout_source_zip")
        image_data_map: dict[str, dict[str, str]] = {}
        zip_file = None
        
        # Map to store image data by path: {image_path: "data:image/...;base64,..."}
        image_data_by_path: dict[str, str] = {}
        
        if zip_bytes:
            try:
                zip_file = zipfile.ZipFile(io.BytesIO(zip_bytes))
                images_bytes_map = BlockProcessor.extract_all_images_from_layout(layout_doc, zip_file)
                # Convert image bytes to base64 data URIs
                for img_path, img_bytes in images_bytes_map.items():
                    mime = mimetypes.guess_type(img_path)[0] or "image/png"
                    data_uri = f"data:{mime};base64,{base64.b64encode(img_bytes).decode('ascii')}"
                    image_data_by_path[img_path] = data_uri
            except Exception as e:
                logger.trace(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Failed to extract images from ZIP: {e}")
                image_data_by_path = {}
        
        # Use LayoutMarkdownBuilder to generate deep-split segments (same as _prepare_layout_preview_from_layout)
        from layout.markdown_builder import LayoutMarkdownBuilder
        
        # Get chunk_size: prioritize query parameter, then task_state, then default
        # This ensures chunks are always regenerated with the current chunk_size setting
        chunk_size_query = chunk_size  # Save original query parameter for logging
        if chunk_size is None:
            chunk_size = st.get("segments_metadata", {}).get("chunk_size")
        if chunk_size is None:
            # Try to get from payload
            payload = st.get("payload")
            if payload:
                if isinstance(payload, dict):
                    chunk_size = payload.get("chunk_size")
                else:
                    chunk_size = getattr(payload, 'chunk_size', None)
        if chunk_size is None:
            chunk_size = 3000  # Default fallback
        
        # Get deep_split from task_state or payload
        deep_split_enabled = True  # Default to True
        source = "default"
        if "deep_split" in st:
            deep_split_enabled = bool(st["deep_split"])
            source = "task_state"
        else:
            payload = st.get("payload")
            if payload:
                if isinstance(payload, dict):
                    deep_split_enabled = bool(payload.get("deep_split", True))
                else:
                    deep_split_enabled = bool(getattr(payload, 'deep_split', True))
                source = "payload"
        
        logger.trace(
            LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Building layout extract with deep_split={deep_split_enabled} "
            f"(from {source}), chunk_size={chunk_size} (chunks will be regenerated)"
            )
        # Get equation_format from payload (default: "text")
        equation_format = "text"
        payload = st.get("payload")
        if payload:
            try:
                if isinstance(payload, dict):
                    equation_format = (payload.get("equation_format") or "text").lower()
                else:
                    equation_format = (getattr(payload, "equation_format", None) or "text").lower()
            except Exception:
                equation_format = "text"
        if equation_format not in ("text", "image"):
            equation_format = "text"
        builder = LayoutMarkdownBuilder(
            max_chunk_chars=chunk_size, 
            deep_split=deep_split_enabled,
            equation_format=equation_format,
            include_structural_blocks=True,  # Include header/footer blocks for Extract phase (user can choose to exclude)
        )
        st["message"] = "Building layout..."
        st["progress"] = 2
        # Run in thread so event loop can serve getStatus/auth (avoid connection timeout during long build)
        layout_result = await asyncio.to_thread(builder.build, layout_doc)
        logger.trace(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: LayoutMarkdownBuilder generated {len(layout_result.chunks)} chunks (deep_split={deep_split_enabled}, chunk_size={chunk_size})")
        
        if not layout_result.chunks:
            raise HTTPException(
                status_code=404,
                detail="No segments generated from layout document."
            )
        
        # Import unified exclusion detection utility
        from utils.translation_segments import _is_image_segment
        
        # Process chunks from LayoutMarkdownBuilder (these are already deep-split)
        all_segments = []
        original_text_parts = []  # Text without headers and footers
        excluded_segment_indices_list = []  # Track excluded segment indices (for translation)
        
        # CRITICAL: Read pre-existing excluded segments using ExclusionManager (single source of truth)
        # This ensures that manually excluded segments from Extract phase are correctly identified
        # ExclusionManager reads from segments_metadata.excluded_segments (new format) and
        # falls back to excluded_segment_indices (legacy format) for backward compatibility
        from exclusion.core import ExclusionManager
        excluded_segments_with_reasons = ExclusionManager.get_excluded_segments(st)
        pre_existing_excluded_indices = set(excluded_segments_with_reasons.keys())
        segments_metadata = st.get("segments_metadata", {})
        
        if pre_existing_excluded_indices:
            logger.trace(
                LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Found {len(pre_existing_excluded_indices)} pre-existing excluded segments "
                f"from ExclusionManager (single source of truth): {sorted(pre_existing_excluded_indices)[:20]}{'...' if len(pre_existing_excluded_indices) > 20 else ''}"
                )
        
        # If clear_user_exclusions is True, filter out ref_text blocks from pre_existing_excluded_indices
        # This allows clearing user-selected exclusions (references) while keeping auto-detected ones (equations)
        if clear_user_exclusions and pre_existing_excluded_indices:
            # We need to check block_type for each excluded index, but we don't have all_segments yet
            # So we'll filter them out during segment processing instead
            # For now, just log that we'll filter them
            logger.trace(LogModule.EXCLUSION, f"[LAYOUT-EXTRACT] Task {task_id}: Will filter out ref_text blocks from {len(pre_existing_excluded_indices)} pre-existing excluded segment indices")
        
        # Map to track which blocks are headers/footers/images for exclusion
        block_type_map = {}  # {block_index: block_type}
        block_image_map = {}  # {block_index: image_path}
        
        # First pass: collect block metadata from layout_doc
        for page in layout_doc.pages:
            for block in page.blocks:
                block_type = getattr(block, "type", "unknown") or "unknown"
                block_index = getattr(block, "index", None)
                image_path = getattr(block, "image_path", None)
                if block_index is not None:
                    block_type_map[block_index] = block_type
                    if image_path:
                        block_image_map[block_index] = image_path
        
        # Process chunks from LayoutMarkdownBuilder (yield to event loop periodically so getStatus/auth can run)
        total_segments_building = len(layout_result.chunks)
        for chunk_idx, chunk in enumerate(layout_result.chunks):
            chunk_text = chunk.text
            chunk_block_indices = chunk.block_indices if hasattr(chunk, 'block_indices') else []
            
            # Determine block type and metadata from chunk's block indices
            is_image = False
            is_header = False
            is_footer = False
            block_type = "text"
            block_index = None
            image_path = None
            placeholder_id = None
            
            # Check if chunk is an image chunk (from LayoutMarkdownBuilder) - check this first
            if chunk.chunk_type == "image":
                is_image = True
                image_path = chunk.image_path
                # Use image_placeholder from chunk if available
                if chunk.image_placeholder:
                    placeholder_id = chunk.image_placeholder
                else:
                    placeholder_id = f"img-{chunk_idx}"
            
            # Check if any block in this chunk is an image/header/footer/table
            # CRITICAL: For table blocks, we need to distinguish between:
            # - table_body: The actual table content (markdown/HTML table) - can be excluded by user
            # - table_caption: Table caption text - should be treated as normal text
            # - table_footnote: Table footnote text - should be treated as normal text
            # We identify table_body by checking if chunk text is a table segment (markdown/HTML table syntax)
            # CRITICAL: For image blocks, we need to distinguish between:
            # - image: Pure image content (placeholder only) - should be excluded
            # - image_caption: Image caption text - should be treated as normal text and NOT excluded as image
            is_table_body = False
            for block_idx in chunk_block_indices:
                if block_idx in block_type_map:
                    block_type = block_type_map[block_idx]
                    block_index = block_idx
                    if block_type == "image":
                        # Only set is_image=True if this is NOT an image_caption block
                        # Image captions contain actual text content and should not be excluded as images
                        is_image = True
                        if not image_path:
                            image_path = block_image_map.get(block_idx)
                    elif block_type in ("image_caption", "caption"):
                        # CRITICAL: image_caption blocks contain actual text content and should NOT be excluded as images
                        # Reset is_image to False if we encounter an image_caption block
                        is_image = False
                        # Keep block_type as image_caption for proper exclusion detection
                    elif block_type == "header":
                        is_header = True
                    elif block_type == "footer":
                        is_footer = True
                    elif block_type == "table":
                        # Prefer layout: chunk_type "table_body" from LayoutMarkdownBuilder (no string check).
                        # Fallback: chunk may be from older path without chunk_type; use _is_table_segment then.
                        chunk_type = getattr(chunk, "chunk_type", None)
                        if chunk_type == "table_body":
                            is_table_body = True
                            block_type = "table_body"
                        else:
                            from utils.translation_segments import _is_table_segment
                            if _is_table_segment(chunk_text):
                                is_table_body = True
                                block_type = "table_body"
                        # If not table_body, treat as normal text (caption or footnote)
                    elif block_type == "chart":
                        # Check if chunk is a chart_body type from LayoutMarkdownBuilder
                        chunk_type = getattr(chunk, "chunk_type", None)
                        if chunk_type == "chart_body":
                            block_type = "chart_body"
            
            # Also check if chunk text contains image placeholder (fallback)
            # CRITICAL: Only treat as image if block_type is NOT image_caption
            # Image captions may contain placeholders but also have actual text content
            placeholder_match = re.search(r'<ph-([^>]+)>', chunk_text)
            if placeholder_match and not is_image and block_type not in ("image_caption", "caption"):
                is_image = True
                placeholder_id = placeholder_match.group(1)
                # Try to find image_path from placeholder_id
                if placeholder_id.startswith("img-"):
                    try:
                        img_block_idx = int(placeholder_id.replace("img-", ""))
                        image_path = block_image_map.get(img_block_idx)
                    except ValueError:
                        pass
                # Check if it's a layoutimg placeholder (from LayoutMarkdownBuilder)
                elif placeholder_id.startswith("layoutimg"):
                    # Try to find image_path from chunk's block_indices
                    for block_idx in chunk_block_indices:
                        if block_idx in block_image_map:
                            image_path = block_image_map[block_idx]
                            break
            
            # Determine if should be excluded
            # Check if chunk text is actually a placeholder (not actual text content)
            # Image captions have actual text content, so they should not be excluded even if block_type is image
            placeholder_pattern = r'^<ph-[^>]+>$'
            is_placeholder_only = re.match(placeholder_pattern, chunk_text.strip())
            
            # CRITICAL: Check if this segment contains actual text content (not just placeholder)
            # Image captions (e.g., "Figure 1: ...") contain actual text and should NOT be excluded
            # even if block_type is "image" (because they share block_indices with image blocks)
            has_actual_text = False
            if chunk_text:
                # Remove placeholder and whitespace, check if there's actual text left
                text_without_placeholder = re.sub(r'<ph-[^>]+>', '', chunk_text).strip()
                # Check if there's meaningful text (not just whitespace or very short)
                if text_without_placeholder and len(text_without_placeholder) > 3:
                    has_actual_text = True
            
            is_excluded = False
            # CRITICAL: Check if user explicitly unexcluded this segment to prevent re-detection
            # This check must be done FIRST, before any auto-detection logic
            user_unexcluded_segments = st.get("segments_metadata", {}).get("user_unexcluded_segments", [])
            if chunk_idx in user_unexcluded_segments:
                # User explicitly chose to unexclude this segment, skip ALL auto-detection
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Segment {chunk_idx} was user-unexcluded, skipping ALL auto-detection "
                    f"(including language match and table detection). user_unexcluded_segments={user_unexcluded_segments}"
                    )
                is_excluded = False
            # CRITICAL: If segment has actual text content (e.g., image caption), do NOT exclude it as image
            # even if block_type is "image" (because image captions share block_indices with image blocks)
            elif has_actual_text and is_image:
                # Image caption: real text shares layout block index with image block — not an image segment
                is_image = False
                is_excluded = False
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Segment {chunk_idx} has actual text content (likely image caption), "
                    f"not excluding as image. block_type={block_type}, text_preview={chunk_text[:50]}..."
                    )
            # CRITICAL: First check if this segment was pre-marked as excluded (e.g., by frontend user choice)
            # This allows frontend to mark segments as excluded (e.g., references) and have them
            # automatically excluded when regenerating chunks
            # NOTE: chunk_idx is the index in layout_result.chunks, which corresponds to the index in all_segments
            # Frontend passes segment indices from segmentsData array, which matches all_segments indices
            elif chunk_idx in pre_existing_excluded_indices:
                # If clear_user_exclusions is True, skip ref_text blocks (user-selected exclusions)
                # but keep other pre-existing exclusions (auto-detected)
                if clear_user_exclusions and block_type == "ref_text":
                    # Skip this exclusion (user cleared it)
                    logger.trace(LogModule.EXCLUSION, f"[LAYOUT-EXTRACT] Task {task_id}: Segment {chunk_idx} (ref_text) exclusion cleared by user, will not be excluded")
                else:
                    is_excluded = True
                    logger.trace(LogModule.EXCLUSION, f"[LAYOUT-EXTRACT] Task {task_id}: Segment {chunk_idx} was pre-marked as excluded (e.g., by user choice for references). block_type={block_type}, chunk_text preview={chunk_text[:50] if chunk_text else 'empty'}...")
            # Exclude interline_equation blocks to preserve original LaTeX format and avoid AI modification
            elif block_type == "interline_equation":
                is_excluded = True
                logger.trace(LogModule.EXCLUSION, f"[LAYOUT-EXTRACT] Task {task_id}: Marking interline_equation block {block_index} as excluded to preserve LaTeX format")
            # NOTE: Tables are NOT automatically excluded by default
            # Most tables can be translated, so we only mark them with block_type="table"
            # User can choose to exclude tables if needed (will be marked as TABLE exclusion reason)
            # Tables will be detected and marked in segment_data.block_type for frontend display
            elif is_image and is_placeholder_only:
                is_excluded = True  # Only exclude if it's actually an image placeholder
            else:
                # Get target language from task_state for language-based exclusion
                # CRITICAL: Only perform auto-detection if segment is NOT in user_unexcluded_segments
                # (already checked above, but double-check for safety)
                if chunk_idx not in user_unexcluded_segments:
                    # IMPORTANT: Do NOT overwrite the get_layout_extract() function parameter `target_lang`.
                    # We only need a local value here for per-segment auto-detection during segment building.
                    payload = st.get("payload")
                    segment_target_lang = self._get_payload_target_lang(payload)
                    # Use detect_exclusion_reason instead of should_exclude_text
                    from exclusion.core.exclusion_detector import detect_exclusion_reason, ExclusionReason
                    # CRITICAL: If segment has actual text content (e.g., image caption), set is_image=False
                    # to prevent it from being detected as IMAGE exclusion reason
                    # Image captions share block_indices with image blocks, but contain actual text content
                    detection_is_image = is_image if not has_actual_text else False
                    # CRITICAL: PDF layout - table takes priority over identifier (strict_table_priority=True).
                    # Otherwise table body with URLs/emails is wrongly detected as IDENTIFIER and excluded.
                    detected_result = detect_exclusion_reason(
                        text=chunk_text,
                        block_type=block_type,
                        target_lang=segment_target_lang,
                        is_image=detection_is_image,
                        is_table=is_table_body,
                        strict_table_priority=True,
                    )
                    if detected_result:
                        detected_reason, _ = detected_result
                        default_excluded = ExclusionReason.get_default_excluded()
                        if detected_reason in default_excluded:
                            is_excluded = True
            
            # Track excluded segment indices for translation
            if is_excluded:
                excluded_segment_indices_list.append(chunk_idx)
            
            # Determine if this is a chart_body segment
            is_chart_body = block_type == "chart_body"
            
            # Build segment data
            # CRITICAL: Store original segment index (chunk_idx) for proper mapping
            # chunk_idx is the index in layout_result.chunks, which is the original segment index
            segment_data = {
                "text": chunk_text,
                "block_type": block_type,  # For table_body, this will be "table_body"
                "block_index": block_index,
                "layout_block_indices": list(chunk_block_indices)
                if chunk_block_indices
                else ([block_index] if block_index is not None else []),
                "chunk_index": chunk_idx,
                "segment_index": chunk_idx,  # CRITICAL: Store original segment index for proper mapping
                "is_image": is_image,
                "is_header": is_header,
                "is_footer": is_footer,
                "is_excluded": is_excluded,
                "is_table_body": is_table_body,  # Mark if this is table_body (for frontend display)
                "is_chart_body": is_chart_body,  # Mark if this is chart_body (for frontend display)
            }
            
            # Add exclusion_reason if segment is excluded
            # This will be populated after we detect exclusion reasons for all segments
            # For now, set to None - it will be updated later in the function
            
            # Add image data if available
            if is_image:
                # Generate placeholder ID if not already set
                if not placeholder_id:
                    placeholder_id = f"img-{block_index or chunk_idx}"
                segment_data["placeholder_id"] = placeholder_id
                segment_data["image_path"] = image_path or "Image" if image_path else "Image"
                
                # Ensure text contains placeholder if it's an image chunk
                if not placeholder_match and chunk.chunk_type == "image":
                    segment_data["text"] = f"<ph-{placeholder_id}>"
                
                # Try to find image data in image_data_by_path
                if image_path:
                    normalized_path = image_path.replace("\\", "/").lstrip("./")
                    image_data = None
                    for key, value in image_data_by_path.items():
                        if key.endswith(normalized_path) or normalized_path.endswith(key):
                            image_data = value
                            break
                    
                    if image_data:
                        segment_data["image_data"] = image_data
                    else:
                        # If image data not found, still include placeholder for frontend to handle
                        segment_data["image_data"] = None
            
            all_segments.append(segment_data)
            # Yield to event loop every 200 segments so getStatus/auth requests can be served (avoid connection timeout)
            if (chunk_idx + 1) % 200 == 0 or (chunk_idx + 1) == total_segments_building:
                await asyncio.sleep(0)
                pct = min(10, int((chunk_idx + 1) / total_segments_building * 10)) if total_segments_building else 0
                st["message"] = f"Building segments: {chunk_idx + 1}/{total_segments_building}"
                st["progress"] = pct
        
        # Build original text from LayoutMarkdownBuilder's markdown_text
        # This includes all segments (text, images, tables) merged together, excluding only headers and footers
        # Use the markdown_text from LayoutMarkdownBuilder which already contains all content in proper order
        original_text = layout_result.markdown_text
        
        # Remove headers and footers from original text if needed
        # We can identify headers/footers by checking if chunks are marked as header/footer
        # But since markdown_text is already built from chunks, we need to filter it differently
        # For now, we'll use the markdown_text as-is, but we can enhance this later if needed
        # The frontend can handle displaying headers/footers separately if needed
        
        # Alternative: Build original text from chunks excluding headers and footers
        # This ensures we have the merged text without headers/footers
        original_text_parts = []
        for chunk in layout_result.chunks:
            chunk_text = chunk.text
            # Check if this chunk contains header/footer blocks
            chunk_block_indices = chunk.block_indices if hasattr(chunk, 'block_indices') else []
            is_header_chunk = False
            is_footer_chunk = False
            for block_idx in chunk_block_indices:
                if block_idx in block_type_map:
                    block_type = block_type_map[block_idx]
                    if block_type == "header":
                        is_header_chunk = True
                    elif block_type == "footer":
                        is_footer_chunk = True
            
            # Include all chunks except headers and footers
            # Include images (as placeholders), tables, and all text content
            if not is_header_chunk and not is_footer_chunk and chunk_text.strip():
                original_text_parts.append(chunk_text)
        
        # Get chunk_size from centralized service (priority: payload → platform config → fallback)
        # Priority 0: Query parameter (if provided, use it directly)
        if chunk_size is None or chunk_size == 0:
            payload = st.get("payload")
            if payload:
                from app.services.translation.chunk_size_service import chunk_size_service
                chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
            if not chunk_size or chunk_size == 0:
                chunk_size = 3000  # Default fallback
                logger.warning(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: chunk_size is None or 0, using fallback value 3000")
        
        # Build chunks from segments (merge segments according to chunk_size)
        # Chunks are used for translation, segments are for proofreading
        # We need to merge segments into chunks based on chunk_size
        # IMPORTANT: This is always regenerated on each API call, never cached
        # CRITICAL: We need to build TWO sets of chunks:
        # 1. all_chunks: Only non-excluded chunks (for translation)
        # 2. all_chunks_with_excluded: All chunks including excluded (for layout_prepared_chunks)
        # This ensures MDTranslator can correctly map chunks to original segments
        logger.info(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: Regenerating chunks from {len(all_segments)} segments with chunk_size={chunk_size}")
        all_chunks = []  # Only non-excluded chunks (for translation)
        all_chunks_with_excluded = []  # All chunks including excluded (for layout_prepared_chunks)
        current_chunk_parts = []
        current_chunk_chars = 0
        current_chunk_segment_indices = []  # Track which segments are in this chunk
        current_chunk_parts_with_excluded = []
        current_chunk_chars_with_excluded = 0
        current_chunk_segment_indices_with_excluded = []  # Track which segments are in this chunk (including excluded)
        
        for seg_idx, segment_data in enumerate(all_segments):
            # Yield to event loop every 200 segments so getStatus/auth requests can be served
            if seg_idx > 0 and seg_idx % 200 == 0:
                await asyncio.sleep(0)
            segment_text = segment_data.get("text", "")
            segment_chars = len(segment_text)
            is_excluded = segment_data.get("is_excluded", False)
            
            # CRITICAL: Use original segment index from segment_data, not enumerate index
            # segment_data["segment_index"] must equal chunk_idx (0..N-1) as set when building all_segments.
            # If it does not match seg_idx, the data is inconsistent (e.g. block_index mistaken for segment_index).
            raw_segment_index = segment_data.get("segment_index", seg_idx)
            if raw_segment_index != seg_idx:
                logger.warning(
                    LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: segment_index mismatch at seg_idx={seg_idx}: "
                    f"segment_data['segment_index']={raw_segment_index} (expected {seg_idx}). "
                    f"Using seg_idx to avoid wrong chunk->segment mapping. "
                    f"segment_data keys: segment_index, block_index={segment_data.get('block_index')}"
                )
                original_segment_index = seg_idx
            else:
                original_segment_index = raw_segment_index
            
            # CRITICAL: Build chunks with excluded segments for layout_prepared_chunks
            # This ensures MDTranslator can correctly map chunks to original segments
            # Excluded segments are added as separate chunks (one segment per chunk)
            if is_excluded:
                # Flush current chunk (if any) before adding excluded segment as separate chunk
                if current_chunk_parts_with_excluded:
                    chunk_text_with_excluded = "\n\n".join(current_chunk_parts_with_excluded)
                    all_chunks_with_excluded.append({
                        "text": chunk_text_with_excluded,
                        "segment_indices": current_chunk_segment_indices_with_excluded.copy(),
                        "chunk_index": len(all_chunks_with_excluded),
                        "is_excluded": False,  # This chunk itself is not excluded, it contains non-excluded segments
                    })
                    current_chunk_parts_with_excluded = []
                    current_chunk_chars_with_excluded = 0
                    current_chunk_segment_indices_with_excluded = []
                
                # Add excluded segment as a separate chunk
                all_chunks_with_excluded.append({
                    "text": segment_text,
                    "segment_indices": [original_segment_index],
                    "chunk_index": len(all_chunks_with_excluded),
                    "is_excluded": True,  # This chunk is excluded
                })
                # Excluded segments are not added to all_chunks (for translation)
                continue
            
            # Check if we should start a new chunk
            # Start new chunk if:
            # 1. Current chunk would exceed chunk_size
            # 2. Current chunk is not empty and adding this segment would exceed chunk_size
            should_start_new_chunk = False
            if current_chunk_chars > 0 and current_chunk_chars + segment_chars > chunk_size:
                should_start_new_chunk = True
            
            if should_start_new_chunk and current_chunk_parts:
                # Flush current chunk (for both all_chunks and all_chunks_with_excluded)
                chunk_text = "\n\n".join(current_chunk_parts)
                # Estimate tokens for this chunk (including system prompt)
                from utils.token_estimator import estimate_chunk_input_tokens
                # Convert chunk segments to JSON format for token estimation
                chunk_dict = {str(i): seg for i, seg in enumerate(current_chunk_parts)}
                chunk_json = json.dumps(chunk_dict, ensure_ascii=False)
                estimated_tokens = estimate_chunk_input_tokens(chunk_json)
                all_chunks.append({
                    "text": chunk_text,
                    "segment_indices": current_chunk_segment_indices.copy(),
                    "chunk_index": len(all_chunks),
                    "estimated_input_tokens": estimated_tokens,
                })
                # Also add to all_chunks_with_excluded
                all_chunks_with_excluded.append({
                    "text": chunk_text,
                    "segment_indices": current_chunk_segment_indices_with_excluded.copy(),
                    "chunk_index": len(all_chunks_with_excluded),
                    "is_excluded": False,
                })
                current_chunk_parts = []
                current_chunk_chars = 0
                current_chunk_segment_indices = []
                current_chunk_parts_with_excluded = []
                current_chunk_chars_with_excluded = 0
                current_chunk_segment_indices_with_excluded = []
            
            # Add segment to current chunk (only non-excluded segments)
            # CRITICAL: Use original_segment_index instead of seg_idx to ensure proper mapping
            if segment_text.strip():
                current_chunk_parts.append(segment_text)
                current_chunk_chars += segment_chars
                current_chunk_segment_indices.append(original_segment_index)
                # Also add to chunks_with_excluded
                current_chunk_parts_with_excluded.append(segment_text)
                current_chunk_chars_with_excluded += segment_chars
                current_chunk_segment_indices_with_excluded.append(original_segment_index)
        
        # Flush last chunk if any (for both all_chunks and all_chunks_with_excluded)
        if current_chunk_parts:
            chunk_text = "\n\n".join(current_chunk_parts)
            # Estimate tokens for this chunk (including system prompt)
            from utils.token_estimator import estimate_chunk_input_tokens
            # Convert chunk segments to JSON format for token estimation
            chunk_dict = {str(i): seg for i, seg in enumerate(current_chunk_parts)}
            chunk_json = json.dumps(chunk_dict, ensure_ascii=False)
            estimated_tokens = estimate_chunk_input_tokens(chunk_json)
            all_chunks.append({
                "text": chunk_text,
                "segment_indices": current_chunk_segment_indices.copy(),
                "chunk_index": len(all_chunks),
                "estimated_input_tokens": estimated_tokens,
            })
            # Also add to all_chunks_with_excluded
            all_chunks_with_excluded.append({
                "text": chunk_text,
                "segment_indices": current_chunk_segment_indices_with_excluded.copy(),
                "chunk_index": len(all_chunks_with_excluded),
                "is_excluded": False,
            })
        elif current_chunk_parts_with_excluded:
            # Flush last chunk for all_chunks_with_excluded (if any)
            chunk_text_with_excluded = "\n\n".join(current_chunk_parts_with_excluded)
            all_chunks_with_excluded.append({
                "text": chunk_text_with_excluded,
                "segment_indices": current_chunk_segment_indices_with_excluded.copy(),
                "chunk_index": len(all_chunks_with_excluded),
                "is_excluded": False,
            })
        
        excluded_count = sum(1 for seg in all_segments if seg.get("is_excluded", False))
        logger.info(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: Regenerated {len(all_chunks)} chunks (for translation) and {len(all_chunks_with_excluded)} chunks (including excluded) from {len(all_segments)} segments (chunk_size={chunk_size}, excluded_segments={excluded_count})")
        
        # CRITICAL: Use new unified exclusion detection architecture
        # This replaces the previous manual detection loop with a unified batch detection service
        # Wrap in try-except to ensure segments are returned even if detection fails
        excluded_segments_with_reasons = {}
        all_detected_reasons = {}
        try:
            from exclusion.extractors.pdf_extractor import PDFMetadataExtractor
            from exclusion.detection.batch_detector import ExclusionDetectionBatch
            from exclusion.core import ExclusionReason
            
            # Get target_lang from multiple sources (priority order):
            # 1. Function parameter (from frontend API call) - highest priority
            # 2. segments_metadata.last_target_lang_for_language_match (from frontend language change in Extract phase)
            # 3. task_state payload (from translation task creation)
            # If target_lang is not available, skip language_match detection to prevent incorrect exclusions
            
            logger.debug(
                LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Received target_lang parameter: {target_lang} (type: {type(target_lang)})"
                )
            
            # Normalize empty string to None
            if target_lang == "":
                target_lang = None
            
            # Get segments_metadata and payload once
            segments_metadata = st.get("segments_metadata", {})
            if not segments_metadata:
                segments_metadata = {}
                st["segments_metadata"] = segments_metadata
            
            payload = st.get("payload")
            stored_target_lang = segments_metadata.get("last_target_lang_for_language_match")
            payload_target_lang = self._get_payload_target_lang(payload)
            
            # Determine target_lang and its source in priority order
            target_lang_source = None
            if target_lang:
                target_lang_source = "parameter (function argument)"
            elif stored_target_lang:
                target_lang = stored_target_lang
                target_lang_source = "stored (segments_metadata.last_target_lang_for_language_match)"
            elif payload_target_lang:
                target_lang = payload_target_lang
                target_lang_source = "payload"
            
            # Log the final target_lang and source
            if target_lang:
                logger.info(
                    LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Using target_lang={target_lang} for language_match detection "
                    f"(source: {target_lang_source})"
                    )
            else:
                logger.warning(
                    LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: target_lang is None. "
                    f"Checked sources: parameter=None, stored={stored_target_lang}, payload={payload_target_lang}. "
                    f"Skipping language_match detection."
                    )
            
            # Prepare segments and format-specific data for batch detection
            segments = [seg.get("text", "") for seg in all_segments]
            format_specific_data = {}
            for seg_idx, segment_data in enumerate(all_segments):
                chunk_idx = segment_data.get("chunk_index", seg_idx)
                if chunk_idx < len(layout_result.chunks):
                    chunk = layout_result.chunks[chunk_idx]
                    chunk_block_indices = chunk.block_indices if hasattr(chunk, 'block_indices') else []
                    format_specific_data[seg_idx] = {
                        "chunk_block_indices": chunk_block_indices,
                        "chunk_type": getattr(chunk, 'chunk_type', None),
                        "image_path": segment_data.get("image_path"),
                        "image_placeholder": segment_data.get("placeholder_id"),
                    }
            
            # Create PDF metadata extractor
            extractor = PDFMetadataExtractor(block_type_map, block_image_map)
            
            # Batch detect exclusions using new architecture
            # Note: We detect for ALL segments, not just currently excluded ones
            # This ensures we can properly classify all segments (including optional TABLE)
            # CRITICAL: Log target_lang before calling detect_exclusions_batch to verify it's correct
            logger.info(
                LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Calling ExclusionDetectionBatch.detect_exclusions_batch "
                f"with target_lang={target_lang} (type: {type(target_lang)})"
                )
            # task_id already set at start of get_layout_extract for progress visibility
            # Run exclusion detection in thread so event loop can handle getStatus polls (frontend Extract progress bar)
            excluded_segments_with_reasons, all_detected_reasons = await asyncio.to_thread(
                ExclusionDetectionBatch.detect_exclusions_batch,
                segments=segments,
                metadata_extractor=extractor,
                task_state=st,
                target_lang=target_lang,
                format_specific_data=format_specific_data,
                preserve_existing=True,
                auto_exclude_optional=False  # TABLE not auto-excluded
            )
        except Exception as e:
            logger.error(
                LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Failed to detect exclusions using new architecture: {e}",
                exc_info=True,
                )
            # Fallback: use existing excluded segments from task_state
            from exclusion.core import ExclusionManager, ExclusionReason
            excluded_segments_with_reasons = ExclusionManager.get_excluded_segments(st)
            logger.warning(
                LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Using fallback exclusion detection, found {len(excluded_segments_with_reasons)} existing excluded segments"
                )
        
        # CRITICAL: Import ExclusionReason if not already imported (for fallback case)
        if 'ExclusionReason' not in locals():
            from exclusion.core import ExclusionReason
        
        # CRITICAL: Filter out optional/user-choice exclusions from excluded_segments_with_reasons.
        # TABLE and LANGUAGE_MATCH are "default not excluded" - only store when user already excluded them.
        from exclusion.core import ExclusionManager
        pre_existing_excluded_indices = set(ExclusionManager.get_excluded_segments(st).keys())
        filtered_excluded_segments = {}
        for seg_idx, reason in excluded_segments_with_reasons.items():
            # Content-based (IMAGE, FORMULA, REFERENCE, IDENTIFIER, STRUCTURAL): always include
            if ExclusionReason.is_content_based(reason):
                filtered_excluded_segments[seg_idx] = reason
            # LANGUAGE_MATCH: only include if user already excluded (in pre_existing); do not auto-store
            elif ExclusionReason.is_language_based(reason):
                if seg_idx in pre_existing_excluded_indices:
                    filtered_excluded_segments[seg_idx] = reason
                else:
                    logger.debug(
                        LogModule.EXCLUSION,
                        f"[LAYOUT-EXTRACT] Task {task_id}: Not storing language_match for segment {seg_idx} "
                        f"(user has not excluded language match, not in pre_existing_excluded_indices)"
                    )
            elif reason == ExclusionReason.TABLE:
                # TABLE is optional - only include if it was in pre_existing_excluded_indices (user choice)
                if seg_idx in pre_existing_excluded_indices:
                    filtered_excluded_segments[seg_idx] = reason
                    logger.debug(
                        LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Including table segment {seg_idx} in excluded_segments "
                        f"(user explicitly excluded it, in pre_existing_excluded_indices)"
                        )
                else:
                    logger.debug(
                        LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Excluding table segment {seg_idx} from excluded_segments "
                        f"(optional exclusion, not auto-excluded, not in pre_existing_excluded_indices)"
                        )
            else:
                # Other reasons (USER_SELECTED, UNKNOWN) - include them
                filtered_excluded_segments[seg_idx] = reason
        
        # Store exclusions using new architecture
        if filtered_excluded_segments:
            segment_metadata = {
                idx: {"block_type": all_segments[idx].get("block_type")}
                for idx in filtered_excluded_segments.keys()
                if idx < len(all_segments)
            }
            ExclusionDetectionBatch.store_exclusions(
                task_state=st,
                excluded_segments=filtered_excluded_segments,
                segment_metadata=segment_metadata,
                source="pdf_layout_extract",
                all_detected_reasons=all_detected_reasons  # Store all detected reasons for frontend display
            )
        
        # Store chunk_size for future re-extract detection
        if chunk_size:
            segments_metadata = st.get("segments_metadata", {})
            segments_metadata["last_chunk_size"] = chunk_size
            logger.debug(
                LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Stored chunk_size={chunk_size} for re-extract detection"
                )
            
            # Refine UNKNOWN reasons using unified method
            # Include target_lang if available to detect language_match
            from exclusion.core import ExclusionManager
            segment_texts = {idx: seg.get("text", "") for idx, seg in enumerate(all_segments)}
            segment_block_types = {idx: seg.get("block_type") for idx, seg in enumerate(all_segments) if seg.get("block_type")}
            updated_reasons = ExclusionManager.refine_exclusion_reasons(
                st,
                segment_indices=[idx for idx, reason in excluded_segments_with_reasons.items() 
                               if reason == ExclusionReason.UNKNOWN],
                segment_texts=segment_texts,
                segment_block_types=segment_block_types if segment_block_types else None,
                target_lang=target_lang,  # Check language match if target_lang is available
                layout_chunk_block_map=None  # Not available in get_layout_extract
            )
            
            if updated_reasons:
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Refined {len(updated_reasons)} exclusion reasons "
                    f"from UNKNOWN to specific reasons"
                    )
            
            logger.debug(
                LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Updated segments_metadata.excluded_segments "
                f"with {len(filtered_excluded_segments)} excluded segments "
                f"(filtered from {len(excluded_segments_with_reasons)} detected exclusions, "
                f"removed {len(excluded_segments_with_reasons) - len(filtered_excluded_segments)} optional TABLE exclusions)"
                )
            
            # CRITICAL: Update all_segments with exclusion_reason and is_excluded for frontend display
            # This ensures frontend can correctly display exclusion statistics
            # CRITICAL: Use filtered_excluded_segments (not excluded_segments_with_reasons) for setting is_excluded
            # This ensures that optional exclusions (TABLE) are not automatically marked as excluded
            # Only segments in filtered_excluded_segments (content-based or user-selected) should be marked as excluded
            total_segments = len(all_segments)
            invalid_excluded_indices = [idx for idx in filtered_excluded_segments.keys() if idx >= total_segments]
            if invalid_excluded_indices:
                logger.warning(
                    LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Found {len(invalid_excluded_indices)} excluded segment indices "
                    f"that are out of range (total_segments={total_segments}): "
                    f"{sorted(invalid_excluded_indices)[:10]}{'...' if len(invalid_excluded_indices) > 10 else ''}. "
                    f"These will be filtered out to prevent frontend state inconsistencies."
                    )
                # Filter out invalid indices to prevent frontend state inconsistencies
                filtered_excluded_segments = {
                    idx: reason for idx, reason in filtered_excluded_segments.items()
                    if idx < total_segments
                }
            
            # CRITICAL: Update all_segments with exclusion_reason, detected_exclusion_reason, and is_excluded for frontend display
            # This ensures frontend can correctly display exclusion statistics and segment labels
            # CRITICAL: Use filtered_excluded_segments (not excluded_segments_with_reasons) for setting is_excluded
            # This ensures that optional exclusions (TABLE) are not automatically marked as excluded
            # Only segments in filtered_excluded_segments (content-based or user-selected) should be marked as excluded
            
            # CRITICAL: Get detected_exclusion_reasons from segments_metadata for frontend display
            # This includes all detected types (identifier, language_match, etc.) even if not excluded
            detected_exclusion_reasons_dict = segments_metadata.get("detected_exclusion_reasons", {})
            
            for seg_idx in range(len(all_segments)):
                # CRITICAL: Set detected_exclusion_reason for ALL segments (including non-excluded ones)
                # This allows frontend to display all detected types (identifier, language_match, etc.)
                # BUT: If user manually retried and successfully translated this segment, don't overlay old reason
                seg_idx_str = str(seg_idx)
                if seg_idx_str in detected_exclusion_reasons_dict:
                    detected_reason_info = detected_exclusion_reasons_dict[seg_idx_str]
                    if isinstance(detected_reason_info, dict):
                        detected_reason_str = detected_reason_info.get("reason", "")
                    else:
                        detected_reason_str = str(detected_reason_info)
                    if detected_reason_str:
                        is_user_translated = (
                            all_segments[seg_idx].get("status") == "translated"
                            and all_segments[seg_idx].get("modified") is True
                        )
                        if not is_user_translated:
                            all_segments[seg_idx]["detected_exclusion_reason"] = detected_reason_str
                
                if seg_idx in filtered_excluded_segments:
                    # Segment is excluded: set is_excluded=True and exclusion_reason
                    all_segments[seg_idx]["is_excluded"] = True
                    reason = filtered_excluded_segments[seg_idx]
                    all_segments[seg_idx]["exclusion_reason"] = reason.value
                    # Also add exclusion_metadata if available
                    block_type = all_segments[seg_idx].get("block_type")
                    if block_type:
                        all_segments[seg_idx]["exclusion_metadata"] = {
                            "block_type": block_type
                        }
                else:
                    # Segment is NOT excluded: ensure is_excluded=False and clear exclusion_reason
                    # This ensures consistency with excluded_segments (single source of truth)
                    all_segments[seg_idx]["is_excluded"] = False
                    if "exclusion_reason" in all_segments[seg_idx]:
                        del all_segments[seg_idx]["exclusion_reason"]
                    if "exclusion_metadata" in all_segments[seg_idx]:
                        del all_segments[seg_idx]["exclusion_metadata"]
        
        # CRITICAL: Update layout_prepared_chunks with ALL chunks (including excluded)
        # This ensures MDTranslator can correctly map chunks to original segments
        # Each chunk in layout_prepared_chunks has segment_indices that reference original all_segments indices
        serialized_chunks = []
        chunk_block_map = []
        chunk_block_texts_map = []
        for chunk in all_chunks_with_excluded:
            chunk_text = chunk.get("text", "")
            chunk_segment_indices = chunk.get("segment_indices", [])
            
            # Build block_indices and block_texts from segment indices
            block_indices = []
            block_texts = chunk.get("block_texts") or []
            for seg_idx in chunk_segment_indices:
                if seg_idx < len(all_segments):
                    seg = all_segments[seg_idx]
                    seg_block_indices = list(seg.get("layout_block_indices") or [])
                    if not seg_block_indices and seg.get("block_index") is not None:
                        try:
                            seg_block_indices = [int(seg["block_index"])]
                        except (TypeError, ValueError):
                            pass
                    for bidx in seg_block_indices:
                        if bidx not in block_indices:
                            block_indices.append(bidx)
            if not block_texts:
                for seg_idx in chunk_segment_indices:
                    if seg_idx < len(all_segments):
                        seg_text = all_segments[seg_idx].get("text", "")
                        if seg_text:
                            block_texts.append(seg_text)
            
            # Estimate tokens for this chunk
            estimated_tokens = None
            if chunk_text.strip():  # Only estimate for non-empty chunks
                from utils.token_estimator import estimate_chunk_input_tokens
                chunk_dict = {str(i): seg for i, seg in enumerate(chunk_text.split("\n\n"))}
                chunk_json = json.dumps(chunk_dict, ensure_ascii=False)
                estimated_tokens = estimate_chunk_input_tokens(chunk_json, system_prompt_approx=2000)
            
            # Determine chunk type based on whether it's excluded
            is_chunk_excluded = chunk.get("is_excluded", False)
            chunk_type = "text"
            
            # CRITICAL: Check for chart_body type from original layout_result.chunks
            # This ensures chart segments are properly identified even when rendered as images
            original_chunk_type = None
            if layout_result and layout_result.chunks and len(layout_result.chunks) > len(serialized_chunks):
                # Try to find matching chunk by segment indices
                for orig_chunk in layout_result.chunks:
                    orig_seg_indices = orig_chunk.block_indices if hasattr(orig_chunk, 'block_indices') else []
                    if orig_seg_indices == chunk_segment_indices:
                        original_chunk_type = orig_chunk.chunk_type if hasattr(orig_chunk, 'chunk_type') else None
                        break
            
            if original_chunk_type == "chart_body":
                chunk_type = "chart_body"  # Preserve chart_body type for exclusion detection
            elif is_chunk_excluded:
                # Check if this excluded chunk is an image segment
                if chunk_segment_indices and chunk_segment_indices[0] < len(all_segments):
                    seg = all_segments[chunk_segment_indices[0]]
                    if seg.get("is_image", False):
                        chunk_type = "image"
            
            chunk_info = {
                "text": chunk_text,
                "chunk_type": chunk_type,
                "block_indices": block_indices,
                "block_texts": block_texts,
                "segment_indices": chunk_segment_indices,  # CRITICAL: Store original segment indices for proper mapping
                "is_image": chunk_type == "image",
                "is_excluded": is_chunk_excluded,  # CRITICAL: Mark excluded chunks for MDTranslator
            }
            
            # CRITICAL: Log excluded chunks for debugging
            if is_chunk_excluded:
                logger.trace(
                    LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Chunk {len(serialized_chunks)} is marked as excluded "
                    f"(segment_indices={chunk_segment_indices}, chunk_type={chunk_type}, "
                    f"text_preview={chunk_text[:50] if chunk_text else 'empty'}...)"
                    )
            if estimated_tokens is not None:
                chunk_info["estimated_input_tokens"] = estimated_tokens
            
            serialized_chunks.append(chunk_info)
            chunk_block_map.append(block_indices)
            chunk_block_texts_map.append(block_texts)
        
        # Also update total_estimated_input_tokens in task_state
        total_estimated_tokens_updated = sum(
            chunk.get("estimated_input_tokens", 0)
            for chunk in serialized_chunks
            if isinstance(chunk, dict) and not chunk.get("is_image", False)
        )
        if total_estimated_tokens_updated > 0:
            st["total_estimated_input_tokens"] = total_estimated_tokens_updated
            logger.debug(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: Updated total_estimated_input_tokens={total_estimated_tokens_updated} in task_state")
        
        # Update layout_prepared_chunks in task_state
        # This ensures translation phase uses chunks that exclude excluded segments
        st["layout_prepared_chunks"] = serialized_chunks
        st["layout_chunk_block_map"] = chunk_block_map
        st["layout_chunk_block_texts"] = chunk_block_texts_map
        from utils.translation_segments import build_segment_layout_block_map
        st["segment_layout_block_map"] = build_segment_layout_block_map(all_segments)
        logger.info(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: Updated layout_prepared_chunks with {len(serialized_chunks)} chunks (excluding {excluded_count} excluded segments)")
        
        # Build chunks text (merged segments, excluding headers and footers)
        # CRITICAL: chunks_text should be a LIST of strings (one per chunk), not a single joined string
        # This matches the format expected by frontend
        chunks_text = []
        for chunk in all_chunks:
            chunk_text = chunk.get("text", "")
            # Check if any segment in this chunk is header/footer
            chunk_segment_indices = chunk.get("segment_indices", [])
            is_header_chunk = False
            is_footer_chunk = False
            for seg_idx in chunk_segment_indices:
                if seg_idx < len(all_segments):
                    seg = all_segments[seg_idx]
                    if seg.get("is_header", False):
                        is_header_chunk = True
                    if seg.get("is_footer", False):
                        is_footer_chunk = True
            
            # Include all chunks except headers and footers
            if not is_header_chunk and not is_footer_chunk and chunk_text.strip():
                chunks_text.append(chunk_text)
        
        # Close ZIP file if opened
        if zip_file:
            try:
                zip_file.close()
            except Exception:
                pass
        
        # Build image_data_map in format expected by frontend: {placeholder_id: {"data": "...", "alt": "..."}}
        image_data_map_for_frontend: dict[str, dict[str, str]] = {}
        
        # Step 1: Add regular image segments (tables, images, etc.)
        for seg in all_segments:
            if seg.get("is_image") and seg.get("placeholder_id") and seg.get("image_data"):
                placeholder_id = seg["placeholder_id"]
                image_path = seg.get("image_path", "Image")
                image_data_map_for_frontend[placeholder_id] = {
                    "data": seg["image_data"],
                    "alt": image_path,
                }
        
        # Step 2: Add equation images from segments containing ![Equation](filename.jpg) format
        # This handles interline_equation blocks that are rendered as images
        equation_image_pattern = re.compile(r'!\[([^\]]*)\]\(([^\)]+\.(jpg|jpeg|png|gif|webp))\)')
        equation_counter = 0
        for seg in all_segments:
            segment_text = seg.get("text", "")
            if not segment_text:
                continue
            
            # Check if segment text contains equation image references
            matches = equation_image_pattern.findall(segment_text)
            for match in matches:
                alt_text = match[0] or "Equation"
                filename = match[1]
                
                # Normalize filename (remove ./ prefix and path)
                normalized_filename = filename.replace("\\", "/").lstrip("./").split("/")[-1]
                
                # Try to find image data in image_data_by_path
                image_data = None
                for key, value in image_data_by_path.items():
                    # Match by filename (end of path or exact match)
                    key_normalized = key.replace("\\", "/").lstrip("./").split("/")[-1]
                    if key_normalized == normalized_filename or key.endswith(normalized_filename) or normalized_filename.endswith(key_normalized):
                        image_data = value
                        break
                
                if image_data:
                    # Generate a unique placeholder ID for this equation image
                    equation_counter += 1
                    placeholder_id = f"equation_{equation_counter}_{normalized_filename}"
                    image_data_map_for_frontend[placeholder_id] = {
                        "data": image_data,
                        "alt": filename,  # Use full filename as alt text
                    }
                    logger.debug(
                        LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Added equation image to image_data_map: "
                        f"placeholder={placeholder_id}, filename={filename}, normalized={normalized_filename}"
                        )
                else:
                    logger.warning(
                        LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Equation image not found in image_data_by_path: "
                        f"filename={filename}, normalized={normalized_filename}, "
                        f"available_keys={list(image_data_by_path.keys())[:5]}"
                        )
        
        # Calculate total estimated input tokens for all chunks
        total_estimated_tokens = sum(
            chunk.get("estimated_input_tokens", 0) 
            for chunk in all_chunks 
            if isinstance(chunk, dict)
        )
        
        # Log detailed token calculation for debugging
        if all_chunks:
            chunk_tokens_detail = [
                (i, chunk.get("estimated_input_tokens", 0) if isinstance(chunk, dict) else 0)
                for i, chunk in enumerate(all_chunks)
            ]
            chunks_with_tokens = [ct for ct in chunk_tokens_detail if ct[1] > 0]
            chunks_without_tokens = [ct for ct in chunk_tokens_detail if ct[1] == 0]
            logger.debug(
                LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Token calculation - "
                f"total_chunks={len(all_chunks)}, "
                f"chunks_with_tokens={len(chunks_with_tokens)}, "
                f"chunks_without_tokens={len(chunks_without_tokens)}, "
                f"total_estimated_tokens={total_estimated_tokens}"
                )
            if chunks_with_tokens:
                logger.debug(
                    LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Sample chunks with tokens (first 5): "
                    f"{chunks_with_tokens[:5]}"
                    )
            if chunks_without_tokens and len(chunks_without_tokens) <= 10:
                logger.debug(
                    LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Chunks without tokens: {chunks_without_tokens}"
                    )
        
        # Build response
        from layout.pdf_renderer.typst_overlay.segment_font_metrics import (
            enrich_segments_font_fields,
        )
        font_override_by_index: Dict[int, float] = {}
        ts_data = st.get("translation_segments")
        if isinstance(ts_data, dict):
            for stored in ts_data.get("segments") or []:
                if not isinstance(stored, dict):
                    continue
                seg_idx = stored.get("segment_index")
                if seg_idx is None:
                    continue
                if stored.get("font_size_pt") is not None:
                    font_override_by_index[int(seg_idx)] = stored.get("font_size_pt")
        for seg in all_segments:
            if not isinstance(seg, dict):
                continue
            seg_idx = seg.get("segment_index", seg.get("chunk_index"))
            if seg_idx is not None and int(seg_idx) in font_override_by_index:
                seg["font_size_pt"] = font_override_by_index[int(seg_idx)]
        enrich_segments_font_fields(
            layout_doc, all_segments, text_field="text", task_state=st,
        )

        response = {
            "task_id": task_id,
            "segments": all_segments,  # Deep split segments (for left panel - Segments)
            "chunks": all_chunks,  # Merged chunks (for right panel - Chunks)
            "chunks_text": chunks_text,  # Merged chunks text (without headers/footers, for display)
            "total_segments": len(all_segments),
            "total_chunks": len(all_chunks),
            "image_data_map": image_data_map_for_frontend,  # Map of placeholder_id to image data
            "ready": True,
            # chunk_size removed from response - frontend should use global settings instead
        }
        
        # Add total estimated input tokens if available
        if total_estimated_tokens > 0:
            response["total_estimated_input_tokens"] = total_estimated_tokens
            logger.info(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: Added total_estimated_input_tokens={total_estimated_tokens} to response")
        else:
            logger.warning(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: total_estimated_tokens is 0 or negative, NOT adding to response")
        
        logger.info(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: Returning {len(all_segments)} segments, {len(all_chunks)} chunks, {len(image_data_map_for_frontend)} images, chunk_size={chunk_size}")
        
        # Update segments_metadata with excluded_segment_indices for translation
        # CRITICAL: Merge with pre-existing excluded_segment_indices to preserve user choices (e.g., references)
        # This allows frontend to mark segments as excluded (e.g., references) and have them
        # automatically excluded when regenerating chunks
        if "segments_metadata" not in st:
            st["segments_metadata"] = {}
        
        # Merge pre-existing excluded_segment_indices with newly detected ones
        existing_excluded = set(st["segments_metadata"].get("excluded_segment_indices", []))
        newly_detected_excluded = set(excluded_segment_indices_list)
        
        # CRITICAL: Remove user_unexcluded_segments from both existing_excluded and newly_detected_excluded
        # This ensures that segments explicitly unexcluded by user are not re-added even if they are detected again
        user_unexcluded_segments = st.get("segments_metadata", {}).get("user_unexcluded_segments", [])
        if user_unexcluded_segments:
            user_unexcluded_set = set(user_unexcluded_segments)
            existing_excluded = existing_excluded.difference(user_unexcluded_set)
            newly_detected_excluded = newly_detected_excluded.difference(user_unexcluded_set)
            logger.info(
                LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Removed {len(user_unexcluded_segments)} user-unexcluded segments "
                f"from existing_excluded and newly_detected_excluded: {sorted(user_unexcluded_segments)}"
                )
        
        # If clear_user_exclusions is True, remove ref_text blocks from existing_excluded
        if clear_user_exclusions:
            # Filter out ref_text blocks from existing_excluded
            filtered_existing = set()
            for idx in existing_excluded:
                if idx < len(all_segments):
                    block_type = all_segments[idx].get('block_type')
                    if block_type != 'ref_text':
                        filtered_existing.add(idx)
            existing_excluded = filtered_existing
            logger.info(LogModule.EXCLUSION, f"[LAYOUT-EXTRACT] Task {task_id}: Filtered out ref_text blocks from existing excluded indices (clear_user_exclusions=True)")
        
        merged_excluded = sorted(list(existing_excluded.union(newly_detected_excluded)))
        
        if merged_excluded:
            # CRITICAL: Use excluded_segments (dict) as single source of truth. Only persist indices that have a reason.
            # Indices in merged_excluded but not in the dict are stale (e.g. language_match cleared on language switch
            # but list was not fully cleared). Do NOT add them to the list or mark as USER_SELECTED.
            current_excluded_dict = st["segments_metadata"].get("excluded_segments", {})
            current_dict_indices = {int(k) for k in current_excluded_dict.keys()}
            synced_list = sorted([i for i in merged_excluded if i in current_dict_indices])
            dropped_count = len(merged_excluded) - len(synced_list)
            if dropped_count > 0:
                logger.info(
                    LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Synced excluded_segment_indices to excluded_segments dict: "
                    f"dropped {dropped_count} indices without reason (e.g. cleared on language switch). "
                    f"list length {len(merged_excluded)} -> {len(synced_list)}"
                )
            st["segments_metadata"]["excluded_segment_indices"] = synced_list
            equation_excluded_count = sum(1 for idx in synced_list if idx < len(all_segments) and all_segments[idx].get('block_type') == 'interline_equation')
            ref_excluded_count = sum(1 for idx in synced_list if idx < len(all_segments) and all_segments[idx].get('block_type') == 'ref_text')
            newly_added_count = len(newly_detected_excluded.difference(existing_excluded))
            logger.info(
                LogModule.EXCLUSION,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Stored {len(synced_list)} excluded segment indices in segments_metadata "
                f"(including {equation_excluded_count} interline_equation blocks, {ref_excluded_count} ref_text blocks, "
                f"{newly_added_count} newly detected). Excluded indices: {sorted(synced_list)[:30]}{'...' if len(synced_list) > 30 else ''}"
                )
        
        # CRITICAL: Update source_chunks_cache with all_segments to ensure correct indexing
        # This ensures that source_segments in record_translation_segments matches actual_segment_index
        # source_chunks_cache["segments"] should be a list of strings, where segments[i] corresponds to segment_index=i
        # all_segments should already be ordered by segment_index (0, 1, 2, ...)
        # Build cache_segments by extracting text from all_segments in order
        cache_segments = []
        max_segment_index = -1
        for seg in all_segments:
            seg_idx = seg.get("segment_index", len(cache_segments))
            seg_text = seg.get("text", "")
            # Ensure cache_segments has enough elements (handle non-contiguous indices)
            while len(cache_segments) <= seg_idx:
                cache_segments.append("")
            cache_segments[seg_idx] = seg_text
            if seg_idx > max_segment_index:
                max_segment_index = seg_idx
        
        # Trim empty trailing elements (if any segments were skipped)
        # But keep all valid segments up to max_segment_index
        cache_segments = cache_segments[:max_segment_index + 1] if max_segment_index >= 0 else []
        
        # CRITICAL: Check if existing source_chunks_cache has more segments than we're about to store
        # This can happen if source_chunks_cache was populated by prepare_source_preview_for_extractor_based
        # (which stores ALL segments) but get_layout_extract only processes a subset (e.g., from layout_result.chunks)
        existing_cache = st.get("source_chunks_cache", {})
        existing_segments = existing_cache.get("segments", [])
        existing_total = existing_cache.get("total_segments")
        
        if existing_segments and existing_total is not None:
            if len(existing_segments) > len(cache_segments) and existing_total > len(cache_segments):
                logger.warning(
                    LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: WARNING - Existing source_chunks_cache has {len(existing_segments)} segments "
                    f"(total_segments={existing_total}), but get_layout_extract only generated {len(cache_segments)} segments. "
                    f"This may indicate that layout_result.chunks is incomplete. Preserving existing cache to avoid data loss."
                    )
                # Preserve existing cache instead of overwriting with incomplete data
                # This ensures that get_source_preview can still access all segments
                # Only update chunk_size and created_at
                st["source_chunks_cache"] = {
                    "content_hash": existing_cache.get("content_hash", ""),
                    "chunk_size": chunk_size,  # Update chunk_size
                    "segments": existing_segments,  # Preserve existing segments
                    "total_segments": existing_total,  # Preserve existing total
                    "created_at": existing_cache.get("created_at", time.time()),
                }
                logger.info(
                    LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Preserved existing source_chunks_cache with {len(existing_segments)} segments "
                    f"instead of overwriting with {len(cache_segments)} segments from layout_result.chunks"
                    )
            else:
                # Update source_chunks_cache normally
                st["source_chunks_cache"] = {
                    "content_hash": existing_cache.get("content_hash", ""),
                    "chunk_size": chunk_size,
                    "segments": cache_segments,  # Indexed by segment_index, matches all_segments
                    "total_segments": len(cache_segments),
                    "created_at": time.time(),
                }
        else:
            # No existing cache, create new one
            st["source_chunks_cache"] = {
                "content_hash": existing_cache.get("content_hash", ""),
                "chunk_size": chunk_size,
                "segments": cache_segments,  # Indexed by segment_index, matches all_segments
                "total_segments": len(cache_segments),
                "created_at": time.time(),
            }

        # Write Extract-phase segments to temp dir in same JSON format as LLM input for diagnosis
        try:
            from utils.extract_segments_debug import write_extract_segments_json
            written = write_extract_segments_json(
                st.get("temp_dir"), cache_segments, task_id=task_id
            )
            if written:
                logger.debug(
                    LogModule.EXTRACT,
                    f"[LAYOUT-EXTRACT] Task {task_id}: Wrote {len(cache_segments)} segments to {written}"
                )
        except Exception as _e:
            logger.debug(LogModule.EXTRACT, f"[LAYOUT-EXTRACT] Task {task_id}: Failed to write extract_segments.json: {_e}")

        return JSONResponse(content=response)
    
    async def update_excluded_segments_for_language(
        self,
        task_id: str,
        target_lang: str,
        auto_exclude: bool = False
    ):
        """
        Update excluded segments based on new target language.
        Re-detects language for all segments and updates excluded status.
        
        Args:
            task_id: Unique task identifier
            target_lang: New target language code (e.g., 'zh', 'en')
            
        Returns:
            JSONResponse with updated excluded segment indices
            
        Raises:
            HTTPException: If task not found
        """
        task_state = self.task_manager.get_task(task_id)
        if task_state is None:
            raise HTTPException(status_code=404, detail=f"Task ID '{task_id}' not found.")
        
        st = task_state
        
        # Check if task has failed
        task_status = st.get("status")
        if task_status == "failed":
            error_message = st.get("error") or st.get("message") or "Task processing failed"
            raise HTTPException(
                status_code=400,
                detail=f"Task '{task_id}' has failed and cannot update excluded segments. Error: {error_message}"
            )
        
        # Get all segments from cache or preview
        cache_info = st.get("source_chunks_cache", {})
        all_segments = cache_info.get("segments", [])
        
        # Fallback to preview segments if cache not available
        if not all_segments:
            preview = st.get("source_preview") or {"segments": [], "total_segments": 0, "ready": False}
            all_segments = preview.get("segments", [])
        
        # CRITICAL: In Extract phase, segments may not be ready yet (e.g., MinerU processing)
        # If segments are not ready, just store the target language for later use
        # This allows MinerU to use the latest target language when processing completes
        if not all_segments:
            # Check if we're in Extract phase (no translation_segments yet)
            translation_segments_data = st.get("translation_segments")
            is_translate_phase = bool(translation_segments_data and isinstance(translation_segments_data, dict))
            
            if not is_translate_phase:
                # Extract phase - segments not ready yet, just store target language
                segments_metadata = st.get("segments_metadata", {})
                if not segments_metadata:
                    segments_metadata = {}
                    st["segments_metadata"] = segments_metadata
                
                # Store target language for MinerU processing
                segments_metadata["last_target_lang_for_language_match"] = target_lang
                logger.info(
                    LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Extract phase - segments not ready yet. "
                    f"Stored target_lang={target_lang} for MinerU processing. "
                    f"Language Match detection will be performed when segments are ready."
                    )
                
                # Return success response indicating target language was stored
                return JSONResponse(content={
                    "task_id": task_id,
                    "target_lang": target_lang,
                    "excluded_segment_indices": [],
                    "total_segments": 0,
                    "excluded_count": 0,
                    "language_matched_count": 0,
                    "message": f"Target language '{target_lang}' stored for MinerU processing. "
                               f"Segments not ready yet - Language Match detection will be performed when extraction completes."
                })
            else:
                # Translate phase - segments should be available
                raise HTTPException(
                    status_code=404,
                    detail="No segments found. Please ensure the task has completed extraction."
                )
        
        # CRITICAL: Check if we're in Translate phase
        # Translate phase: translation_segments exists (translation has started)
        # Extract phase: translation_segments does not exist (only extraction completed)
        translation_segments_data = st.get("translation_segments")
        is_translate_phase = bool(translation_segments_data and isinstance(translation_segments_data, dict))
        
        # Import unified exclusion detection utility
        from utils.translation_segments import _is_image_segment
        from exclusion.core import ExclusionManager, ExclusionReason, detect_exclusion_reason
        
        # CRITICAL: Use ExclusionManager to get current excluded segments with their reasons
        # This provides unified exclusion management and preserves exclusion reasons
        current_excluded_segments = ExclusionManager.get_excluded_segments(st)
        current_excluded_indices = set(current_excluded_segments.keys())
        
        # Get segments_metadata for backward compatibility
        segments_metadata = st.get("segments_metadata", {})
        if not segments_metadata:
            segments_metadata = {}
            st["segments_metadata"] = segments_metadata
        
        # CRITICAL: Check if target language has changed
        # If target language hasn't changed, skip re-detection to avoid unnecessary work
        stored_target_lang = segments_metadata.get("last_target_lang_for_language_match")
        if stored_target_lang == target_lang:
            logger.info(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Target language '{target_lang}' has not changed "
                f"(stored: '{stored_target_lang}'). auto_exclude={auto_exclude}"
                )
            
            # CRITICAL: If auto_exclude=False, we need to clear language_match exclusions
            # even if target_lang hasn't changed (user explicitly unchecking the checkbox)
            if not auto_exclude:
                logger.info(
                    LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: auto_exclude=False, clearing language_match exclusions "
                    f"even though target_lang unchanged"
                    )
                # Clear language_match exclusions, keep content-based and user-based exclusions
                from exclusion.core import ExclusionReason
                cleared_excluded_segments = {}
                language_match_cleared_count = 0
                for idx, reason in current_excluded_segments.items():
                    if ExclusionReason.is_language_based(reason):
                        language_match_cleared_count += 1
                        # Skip language_match exclusions
                    else:
                        # Keep content-based and user-based exclusions
                        cleared_excluded_segments[idx] = reason
                
                if language_match_cleared_count > 0:
                    logger.info(
                        LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Cleared {language_match_cleared_count} language_match exclusions. "
                        f"Remaining excluded: {len(cleared_excluded_segments)} (content-based and user-based only)"
                        )
                    # Update exclusion data using ExclusionManager
                    excluded_segments_for_manager = {}
                    for idx, reason in cleared_excluded_segments.items():
                        if isinstance(reason, ExclusionReason):
                            excluded_segments_for_manager[idx] = reason
                        else:
                            try:
                                excluded_segments_for_manager[idx] = ExclusionReason(reason)
                            except (ValueError, TypeError):
                                logger.warning(
                                    LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Invalid exclusion reason '{reason}' for segment {idx}, skipping"
                                    )
                                continue
                    
                    ExclusionManager.update_excluded_segments(st, excluded_segments_for_manager)
                    excluded_segment_indices = sorted(cleared_excluded_segments.keys())
                else:
                    # No language_match exclusions to clear, return current state
                    excluded_segment_indices = sorted(current_excluded_segments.keys())
            else:
                # auto_exclude=True, return current excluded segments without re-detection
                excluded_segment_indices = sorted(current_excluded_segments.keys())
                logger.info(
                    LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: auto_exclude=True, returning current excluded segments "
                    f"without re-detection"
                    )
            
            return JSONResponse(content={
                "task_id": task_id,
                "target_lang": target_lang,
                "excluded_segment_indices": excluded_segment_indices,
                "total_segments": len(all_segments),
                "excluded_count": len(excluded_segment_indices),
                "message": f"Target language unchanged, {'cleared language_match exclusions' if not auto_exclude else 'no re-detection performed'}"
            })
        
        # Target language has changed - proceed with re-detection
        logger.info(
            LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Target language changed from '{stored_target_lang}' to '{target_lang}'. "
            f"Proceeding with Language Match re-detection."
            )
        
        # Store new target language for future comparison
        segments_metadata["last_target_lang_for_language_match"] = target_lang
        
        if is_translate_phase:
            # CRITICAL: Translate phase - Only re-detect Language Match, preserve all content-based exclusions (including Identifier)
            # Identifier has higher priority than Language Match
            logger.info(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Translate phase detected. "
                f"Re-detecting Language Match exclusions only, preserving all content-based exclusions (including Identifier)."
                )
            
            # Step 1: Identify and preserve content-based exclusions (including Identifier)
            # Content-based exclusions should always be excluded regardless of target language
            content_based_excluded = {}
            language_based_excluded = {}
            user_based_excluded = {}
            
            # CRITICAL: Clear existing Language Match exclusions when target language changes
            # This ensures old Language Match markers are removed before re-detection
            language_match_cleared_count = 0
            for idx, reason in current_excluded_segments.items():
                if ExclusionReason.is_content_based(reason):
                    content_based_excluded[idx] = reason
                elif ExclusionReason.is_language_based(reason):
                    # Clear Language Match exclusions - they will be re-detected with new target language
                    language_match_cleared_count += 1
                    logger.debug(
                        LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Clearing Language Match exclusion for segment {idx} "
                        f"(target language changed from '{stored_target_lang}' to '{target_lang}')"
                        )
                elif ExclusionReason.is_user_based(reason):
                    user_based_excluded[idx] = reason
            
            if language_match_cleared_count > 0:
                logger.info(
                    LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Cleared {language_match_cleared_count} Language Match exclusions "
                    f"before re-detection with new target language '{target_lang}'"
                    )
            
            # Step 2: Re-detect Language Match ONLY for segments that are NOT content-based excluded
            # This ensures Identifier and other content-based exclusions are preserved
            final_excluded_segments = {}
            all_detected_reasons = {}  # segment_index -> ExclusionReason
            
            # First, preserve all content-based and user-based exclusions
            final_excluded_segments.update(content_based_excluded)
            final_excluded_segments.update(user_based_excluded)
            
            # Get user_unexcluded_segments to prevent re-detection of user-unexcluded segments
            user_unexcluded_segments = segments_metadata.get("user_unexcluded_segments", [])
            user_unexcluded_set = set(user_unexcluded_segments) if user_unexcluded_segments else set()
            
            # Re-detect Language Match for segments that are NOT content-based excluded
            # CRITICAL: Also allow detection of new Identifier segments (even if previously Language Match)
            # Identifier has higher priority than Language Match
            
            # Prepare segments for concurrent processing
            segments_to_process = []
            segment_index_map = {}  # Map from list index to original segment index
            for idx, seg in enumerate(all_segments):
                # Skip if already excluded as content-based (including Identifier)
                if idx in content_based_excluded:
                    all_detected_reasons[idx] = content_based_excluded[idx]
                    continue
                
                # Skip if already excluded as user-based
                if idx in user_based_excluded:
                    continue
                
                # Skip if user explicitly unexcluded
                if idx in user_unexcluded_set:
                    continue
                
                # Store segment and map list index to original segment index
                list_idx = len(segments_to_process)
                segments_to_process.append(seg)
                segment_index_map[list_idx] = idx
            
            # Process segments concurrently with progress updates
            def process_segment_detection(list_idx: int, seg: Any) -> Tuple[int, Optional[Tuple[ExclusionReason, dict]]]:
                """Process a single segment for exclusion detection."""
                # Get original segment index from map
                idx = segment_index_map.get(list_idx, list_idx)
                
                # Handle both string and dict formats
                if isinstance(seg, dict):
                    seg_text = seg.get("text", str(seg.get("source_text", "")))
                    block_type = seg.get("block_type")
                    is_table_body = seg.get("is_table_body", False)
                else:
                    seg_text = str(seg)
                    block_type = None
                    is_table_body = False
                
                # Check if this is a table segment
                is_table = (block_type == "table_body" or block_type == "table" or is_table_body)
                if not is_table:
                    from utils.translation_segments import _is_table_segment
                    is_table = _is_table_segment(seg_text)
                
                # CRITICAL: Detect exclusions, but prioritize Identifier over Language Match
                # If Identifier is detected, it takes priority and Language Match is ignored
                detected_result = detect_exclusion_reason(
                    text=seg_text,
                    block_type=block_type,
                    target_lang=target_lang,
                    is_image=False,
                    is_table=is_table
                )
                
                return (idx, detected_result)
            
            # Process segments concurrently (50% progress for detection phase)
            detection_results = self._process_segments_concurrently(
                segments_to_process,
                process_segment_detection,
                task_id,
                "detecting_exclusions",
                base_progress=0,
                progress_range=50,
                progress_interval=500
            )
            
            # Process detection results
            for idx, detected_result in detection_results.items():
                if detected_result:
                    detected_reason, detected_metadata = detected_result
                    
                    # Store all detected reasons for frontend display
                    all_detected_reasons[idx] = detected_reason
                    
                    # CRITICAL: Identifier has higher priority than Language Match
                    # If Identifier is detected, add it and skip Language Match
                    if ExclusionReason.is_content_based(detected_reason):
                        # Content-based exclusion detected (e.g., Identifier) - add it with highest priority
                        final_excluded_segments[idx] = detected_reason
                        logger.debug(
                            LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Translate phase - Segment {idx} detected as content-based "
                            f"(reason={detected_reason.value}), adding with highest priority"
                            )
                    elif ExclusionReason.is_language_based(detected_reason):
                        # Language Match detected - add to excluded segments
                        # But only if it's not already excluded as content-based (Identifier has priority)
                        final_excluded_segments[idx] = detected_reason
                        # NOTE: Removed verbose debug log for each detected language match to reduce log noise
                    # Other exclusion types (TABLE, etc.) are handled by auto_exclude flag
            
            # Step 3: Update exclusion data using ExclusionManager
            # CRITICAL: ExclusionManager.update_excluded_segments expects Dict[int, ExclusionReason]
            # Convert final_excluded_segments to the format expected by ExclusionManager
            excluded_segments_for_manager = {}
            for idx, reason in final_excluded_segments.items():
                if isinstance(reason, ExclusionReason):
                    excluded_segments_for_manager[idx] = reason
                else:
                    # If reason is a string, try to convert it to ExclusionReason
                    try:
                        excluded_segments_for_manager[idx] = ExclusionReason(reason)
                    except (ValueError, TypeError):
                        # If conversion fails, skip this segment
                        logger.warning(
                            LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Invalid exclusion reason '{reason}' for segment {idx}, skipping"
                            )
                        continue
            
            # Update using ExclusionManager (expects Dict[int, ExclusionReason])
            ExclusionManager.update_excluded_segments(st, excluded_segments_for_manager)
            
            # Also update detected_exclusion_reasons for frontend display
            detected_reasons_dict = {str(idx): reason.value if isinstance(reason, ExclusionReason) else str(reason)
                                    for idx, reason in all_detected_reasons.items()}
            segments_metadata["detected_exclusion_reasons"] = detected_reasons_dict
            
            # CRITICAL: excluded_segments is already updated by ExclusionManager.update_excluded_segments
            # No need to manually update segments_metadata["excluded_segments"] here
            
            # Count by reason type for frontend display
            reason_counts = {}
            for idx, reason in final_excluded_segments.items():
                reason_str = reason.value if isinstance(reason, ExclusionReason) else str(reason)
                reason_counts[reason_str] = reason_counts.get(reason_str, 0) + 1
            
            excluded_indices = sorted(final_excluded_segments.keys())
            
            logger.info(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Translate phase - Language Match re-detection completed. "
                f"Preserved {len(content_based_excluded)} content-based exclusions (including Identifier), "
                f"detected {reason_counts.get('language_match', 0)} Language Match exclusions."
                )
            
            return JSONResponse({
                "excluded_segment_indices": excluded_indices,
                "excluded_count": len(excluded_indices),
                "reason_counts": reason_counts,
                "language_matched_count": reason_counts.get("language_match", 0),
                "message": "Translate phase: Re-detected Language Match exclusions, preserved all content-based exclusions (including Identifier)."
            })
        
        # Import unified exclusion detection utility
        from utils.translation_segments import _is_image_segment
        from exclusion.core import ExclusionManager, ExclusionReason, detect_exclusion_reason
        
        # CRITICAL: Use ExclusionManager to get current excluded segments with their reasons
        # This provides unified exclusion management and preserves exclusion reasons
        current_excluded_segments = ExclusionManager.get_excluded_segments(st)
        current_excluded_indices = set(current_excluded_segments.keys())
        
        # Get segments_metadata for backward compatibility
        segments_metadata = st.get("segments_metadata", {})
        if not segments_metadata:
            segments_metadata = {}
            st["segments_metadata"] = segments_metadata
        
        # Debug: Log current excluded indices for troubleshooting
        logger.debug(
            LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Current excluded indices from segments_metadata: {sorted(current_excluded_indices)} "
            f"(total segments: {len(all_segments)})"
        )
        # Debug: Log sample segments to understand format
        if all_segments:
            sample_seg = all_segments[0] if len(all_segments) > 0 else None
            logger.debug(
                LogModule.EXCLUSION,
                f"[UPDATE-EXCLUDED] Task {task_id}: Sample segment format: type={type(sample_seg)}, "
                f"value={repr(str(sample_seg)[:50]) if sample_seg else 'None'}"
            )
        
        # Step 2: Identify content-based exclusions (these should be preserved)
        # Content-based exclusions should always be excluded regardless of target language
        content_based_excluded = {}
        language_based_excluded = {}
        user_based_excluded = {}
        
        # CRITICAL: Clear existing Language Match exclusions when target language changes
        # This ensures old Language Match markers are removed before re-detection
        language_match_cleared_count = 0
        for idx, reason in current_excluded_segments.items():
            if ExclusionReason.is_content_based(reason):
                content_based_excluded[idx] = reason
            elif ExclusionReason.is_language_based(reason):
                # Clear Language Match exclusions - they will be re-detected with new target language
                language_match_cleared_count += 1
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Clearing Language Match exclusion for segment {idx} "
                    f"(target language changed from '{stored_target_lang}' to '{target_lang}')"
                    )
            elif ExclusionReason.is_user_based(reason):
                user_based_excluded[idx] = reason
            else:
                # UNKNOWN reason - try to detect actual reason
                if idx < len(all_segments):
                    seg = all_segments[idx]
                    if isinstance(seg, dict):
                        seg_text = seg.get("text", str(seg.get("source_text", "")))
                    else:
                        seg_text = str(seg)
                    
                    detected_result = detect_exclusion_reason(
                        text=seg_text,
                        block_type=None,
                        target_lang=None,  # Check without language to determine if content-based
                        is_image=False
                    )
                    if detected_result:
                        detected_reason, _ = detected_result
                        if ExclusionReason.is_content_based(detected_reason):
                            content_based_excluded[idx] = detected_reason
                        else:
                            # Keep as UNKNOWN for now
                            content_based_excluded[idx] = reason
                    else:
                        # Not excluded without language check, treat as language-based
                        language_based_excluded[idx] = reason
        
        if language_match_cleared_count > 0:
            logger.info(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Cleared {language_match_cleared_count} Language Match exclusions "
                f"before re-detection with new target language '{target_lang}'"
                )
        
        # Step 3: Re-evaluate ALL segments with new target language using detect_exclusion_reason
        # This will correctly identify both content-based and language-based exclusions
        final_excluded_segments = {}
        exclusion_metadata_dict = {}
        language_matched_segments = []  # Segments that match target language (need user confirmation)
        # CRITICAL: Also collect all detected reasons (including non-excluded) for frontend display
        all_detected_reasons = {}  # segment_index -> ExclusionReason
        
        # First, preserve all content-based and user-based exclusions
        final_excluded_segments.update(content_based_excluded)
        final_excluded_segments.update(user_based_excluded)
        
        # CRITICAL: Also add content-based exclusions to all_detected_reasons for frontend display
        # This ensures frontend can show identifier, formula, etc. labels even after target_lang changes
        for idx, reason in content_based_excluded.items():
            all_detected_reasons[idx] = reason
        
        # Get block_type information if available (for better exclusion detection)
        # This is a simplified approach - full implementation would require layout_document
        block_type_map = {}  # segment_index -> block_type mapping (if available)
        
        # CRITICAL: Get user_unexcluded_segments to prevent re-detection of user-unexcluded segments
        user_unexcluded_segments = segments_metadata.get("user_unexcluded_segments", [])
        user_unexcluded_set = set(user_unexcluded_segments) if user_unexcluded_segments else set()
        
        if user_unexcluded_segments:
            logger.debug(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Found {len(user_unexcluded_segments)} user-unexcluded segments: "
                f"{sorted(user_unexcluded_segments)}. These will be skipped during re-detection."
                )
        
        for idx, seg in enumerate(all_segments):
            # CRITICAL: Skip content-based excluded segments (Identifier, Image, Formula, etc.)
            # These have higher priority and should NOT be re-detected for Language Match
            # Only detect Language Match for segments that are NOT content-based excluded
            is_content_based_excluded = idx in content_based_excluded
            
            # Skip content-based excluded segments - they have higher priority
            if is_content_based_excluded:
                # Already in all_detected_reasons from earlier, skip re-detection
                continue
            
            # Skip if already excluded (user-based) - but still detect for all_detected_reasons
            # CRITICAL: We still need to detect to update all_detected_reasons for Language Match
            # This ensures frontend can display updated Language Match labels
            if idx in final_excluded_segments and not is_content_based_excluded:
                # Skip non-content-based excluded segments (they're already handled)
                continue
            
            # CRITICAL: Skip segments that user explicitly unexcluded
            # This prevents re-detection of optional exclusions (TABLE) or language-based exclusions (LANGUAGE_MATCH)
            # that the user has explicitly chosen to unexclude
            if idx in user_unexcluded_set:
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Skipping segment {idx} in re-detection "
                    f"(user explicitly unexcluded it, in user_unexcluded_segments)"
                    )
                continue
            
            # Handle both string and dict formats
            if isinstance(seg, dict):
                seg_text = seg.get("text", str(seg.get("source_text", "")))
                block_type = seg.get("block_type")
                is_table_body = seg.get("is_table_body", False)
            else:
                seg_text = str(seg)
                block_type = None
                is_table_body = False
            
            # Check if this is a table segment (for better detection)
            is_table = (block_type == "table_body" or block_type == "table" or is_table_body)
            if not is_table:
                from utils.translation_segments import _is_table_segment
                is_table = _is_table_segment(seg_text)
            
            # Detect exclusion reason with new target language
            detected_result = detect_exclusion_reason(
                text=seg_text,
                block_type=block_type,
                target_lang=target_lang,
                is_image=False,
                is_table=is_table  # Pass is_table for better TABLE detection
            )
            
            if detected_result:
                detected_reason, detected_metadata = detected_result
            
                # CRITICAL: Store all detected reasons (including non-excluded) for frontend display
                # This ensures frontend can show identifier, language_match, etc. even if not excluded
                # CRITICAL: For content-based excluded segments, preserve the original reason (e.g., identifier)
                # but also detect for Language Match to update all_detected_reasons
                if is_content_based_excluded:
                    # Content-based exclusion - preserve original reason in all_detected_reasons
                    # This ensures frontend shows the correct label (e.g., identifier, not language_match)
                    original_reason = content_based_excluded.get(idx)
                    if original_reason:
                        # Preserve original content-based reason (e.g., identifier)
                        all_detected_reasons[idx] = original_reason
                        # If also detected as language_match, log it but don't override identifier
                        if detected_reason == ExclusionReason.LANGUAGE_MATCH:
                            logger.debug(
                                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Segment {idx} is content-based excluded ({original_reason.value}) "
                                f"but also matches target language. Preserving {original_reason.value} in all_detected_reasons."
                                )
                else:
                    # Not content-based excluded - use detected reason
                    all_detected_reasons[idx] = detected_reason
            
                # CRITICAL: For optional exclusions (TABLE), only add if not in user_unexcluded_segments
                # This ensures user choice is respected
                if ExclusionReason.is_optional(detected_reason) and idx in user_unexcluded_set:
                    logger.debug(
                        LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Skipping optional exclusion (TABLE) for segment {idx} "
                        f"(user explicitly unexcluded it)"
                        )
                    continue
            
                # Check if this is a language-based exclusion
                if ExclusionReason.is_language_based(detected_reason):
                    # Language-based exclusion - need user confirmation if auto_exclude=False
                    if auto_exclude:
                        final_excluded_segments[idx] = detected_reason
                        exclusion_metadata_dict[idx] = detected_metadata
                    else:
                        # Don't exclude automatically, but add to language_matched_segments for user confirmation
                        language_matched_segments.append({
                            "index": idx,
                            "text": seg_text[:100] + ("..." if len(seg_text) > 100 else ""),
                            "preview": seg_text[:50] + ("..." if len(seg_text) > 50 else "")
                        })
                elif ExclusionReason.is_optional(detected_reason):
                    # Optional exclusion (TABLE) - do NOT automatically exclude
                    # Only add if it was already excluded (user choice)
                    if idx in current_excluded_indices:
                        # Was excluded before - keep it excluded (user choice)
                        final_excluded_segments[idx] = detected_reason
                        exclusion_metadata_dict[idx] = detected_metadata
                        logger.debug(
                            LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Preserving optional exclusion (TABLE) for segment {idx} "
                            f"(was already excluded, user choice)"
                            )
                    else:
                        # Was not excluded before - do NOT add (optional exclusion, default: not excluded)
                        logger.debug(
                            LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Skipping optional exclusion (TABLE) for segment {idx} "
                            f"(was not excluded before, default: not excluded)"
                            )
                else:
                    # Content-based exclusion - always exclude
                    final_excluded_segments[idx] = detected_reason
                    exclusion_metadata_dict[idx] = detected_metadata
        
        # Step 4: Log the changes for debugging
        final_excluded_set = set(final_excluded_segments.keys())
        preserved_content_based = len(final_excluded_set.intersection(set(content_based_excluded.keys())))
        preserved_user_based = len(final_excluded_set.intersection(set(user_based_excluded.keys())))
        # Language-based exclusions that were cleared (were excluded before but not excluded now)
        cleared_lang_based = set(language_based_excluded.keys()) - final_excluded_set
        # New exclusions added (not in current excluded list)
        newly_added = final_excluded_set - current_excluded_indices
        
        # Group by reason for logging
        reason_counts = {}
        for reason in final_excluded_segments.values():
            reason_counts[reason] = reason_counts.get(reason, 0) + 1
        
        logger.info(
            LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Language-based exclusion update: "
            f"preserved {preserved_content_based} content-based exclusions (out of {len(content_based_excluded)}), "
            f"preserved {preserved_user_based} user-based exclusions (out of {len(user_based_excluded)}), "
            f"cleared {len(cleared_lang_based)} language-based exclusions, "
            f"newly added {len(newly_added)} exclusions, "
            f"final excluded count: {len(final_excluded_segments)} (was {len(current_excluded_segments)}). "
            f"Breakdown: {', '.join(f'{count} {reason.value}' for reason, count in sorted(reason_counts.items()))}"
            )
        # Additional debug logging
        if content_based_excluded:
            logger.debug(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Content-based excluded indices: {sorted(content_based_excluded.keys())}"
                )
        if cleared_lang_based:
            logger.debug(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Cleared language-based excluded indices: {sorted(cleared_lang_based)}"
                )
        if newly_added:
            logger.debug(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Newly added excluded indices: {sorted(newly_added)}"
                )
        
        # CRITICAL: Store all detected reasons (including non-excluded) for frontend display
        # This ensures frontend can show identifier, language_match, etc. even if not excluded
        # Use ExclusionDetectionBatch.store_exclusions to update detected_exclusion_reasons
        if all_detected_reasons:
            from exclusion.detection.batch_detector import ExclusionDetectionBatch
            # Get segment metadata for detected reasons
            segment_metadata_for_detected = {}
            for idx in all_detected_reasons.keys():
                if idx < len(all_segments):
                    seg = all_segments[idx]
                    if isinstance(seg, dict):
                        block_type = seg.get("block_type")
                        if block_type:
                            segment_metadata_for_detected[idx] = {"block_type": block_type}
            
            ExclusionDetectionBatch.store_exclusions(
                task_state=st,
                excluded_segments=final_excluded_segments,  # Only excluded segments
                segment_metadata=segment_metadata_for_detected if segment_metadata_for_detected else None,
                source="update_excluded_segments_for_language",
                all_detected_reasons=all_detected_reasons  # All detected reasons (including non-excluded)
            )
            logger.debug(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Stored {len(all_detected_reasons)} detected exclusion reasons "
                f"(including {len(final_excluded_segments)} excluded) for frontend display"
                )
        
        # Update segments_metadata with new excluded segments using ExclusionManager
        # CRITICAL: Always update to ensure state is consistent
        # If auto_exclude=False, this will only contain content-based and user-based exclusions
        # If auto_exclude=True, this will contain all exclusions (including language-based)
        ExclusionManager.update_excluded_segments(
            task_state=st,
            excluded_segments=final_excluded_segments,
            metadata=exclusion_metadata_dict if exclusion_metadata_dict else None
        )
        
        # Also update excluded_segment_indices for backward compatibility
        excluded_segment_indices = sorted(final_excluded_segments.keys())
        segments_metadata["excluded_segment_indices"] = excluded_segment_indices
        
        # CRITICAL: Also update translation_segments to reflect excluded status
        # This ensures frontend can see updated excluded labels immediately
        # Always update to ensure state is consistent (whether auto_exclude is True or False)
        translation_segments_data = st.get("translation_segments")
        if translation_segments_data and isinstance(translation_segments_data, dict):
            segments_list = translation_segments_data.get("segments", [])
            excluded_set = set(excluded_segment_indices)
            
            updated_count = 0
            for segment in segments_list:
                if isinstance(segment, dict):
                    segment_index = segment.get("segment_index")
                    if segment_index is not None:
                        is_excluded = segment_index in excluded_set
                        old_excluded = segment.get("is_excluded", False)
                        segment["is_excluded"] = is_excluded
                        
                        # If segment is now excluded and wasn't before, set target_text to source_text
                        if is_excluded and not old_excluded:
                            source_text = segment.get("source_text", "")
                            if source_text:
                                segment["target_text"] = source_text
                                segment["modified_text"] = source_text
                        
                        if is_excluded != old_excluded:
                            updated_count += 1
            
            logger.info(
                LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Updated {updated_count} segments in translation_segments "
                f"to reflect excluded status changes"
            )
        
        # CRITICAL: Clear detected_exclusion_reasons for LANGUAGE_MATCH when target_lang changes
        # This prevents incorrect Language Match statistics from persisting after target_lang update
        segments_metadata = st.get("segments_metadata", {})
        detected_exclusion_reasons = segments_metadata.get("detected_exclusion_reasons", {})
        if detected_exclusion_reasons and isinstance(detected_exclusion_reasons, dict):
            # Remove LANGUAGE_MATCH entries from detected_exclusion_reasons
            # They will be re-detected with the new target_lang
            language_match_removed = 0
            updated_detected_reasons = {}
            for seg_idx_str, detection_info in detected_exclusion_reasons.items():
                if isinstance(detection_info, dict):
                    reason = detection_info.get("reason", "unknown")
                elif isinstance(detection_info, str):
                    reason = detection_info
                else:
                    reason = "unknown"
                
                # Skip LANGUAGE_MATCH entries (they will be re-detected with new target_lang)
                if reason == "language_match":
                    language_match_removed += 1
                    logger.debug(
                        LogModule.EXCLUSION,
                        f"[UPDATE-EXCLUDED] Task {task_id}: Removing LANGUAGE_MATCH from detected_exclusion_reasons "
                        f"for segment {seg_idx_str} (will be re-detected with new target_lang={target_lang})"
                        )
                    continue
                
                # Keep other detected reasons (identifier, formula, etc.)
                updated_detected_reasons[seg_idx_str] = detection_info
            
            if language_match_removed > 0:
                segments_metadata["detected_exclusion_reasons"] = updated_detected_reasons
                logger.info(
                    LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Removed {language_match_removed} LANGUAGE_MATCH entries "
                    f"from detected_exclusion_reasons (will be re-detected with new target_lang={target_lang})"
                    )
        
        # Update payload with new target language (if available)
        payload = st.get("payload")
        if payload:
            if isinstance(payload, dict):
                payload["to_lang"] = target_lang
                payload["target_lang"] = target_lang
            else:
                if hasattr(payload, 'to_lang'):
                    payload.to_lang = target_lang
                if hasattr(payload, 'target_lang'):
                    payload.target_lang = target_lang
        
        logger.info(
            LogModule.EXCLUSION,
                    f"[UPDATE-EXCLUDED] Task {task_id}: Updated excluded segments for target_lang={target_lang}. "
            f"Found {len(excluded_segment_indices)} excluded segments out of {len(all_segments)} total segments."
            )
        
        # Return updated excluded segment indices
        response_data = {
            "task_id": task_id,
            "target_lang": target_lang,
            "excluded_segment_indices": excluded_segment_indices,
            "total_segments": len(all_segments),
            "excluded_count": len(excluded_segment_indices),
            "message": f"Updated excluded segments for target language '{target_lang}'"
        }
        
        # If auto_exclude is False and there are language-matched segments, include them in response
        if not auto_exclude and language_matched_segments:
            response_data["language_matched_segments"] = language_matched_segments
            response_data["language_matched_count"] = len(language_matched_segments)
            response_data["requires_confirmation"] = True
        
        return JSONResponse(content=response_data)
    
    def _extract_page_count_from_document(
        self, task_state: Dict[str, Any], original_filename: str
    ) -> int:
        """
        Extract page count from document based on file format.
        
        Args:
            task_state: Task state dictionary
            original_filename: Original filename with extension
            
        Returns:
            Page count (0 if unable to determine)
        """
        if not original_filename:
            return 0
        
        file_ext = Path(original_filename).suffix.lower()
        page_count = 0
        
        try:
            # Get workflow type to determine document format
            payload = task_state.get("payload")
            workflow_type = None
            if payload:
                if isinstance(payload, dict):
                    workflow_type = payload.get("workflow_type")
                elif hasattr(payload, 'workflow_type'):
                    workflow_type = getattr(payload, 'workflow_type')
            
            # If workflow_type not available, infer from file extension
            if not workflow_type:
                ext_to_workflow = {
                    '.docx': 'docx',
                    '.doc': 'docx',
                    '.pptx': 'pptx',
                    '.ppt': 'pptx',
                    '.xlsx': 'xlsx',
                    '.xls': 'xlsx',
                    '.pdf': 'pdf',
                }
                workflow_type = ext_to_workflow.get(file_ext, 'markdown_based')
            
            # Extract page count based on format
            if workflow_type == 'docx':
                page_count = self._get_docx_page_count(task_state)
            elif workflow_type == 'pptx':
                page_count = self._get_pptx_page_count(task_state)
            elif workflow_type == 'xlsx':
                page_count = self._get_xlsx_page_count(task_state)
            elif workflow_type == 'pdf':
                # For PDF, page_count may have been stored early in task_state
                page_count = task_state.get("page_count", 0)
            elif workflow_type in ['txt', 'markdown_based', 'html', 'md']:
                # For text-based formats, estimate pages based on content length
                page_count = self._estimate_text_pages(task_state)
            
            if page_count > 0:
                logger.trace(LogModule.WORKFLOW, f"Extracted page_count={page_count} for {workflow_type} file: {original_filename}")
        
        except Exception as e:
            logger.warning(LogModule.WORKFLOW, f"[STATUS] Failed to extract page count for {original_filename}: {e}")
        
        return page_count
    
    def _get_docx_page_count(self, task_state: Dict[str, Any]) -> int:
        """Get page count from DOCX document."""
        try:
            # Try to get from segments metadata or source preview
            segments_metadata = task_state.get("segments_metadata", {})
            page_count = segments_metadata.get("page_count")
            if page_count:
                return int(page_count)
            
            # Try to estimate from document structure
            # DOCX doesn't have direct page count, but we can estimate from sections
            # For now, try to get from source_preview or segments
            source_preview = task_state.get("source_preview", {})
            segments = source_preview.get("segments", [])
            if segments:
                # CRITICAL: Handle both string list (DOCX, HTML, TXT) and dict list formats
                total_chars = 0
                for seg in segments:
                    if isinstance(seg, str):
                        # String format (DOCX, HTML, TXT workflows)
                        total_chars += len(seg)
                    elif isinstance(seg, dict):
                        # Dictionary format (some workflows)
                        total_chars += len(seg.get("text", ""))
                    else:
                        # Fallback: try to convert to string
                        total_chars += len(str(seg))
                
                # Estimate: roughly 1 page per 500 words or 3000 characters
                estimated_pages = max(1, int(total_chars / 3000))
                return estimated_pages
            
            # Try to load document if file path is available
            original_file_path = task_state.get("original_file_path")
            if original_file_path and Path(original_file_path).exists():
                try:
                    from docx import Document
                    doc = Document(original_file_path)
                    # Estimate pages: count paragraphs and sections
                    # Rough estimate: 1 page = ~30 paragraphs or 1 section
                    para_count = len(doc.paragraphs)
                    section_count = len(doc.sections)
                    estimated_pages = max(1, max(para_count // 30, section_count))
                    return estimated_pages
                except Exception:
                    pass
        
        except Exception as e:
            logger.debug(LogModule.WORKFLOW, f"[STATUS] Failed to get DOCX page count: {e}")
        
        return 0
    
    def _get_pptx_page_count(self, task_state: Dict[str, Any]) -> int:
        """Get page count (slide count) from PPTX document."""
        try:
            # Try to get from segments metadata
            segments_metadata = task_state.get("segments_metadata", {})
            slide_count = segments_metadata.get("slide_count") or segments_metadata.get("page_count")
            if slide_count:
                return int(slide_count)
            
            # Try to get from source_preview total_segments (but this is not accurate for PPTX)
            # PPTX can have multiple segments per slide, so we should not use segment count
            # Instead, we should rely on slide_count from segments_metadata or file loading
            # This fallback is kept only for backward compatibility
            source_preview = task_state.get("source_preview", {})
            total_segments = source_preview.get("total_segments", 0)
            if total_segments > 0:
                # Log warning that we're using an inaccurate method
                logger.debug(LogModule.WORKFLOW, f"[STATUS] PPTX page count: Using total_segments ({total_segments}) as fallback, but this may be inaccurate (multiple segments per slide)")
                # Don't use segment count directly - it's not reliable for PPTX
                # Return 0 to trigger file-based extraction instead
                pass
            
            # Try to load presentation if file path is available
            original_file_path = task_state.get("original_file_path")
            if original_file_path and Path(original_file_path).exists():
                try:
                    from pptx import Presentation
                    prs = Presentation(original_file_path)
                    return len(prs.slides)
                except ImportError:
                    logger.debug(LogModule.WORKFLOW, "[STATUS] python-pptx not available for PPTX page count")
                except Exception:
                    pass
        
        except Exception as e:
            logger.debug(LogModule.WORKFLOW, f"[STATUS] Failed to get PPTX page count: {e}")
        
        return 0
    
    def _get_xlsx_page_count(self, task_state: Dict[str, Any]) -> int:
        """Get page count (worksheet count) from XLSX document."""
        try:
            # Try to get from segments metadata
            segments_metadata = task_state.get("segments_metadata", {})
            sheet_count = segments_metadata.get("sheet_count") or segments_metadata.get("page_count")
            if sheet_count:
                return int(sheet_count)
            
            # Try to load workbook if file path is available
            original_file_path = task_state.get("original_file_path")
            if original_file_path and Path(original_file_path).exists():
                try:
                    from openpyxl import load_workbook
                    wb = load_workbook(original_file_path, read_only=True)
                    return len(wb.worksheets)
                except ImportError:
                    logger.debug(LogModule.WORKFLOW, "[STATUS] openpyxl not available for XLSX page count")
                except Exception:
                    pass
        
        except Exception as e:
            logger.debug(LogModule.WORKFLOW, f"[STATUS] Failed to get XLSX page count: {e}")
        
        return 0
    
    def _estimate_text_pages(self, task_state: Dict[str, Any]) -> int:
        """Estimate page count for text-based formats."""
        try:
            # Estimate based on content length
            source_preview = task_state.get("source_preview", {})
            segments = source_preview.get("segments", [])
            if segments:
                # CRITICAL: Handle both string list and dict list formats
                # source_preview["segments"] can be either:
                # 1. List[str] - list of strings
                # 2. List[Dict] - list of dicts with "text" field
                total_chars = 0
                for seg in segments:
                    if isinstance(seg, str):
                        total_chars += len(seg)
                    elif isinstance(seg, dict):
                        # Extract text from dict (could be "text" or "source_text")
                        text = seg.get("text", "") or seg.get("source_text", "")
                        total_chars += len(text)
                    else:
                        # Fallback: convert to string
                        total_chars += len(str(seg))
                
                # Rough estimate: 1 page = ~3000 characters (A4 page with normal formatting)
                estimated_pages = max(1, int(total_chars / 3000))
                return estimated_pages
            
            # Try cache
            cache_info = task_state.get("source_chunks_cache", {})
            cache_segments = cache_info.get("segments", [])
            if cache_segments:
                total_chars = sum(len(seg) for seg in cache_segments if isinstance(seg, str))
                estimated_pages = max(1, int(total_chars / 3000))
                return estimated_pages
        
        except Exception as e:
            logger.debug(LogModule.WORKFLOW, f"[STATUS] Failed to estimate text pages: {e}")
        
        return 0
    
    def get_format_settings(self, task_id: str) -> Dict[str, Any]:
        """
        Get format settings from task state.
        
        Args:
            task_id: Unique task identifier
            
        Returns:
            Dictionary with format settings
            
        Raises:
            HTTPException: If task not found
        """
        task_state = self.task_manager.get_task(task_id)
        if task_state is None:
            raise HTTPException(status_code=404, detail=f"Task ID '{task_id}' not found.")
        
        # Get format settings from task_state (stored in payload or directly in task_state)
        table_body_format = None
        equation_format = None
        chart_body_format = None
        
        # Check task_state directly first
        if "table_body_format" in task_state:
            table_body_format = task_state["table_body_format"]
        if "equation_format" in task_state:
            equation_format = task_state["equation_format"]
        if "chart_body_format" in task_state:
            chart_body_format = task_state["chart_body_format"]
        
        # Fallback to payload if not in task_state
        if table_body_format is None or equation_format is None or chart_body_format is None:
            payload = task_state.get("payload")
            if payload:
                if isinstance(payload, dict):
                    if table_body_format is None:
                        table_body_format = payload.get("table_body_format")
                    if equation_format is None:
                        equation_format = payload.get("equation_format")
                    if chart_body_format is None:
                        chart_body_format = payload.get("chart_body_format")
                elif hasattr(payload, 'table_body_format'):
                    if table_body_format is None:
                        table_body_format = getattr(payload, 'table_body_format', None)
                    if equation_format is None:
                        equation_format = getattr(payload, 'equation_format', None)
                    if chart_body_format is None:
                        chart_body_format = getattr(payload, 'chart_body_format', None)
        
        return {
            "task_id": task_id,
            "table_body_format": table_body_format,
            "equation_format": equation_format,
            "chart_body_format": chart_body_format,
        }
    
    def update_format_settings(
        self,
        task_id: str,
        table_body_format: Optional[str] = None,
        equation_format: Optional[str] = None,
        chart_body_format: Optional[str] = None,
        bilingual_export: Optional[bool] = None,
        bilingual_order: Optional[str] = None,
        source_text_italic: Optional[bool] = None,
        source_text_color: Optional[str] = None,
        target_text_italic: Optional[bool] = None,
        target_text_color: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Update format settings in task state.
        
        Args:
            task_id: Unique task identifier
            table_body_format: Table format ('html' or 'image')
            equation_format: Equation format ('text' or 'image')
            chart_body_format: Chart format ('html' or 'image', default: 'image')
            bilingual_export: Enable bilingual export (True/False)
            bilingual_order: Bilingual order ('target_after_source' or 'target_before_source')
            source_text_italic: Source text italic (True/False)
            source_text_color: Source text color ('gray', 'blue', 'red', 'green', 'orange', 'black')
            target_text_italic: Target text italic (True/False)
            target_text_color: Target text color ('gray', 'blue', 'red', 'green', 'orange', 'black')
            
        Returns:
            Dictionary with updated format settings
            
        Raises:
            HTTPException: If task not found or invalid format values
        """
        task_state = self.task_manager.get_task(task_id)
        if task_state is None:
            raise HTTPException(status_code=404, detail=f"Task ID '{task_id}' not found.")
        
        # Validate format values (table: html|image; equation: text|latex|image; chart: html|image for PDF export/preview)
        if table_body_format is not None and table_body_format not in ("html", "image"):
            raise HTTPException(status_code=400, detail=f"Invalid table_body_format: {table_body_format}. Must be 'html' or 'image'.")
        if equation_format is not None and equation_format not in ("text", "latex", "image"):
            raise HTTPException(status_code=400, detail=f"Invalid equation_format: {equation_format}. Must be 'text', 'latex', or 'image'.")
        if chart_body_format is not None and chart_body_format not in ("html", "image"):
            raise HTTPException(status_code=400, detail=f"Invalid chart_body_format: {chart_body_format}. Must be 'html' or 'image'.")
        if bilingual_order is not None and bilingual_order not in ("target_after_source", "target_before_source"):
            raise HTTPException(status_code=400, detail=f"Invalid bilingual_order: {bilingual_order}. Must be 'target_after_source' or 'target_before_source'.")
        if source_text_color is not None and source_text_color not in ("gray", "blue", "red", "green", "orange", "black"):
            raise HTTPException(status_code=400, detail=f"Invalid source_text_color: {source_text_color}. Must be 'gray', 'blue', 'red', 'green', 'orange', or 'black'.")
        if target_text_color is not None and target_text_color not in ("gray", "blue", "red", "green", "orange", "black"):
            raise HTTPException(status_code=400, detail=f"Invalid target_text_color: {target_text_color}. Must be 'gray', 'blue', 'red', 'green', 'orange', or 'black'.")
        
        # Update task_state directly
        updates = {}
        if table_body_format is not None:
            updates["table_body_format"] = table_body_format
        if equation_format is not None:
            updates["equation_format"] = equation_format
        if chart_body_format is not None:
            updates["chart_body_format"] = chart_body_format
        if bilingual_export is not None:
            updates["bilingual_export"] = bilingual_export
        if bilingual_order is not None:
            updates["bilingual_order"] = bilingual_order
        if source_text_italic is not None:
            updates["source_text_italic"] = source_text_italic
        if source_text_color is not None:
            updates["source_text_color"] = source_text_color
        if target_text_italic is not None:
            updates["target_text_italic"] = target_text_italic
        if target_text_color is not None:
            updates["target_text_color"] = target_text_color
        
        if updates:
            self.task_manager.update_task(task_id, updates)
            logger.info(LogModule.WORKFLOW, f"[STATUS] Updated format settings for task {task_id}: {updates}")
        
        # Also update payload if it exists
        payload = task_state.get("payload")
        if payload:
            if isinstance(payload, dict):
                if table_body_format is not None:
                    payload["table_body_format"] = table_body_format
                if equation_format is not None:
                    payload["equation_format"] = equation_format
                if chart_body_format is not None:
                    payload["chart_body_format"] = chart_body_format
                if bilingual_export is not None:
                    payload["bilingual_export"] = bilingual_export
                if bilingual_order is not None:
                    payload["bilingual_order"] = bilingual_order
                if source_text_italic is not None:
                    payload["source_text_italic"] = source_text_italic
                if source_text_color is not None:
                    payload["source_text_color"] = source_text_color
                if target_text_italic is not None:
                    payload["target_text_italic"] = target_text_italic
                if target_text_color is not None:
                    payload["target_text_color"] = target_text_color
            elif hasattr(payload, 'table_body_format') or hasattr(payload, 'equation_format'):
                # For object payload, update attributes if possible
                try:
                    if table_body_format is not None:
                        setattr(payload, 'table_body_format', table_body_format)
                    if equation_format is not None:
                        setattr(payload, 'equation_format', equation_format)
                    if bilingual_export is not None:
                        setattr(payload, 'bilingual_export', bilingual_export)
                    if bilingual_order is not None:
                        setattr(payload, 'bilingual_order', bilingual_order)
                    if source_text_italic is not None:
                        setattr(payload, 'source_text_italic', source_text_italic)
                    if source_text_color is not None:
                        setattr(payload, 'source_text_color', source_text_color)
                    if target_text_italic is not None:
                        setattr(payload, 'target_text_italic', target_text_italic)
                    if target_text_color is not None:
                        setattr(payload, 'target_text_color', target_text_color)
                    if chart_body_format is not None:
                        setattr(payload, 'chart_body_format', chart_body_format)
                except Exception as e:
                    logger.debug(LogModule.WORKFLOW, f"[STATUS] Failed to update payload format settings: {e}")
        
        return {
            "task_id": task_id,
            "table_body_format": table_body_format or task_state.get("table_body_format"),
            "equation_format": equation_format or task_state.get("equation_format"),
            "chart_body_format": chart_body_format or task_state.get("chart_body_format"),
            "bilingual_export": bilingual_export if bilingual_export is not None else task_state.get("bilingual_export"),
            "bilingual_order": bilingual_order or task_state.get("bilingual_order"),
            "source_text_italic": source_text_italic if source_text_italic is not None else task_state.get("source_text_italic"),
            "source_text_color": source_text_color or task_state.get("source_text_color"),
            "target_text_italic": target_text_italic if target_text_italic is not None else task_state.get("target_text_italic"),
            "target_text_color": target_text_color or task_state.get("target_text_color"),
            "message": "Format settings updated successfully"
        }

