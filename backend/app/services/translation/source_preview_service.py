# SPDX-FileCopyrightText: 2026 Zampher
# SPDX-License-Identifier: MPL-2.0

"""
Source Preview Service

Handles source text extraction and preview generation for different workflow types.
This service migrates logic from _process_translation_task related to:
- Source text extraction (workflow-specific)
- Preview generation (source_preview, source_chunks_cache)
- Layout document preview generation (for PDF workflows)
"""

import hashlib
import os
import time
from typing import Any, Dict, Optional, List
from pathlib import Path

from logger import unified_logger as logger
from logger.logger import LogModule
from backend.app.services.task import TaskManager
from backend.app.services.translation.chunk_size_service import chunk_size_service
from backend.app.utils.encoding_utils import decode_with_detection
from backend.app.config.pagination_config import SOURCE_PREVIEW_SEGMENTS_LIMIT


class SourcePreviewService:
    """Service for extracting source text and generating previews."""
    
    def __init__(self, task_manager: TaskManager):
        """
        Initialize source preview service.
        
        Args:
            task_manager: Task manager instance
        """
        self.task_manager = task_manager
    
    def _get_target_lang_for_exclusion_detection(
        self,
        task_id: str,
        task_state: Dict[str, Any],
        payload: Any
    ) -> Optional[str]:
        """
        Get target language for exclusion detection from multiple sources (priority order):
        1. segments_metadata.last_target_lang_for_language_match (from frontend language change in Extract phase) - highest priority
        2. payload (from translation task creation)
        
        Args:
            task_id: Task identifier
            task_state: Task state dictionary
            payload: Payload object from translation task creation
            
        Returns:
            Target language code (e.g., 'zh', 'en') or None if not available
        """
        target_lang = None
        
        # Priority 1: Check segments_metadata for stored target language (from frontend language change)
        segments_metadata = task_state.get("segments_metadata", {})
        if not segments_metadata:
            segments_metadata = {}
            task_state["segments_metadata"] = segments_metadata
        
        stored_target_lang = segments_metadata.get("last_target_lang_for_language_match")
        if stored_target_lang:
            target_lang = stored_target_lang
            logger.info(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Using stored target_lang={target_lang} "
                f"from segments_metadata.last_target_lang_for_language_match (from frontend language change in Extract phase)"
            )
            return target_lang
        
        # Priority 2: Fallback to payload if not stored
        target_lang = getattr(payload, 'to_lang', None) or getattr(payload, 'target_lang', None)
        
        # CRITICAL: If target_lang is 'en' and it's the default value (not explicitly set by user),
        # treat it as None to prevent incorrect language_match detection
        # We can't distinguish between user-set 'en' and default 'en', so we'll be conservative
        # and only use it if it's explicitly set (not the default from format conversion)
        if target_lang == 'en':
            # Check if this is from format conversion (which sets default 'en')
            # If so, treat as None to skip language_match detection
            is_format_conversion = task_state.get("is_format_conversion", False) or task_state.get("convert_only", False)
            if is_format_conversion:
                # This is format conversion, 'en' is likely the default value
                # Check if payload has to_lang explicitly set (not None)
                payload_to_lang = getattr(payload, 'to_lang', None)
                if payload_to_lang == 'en':
                    # This is the default 'en' from format conversion, treat as None
                    logger.info(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: target_lang='en' from format conversion (default), treating as None to skip language_match detection")
                    target_lang = None
                else:
                    # to_lang is None, which means it wasn't set, so treat as None
                    target_lang = None
        
        if target_lang:
            logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Using target_lang={target_lang} from payload for exclusion detection")
        else:
            logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: target_lang is None or default 'en', skipping language_match detection during extraction")
        
        return target_lang
    
    def prepare_source_preview_for_docx(
        self,
        task_id: str,
        file_contents: bytes,
        payload: Any,
        task_state: Dict[str, Any]
    ) -> bool:
        """
        Prepare source preview for DOCX workflow.
        
        Args:
            task_id: Task identifier
            file_contents: DOCX file content bytes
            payload: Task payload
            task_state: Task state dictionary
            
        Returns:
            True if preview was prepared successfully, False otherwise
        """
        try:
            from extractor.docx_extractor import DocxExtractor
            
            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
            
            # Clear old chunk_to_segment_map if chunk_size changed
            existing_chunk_size = task_state.get("segments_metadata", {}).get("chunk_size")
            if existing_chunk_size and existing_chunk_size != chunk_size:
                logger.info(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: chunk_size changed from {existing_chunk_size} to {chunk_size}, clearing old chunk_to_segment_map")
                task_state.pop("chunk_to_segment_map", None)
                task_state.pop("chunk_tokens_info", None)
                task_state.pop("total_estimated_input_tokens", None)
            
            result = DocxExtractor(file_contents, chunk_size=chunk_size).extract()
            if result.total_segments == 0:
                return False
            
            # Mark excluded segments during extraction using unified exclusion detection
            # CRITICAL: Get target language using unified method that checks stored value first
            target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
            
            # Use new unified exclusion detection architecture
            from exclusion.extractors.docx_extractor import DOCXMetadataExtractor
            from exclusion.detection.batch_detector import ExclusionDetectionBatch
            
            # Create DOCX metadata extractor
            extractor = DOCXMetadataExtractor(result.segment_info)
            
            # CRITICAL: Set task_id in task_state so _update_progress can persist to task_manager (frontend progress polling)
            task_state["task_id"] = task_id
            # Batch detect exclusions using new architecture
            excluded_segments_with_reasons, all_detected_reasons = ExclusionDetectionBatch.detect_exclusions_batch(
                segments=result.segments,
                metadata_extractor=extractor,
                task_state=task_state,
                target_lang=target_lang,
                preserve_existing=True,
                auto_exclude_optional=False  # TABLE not auto-excluded
            )
            
            # Store exclusions using new architecture
            # CRITICAL: Always call store_exclusions even if excluded_segments_with_reasons is empty
            # This ensures all_detected_reasons are stored for frontend display
            segment_metadata = {
                idx: {"block_type": "table" if result.segment_info[idx].get("is_table_cell", False) else None}
                for idx in (excluded_segments_with_reasons.keys() if excluded_segments_with_reasons else [])
                if idx < len(result.segment_info)
            }
            ExclusionDetectionBatch.store_exclusions(
                task_state=task_state,
                excluded_segments=excluded_segments_with_reasons,
                segment_metadata=segment_metadata if segment_metadata else None,
                source="docx_extraction",
                all_detected_reasons=all_detected_reasons  # Store all detected reasons for frontend display
            )
            
            # Build chunk_to_segment_map for chunks generation
            # Get excluded_segment_indices from excluded_segments_with_reasons (for backward compatibility)
            excluded_segment_indices = sorted(excluded_segments_with_reasons.keys()) if excluded_segments_with_reasons else []
            self._build_chunk_to_segment_map(
                task_id=task_id,
                segments=result.segments,
                chunk_size=chunk_size,
                excluded_segment_indices=excluded_segment_indices,
                task_state=task_state
            )
            
            # Store preview and cache
            content_hash = hashlib.sha1(file_contents).hexdigest()
            task_state["source_preview"] = {
                "segments": result.segments[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                "total_segments": result.total_segments,
                "ready": True,
            }
            task_state["source_chunks_cache"] = {
                "content_hash": content_hash,
                "chunk_size": chunk_size,
                "segments": result.segments,
                "total_segments": result.total_segments,
                "created_at": time.time(),
            }
            
            # CRITICAL: Get excluded_segments AFTER ExclusionManager.update_excluded_segments
            # This ensures we get the newly stored exclusion information
            # ExclusionManager.update_excluded_segments has already updated task_state["segments_metadata"]["excluded_segments"]
            current_segments_metadata = task_state.get("segments_metadata", {})
            current_excluded_segments = current_segments_metadata.get("excluded_segments", {})
            
            # DEBUG: Log current_excluded_segments structure for troubleshooting
            if current_excluded_segments:
                # Count reasons in current_excluded_segments
                reason_counts_in_storage = {}
                for seg_idx_str, exclusion_info in current_excluded_segments.items():
                    if isinstance(exclusion_info, dict):
                        reason = exclusion_info.get("reason", "unknown")
                    elif isinstance(exclusion_info, str):
                        reason = exclusion_info
                    else:
                        reason = "unknown"
                    reason_counts_in_storage[reason] = reason_counts_in_storage.get(reason, 0) + 1
                reason_summary_in_storage = ', '.join(f'{count} {reason}' for reason, count in sorted(reason_counts_in_storage.items()))
                logger.info(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Current excluded_segments in storage: {len(current_excluded_segments)} segments ({reason_summary_in_storage})"
                )
            
            # Create new segments_metadata
            new_segments_metadata = {
                "source": "docx",
                "workflow_type": getattr(payload, 'workflow_type', 'docx'),
                "chunk_size": chunk_size,
                "content_hash": content_hash,
                "separators_after": result.separators_after,
                "segment_info": result.segment_info,
                "excluded_segment_indices": excluded_segment_indices,  # Legacy format for backward compatibility
            }
            
            # CRITICAL: Preserve excluded_segments from ExclusionManager (includes newly detected exclusions)
            # This is the source of truth for exclusion information
            if current_excluded_segments:
                new_segments_metadata["excluded_segments"] = current_excluded_segments
                logger.info(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Preserved {len(current_excluded_segments)} excluded_segments from ExclusionManager "
                    f"(includes {len(excluded_segments_with_reasons) if excluded_segments_with_reasons else 0} newly detected exclusions)"
                )
            elif excluded_segments_with_reasons:
                # Fallback: if ExclusionManager didn't store it (shouldn't happen), log a warning
                logger.warning(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: ExclusionManager.update_excluded_segments was called with {len(excluded_segments_with_reasons)} exclusions, "
                    f"but current_excluded_segments is empty. This may indicate a storage issue."
                )
            
            task_state["segments_metadata"] = new_segments_metadata
            
            # Verify chunk_to_segment_map was created
            if "chunk_to_segment_map" not in task_state or task_state.get("chunk_to_segment_map") is None:
                logger.warning(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: chunk_to_segment_map not created for DOCX, chunks may not be available")
            else:
                logger.info(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: DOCX preview prepared successfully with {len(task_state.get('chunk_to_segment_map', []))} chunks")
                
                # Output segments and chunks to temporary folder for debugging
                self._output_extract_debug_files(
                    task_id=task_id,
                    task_state=task_state,
                    segments=result.segments,
                    chunk_to_segment_map=task_state.get("chunk_to_segment_map", [])
                )
            
            return True
        except ValueError as e:
            # ValueError from DocxExtractor indicates file format issue (e.g., WPS format)
            error_msg = str(e)
            logger.error(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Failed to prepare DOCX source preview - file format issue: {error_msg}. "
                f"This may be a WPS format file (.wps.docx) which is not supported. File size: {len(file_contents)} bytes",
                exc_info=True
            )
            return False
        except Exception as e:
            error_msg = str(e)
            logger.error(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Failed to prepare DOCX source preview: {error_msg}. "
                f"File size: {len(file_contents)} bytes, error type: {type(e).__name__}",
                exc_info=True
            )
            return False
    
    def prepare_source_preview_for_pptx(
        self,
        task_id: str,
        file_contents: bytes,
        original_filename: str,
        payload: Any,
        task_state: Dict[str, Any],
        temp_dir: Optional[str] = None
    ) -> bool:
        """
        Prepare source preview for PPTX workflow.
        
        Args:
            task_id: Task identifier
            file_contents: PPTX file content bytes
            original_filename: Original filename
            payload: Task payload
            task_state: Task state dictionary
            temp_dir: Temporary directory path (optional)
            
        Returns:
            True if preview was prepared successfully, False otherwise
        """
        try:
            # Check if python-pptx is available
            try:
                from pptx import Presentation
                pptx_available = True
            except ImportError:
                pptx_available = False
                logger.error(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: python-pptx is not installed. Cannot extract segments from PPTX file.")
                return False
            
            if not pptx_available:
                return False
            
            from translator.ai_translator.pptx_translator import PptxTranslatorConfig, PptxTranslator
            from ir.document import Document
            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
            
            # Clear old chunk_to_segment_map if chunk_size changed
            existing_chunk_size = task_state.get("segments_metadata", {}).get("chunk_size")
            if existing_chunk_size and existing_chunk_size != chunk_size:
                logger.info(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: chunk_size changed from {existing_chunk_size} to {chunk_size}, clearing old chunk_to_segment_map")
                task_state.pop("chunk_to_segment_map", None)
                task_state.pop("chunk_tokens_info", None)
                task_state.pop("total_estimated_input_tokens", None)
            
            # Create a temporary translator config for extraction (skip_translate=True)
            extractor_config = PptxTranslatorConfig(
                skip_translate=True,
                translate_notes=getattr(payload, 'translate_notes', False),
                translate_master=getattr(payload, 'translate_master', False),
                translate_tables=getattr(payload, 'translate_tables', True),
                translate_textboxes=getattr(payload, 'translate_textboxes', True),
                logger=logger
            )
            extractor = PptxTranslator(extractor_config)
            
            # Create Document object from file contents
            file_stem = Path(original_filename).stem
            file_suffix = Path(original_filename).suffix
            document = Document.from_bytes(content=file_contents, stem=file_stem, suffix=file_suffix)
            
            # Extract segments using _pre_translate
            # CRITICAL: Get elements_to_translate to access element type information (table_cell, notes, etc.)
            prs, elements_to_translate, original_texts, _ = extractor._pre_translate(document, temp_dir=None)
            
            if not original_texts or len(original_texts) == 0:
                logger.warning(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: No segments extracted from PPTX file")
                return False
            
            # Extract slide count from presentation
            slide_count = len(prs.slides) if prs else 0
            if slide_count > 0:
                logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Extracted {slide_count} slides from PPTX file")
            
            # Build element type map for exclusion detection
            # Map segment index to element type (table_cell, notes, master, text_frame, etc.)
            element_type_map = {}  # {segment_index: element_type}
            for idx, element_info in enumerate(elements_to_translate):
                if idx < len(original_texts):
                    element_type = element_info.get("type", "text_frame")
                    element_type_map[idx] = element_type
            
            # Use new unified exclusion detection architecture
            # CRITICAL: Get target language using unified method that checks stored value first
            target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
            
            from exclusion.extractors.pptx_extractor import PPTXMetadataExtractor
            from exclusion.detection.batch_detector import ExclusionDetectionBatch
            
            # Create PPTX metadata extractor
            extractor = PPTXMetadataExtractor(element_type_map)
            
            # CRITICAL: Set task_id in task_state so _update_progress can persist to task_manager (frontend progress polling)
            task_state["task_id"] = task_id
            # Batch detect exclusions using new architecture
            excluded_segments_with_reasons, all_detected_reasons = ExclusionDetectionBatch.detect_exclusions_batch(
                segments=original_texts,
                metadata_extractor=extractor,
                task_state=task_state,
                target_lang=target_lang,
                preserve_existing=True,
                auto_exclude_optional=False  # TABLE not auto-excluded
            )
            
            # Store exclusions using new architecture
            # CRITICAL: Always call store_exclusions even if excluded_segments_with_reasons is empty
            # This ensures all_detected_reasons are stored for frontend display
            segment_metadata = {
                idx: {
                    "block_type": "table" if element_type_map.get(idx) == "table_cell" else None,
                    "element_type": element_type_map.get(idx, "text_frame")
                }
                for idx in (excluded_segments_with_reasons.keys() if excluded_segments_with_reasons else [])
                if idx < len(original_texts)
            }
            ExclusionDetectionBatch.store_exclusions(
                task_state=task_state,
                excluded_segments=excluded_segments_with_reasons,
                segment_metadata=segment_metadata if segment_metadata else None,
                source="pptx_extraction",
                all_detected_reasons=all_detected_reasons  # Store all detected reasons for frontend display
            )
            
            # Build chunk_to_segment_map for chunks generation
            # Get excluded_segment_indices from excluded_segments_with_reasons (for backward compatibility)
            excluded_segment_indices = sorted(excluded_segments_with_reasons.keys()) if excluded_segments_with_reasons else []
            self._build_chunk_to_segment_map(
                task_id=task_id,
                segments=original_texts,
                chunk_size=chunk_size,
                excluded_segment_indices=excluded_segment_indices,
                task_state=task_state
            )
            
            # Store preview and cache
            content_hash = hashlib.sha1(file_contents).hexdigest()
            task_state["source_preview"] = {
                "segments": original_texts[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                "total_segments": len(original_texts),
                "ready": True,
            }
            task_state["source_chunks_cache"] = {
                "content_hash": content_hash,
                "chunk_size": chunk_size,
                "segments": original_texts,
                "total_segments": len(original_texts),
                "created_at": time.time(),
            }
            # CRITICAL: Get excluded_segments AFTER ExclusionManager.update_excluded_segments
            # This ensures we get the newly stored exclusion information
            # ExclusionManager.update_excluded_segments has already updated task_state["segments_metadata"]["excluded_segments"]
            current_segments_metadata = task_state.get("segments_metadata", {})
            current_excluded_segments = current_segments_metadata.get("excluded_segments", {})
            
            # Create new segments_metadata
            new_segments_metadata = {
                "source": "pptx",
                "workflow_type": getattr(payload, 'workflow_type', 'pptx'),
                "chunk_size": chunk_size,
                "content_hash": content_hash,
                "excluded_segment_indices": excluded_segment_indices,  # Legacy format for backward compatibility
                "slide_count": slide_count,  # Store slide count for page count extraction
                "segment_info": elements_to_translate,  # Store element info for later use
            }
            
            # CRITICAL: Preserve excluded_segments from ExclusionManager (includes newly detected exclusions)
            # This is the source of truth for exclusion information
            if current_excluded_segments:
                new_segments_metadata["excluded_segments"] = current_excluded_segments
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Preserved {len(current_excluded_segments)} excluded_segments from ExclusionManager "
                    f"(includes {len(excluded_segments_with_reasons) if excluded_segments_with_reasons else 0} newly detected exclusions)"
                )
            elif excluded_segments_with_reasons:
                # Fallback: if ExclusionManager didn't store it (shouldn't happen), log a warning
                logger.warning(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: ExclusionManager.update_excluded_segments was called with {len(excluded_segments_with_reasons)} exclusions, "
                    f"but current_excluded_segments is empty. This may indicate a storage issue."
                )
            
            task_state["segments_metadata"] = new_segments_metadata
            
            # Verify chunk_to_segment_map was created and output debug files
            if "chunk_to_segment_map" in task_state and task_state.get("chunk_to_segment_map") is not None:
                # Output segments and chunks to temporary folder for debugging
                self._output_extract_debug_files(
                    task_id=task_id,
                    task_state=task_state,
                    segments=original_texts,
                    chunk_to_segment_map=task_state.get("chunk_to_segment_map", [])
                )
            
            logger.info(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Prepared source preview for PPTX: {len(original_texts)} segments")
            return True
        except Exception as e:
            logger.error(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Failed to prepare PPTX source preview: {e}", exc_info=True)
            return False
    
    def prepare_source_preview_for_extractor_based(
        self,
        task_id: str,
        file_contents: bytes,
        payload: Any,
        task_state: Dict[str, Any],
        extractor_class: Any,
        workflow_type: str
    ) -> bool:
        """
        Prepare source preview for extractor-based workflows (EPUB, MOBI, Qt TS).
        
        Args:
            task_id: Task identifier
            file_contents: File content bytes
            payload: Task payload
            task_state: Task state dictionary
            extractor_class: Extractor class to use
            workflow_type: Workflow type name
            
        Returns:
            True if preview was prepared successfully, False otherwise
        """
        try:
            # Reuse Extract-phase result when Translate task inherited from Convert (skip re-extraction and re-detection)
            convert_task_id = task_state.get("convert_task_id")
            cache = task_state.get("source_chunks_cache") or {}
            segments_list = cache.get("segments") if isinstance(cache.get("segments"), list) else None
            total_segments = cache.get("total_segments")
            segments_metadata = task_state.get("segments_metadata") or {}

            if convert_task_id and segments_list is not None and total_segments is not None and len(segments_list) == total_segments:
                logger.info(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: Reusing Extract result from convert_task_id={convert_task_id} "
                    f"(segments={total_segments}), skipping re-extraction and exclusion detection for {workflow_type.upper()}"
                )
                chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
                # Ensure source_preview is set from cache
                if not task_state.get("source_preview", {}).get("ready"):
                    task_state["source_preview"] = {
                        "segments": segments_list[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                        "total_segments": total_segments,
                        "ready": True,
                    }
                # Build chunk_to_segment_map if missing (Translate phase needs it; Convert skips it)
                convert_only = task_state.get("convert_only", False)
                if not convert_only and (not task_state.get("chunk_to_segment_map") or task_state.get("chunk_to_segment_map") is None):
                    excluded_segment_indices = segments_metadata.get("excluded_segment_indices")
                    if excluded_segment_indices is None and isinstance(segments_metadata.get("excluded_segments"), dict):
                        try:
                            excluded_segment_indices = sorted(int(k) for k in segments_metadata["excluded_segments"].keys() if str(k).isdigit())
                        except (ValueError, TypeError):
                            excluded_segment_indices = []
                    if excluded_segment_indices is None:
                        excluded_segment_indices = []
                    self._build_chunk_to_segment_map(
                        task_id=task_id,
                        segments=segments_list,
                        chunk_size=chunk_size,
                        excluded_segment_indices=excluded_segment_indices,
                        task_state=task_state
                    )
                    if task_state.get("chunk_to_segment_map"):
                        logger.debug(
                            LogModule.EXTRACT,
                            f"[PREVIEW] Task {task_id}: [{workflow_type.upper()}] Built chunk_to_segment_map from inherited cache: {len(task_state['chunk_to_segment_map'])} chunks"
                        )
                return True

            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
            result = extractor_class(file_contents, chunk_size=chunk_size).extract()
            
            if result.total_segments == 0:
                logger.warning(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: {workflow_type.upper()} extraction returned 0 segments.")
                return False
            
            # Mark excluded segments using unified exclusion detection
            from utils.translation_segments import _is_image_segment
            from exclusion.core import detect_exclusion_reason, ExclusionReason, ExclusionManager
            
            # CRITICAL: Get target language using unified method that checks stored value first
            target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
            excluded_segments_with_reasons = {}
            excluded_segment_indices = []
            # CRITICAL: Store all detected reasons (including non-excluded ones) for frontend display
            # This allows frontend to display identifier, language_match, etc. even if not excluded
            all_detected_reasons = {}
            
            for idx, seg_text in enumerate(result.segments):
                is_image = _is_image_segment(seg_text)
                detected_result = detect_exclusion_reason(
                    text=seg_text,
                    block_type=None,
                    target_lang=target_lang,
                    is_image=is_image,
                    is_table=False
                )
                if detected_result:
                    detected_reason, _ = detected_result
                    # CRITICAL: Store all detected reasons (including non-excluded ones)
                    all_detected_reasons[idx] = detected_reason
                    # Only exclude non-optional types (TABLE is optional and not auto-excluded)
                    if not ExclusionReason.is_optional(detected_reason):
                        excluded_segment_indices.append(idx)
                        excluded_segments_with_reasons[idx] = detected_reason
            
            if excluded_segments_with_reasons:
                logger.debug(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: Marked {len(excluded_segments_with_reasons)} segments as excluded during {workflow_type.upper()} extraction"
                )
                # Store excluded segments using ExclusionManager
                ExclusionManager.update_excluded_segments(
                    task_state=task_state,
                    excluded_segments=excluded_segments_with_reasons,
                    metadata=None
                )
            
            # CRITICAL: Get excluded_segments AFTER ExclusionManager.update_excluded_segments
            current_segments_metadata = task_state.get("segments_metadata", {})
            current_excluded_segments = current_segments_metadata.get("excluded_segments", {})
            
            content_hash = hashlib.sha1(file_contents).hexdigest()
            task_state["source_preview"] = {
                "segments": result.segments[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                "total_segments": result.total_segments,
                "ready": True,
            }
            # CRITICAL: Store ALL segments in source_chunks_cache (not truncated)
            # This ensures get_source_preview can access all segments for pagination
            # Log segment count for debugging
            logger.info(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Storing {workflow_type.upper()} extraction result - "
                f"result.segments count: {len(result.segments)}, "
                f"result.total_segments: {result.total_segments}, "
                f"segments match: {len(result.segments) == result.total_segments}, "
                f"result.segments type: {type(result.segments).__name__}, "
                f"result.segments is list: {isinstance(result.segments, list)}"
            )
            
            # CRITICAL: Check if result.segments is actually a list and has the expected length
            if not isinstance(result.segments, list):
                logger.error(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: CRITICAL ERROR - result.segments is not a list! "
                    f"type: {type(result.segments).__name__}, value: {result.segments}"
                )
                return False
            
            if len(result.segments) != result.total_segments:
                logger.error(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: CRITICAL ERROR - result.segments length ({len(result.segments)}) "
                    f"does not match result.total_segments ({result.total_segments})! "
                    f"This indicates a bug in the extractor."
                )
                # Still proceed, but log the mismatch
            
            # CRITICAL: Create a copy of segments list to avoid any reference issues
            segments_to_store = list(result.segments)  # Create a new list copy
            
            logger.info(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: About to save source_chunks_cache - "
                f"segments_to_store count: {len(segments_to_store)}, "
                f"task_state type: {type(task_state).__name__}, "
                f"task_state id: {id(task_state)}, "
                f"task_state has source_chunks_cache before: {'source_chunks_cache' in task_state}"
            )
            
            try:
                task_state["source_chunks_cache"] = {
                    "content_hash": content_hash,
                    "chunk_size": chunk_size,
                    "segments": segments_to_store,  # CRITICAL: Store ALL segments, not truncated
                    "total_segments": result.total_segments,
                    "created_at": time.time(),
                }
                logger.info(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: Successfully assigned source_chunks_cache - "
                    f"task_state has source_chunks_cache after: {'source_chunks_cache' in task_state}"
                )
            except Exception as save_error:
                logger.error(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: CRITICAL ERROR - Failed to save source_chunks_cache: {save_error}",
                    exc_info=True
                )
                raise  # Re-raise to be caught by outer try-except
            
            # CRITICAL: Verify that segments were stored correctly immediately after assignment
            stored_segments = task_state["source_chunks_cache"].get("segments", [])
            stored_total = task_state["source_chunks_cache"].get("total_segments")
            logger.info(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Verification - stored segments count: {len(stored_segments)}, "
                f"stored total_segments: {stored_total}, "
                f"expected count: {len(segments_to_store)}, "
                f"expected total: {result.total_segments}, "
                f"match: {len(stored_segments) == len(segments_to_store)}, "
                f"stored_segments type: {type(stored_segments).__name__}, "
                f"stored_segments is list: {isinstance(stored_segments, list)}, "
                f"task_state has source_chunks_cache: {'source_chunks_cache' in task_state}"
            )
            
            # CRITICAL: If stored segments count doesn't match, log error
            if len(stored_segments) != len(segments_to_store):
                logger.error(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: CRITICAL ERROR - Stored segments count ({len(stored_segments)}) "
                    f"does not match expected count ({len(segments_to_store)})! "
                    f"This indicates a bug in task_state storage or serialization. "
                    f"stored_segments type: {type(stored_segments).__name__}, "
                    f"stored_segments[:5]: {stored_segments[:5] if stored_segments else 'empty'}"
                )
            
            # Create new segments_metadata
            new_segments_metadata = {
                "source": workflow_type,
                "workflow_type": workflow_type,
                "chunk_size": chunk_size,
                "content_hash": content_hash,
                "segment_info": getattr(result, 'segment_info', []),
                "excluded_segment_indices": excluded_segment_indices,  # Legacy format for backward compatibility
            }
            
            # CRITICAL: Preserve excluded_segments from ExclusionManager
            if current_excluded_segments:
                new_segments_metadata["excluded_segments"] = current_excluded_segments
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Preserved {len(current_excluded_segments)} excluded_segments from ExclusionManager for {workflow_type.upper()}"
                )
            
            # CRITICAL: Store all detected exclusion reasons (including non-excluded ones) for frontend display
            # This allows frontend to display identifier, language_match, etc. even if not excluded
            if all_detected_reasons:
                # NOTE: time module is already imported at the top of the file, no need to import again
                detected_reasons_dict = {}
                for idx, reason in all_detected_reasons.items():
                    detected_reasons_dict[str(idx)] = {
                        "reason": reason.value,
                        "detected_at": time.time()
                    }
                new_segments_metadata["detected_exclusion_reasons"] = detected_reasons_dict
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Stored {len(all_detected_reasons)} detected exclusion reasons "
                    f"(including {len(excluded_segments_with_reasons)} excluded) for {workflow_type.upper()}"
                )
            
            task_state["segments_metadata"] = new_segments_metadata
            
            # Build chunk_to_segment_map for chunks generation (if not convert_only)
            # CRITICAL: This prevents segments from being merged into a single chunk in get_source_preview
            convert_only = task_state.get("convert_only", False)
            logger.debug(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: [{workflow_type.upper()}] convert_only={convert_only}, will build chunk_to_segment_map: {not convert_only}"
            )
            if not convert_only:
                try:
                    self._build_chunk_to_segment_map(
                        task_id=task_id,
                        segments=result.segments,
                        chunk_size=chunk_size,
                        excluded_segment_indices=excluded_segment_indices,
                        task_state=task_state
                    )
                    # Verify chunk_to_segment_map was created
                    if "chunk_to_segment_map" not in task_state or task_state.get("chunk_to_segment_map") is None:
                        logger.error(
                            LogModule.EXTRACT,
                            f"[PREVIEW] Task {task_id}: [{workflow_type.upper()}] CRITICAL - _build_chunk_to_segment_map completed but chunk_to_segment_map not found in task_state"
                        )
                    else:
                        logger.debug(
                            LogModule.EXTRACT,
                            f"[PREVIEW] Task {task_id}: [{workflow_type.upper()}] Successfully built and stored chunk_to_segment_map: {len(task_state.get('chunk_to_segment_map', []))} chunks"
                        )
                except Exception as e:
                    logger.error(
                        LogModule.EXTRACT,
                        f"[PREVIEW] Task {task_id}: [{workflow_type.upper()}] Failed to build chunk_to_segment_map: {e}",
                        exc_info=True
                    )
                    # Do not re-raise - allow preview to continue without chunks
                    # This allows frontend to display segments even if chunks generation fails
            else:
                logger.debug(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: [{workflow_type.upper()}] Skipping chunk_to_segment_map build (convert_only=True)"
                )
            
            # For convert_only mode, create translation_segments
            if task_state.get("convert_only", False):
                ts_segments = []
                # CRITICAL: Reuse exclusion detection results from above instead of re-detecting
                # This avoids duplicate detection and improves performance
                # Get excluded segments from segments_metadata (already computed above)
                current_segments_metadata = task_state.get("segments_metadata", {})
                excluded_reasons_map = {}
                if current_segments_metadata.get("excluded_segments"):
                    for seg_idx_str, exclusion_info in current_segments_metadata["excluded_segments"].items():
                        seg_idx = int(seg_idx_str)
                        if isinstance(exclusion_info, dict):
                            reason = exclusion_info.get("reason", "unknown")
                        elif isinstance(exclusion_info, str):
                            reason = exclusion_info
                        else:
                            reason = "unknown"
                        excluded_reasons_map[seg_idx] = reason
                
                for i, s in enumerate(result.segments):
                    is_image = _is_image_segment(s)
                    # Reuse exclusion detection result from above instead of re-detecting
                    is_excluded = i in excluded_reasons_map
                    reason_val = excluded_reasons_map.get(i) if is_excluded else None
                    seg = {
                        "segment_index": i,
                        "source_text": s,
                        "target_text": s,
                        "modified": False,
                        "separator_after": "",
                        "is_image": is_image,
                        "is_excluded": is_excluded,
                    }
                    if is_excluded and reason_val:
                        seg["exclusion_reason"] = reason_val
                    ts_segments.append(seg)
                task_state["translation_segments"] = {
                    "segments": ts_segments,
                    "metadata": task_state.get("segments_metadata", {})
                }
            
            self.task_manager.add_log(task_id, "success", f"Source preview prepared from {workflow_type.upper()}: {min(result.total_segments, SOURCE_PREVIEW_SEGMENTS_LIMIT)}/{result.total_segments} segments")
            return True
        except Exception as e:
            logger.error(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Failed to prepare {workflow_type} source preview: {e}", exc_info=True)
            return False
    
    def prepare_source_preview_for_html(
        self,
        task_id: str,
        file_contents: bytes,
        payload: Any,
        task_state: Dict[str, Any]
    ) -> bool:
        """
        Prepare source preview for HTML workflow.
        
        Args:
            task_id: Task identifier
            file_contents: HTML file content bytes
            payload: Task payload
            task_state: Task state dictionary
            
        Returns:
            True if preview was prepared successfully, False otherwise
        """
        try:
            from extractor.html_extractor import HtmlExtractor
            
            decoded = decode_with_detection(file_contents)
            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
            deep_split_enabled = bool(task_state.get("deep_split") or getattr(payload, 'deep_split', True))
            source = "task_state" if task_state.get("deep_split") is not None else "payload"
            
            logger.info(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Using deep_split={deep_split_enabled} "
                f"(from {source}) for html preview (chunk_size={chunk_size})"
            )
            
            result = HtmlExtractor(decoded, chunk_size=chunk_size, deep_split=deep_split_enabled).extract()
            if result.total_segments == 0:
                return False
            
            content_hash = hashlib.sha1(file_contents).hexdigest()
            task_state["source_preview"] = {
                "segments": result.segments[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                "total_segments": result.total_segments,
                "ready": True,
            }
            task_state["source_chunks_cache"] = {
                "content_hash": content_hash,
                "chunk_size": chunk_size,
                "segments": result.segments,
                "total_segments": result.total_segments,
                "created_at": time.time(),
            }
            
            # CRITICAL: Preserve existing exclusion data from convert phase before overwriting segments_metadata.
            # When a translate task inherits segments_metadata from its convert task, user-selected exclusions
            # (e.g. exclude-all) must not be wiped out by re-running HTML extraction.
            existing_segments_metadata = task_state.get("segments_metadata", {})
            existing_excluded_segments = existing_segments_metadata.get("excluded_segments")
            existing_excluded_segment_indices = existing_segments_metadata.get("excluded_segment_indices")
            existing_user_unexcluded_segments = existing_segments_metadata.get("user_unexcluded_segments")
            
            task_state["segments_metadata"] = {
                "source": "html",
                "workflow_type": getattr(payload, 'workflow_type', 'html'),
                "chunk_size": chunk_size,
                "content_hash": content_hash,
                "separators_after": result.separators_after,
                "segment_info": result.segment_info,
            }
            
            # Restore inherited exclusion data so it survives the metadata rebuild
            if existing_excluded_segments is not None:
                task_state["segments_metadata"]["excluded_segments"] = existing_excluded_segments
            if existing_excluded_segment_indices is not None:
                task_state["segments_metadata"]["excluded_segment_indices"] = existing_excluded_segment_indices
            if existing_user_unexcluded_segments is not None:
                task_state["segments_metadata"]["user_unexcluded_segments"] = existing_user_unexcluded_segments
            
            # Mark excluded segments during extraction (same as Markdown workflow)
            from utils.translation_segments import _is_image_segment, _is_table_segment
            from exclusion.core import ExclusionManager, ExclusionReason, detect_exclusion_reason
            
            # CRITICAL: Get target language using unified method that checks stored value first
            target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
            excluded_segments_with_reasons = {}
            excluded_segment_indices = []
            
            for idx, seg_text in enumerate(result.segments):
                is_image = _is_image_segment(seg_text)
                is_table = _is_table_segment(seg_text)
                
                # CRITICAL: Tables are detected but NOT automatically excluded by default
                # Most tables can be translated, so we only exclude non-table segments automatically
                # If user manually excludes a table, it will be marked with TABLE exclusion reason
                # Tables are detected and can be identified by frontend (via block_type or is_table flag)
                
                # Use detect_exclusion_reason to get unified exclusion reason
                # Pass is_table=True to detect TABLE, but only exclude if it's not a table
                # (tables are detected but not auto-excluded)
                detected_result = detect_exclusion_reason(
                    text=seg_text,
                    block_type="table" if is_table else None,  # Mark table segments with block_type
                    target_lang=target_lang,
                    is_image=is_image,
                    is_table=is_table  # Pass table detection to identify TABLE exclusion reason
                )
                
                # Only exclude if detected_result is not TABLE (tables are not auto-excluded)
                # OR if it's TABLE but user has already manually excluded it (handled elsewhere)
                if detected_result:
                    detected_reason, detected_metadata = detected_result
                    # Tables are detected but NOT automatically excluded
                    if detected_reason == ExclusionReason.TABLE:
                        # Table detected but not excluded - will be marked in segment_data for frontend
                        # Frontend can show table indicator and provide exclude option
                        logger.trace(
                            LogModule.EXTRACT,
                            f"[PREVIEW] Task {task_id}: Detected table segment {idx} (not auto-excluded): '{seg_text[:50]}...'"
                        )
                    else:
                        # Non-table exclusion (image, formula, etc.) - auto-exclude
                        excluded_segments_with_reasons[idx] = detected_reason
                        excluded_segment_indices.append(idx)
                # Note: No fallback needed - detect_exclusion_reason should handle all cases
            
            # Update segments_metadata using ExclusionManager
            if excluded_segments_with_reasons:
                ExclusionManager.update_excluded_segments(
                    task_state=task_state,
                    excluded_segments=excluded_segments_with_reasons,
                    metadata=None
                )
                logger.debug(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: Marked {len(excluded_segments_with_reasons)} segments as excluded "
                    f"during HTML extraction (using ExclusionManager)"
                )
            
            # CRITICAL: Get excluded_segments AFTER ExclusionManager.update_excluded_segments
            # This ensures we get the newly stored exclusion information
            current_segments_metadata = task_state.get("segments_metadata", {})
            current_excluded_segments = current_segments_metadata.get("excluded_segments", {})
            
            # Preserve excluded_segments from ExclusionManager
            if current_excluded_segments:
                task_state["segments_metadata"]["excluded_segments"] = current_excluded_segments
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Preserved {len(current_excluded_segments)} excluded_segments from ExclusionManager"
                )
            
            # For convert_only mode, create translation_segments
            if task_state.get("convert_only", False):
                ts_segments = []
                separators_after = result.separators_after or []
                for i, s in enumerate(result.segments):
                    is_image = _is_image_segment(s)
                    is_excluded = i in excluded_segment_indices
                    ts_segments.append({
                        "segment_index": i,
                        "source_text": s,
                        "target_text": s,
                        "modified": False,
                        "separator_after": separators_after[i] if i < len(separators_after) else "",
                        "is_image": is_image,
                        "is_excluded": is_excluded,
                    })
                task_state["translation_segments"] = {
                    "segments": ts_segments,
                    "metadata": task_state.get("segments_metadata", {})
                }
            
            self.task_manager.add_log(task_id, "success", f"Source preview prepared from HTML: {min(result.total_segments, 200)}/{result.total_segments} segments")
            return True
        except Exception as e:
            logger.error(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Failed to prepare HTML source preview: {e}", exc_info=True)
            return False
    
    def prepare_source_preview_for_srt(
        self,
        task_id: str,
        file_contents: bytes,
        payload: Any,
        task_state: Dict[str, Any]
    ) -> bool:
        """
        Prepare source preview for SRT workflow.
        
        Args:
            task_id: Task identifier
            file_contents: SRT file content bytes
            payload: Task payload
            task_state: Task state dictionary
            
        Returns:
            True if preview was prepared successfully, False otherwise
        """
        try:
            from extractor.srt_extractor import SrtExtractor
            
            try:
                decoded = file_contents.decode('utf-8')
            except UnicodeDecodeError:
                decoded = file_contents.decode('utf-8', errors='replace')
            
            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
            result = SrtExtractor(decoded, chunk_size=chunk_size).extract()
            
            if result.total_segments == 0:
                return False
            
            content_hash = hashlib.sha1(file_contents).hexdigest()
            task_state["source_preview"] = {
                "segments": result.segments[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                "total_segments": result.total_segments,
                "ready": True,
            }
            task_state["source_chunks_cache"] = {
                "content_hash": content_hash,
                "chunk_size": chunk_size,
                "segments": result.segments,
                "total_segments": result.total_segments,
                "created_at": time.time(),
            }
            
            # Mark excluded segments using unified exclusion detection
            from utils.translation_segments import _is_image_segment
            from exclusion.core import detect_exclusion_reason, ExclusionReason, ExclusionManager
            excluded_segment_indices = []
            excluded_segments_with_reasons = {}
            # CRITICAL: Get target language using unified method that checks stored value first
            target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
            for idx, seg_text in enumerate(result.segments):
                is_image = _is_image_segment(seg_text)
                detected_result = detect_exclusion_reason(
                    text=seg_text,
                    block_type=None,
                    target_lang=target_lang,
                    is_image=is_image,
                    is_table=False
                )
                if detected_result:
                    detected_reason, _ = detected_result
                    # Only exclude non-optional types (TABLE is optional and not auto-excluded)
                    if not ExclusionReason.is_optional(detected_reason):
                        excluded_segment_indices.append(idx)
                        excluded_segments_with_reasons[idx] = detected_reason
            
            if excluded_segment_indices:
                logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Marked {len(excluded_segment_indices)} segments as excluded during SRT extraction")
                # Store excluded segments using ExclusionManager
                ExclusionManager.update_excluded_segments(
                    task_state=task_state,
                    excluded_segments=excluded_segments_with_reasons,
                    metadata=None
                )
            
            # CRITICAL: Get excluded_segments AFTER ExclusionManager.update_excluded_segments
            current_segments_metadata = task_state.get("segments_metadata", {})
            current_excluded_segments = current_segments_metadata.get("excluded_segments", {})
            
            # Create new segments_metadata
            new_segments_metadata = {
                "source": "srt",
                "workflow_type": getattr(payload, 'workflow_type', 'srt'),
                "chunk_size": chunk_size,
                "content_hash": content_hash,
                "separators_after": result.separators_after,
                "segment_info": result.segment_info,
                "excluded_segment_indices": excluded_segment_indices,  # Legacy format for backward compatibility
            }
            
            # CRITICAL: Preserve excluded_segments from ExclusionManager
            if current_excluded_segments:
                new_segments_metadata["excluded_segments"] = current_excluded_segments
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Preserved {len(current_excluded_segments)} excluded_segments from ExclusionManager"
                )
            
            task_state["segments_metadata"] = new_segments_metadata
            
            # For convert_only mode, create translation_segments
            if task_state.get("convert_only", False):
                from utils.translation_segments import _is_image_segment
                from exclusion.core.exclusion_detector import detect_exclusion_reason
                ts_segments = []
                separators_after = result.separators_after or []
                # Get target language from payload for language-based exclusion
                # CRITICAL: Get target language using unified method that checks stored value first
                target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
                for i, s in enumerate(result.segments):
                    is_image = _is_image_segment(s)
                    detected_result = detect_exclusion_reason(
                        text=s,
                        block_type=None,
                        target_lang=target_lang,
                        is_image=is_image,
                        is_table=False
                    )
                    is_excluded = bool(detected_result)
                    reason_val = detected_result[0].value if detected_result else None
                    meta = detected_result[1] if detected_result and len(detected_result) > 1 else {}
                    seg = {
                        "segment_index": i,
                        "source_text": s,
                        "target_text": s,
                        "modified": False,
                        "separator_after": separators_after[i] if i < len(separators_after) else "",
                        "is_image": is_image,
                        "is_excluded": is_excluded,
                    }
                    if is_excluded and reason_val:
                        seg["exclusion_reason"] = reason_val
                        if meta:
                            seg["exclusion_metadata"] = meta
                    ts_segments.append(seg)
                task_state["translation_segments"] = {
                    "segments": ts_segments,
                    "metadata": task_state.get("segments_metadata", {})
                }
            
            self.task_manager.add_log(task_id, "success", f"Source preview prepared from SRT: {min(result.total_segments, 200)}/{result.total_segments} segments")
            return True
        except Exception as e:
            logger.error(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Failed to prepare SRT source preview: {e}", exc_info=True)
            return False
    
    def prepare_source_preview_for_txt(
        self,
        task_id: str,
        file_contents: bytes,
        payload: Any,
        task_state: Dict[str, Any]
    ) -> bool:
        """
        Prepare source preview for TXT workflow.
        
        Args:
            task_id: Task identifier
            file_contents: TXT file content bytes
            payload: Task payload
            task_state: Task state dictionary
            
        Returns:
            True if preview was prepared successfully, False otherwise
        """
        try:
            from utils.markdown_splitter import split_text_into_paragraphs

            decoded = decode_with_detection(file_contents)

            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)

            logger.info(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Using paragraph-first segmentation for txt preview (chunk_size={chunk_size})"
            )

            segments = split_text_into_paragraphs(decoded, max_block_size=chunk_size)
            
            if not segments:
                return False
            
            # Mark excluded segments using unified exclusion detection
            from utils.translation_segments import _is_image_segment
            from exclusion.core import detect_exclusion_reason, ExclusionReason, ExclusionManager
            excluded_segment_indices = []
            excluded_segments_with_reasons = {}
            # Get target language from payload for language-based exclusion
            target_lang = getattr(payload, 'to_lang', None) or getattr(payload, 'target_lang', None)
            for idx, seg_text in enumerate(segments):
                is_image = _is_image_segment(seg_text)
                detected_result = detect_exclusion_reason(
                    text=seg_text,
                    block_type=None,
                    target_lang=target_lang,
                    is_image=is_image,
                    is_table=False
                )
                if detected_result:
                    detected_reason, _ = detected_result
                    # Only exclude non-optional types (TABLE is optional and not auto-excluded)
                    if not ExclusionReason.is_optional(detected_reason):
                        excluded_segment_indices.append(idx)
                        excluded_segments_with_reasons[idx] = detected_reason
            
            if excluded_segment_indices:
                logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Marked {len(excluded_segment_indices)} segments as excluded during TXT extraction")
                # Store excluded segments using ExclusionManager
                ExclusionManager.update_excluded_segments(
                    task_state=task_state,
                    excluded_segments=excluded_segments_with_reasons,
                    metadata=None
                )
            
            # CRITICAL: Get excluded_segments AFTER ExclusionManager.update_excluded_segments
            current_segments_metadata = task_state.get("segments_metadata", {})
            current_excluded_segments = current_segments_metadata.get("excluded_segments", {})
            
            content_hash = hashlib.sha1(file_contents).hexdigest()
            task_state["source_preview"] = {
                "segments": segments[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                "total_segments": len(segments),
                "ready": True,
            }
            task_state["source_chunks_cache"] = {
                "content_hash": content_hash,
                "chunk_size": chunk_size,
                "segments": segments,
                "total_segments": len(segments),
                "created_at": time.time(),
            }
            
            # Create new segments_metadata
            new_segments_metadata = {
                "source": "txt",
                "workflow_type": getattr(payload, 'workflow_type', 'txt'),
                "chunk_size": chunk_size,
                "content_hash": content_hash,
                "excluded_segment_indices": excluded_segment_indices,  # Legacy format for backward compatibility
            }
            
            # CRITICAL: Preserve excluded_segments from ExclusionManager
            if current_excluded_segments:
                new_segments_metadata["excluded_segments"] = current_excluded_segments
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Preserved {len(current_excluded_segments)} excluded_segments from ExclusionManager"
                )
            
            task_state["segments_metadata"] = new_segments_metadata
            
            # For convert_only mode, create translation_segments
            if task_state.get("convert_only", False):
                from exclusion.core.exclusion_detector import detect_exclusion_reason
                ts_segments = []
                # CRITICAL: Get target language using unified method that checks stored value first
                target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
                for i, s in enumerate(segments):
                    is_image = _is_image_segment(s)
                    detected_result = detect_exclusion_reason(
                        text=s,
                        block_type=None,
                        target_lang=target_lang,
                        is_image=is_image,
                        is_table=False
                    )
                    is_excluded = bool(detected_result)
                    reason_val = detected_result[0].value if detected_result else None
                    meta = detected_result[1] if detected_result and len(detected_result) > 1 else {}
                    seg = {
                        "segment_index": i,
                        "source_text": s,
                        "target_text": s,
                        "modified": False,
                        "separator_after": "",
                        "is_image": is_image,
                        "is_excluded": is_excluded,
                    }
                    if is_excluded and reason_val:
                        seg["exclusion_reason"] = reason_val
                        if meta:
                            seg["exclusion_metadata"] = meta
                    ts_segments.append(seg)
                task_state["translation_segments"] = {
                    "segments": ts_segments,
                    "metadata": task_state.get("segments_metadata", {})
                }
            
            self.task_manager.add_log(task_id, "success", f"Source preview prepared from TXT: {min(len(segments), 200)}/{len(segments)} segments")
            return True
        except Exception as e:
            logger.error(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Failed to prepare TXT source preview: {e}", exc_info=True)
            return False
    
    def prepare_source_preview_for_json(
        self,
        task_id: str,
        file_contents: bytes,
        payload: Any,
        task_state: Dict[str, Any]
    ) -> bool:
        """
        Prepare source preview for JSON workflow.
        
        Args:
            task_id: Task identifier
            file_contents: JSON file content bytes
            payload: Task payload
            task_state: Task state dictionary
            
        Returns:
            True if preview was prepared successfully, False otherwise
        """
        try:
            from extractor.json_extractor import JsonExtractor
            
            try:
                decoded = file_contents.decode('utf-8')
            except UnicodeDecodeError:
                decoded = file_contents.decode('utf-8', errors='replace')
            
            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
            json_paths = getattr(payload, 'json_paths', None) or []
            result = JsonExtractor(decoded, json_paths=json_paths, chunk_size=chunk_size).extract()
            
            if result.total_segments == 0:
                return False

            # Preserve JSON layout (key order and newlines) for export
            try:
                from utils.json_layout import parse_json_layout
                json_layout = parse_json_layout(decoded)
            except Exception as layout_err:
                logger.debug(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: Could not parse JSON layout: {layout_err}"
                )
                json_layout = []
            
            # Mark excluded segments using unified exclusion detection
            from utils.translation_segments import _is_image_segment
            from exclusion.core import detect_exclusion_reason, ExclusionReason, ExclusionManager
            excluded_segment_indices = []
            excluded_segments_with_reasons = {}
            # CRITICAL: Get target language using unified method that checks stored value first
            target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
            for idx, seg_text in enumerate(result.segments):
                is_image = _is_image_segment(seg_text)
                detected_result = detect_exclusion_reason(
                    text=seg_text,
                    block_type=None,
                    target_lang=target_lang,
                    is_image=is_image,
                    is_table=False
                )
                if detected_result:
                    detected_reason, _ = detected_result
                    # Only exclude non-optional types (TABLE is optional and not auto-excluded)
                    if not ExclusionReason.is_optional(detected_reason):
                        excluded_segment_indices.append(idx)
                        excluded_segments_with_reasons[idx] = detected_reason
            
            if excluded_segment_indices:
                logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Marked {len(excluded_segment_indices)} segments as excluded during JSON extraction")
                # Store excluded segments using ExclusionManager
                ExclusionManager.update_excluded_segments(
                    task_state=task_state,
                    excluded_segments=excluded_segments_with_reasons,
                    metadata=None
                )
            
            content_hash = hashlib.sha1(file_contents).hexdigest()
            task_state["source_preview"] = {
                "segments": result.segments[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                "total_segments": result.total_segments,
                "ready": True,
            }
            task_state["source_chunks_cache"] = {
                "content_hash": content_hash,
                "chunk_size": chunk_size,
                "segments": result.segments,
                "total_segments": result.total_segments,
                "created_at": time.time(),
            }
            
            # CRITICAL: Get excluded_segments AFTER ExclusionManager.update_excluded_segments
            # This ensures we get the newly stored exclusion information
            # ExclusionManager.update_excluded_segments has already updated task_state["segments_metadata"]["excluded_segments"]
            current_segments_metadata = task_state.get("segments_metadata", {})
            current_excluded_segments = current_segments_metadata.get("excluded_segments", {})
            
            # Create new segments_metadata
            new_segments_metadata = {
                "source": "json",
                "workflow_type": getattr(payload, 'workflow_type', 'json'),
                "chunk_size": chunk_size,
                "content_hash": content_hash,
                "separators_after": result.separators_after,
                "segment_info": result.segment_info,
                "excluded_segment_indices": excluded_segment_indices,  # Legacy format for backward compatibility
            }
            if json_layout:
                new_segments_metadata["json_layout"] = json_layout
            
            # CRITICAL: Preserve excluded_segments from ExclusionManager (includes newly detected exclusions)
            # This is the source of truth for exclusion information
            if current_excluded_segments:
                new_segments_metadata["excluded_segments"] = current_excluded_segments
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Preserved {len(current_excluded_segments)} excluded_segments from ExclusionManager "
                    f"(includes {len(excluded_segments_with_reasons) if excluded_segments_with_reasons else 0} newly detected exclusions)"
                )
            elif excluded_segments_with_reasons:
                # Fallback: if ExclusionManager didn't store it (shouldn't happen), log a warning
                logger.warning(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: ExclusionManager.update_excluded_segments was called with {len(excluded_segments_with_reasons)} exclusions, "
                    f"but current_excluded_segments is empty. This may indicate a storage issue."
                )
            
            task_state["segments_metadata"] = new_segments_metadata
            
            # Build chunk_to_segment_map for chunks generation (if not convert_only)
            # CRITICAL: This prevents segments from being merged into a single chunk in get_source_preview
            if not task_state.get("convert_only", False):
                self._build_chunk_to_segment_map(
                    task_id=task_id,
                    segments=result.segments,
                    chunk_size=chunk_size,
                    excluded_segment_indices=excluded_segment_indices,
                    task_state=task_state
                )
            
            # For convert_only mode, create translation_segments
            if task_state.get("convert_only", False):
                from utils.translation_segments import _is_image_segment
                from exclusion.core.exclusion_detector import detect_exclusion_reason
                ts_segments = []
                separators_after = result.separators_after or []
                # Get target language from payload for language-based exclusion
                # CRITICAL: Get target language using unified method that checks stored value first
                target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
                for i, s in enumerate(result.segments):
                    is_image = _is_image_segment(s)
                    detected_result = detect_exclusion_reason(
                        text=s,
                        block_type=None,
                        target_lang=target_lang,
                        is_image=is_image,
                        is_table=False
                    )
                    is_excluded = bool(detected_result)
                    reason_val = detected_result[0].value if detected_result else None
                    meta = detected_result[1] if detected_result and len(detected_result) > 1 else {}
                    seg = {
                        "segment_index": i,
                        "source_text": s,
                        "target_text": s,
                        "modified": False,
                        "separator_after": separators_after[i] if i < len(separators_after) else "",
                        "is_image": is_image,
                        "is_excluded": is_excluded,
                    }
                    if is_excluded and reason_val:
                        seg["exclusion_reason"] = reason_val
                        if meta:
                            seg["exclusion_metadata"] = meta
                    ts_segments.append(seg)
                task_state["translation_segments"] = {
                    "segments": ts_segments,
                    "metadata": task_state.get("segments_metadata", {})
                }
            
            self.task_manager.add_log(task_id, "success", f"Source preview prepared from JSON: {min(result.total_segments, 200)}/{result.total_segments} segments")
            return True
        except Exception as e:
            logger.error(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Failed to prepare JSON source preview: {e}", exc_info=True)
            return False
    
    def prepare_source_preview_for_xlsx(
        self,
        task_id: str,
        file_contents: bytes,
        payload: Any,
        task_state: Dict[str, Any]
    ) -> bool:
        """
        Prepare source preview for XLSX workflow.
        
        Args:
            task_id: Task identifier
            file_contents: XLSX file content bytes
            payload: Task payload
            task_state: Task state dictionary
            
        Returns:
            True if preview was prepared successfully, False otherwise
        """
        try:
            from extractor.xlsx_extractor import XlsxExtractor
            
            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
            translate_regions = getattr(payload, 'translate_regions', None) or []
            result = XlsxExtractor(file_contents, translate_regions=translate_regions, chunk_size=chunk_size).extract()
            
            if result.total_segments == 0:
                return False
            
            # Mark excluded segments using unified exclusion detection
            from utils.translation_segments import _is_image_segment
            from exclusion.core import detect_exclusion_reason, ExclusionReason, ExclusionManager
            excluded_segment_indices = []
            excluded_segments_with_reasons = {}
            # CRITICAL: Get target language using unified method that checks stored value first
            target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
            for idx, seg_text in enumerate(result.segments):
                is_image = _is_image_segment(seg_text)
                detected_result = detect_exclusion_reason(
                    text=seg_text,
                    block_type=None,
                    target_lang=target_lang,
                    is_image=is_image,
                    is_table=False
                )
                if detected_result:
                    detected_reason, _ = detected_result
                    # Only exclude non-optional types (TABLE is optional and not auto-excluded)
                    if not ExclusionReason.is_optional(detected_reason):
                        excluded_segment_indices.append(idx)
                        excluded_segments_with_reasons[idx] = detected_reason
            
            if excluded_segment_indices:
                logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Marked {len(excluded_segment_indices)} segments as excluded during XLSX extraction")
                # Store excluded segments using ExclusionManager
                ExclusionManager.update_excluded_segments(
                    task_state=task_state,
                    excluded_segments=excluded_segments_with_reasons,
                    metadata=None
                )
            
            # CRITICAL: Get excluded_segments AFTER ExclusionManager.update_excluded_segments
            current_segments_metadata = task_state.get("segments_metadata", {})
            current_excluded_segments = current_segments_metadata.get("excluded_segments", {})
            
            content_hash = hashlib.sha1(file_contents).hexdigest()
            task_state["source_preview"] = {
                "segments": result.segments[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                "total_segments": result.total_segments,
                "ready": True,
            }
            task_state["source_chunks_cache"] = {
                "content_hash": content_hash,
                "chunk_size": chunk_size,
                "segments": result.segments,
                "total_segments": result.total_segments,
                "created_at": time.time(),
            }
            
            # Create new segments_metadata
            new_segments_metadata = {
                "source": "xlsx",
                "workflow_type": getattr(payload, 'workflow_type', 'xlsx'),
                "chunk_size": chunk_size,
                "content_hash": content_hash,
                "separators_after": result.separators_after,
                "segment_info": result.segment_info,
                "excluded_segment_indices": excluded_segment_indices,  # Legacy format for backward compatibility
            }
            
            # CRITICAL: Preserve excluded_segments from ExclusionManager
            if current_excluded_segments:
                new_segments_metadata["excluded_segments"] = current_excluded_segments
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Preserved {len(current_excluded_segments)} excluded_segments from ExclusionManager"
                )
            
            task_state["segments_metadata"] = new_segments_metadata
            
            # Build chunk_to_segment_map for chunks generation (if not convert_only)
            convert_only = task_state.get("convert_only", False)
            logger.debug(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: convert_only={convert_only}, will build chunk_to_segment_map: {not convert_only}"
            )
            if not convert_only:
                try:
                    self._build_chunk_to_segment_map(
                        task_id=task_id,
                        segments=result.segments,
                        chunk_size=chunk_size,
                        excluded_segment_indices=excluded_segment_indices,
                        task_state=task_state
                    )
                    # Verify chunk_to_segment_map was created
                    if "chunk_to_segment_map" not in task_state or task_state.get("chunk_to_segment_map") is None:
                        logger.error(
                            LogModule.EXTRACT,
                            f"[PREVIEW] Task {task_id}: CRITICAL - _build_chunk_to_segment_map completed but chunk_to_segment_map not found in task_state"
                        )
                    else:
                        logger.debug(
                            LogModule.EXTRACT,
                            f"[PREVIEW] Task {task_id}: Successfully built and stored chunk_to_segment_map: {len(task_state.get('chunk_to_segment_map', []))} chunks"
                        )
                except Exception as e:
                    logger.error(
                        LogModule.EXTRACT,
                        f"[PREVIEW] Task {task_id}: Failed to build chunk_to_segment_map: {e}",
                        exc_info=True
                    )
                    # Do not re-raise - allow preview to continue without chunks
                    # This allows frontend to display segments even if chunks generation fails
            else:
                logger.debug(
                    LogModule.EXTRACT,
                    f"[PREVIEW] Task {task_id}: Skipping chunk_to_segment_map build (convert_only=True)"
                )
            
            # For convert_only mode, create translation_segments
            if task_state.get("convert_only", False):
                from utils.translation_segments import _is_image_segment
                from exclusion.core.exclusion_detector import detect_exclusion_reason
                ts_segments = []
                separators_after = result.separators_after or []
                for i, s in enumerate(result.segments):
                    is_image = _is_image_segment(s)
                    detected_result = detect_exclusion_reason(
                        text=s,
                        block_type=None,
                        target_lang=None,  # No target_lang for convert_only
                        is_image=is_image,
                        is_table=False
                    )
                    is_excluded = bool(detected_result)
                    reason_val = detected_result[0].value if detected_result else None
                    meta = detected_result[1] if detected_result and len(detected_result) > 1 else {}
                    separator_after = separators_after[i] if i < len(separators_after) else ""
                    seg = {
                        "segment_index": i,
                        "source_text": s,
                        "target_text": s,
                        "modified": False,
                        "separator_after": separator_after,
                        "is_image": is_image,
                        "is_excluded": is_excluded,
                    }
                    if is_excluded and reason_val:
                        seg["exclusion_reason"] = reason_val
                        if meta:
                            seg["exclusion_metadata"] = meta
                    ts_segments.append(seg)
                task_state["translation_segments"] = {
                    "segments": ts_segments,
                    "metadata": task_state.get("segments_metadata", {})
                }
            
            # Verify chunk_to_segment_map was created and output debug files
            if "chunk_to_segment_map" in task_state and task_state.get("chunk_to_segment_map") is not None:
                # Output segments and chunks to temporary folder for debugging
                self._output_extract_debug_files(
                    task_id=task_id,
                    task_state=task_state,
                    segments=result.segments,
                    chunk_to_segment_map=task_state.get("chunk_to_segment_map", [])
                )
            
            self.task_manager.add_log(task_id, "success", f"Source preview prepared from XLSX: {min(result.total_segments, 200)}/{result.total_segments} segments")
            return True
        except Exception as e:
            logger.error(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Failed to prepare XLSX source preview: {e}", exc_info=True)
            return False
    
    def prepare_layout_preview_from_layout(
        self,
        task_id: str,
        layout_doc: Any,
        payload: Any,
        task_state: Dict[str, Any],
        reason: str = "auto"
    ) -> bool:
        """
        Build layout-based preview/cache for PDF workflows.
        
        Args:
            task_id: Task identifier
            layout_doc: Layout document instance
            payload: Task payload
            task_state: Task state dictionary
            reason: Reason for preview generation (for logging)
            
        Returns:
            True if preview was prepared successfully, False otherwise
        """
        if layout_doc is None:
            return False
        
        # Translation task inherited from Extract (convert_task_id): keep inherited layout_prepared_chunks
        # and segments_metadata (e.g. after "clear all except image"). Do not rebuild layout or re-run
        # exclusion detection, which would overwrite with content-based exclusions and undo the user's choice.
        convert_task_id = task_state.get("convert_task_id")
        existing_chunks = task_state.get("layout_prepared_chunks")
        if convert_task_id and existing_chunks and isinstance(existing_chunks, list) and len(existing_chunks) > 0:
            if not task_state.get("layout_document"):
                task_state["layout_document"] = layout_doc
            logger.info(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Skipping layout rebuild (inherited from convert_task_id={convert_task_id}, "
                f"{len(existing_chunks)} chunks), preserving layout_prepared_chunks and exclusion state"
            )
            # IMPORTANT: For translation tasks (convert_only=False), we still must ensure that
            # chunk_to_segment_map exists so that get_source_preview can build chunks/token stats.
            convert_only = task_state.get("convert_only", False)
            if not convert_only and (not task_state.get("chunk_to_segment_map") or task_state.get("chunk_to_segment_map") is None):
                cache = task_state.get("source_chunks_cache") or {}
                segments_list = cache.get("segments") if isinstance(cache.get("segments"), list) else None
                if segments_list:
                    try:
                        chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
                        segments_metadata = task_state.get("segments_metadata") or {}
                        excluded_segment_indices = segments_metadata.get("excluded_segment_indices")
                        if excluded_segment_indices is None and isinstance(segments_metadata.get("excluded_segments"), dict):
                            try:
                                excluded_segment_indices = sorted(
                                    int(k) for k in segments_metadata["excluded_segments"].keys() if str(k).isdigit()
                                )
                            except (ValueError, TypeError):
                                excluded_segment_indices = []
                        if excluded_segment_indices is None:
                            excluded_segment_indices = []
                        self._build_chunk_to_segment_map(
                            task_id=task_id,
                            segments=segments_list,
                            chunk_size=chunk_size,
                            excluded_segment_indices=excluded_segment_indices,
                            task_state=task_state,
                        )
                        if task_state.get("chunk_to_segment_map"):
                            logger.debug(
                                LogModule.EXTRACT,
                                f"[PREVIEW] Task {task_id}: [LAYOUT_PDF] Built chunk_to_segment_map from inherited cache: "
                                f"{len(task_state['chunk_to_segment_map'])} chunks"
                            )
                        else:
                            logger.error(
                                LogModule.EXTRACT,
                                f"[PREVIEW] Task {task_id}: [LAYOUT_PDF] _build_chunk_to_segment_map completed but "
                                f"chunk_to_segment_map not found in task_state"
                            )
                    except Exception as e:
                        logger.error(
                            LogModule.EXTRACT,
                            f"[PREVIEW] Task {task_id}: [LAYOUT_PDF] Failed to build chunk_to_segment_map from inherited cache: {e}",
                            exc_info=True,
                        )
                else:
                    logger.warning(
                        LogModule.EXTRACT,
                        f"[PREVIEW] Task {task_id}: [LAYOUT_PDF] source_chunks_cache.segments missing or invalid; "
                        f"cannot build chunk_to_segment_map from inherited cache"
                    )
            return True
        
        try:
            from layout.markdown_builder import LayoutMarkdownBuilder
            import base64
            import io
            import mimetypes
            import zipfile
        except ImportError:
            return False
        
        try:
            chunk_size = chunk_size_service.get_chunk_size(payload, task_id)
            
            # Clear old chunk_to_segment_map if chunk_size changed
            existing_chunk_size = task_state.get("segments_metadata", {}).get("chunk_size")
            if existing_chunk_size and existing_chunk_size != chunk_size:
                logger.info(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: chunk_size changed from {existing_chunk_size} to {chunk_size}, clearing old chunk_to_segment_map")
                task_state.pop("chunk_to_segment_map", None)
                task_state.pop("chunk_tokens_info", None)
                task_state.pop("total_estimated_input_tokens", None)
            
            # Get deep_split from task_state first, then payload, default to True
            deep_split_enabled = True  # Default
            source = "default"
            if "deep_split" in task_state:
                deep_split_enabled = bool(task_state["deep_split"])
                source = "task_state"
            elif payload:
                deep_split_enabled = bool(getattr(payload, 'deep_split', True))
                source = "payload"
            
            # Get equation_format and table_body_format from task_state then payload (for PDF preview: same as export)
            equation_format = (task_state.get("equation_format") or "text")
            if isinstance(equation_format, str):
                equation_format = equation_format.lower().strip()
            if equation_format not in ("text", "latex", "image"):
                equation_format = "text"
            table_body_format = (task_state.get("table_body_format") or "html")
            if isinstance(table_body_format, str):
                table_body_format = table_body_format.lower().strip()
            if table_body_format not in ("html", "image"):
                table_body_format = "html"
            payload_obj = task_state.get("payload")
            try:
                if payload_obj:
                    if isinstance(payload_obj, dict):
                        equation_format = (payload_obj.get("equation_format") or equation_format)
                        table_body_format = (payload_obj.get("table_body_format") or table_body_format)
                    else:
                        equation_format = (getattr(payload_obj, "equation_format", None) or equation_format)
                        table_body_format = (getattr(payload_obj, "table_body_format", None) or table_body_format)
                    if isinstance(equation_format, str):
                        equation_format = equation_format.lower().strip()
                    if equation_format not in ("text", "latex", "image"):
                        equation_format = "text"
                    if isinstance(table_body_format, str):
                        table_body_format = table_body_format.lower().strip()
                    if table_body_format not in ("html", "image"):
                        table_body_format = "html"
            except Exception:
                pass
            logger.debug(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Building layout preview with deep_split={deep_split_enabled} "
                f"(from {source}), chunk_size={chunk_size}, equation_format={equation_format}, table_body_format={table_body_format}, "
                f"total_blocks={len(list(layout_doc.iter_blocks()))}"
            )
            builder = LayoutMarkdownBuilder(
                max_chunk_chars=chunk_size,
                deep_split=deep_split_enabled,
                equation_format=equation_format,
                table_body_format=table_body_format,
                include_structural_blocks=True  # CRITICAL: Must match Extract phase to ensure same segment count (176 vs 164)
            )
            layout_result = builder.build(layout_doc)
            logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: LayoutMarkdownBuilder generated {len(layout_result.chunks)} chunks (deep_split={deep_split_enabled})")
            
            if not layout_result.chunks:
                return False
            
            preview_segments = [chunk.text for chunk in layout_result.chunks]
            content_hash = hashlib.sha1(layout_result.markdown_text.encode("utf-8")).hexdigest()
            
            task_state["source_preview"] = {
                "segments": preview_segments[:SOURCE_PREVIEW_SEGMENTS_LIMIT],
                "total_segments": len(preview_segments),
                "ready": True,
            }
            task_state["source_chunks_cache"] = {
                "content_hash": content_hash,
                "chunk_size": chunk_size,
                "segments": preview_segments,
                "total_segments": len(preview_segments),
                "created_at": time.time(),
            }
            # Write Extract-phase segments to temp dir in same JSON format as LLM input for diagnosis
            try:
                from utils.extract_segments_debug import write_extract_segments_json
                written = write_extract_segments_json(
                    task_state.get("temp_dir"), preview_segments, task_id=task_id
                )
                if written:
                    logger.debug(
                        LogModule.EXTRACT,
                        f"[PREVIEW] Task {task_id}: Wrote {len(preview_segments)} segments to {written}"
                    )
            except Exception as _e:
                logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Failed to write extract_segments.json: {_e}")
            # CRITICAL: Update segments_metadata instead of completely overwriting it
            # This preserves important data from Extract phase (excluded_segments, excluded_segment_indices, etc.)
            # Only update fields that need to change (content_hash, chunk_size) or are specific to this preview generation
            existing_segments_metadata = task_state.get("segments_metadata", {})
            
            # Initialize segments_metadata if it doesn't exist, otherwise update it
            if not existing_segments_metadata:
                task_state["segments_metadata"] = {
                    "source": "layout_pdf",
                    "workflow_type": getattr(payload, 'workflow_type', None),
                    "chunk_size": chunk_size,
                    "content_hash": content_hash,
                    "separators_after": [],
                    "segment_info": [],
                }
            else:
                # Update only the fields that need to change
                task_state["segments_metadata"].update({
                    "source": "layout_pdf",  # Update source to indicate this is from layout preview
                    "chunk_size": chunk_size,  # Update chunk_size if it changed
                    "content_hash": content_hash,  # Update content_hash based on new markdown
                })
                # Preserve workflow_type if it exists, otherwise set it
                if "workflow_type" not in task_state["segments_metadata"]:
                    task_state["segments_metadata"]["workflow_type"] = getattr(payload, 'workflow_type', None)
                # Only set separators_after and segment_info to empty if they don't exist
                # (they may have been set in Extract phase and should be preserved)
                if "separators_after" not in task_state["segments_metadata"]:
                    task_state["segments_metadata"]["separators_after"] = []
                if "segment_info" not in task_state["segments_metadata"]:
                    task_state["segments_metadata"]["segment_info"] = []
            
            # Log preservation of excluded data
            excluded_segments_count = len(task_state["segments_metadata"].get("excluded_segments", {}))
            excluded_indices_count = len(task_state["segments_metadata"].get("excluded_segment_indices", []))
            if excluded_segments_count > 0 or excluded_indices_count > 0:
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Preserved exclusion data from Extract phase: "
                    f"{excluded_segments_count} excluded_segments, {excluded_indices_count} excluded_segment_indices"
                )
            
            # Build chunk_to_segment_map for layout-based PDF preview when not convert_only
            convert_only = task_state.get("convert_only", False)
            logger.debug(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: [LAYOUT_PDF] convert_only={convert_only}, will build chunk_to_segment_map: {not convert_only}"
            )
            if not convert_only:
                try:
                    segments_metadata = task_state.get("segments_metadata") or {}
                    excluded_segment_indices = segments_metadata.get("excluded_segment_indices")
                    if excluded_segment_indices is None and isinstance(segments_metadata.get("excluded_segments"), dict):
                        try:
                            excluded_segment_indices = sorted(
                                int(k) for k in segments_metadata["excluded_segments"].keys() if str(k).isdigit()
                            )
                        except (ValueError, TypeError):
                            excluded_segment_indices = []
                    if excluded_segment_indices is None:
                        excluded_segment_indices = []
                    self._build_chunk_to_segment_map(
                        task_id=task_id,
                        segments=preview_segments,
                        chunk_size=chunk_size,
                        excluded_segment_indices=excluded_segment_indices,
                        task_state=task_state,
                    )
                    if "chunk_to_segment_map" not in task_state or task_state.get("chunk_to_segment_map") is None:
                        logger.error(
                            LogModule.EXTRACT,
                            f"[PREVIEW] Task {task_id}: [LAYOUT_PDF] CRITICAL - _build_chunk_to_segment_map completed "
                            f"but chunk_to_segment_map not found in task_state"
                        )
                    else:
                        logger.debug(
                            LogModule.EXTRACT,
                            f"[PREVIEW] Task {task_id}: [LAYOUT_PDF] Successfully built and stored chunk_to_segment_map: "
                            f"{len(task_state.get('chunk_to_segment_map', []))} chunks"
                        )
                except Exception as e:
                    logger.error(
                        LogModule.EXTRACT,
                        f"[PREVIEW] Task {task_id}: [LAYOUT_PDF] Failed to build chunk_to_segment_map: {e}",
                        exc_info=True,
                    )
            
            # Process chunks and build image data map
            serialized_chunks = []
            chunk_block_map = []
            image_data_map: Dict[str, Dict[str, str]] = {}
            existing_image_map = task_state.get("image_data_map")
            if isinstance(existing_image_map, dict):
                image_data_map.update({
                    str(k): {
                        "data": (v or {}).get("data", ""),
                        "alt": (v or {}).get("alt", ""),
                    }
                    for k, v in existing_image_map.items()
                })

            from utils.mineru_image_data_map import populate_image_data_map_from_mineru_zip

            populate_image_data_map_from_mineru_zip(
                image_data_map,
                task_state,
                layout_doc=layout_doc,
            )
            
            zip_bytes = task_state.get("layout_source_zip")
            zip_file = None
            zip_entries: List[str] = []
            zip_entry_map: Dict[str, str] = {}
            if zip_bytes:
                try:
                    zip_file = zipfile.ZipFile(io.BytesIO(zip_bytes))
                    zip_entries = zip_file.namelist()
                    zip_entry_map = {
                        name.replace("\\", "/"): name for name in zip_entries
                    }
                except Exception as zip_error:
                    logger.debug(LogModule.EXTRACT, f"[LAYOUT] Failed to open MinerU ZIP for images: {zip_error}")
                    zip_file = None
            
            def _normalize_image_path(path: Optional[str]) -> Optional[str]:
                if not path:
                    return None
                return path.replace("\\", "/").lstrip("./")
            
            placeholder_cache: Dict[str, str] = {}
            
            def _read_image_data_uri(image_path: Optional[str]) -> Optional[str]:
                if not image_path or zip_file is None:
                    return None
                normalized = _normalize_image_path(image_path)
                if not normalized:
                    return None
                if normalized in placeholder_cache:
                    return placeholder_cache[normalized]
                
                # Try exact match first
                candidate = zip_entry_map.get(normalized)
                if candidate is None:
                    # Try matching by filename (basename)
                    filename_only = os.path.basename(normalized)
                    for name, original in zip_entry_map.items():
                        if name == filename_only or name.endswith('/' + filename_only) or name.endswith('\\' + filename_only):
                            candidate = original
                            logger.debug(LogModule.EXTRACT, f"[LAYOUT] Matched image '{image_path}' to ZIP entry '{candidate}' by filename")
                            break
                        if name.endswith(normalized):
                            candidate = original
                            logger.debug(LogModule.EXTRACT, f"[LAYOUT] Matched image '{image_path}' to ZIP entry '{candidate}' by path ending")
                            break
                
                if not candidate:
                    logger.warning(LogModule.EXTRACT, f"[LAYOUT] Image path '{image_path}' (normalized: '{normalized}') not found in ZIP entries.")
                    return None
                
                try:
                    raw_bytes = zip_file.read(candidate)
                    logger.debug(LogModule.EXTRACT, f"[LAYOUT] Successfully read image '{image_path}' from ZIP entry '{candidate}' ({len(raw_bytes)} bytes)")
                except KeyError:
                    logger.warning(LogModule.EXTRACT, f"[LAYOUT] Failed to read image '{candidate}' from ZIP (KeyError).")
                    return None
                
                mime = mimetypes.guess_type(candidate)[0] or "image/png"
                data_uri = f"data:{mime};base64,{base64.b64encode(raw_bytes).decode('ascii')}"
                placeholder_cache[normalized] = data_uri
                return data_uri
            
            image_segment_indices: List[int] = []
            chunk_block_texts_map: List[List[str]] = []
            
            # Build block_type_map once so we can mark formula/identifier exclusions per chunk (PDF 1:1 chunk-segment)
            block_type_map: Dict[int, str] = {}
            for page in layout_doc.pages:
                for block in page.blocks:
                    block_type = getattr(block, "type", "unknown") or "unknown"
                    block_index = getattr(block, "index", None)
                    if block_index is not None:
                        block_type_map[block_index] = block_type
            
            from exclusion.core import ExclusionManager
            from exclusion.core.exclusion_reason import ExclusionReason
            from exclusion.core.exclusion_detector import detect_exclusion_reason
            from utils.translation_segments import _is_image_segment
            
            target_lang = self._get_target_lang_for_exclusion_detection(task_id, task_state, payload)
            # Optional exclusions (LANGUAGE_MATCH, TABLE): only mark as excluded if user already excluded them
            pre_existing_excluded_indices = set(ExclusionManager.get_excluded_segments(task_state).keys())

            for idx, chunk in enumerate(layout_result.chunks):
                is_chart_body = chunk.chunk_type == "chart_body"
                is_image = chunk.chunk_type == "image" or (
                    is_chart_body and bool(chunk.image_path)
                )
                placeholder_id = None
                chunk_text = chunk.text
                
                if is_image or is_chart_body:
                    if chunk.image_path or is_chart_body:
                        placeholder_id = chunk.image_placeholder or f"layoutimg{idx}"
                    if is_image and not chunk_text:
                        chunk_text = chunk_text or f"<ph-{placeholder_id}>"
                    data_uri = _read_image_data_uri(chunk.image_path) if chunk.image_path else None
                    if placeholder_id and data_uri:
                        from utils.mineru_image_data_map import register_image_data_uri

                        alt_text = chunk.image_alt or (chunk.image_path or "Image")
                        image_data_map[placeholder_id] = {
                            "data": data_uri or "",
                            "alt": alt_text or "Image",
                        }
                        if chunk.image_path:
                            register_image_data_uri(
                                image_data_map,
                                chunk.image_path,
                                data_uri,
                                alt=alt_text,
                            )
                    if is_image or (is_chart_body and placeholder_id):
                        image_segment_indices.append(idx)
                else:
                    # Check if this is an equation image
                    import re
                    equation_image_pattern = r'!\[Equation\]\(([^)]+\.(jpg|jpeg|png|gif|webp))\)'
                    equation_match = re.search(equation_image_pattern, chunk_text)
                    if equation_match:
                        equation_filename = equation_match.group(1)
                        logger.debug(LogModule.EXTRACT, f"[LAYOUT] Found equation image in text chunk: {equation_filename}")
                        data_uri = _read_image_data_uri(equation_filename)
                        if data_uri:
                            filename_key = equation_filename.split('/')[-1].split('\\')[-1] if '/' in equation_filename or '\\' in equation_filename else equation_filename
                            image_data_map[filename_key] = {
                                "data": data_uri,
                                "alt": equation_filename,
                            }
                            logger.debug(LogModule.EXTRACT, f"[LAYOUT] Added equation image to image_data_map: {filename_key}")
                
                block_texts = getattr(chunk, "block_texts", None) or []
                
                # Estimate tokens for this chunk
                estimated_tokens = None
                if not is_image and not is_chart_body:
                    from utils.token_estimator import estimate_chunk_input_tokens
                    estimated_tokens = estimate_chunk_input_tokens(
                        chunk_text,
                        system_prompt_approx=2000  # PDF uses longer system prompt
                    )
                
                # CRITICAL: Mark excluded segments (image, formula, identifier) so Translate phase does not send them or mark as failed
                segment_indices = [idx]  # PDF layout: 1 chunk = 1 segment
                is_excluded = False
                exclusion_reason_val: Optional[str] = None
                if is_image:
                    is_excluded = True
                    exclusion_reason_val = ExclusionReason.IMAGE.value
                elif is_chart_body:
                    is_excluded = True
                    exclusion_reason_val = ExclusionReason.IMAGE.value
                else:
                    is_interline_equation = bool(
                        chunk.block_indices
                        and any(
                            block_type_map.get(b) == "interline_equation"
                            for b in chunk.block_indices
                        )
                    )
                    if is_interline_equation:
                        is_excluded = True
                        exclusion_reason_val = ExclusionReason.FORMULA.value
                    else:
                        first_block_type = None
                        if chunk.block_indices:
                            first_block_type = block_type_map.get(chunk.block_indices[0])
                        detected = detect_exclusion_reason(
                            text=chunk_text or "",
                            block_type=first_block_type,
                            target_lang=target_lang,
                            is_image=_is_image_segment(chunk_text or ""),
                            is_table=False,
                        )
                        if detected:
                            r = detected[0]
                            if ExclusionReason.is_content_based(r):
                                is_excluded = True
                                exclusion_reason_val = r.value
                            elif ExclusionReason.is_language_based(r) or r == ExclusionReason.TABLE:
                                # Only exclude when user already excluded (e.g. checkbox in Extract)
                                if idx in pre_existing_excluded_indices:
                                    is_excluded = True
                                    exclusion_reason_val = r.value
                            else:
                                is_excluded = True
                                exclusion_reason_val = r.value
                
                chunk_info = {
                    "text": chunk_text,
                    "chunk_type": chunk.chunk_type,
                    "block_indices": chunk.block_indices,
                    "block_texts": block_texts,
                    "image_path": chunk.image_path,
                    "placeholder_id": placeholder_id,
                    "is_image": is_image,
                    "segment_indices": segment_indices,
                    "is_excluded": is_excluded,
                }
                if is_excluded and exclusion_reason_val:
                    chunk_info["exclusion_reason"] = exclusion_reason_val
                if estimated_tokens is not None:
                    chunk_info["estimated_input_tokens"] = estimated_tokens
                
                serialized_chunks.append(chunk_info)
                chunk_block_map.append(chunk.block_indices)
                chunk_block_texts_map.append(block_texts)
            
            if zip_file:
                try:
                    zip_file.close()
                except Exception:
                    pass
            
            if image_data_map:
                task_state["image_data_map"] = image_data_map
            if image_segment_indices:
                task_state["layout_image_segment_indices"] = image_segment_indices
            
            task_state["layout_prepared_chunks"] = serialized_chunks
            task_state["layout_chunk_block_map"] = chunk_block_map
            task_state["layout_chunk_block_texts"] = chunk_block_texts_map
            task_state["layout_markdown_source"] = layout_result.markdown_text
            
            # CRITICAL: Build excluded_segments from layout chunks and store in segments_metadata
            # so Translate phase (and record_segments) sees them and does not mark excluded segments as translation failed.
            # Only store content-based exclusions + optional (LANGUAGE_MATCH, TABLE) when user already excluded them.
            excluded_segments_with_reasons_layout: Dict[int, Any] = {}
            for ch in serialized_chunks:
                if not ch.get("is_excluded"):
                    continue
                reason_val = ch.get("exclusion_reason", ExclusionReason.UNKNOWN.value)
                try:
                    reason = ExclusionReason(reason_val)
                except ValueError:
                    reason = ExclusionReason.UNKNOWN
                for seg_idx in ch.get("segment_indices") or []:
                    excluded_segments_with_reasons_layout[int(seg_idx)] = reason
            # Filter: do not auto-store LANGUAGE_MATCH or TABLE unless user already excluded (pre_existing_excluded_indices)
            pre_existing = set(ExclusionManager.get_excluded_segments(task_state).keys())
            filtered_excluded_layout: Dict[int, Any] = {}
            for seg_idx, reason in excluded_segments_with_reasons_layout.items():
                if ExclusionReason.is_content_based(reason):
                    filtered_excluded_layout[seg_idx] = reason
                elif ExclusionReason.is_language_based(reason) or reason == ExclusionReason.TABLE:
                    if seg_idx in pre_existing:
                        filtered_excluded_layout[seg_idx] = reason
                else:
                    filtered_excluded_layout[seg_idx] = reason
            if filtered_excluded_layout:
                if "segments_metadata" not in task_state:
                    task_state["segments_metadata"] = {}
                ExclusionManager.update_excluded_segments(
                    task_state=task_state,
                    excluded_segments=filtered_excluded_layout,
                )
                task_state["segments_metadata"]["excluded_segment_indices"] = sorted(filtered_excluded_layout.keys())
                logger.debug(
                    LogModule.EXCLUSION,
                    f"[PREVIEW] Task {task_id}: Stored {len(filtered_excluded_layout)} excluded segments from layout chunks "
                    f"(content-based + user-chosen optional) so they are not sent for translation or marked as failed"
                )
            
            # Calculate and store total estimated input tokens
            total_estimated_tokens = sum(
                chunk.get("estimated_input_tokens", 0)
                for chunk in serialized_chunks
                if isinstance(chunk, dict) and not chunk.get("is_image", False)
            )
            if total_estimated_tokens > 0:
                task_state["total_estimated_input_tokens"] = total_estimated_tokens
                logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Calculated total estimated tokens for PDF: {total_estimated_tokens}")
            
            # For convert_only mode, create translation_segments
            if task_state.get("convert_only", False):
                ts_segments = []
                for idx, chunk in enumerate(layout_result.chunks):
                    sc = serialized_chunks[idx] if idx < len(serialized_chunks) else {}
                    seg = {
                        "segment_index": idx,
                        "source_text": sc.get("text", chunk.text),
                        "target_text": sc.get("text", chunk.text),
                        "modified": False,
                        "separator_after": "",
                    }
                    if chunk.block_indices:
                        seg["layout_block_indices"] = chunk.block_indices
                    if sc.get("is_image"):
                        seg["is_image"] = True
                    if sc.get("is_excluded"):
                        seg["is_excluded"] = True
                        seg["exclusion_reason"] = sc.get("exclusion_reason", "unknown")
                    
                    ts_segments.append(seg)
                
                task_state["translation_segments"] = {
                    "segments": ts_segments,
                    "metadata": task_state.get("segments_metadata", {})
                }
            
            return True
        except Exception as e:
            logger.error(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Failed to prepare layout preview: {e}", exc_info=True)
            return False
    
    def _build_chunk_to_segment_map(
        self,
        task_id: str,
        segments: List[str],
        chunk_size: int,
        excluded_segment_indices: List[int],
        task_state: Dict[str, Any]
    ) -> None:
        """
        Build chunk_to_segment_map for chunks generation.
        
        Args:
            task_id: Task identifier
            segments: List of segments
            chunk_size: Chunk size in tokens
            excluded_segment_indices: List of excluded segment indices
            task_state: Task state dictionary
        """
        try:
            from utils.json_utils import segments2json_chunks
            from utils.chunk_size_converter import get_text_content_token_limit
            from utils.translation_segments import _is_image_segment
            
            # Calculate text content token limit
            text_token_limit = get_text_content_token_limit(chunk_size)
            logger.debug(LogModule.EXTRACT, f"[PREVIEW] Task {task_id}: Using text_token_limit={text_token_limit} (from total chunk_size={chunk_size} tokens) for chunking")
            
            # Filter out excluded segments before chunking
            # CRITICAL: Use excluded_segment_indices parameter (from Extract phase) instead of re-detecting
            # This ensures consistency between Extract and Translation phases
            excluded_set = set(excluded_segment_indices) if excluded_segment_indices else set()
            filtered_segments = []
            segment_index_mapping = {}  # Maps filtered index to original index
            for orig_idx, seg_text in enumerate(segments):
                # Check if segment is in excluded_segment_indices (from Extract phase)
                # This includes images, formulas, tables, etc. that were detected in Extract phase
                if orig_idx not in excluded_set:
                    filtered_segments.append(seg_text)
                    segment_index_mapping[len(filtered_segments) - 1] = orig_idx
            
            indexed_originals, chunks, merged_indices_list, chunk_tokens = segments2json_chunks(
                filtered_segments, text_token_limit, estimate_tokens=True
            )
            
            # Store chunk tokens info
            chunk_tokens_info = []
            total_estimated_tokens = 0
            chunk_to_segment_map = []
            
            for chunk_dict in chunks:
                # Map filtered indices back to original segment indices
                filtered_indices = [int(k) for k in sorted(chunk_dict.keys(), key=int)]
                original_indices = [segment_index_mapping[filtered_idx] for filtered_idx in filtered_indices if filtered_idx in segment_index_mapping]
                if original_indices:  # Only add non-empty chunks
                    chunk_to_segment_map.append(original_indices)
                    
                    # Store token estimate for this chunk
                    chunk_idx = len(chunk_tokens_info)
                    if chunk_tokens and chunk_idx < len(chunk_tokens):
                        estimated_tokens = chunk_tokens[chunk_idx]
                        chunk_tokens_info.append(estimated_tokens)
                        total_estimated_tokens += estimated_tokens
            
            task_state["chunk_to_segment_map"] = chunk_to_segment_map
            if chunk_tokens_info:
                task_state["chunk_tokens_info"] = chunk_tokens_info
                task_state["total_estimated_input_tokens"] = total_estimated_tokens
            
            logger.debug(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Built chunk_to_segment_map: {len(chunk_to_segment_map)} chunks "
                f"(excluded {len(excluded_segment_indices)} segments from chunking, total estimated tokens: {total_estimated_tokens})"
            )
            
            # CRITICAL: Handle case where all segments are excluded
            # This can happen when all segments are identifiers, images, etc.
            # In this case, we still need to create an empty chunk_to_segment_map
            # to avoid errors in get_source_preview, but translation will have no chunks
            if not chunk_to_segment_map:
                if len(excluded_segment_indices) == len(segments):
                    # All segments are excluded - this is valid, create empty map
                    logger.warning(
                        LogModule.EXTRACT,
                        f"[PREVIEW] Task {task_id}: All {len(segments)} segments are excluded. "
                        f"Creating empty chunk_to_segment_map. Translation will have no chunks. "
                        f"Excluded reasons: {excluded_segment_indices}"
                    )
                else:
                    # This should not happen - some segments should be available
                    logger.error(
                        LogModule.EXTRACT,
                        f"[PREVIEW] Task {task_id}: CRITICAL - chunk_to_segment_map is empty but not all segments are excluded. "
                        f"This indicates a bug. "
                        f"segments count: {len(segments)}, "
                        f"filtered_segments count: {len(filtered_segments)}, "
                        f"excluded_segment_indices: {len(excluded_segment_indices)}, "
                        f"chunk_size: {chunk_size}, text_token_limit: {text_token_limit}"
                    )
                    raise ValueError(
                        f"chunk_to_segment_map is empty but not all segments are excluded. "
                        f"This should not happen. "
                        f"segments: {len(segments)}, filtered: {len(filtered_segments)}, "
                        f"excluded: {len(excluded_segment_indices)}"
                    )
            
            # Store in task_state
            task_state["chunk_to_segment_map"] = chunk_to_segment_map
            logger.debug(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Stored chunk_to_segment_map in task_state: {len(chunk_to_segment_map)} chunks"
            )
        except Exception as e:
            logger.error(
                LogModule.EXTRACT,
                f"[PREVIEW] Task {task_id}: Failed to build chunk_to_segment_map: {e}. "
                f"segments: {len(segments) if segments else 0}, "
                f"chunk_size: {chunk_size}, excluded: {len(excluded_segment_indices) if excluded_segment_indices else 0}",
                exc_info=True
            )
            raise  # Re-raise to prevent silent failure
            # Ensure chunk_to_segment_map is set to None so fallback can be used
            task_state["chunk_to_segment_map"] = None
    
    def _output_extract_debug_files(
        self,
        task_id: str,
        task_state: Dict[str, Any],
        segments: List[str],
        chunk_to_segment_map: List[List[int]]
    ) -> None:
        """
        Output segments and chunks to temporary folder for debugging during extract phase.
        
        Args:
            task_id: Task identifier
            task_state: Task state dictionary
            segments: List of extracted segments
            chunk_to_segment_map: List mapping chunk indices to segment indices
        """
        try:
            # Get temp_dir from task_state
            temp_dir = task_state.get("temp_dir")
            if not temp_dir or not os.path.isdir(temp_dir):
                logger.debug(LogModule.EXTRACT, f"[EXTRACT-DEBUG] Task {task_id}: temp_dir not available, skipping debug file output")
                return
            
            # Create debug directory following the unified output folder rule
            # Rule: temp_dir/debug/extract/
            debug_dir = os.path.join(temp_dir, "debug", "extract")
            os.makedirs(debug_dir, exist_ok=True)
            
            # Store debug directory path in task_state
            if "debug_files" not in task_state:
                task_state["debug_files"] = {}
            task_state["debug_files"]["extract_debug_dir"] = debug_dir
            
            # Output segments to file (each segment on a line with index)
            segments_file = os.path.join(debug_dir, "segments.txt")
            with open(segments_file, 'w', encoding='utf-8') as f:
                f.write(f"Total segments: {len(segments)}\n")
                f.write("=" * 80 + "\n\n")
                for idx, seg_text in enumerate(segments):
                    f.write(f"Segment {idx}:\n")
                    f.write(f"{seg_text}\n")
                    f.write("-" * 80 + "\n\n")
            
            logger.debug(LogModule.EXTRACT, f"[EXTRACT-DEBUG] Task {task_id}: Saved {len(segments)} segments to {segments_file}")
            
            # Output chunks to file (from chunk_to_segment_map)
            chunks_file = os.path.join(debug_dir, "chunks.txt")
            with open(chunks_file, 'w', encoding='utf-8') as f:
                if chunk_to_segment_map is None:
                    f.write(f"Total chunks: 0 (chunk_to_segment_map is None)\n")
                    f.write("=" * 80 + "\n\n")
                    f.write("No chunks available (chunk_to_segment_map is None)\n")
                else:
                    f.write(f"Total chunks: {len(chunk_to_segment_map)}\n")
                    f.write("=" * 80 + "\n\n")
                    for chunk_idx, segment_indices in enumerate(chunk_to_segment_map):
                        f.write(f"Chunk {chunk_idx} (contains segments: {segment_indices}):\n")
                        # Combine segments in this chunk
                        chunk_text_parts = []
                        for seg_idx in segment_indices:
                            if seg_idx < len(segments):
                                chunk_text_parts.append(f"[Segment {seg_idx}]\n{segments[seg_idx]}")
                        chunk_text = "\n\n".join(chunk_text_parts)
                        f.write(f"{chunk_text}\n")
                        f.write("=" * 80 + "\n\n")
            
            if chunk_to_segment_map is not None:
                logger.debug(LogModule.EXTRACT, f"[EXTRACT-DEBUG] Task {task_id}: Saved {len(chunk_to_segment_map)} chunks to {chunks_file}")
            else:
                logger.debug(LogModule.EXTRACT, f"[EXTRACT-DEBUG] Task {task_id}: chunk_to_segment_map is None, saved empty chunks file to {chunks_file}")
            
            # Output chunk_to_segment_map as JSON for easy parsing
            import json
            map_file = os.path.join(debug_dir, "chunk_to_segment_map.json")
            with open(map_file, 'w', encoding='utf-8') as f:
                if chunk_to_segment_map is None:
                    json.dump([], f, indent=2, ensure_ascii=False)
                else:
                    json.dump(chunk_to_segment_map, f, indent=2, ensure_ascii=False)
            
            logger.debug(LogModule.EXTRACT, f"[EXTRACT-DEBUG] Task {task_id}: Saved chunk_to_segment_map to {map_file}")
            
        except Exception as e:
            logger.warning(LogModule.EXTRACT, f"[EXTRACT-DEBUG] Task {task_id}: Failed to output extract debug files: {e}", exc_info=True)

