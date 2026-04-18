# SPDX-FileCopyrightText: 2026 Zampher
# SPDX-License-Identifier: MPL-2.0
import asyncio
import os
import shutil
import sys
import tempfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Self, Literal, List, Dict, Any, Tuple, Optional

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.shapes.group import GroupShape
    from pptx.shapes.autoshape import Shape
    from pptx.shapes.table import Table
    from pptx.shapes.freeform import FreeformBuilder
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.text.text import TextFrame
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

from agents.segments_agent import SegmentsTranslateAgentConfig, SegmentsTranslateAgent
from ir.document import Document
from translator.ai_translator.base import AiTranslatorConfig, AiTranslator
from logger.logger import LogModule


def get_font_for_language_pptx(target_language: str) -> str:
    """
    Return appropriate font based on target language and operating system for PPTX.
    Same logic as DOCX translator.
    
    Args:
        target_language: Target language name
        
    Returns:
        Recommended font name
    """
    # 根据操作系统选择字体
    is_macos = sys.platform == "darwin"
    is_linux = sys.platform.startswith("linux")
    is_windows = sys.platform == "win32"
    
    language_font_map = {
        # Chinese
        "Chinese": "Microsoft YaHei" if is_windows else ("PingFang SC" if is_macos else "Noto Sans CJK SC"),
        "Simplified Chinese": "Microsoft YaHei" if is_windows else ("PingFang SC" if is_macos else "Noto Sans CJK SC"),
        "Traditional Chinese": "Microsoft JhengHei" if is_windows else ("PingFang TC" if is_macos else "Noto Sans CJK TC"),
        
        # English
        "English": "Calibri" if is_windows else ("Helvetica Neue" if is_macos else "DejaVu Sans"),
        
        # Japanese
        "Japanese": "Yu Gothic" if is_windows else ("Hiragino Sans" if is_macos else "Noto Sans CJK JP"),
        
        # Korean
        "Korean": "Malgun Gothic" if is_windows else ("AppleGothic" if is_macos else "Noto Sans CJK KR"),
        
        # Russian
        "Russian": "Times New Roman",
        
        # Arabic
        "Arabic": "Arial Unicode MS",
        
        # Other European languages
        "Spanish": "Calibri" if is_windows else ("Helvetica Neue" if is_macos else "DejaVu Sans"),
        "French": "Calibri" if is_windows else ("Helvetica Neue" if is_macos else "DejaVu Sans"),
        "German": "Calibri" if is_windows else ("Helvetica Neue" if is_macos else "DejaVu Sans"),
        "Portuguese": "Calibri" if is_windows else ("Helvetica Neue" if is_macos else "DejaVu Sans"),
        
        # Vietnamese
        "Vietnamese": "Arial Unicode MS",
        
        # Hebrew
        "Hebrew": "Arial Unicode MS",
        
        # Thai
        "Thai": "Arial Unicode MS",
        
        # Hindi
        "Hindi": "Arial Unicode MS",
    }
    
    # Find matching font
    font = language_font_map.get(target_language)
    
    # If no matching font found, return default font
    if not font:
        # 根据操作系统选择默认字体
        if is_macos:
            font = "Helvetica Neue"
        elif is_linux:
            font = "DejaVu Sans"
        else:
            font = "Calibri"
    
    return font


@dataclass
class PptxTranslatorConfig(AiTranslatorConfig):
    """
    Configuration class for PptxTranslator.
    """
    insert_mode: Literal["replace", "append", "prepend"] = "replace"
    separator: str = "\n"
    translate_notes: bool = False  # Whether to translate notes pages
    translate_master: bool = False  # Whether to translate master slides (usually not recommended)
    translate_tables: bool = True  # Whether to translate tables
    translate_textboxes: bool = True  # Whether to translate text boxes


class PptxTranslator(AiTranslator):
    """
    Translator for .pptx files.
    Preserves template, formatting, layout, and styles while translating text content.
    """

    def __init__(self, config: PptxTranslatorConfig):
        # Re-check PPTX availability at runtime, in case it was installed after module load
        try:
            from pptx import Presentation
            _pptx_available = True
        except ImportError:
            _pptx_available = False
        
        if not _pptx_available:
            raise ImportError("python-pptx is not installed. Please install it with: pip install python-pptx")
        
        super().__init__(config=config)
        self.chunk_size = config.chunk_size
        self.translate_agent = None
        if not self.skip_translate:
            agent_config = SegmentsTranslateAgentConfig(
                custom_prompt=config.custom_prompt,
                to_lang=config.to_lang,
                base_url=config.base_url,
                api_key=config.api_key,
                model_id=config.model_id,
                api_type=getattr(config, 'api_type', None) or getattr(config, 'api_protocol', None) or 'openai',
                temperature=config.temperature,
                thinking=config.thinking,
                concurrent=config.concurrent,
                connect_timeout=getattr(config, 'connect_timeout', 15),
                timeout=config.timeout,
                logger=self.logger,
                glossary_dict=config.glossary_dict,
                retry=config.retry,
                max_tokens=getattr(config, 'max_tokens', None),
                use_seg_tags=True,  # Use SEG-tag format for PPTX segments
            )
            self.translate_agent = SegmentsTranslateAgent(agent_config)
        self.insert_mode = config.insert_mode
        self.separator = config.separator
        self.translate_notes = config.translate_notes
        self.translate_master = config.translate_master
        self.translate_tables = config.translate_tables
        self.translate_textboxes = config.translate_textboxes
        self.target_language = config.to_lang
        self.temp_pptx_path = None  # Store temporary file path for export

    def _pre_translate(self, document: Document, temp_dir: Optional[str] = None) -> Tuple[Any, List[Dict[str, Any]], List[str], Optional[str]]:
        """
        Extract all text elements from PPTX file, preserving format information.
        Copies source file to temporary directory for safe modification.
        
        Args:
            document: Document object containing .pptx file content
            temp_dir: Optional temporary directory path. If provided, source file will be copied there.
            
        Returns:
            A tuple containing:
            - Presentation object
            - A list containing text element information (preserving format)
            - A list containing all original texts to be translated
            - Path to temporary PPTX file (or None if not using temp file)
        """
        temp_pptx_path = None
        
        # Re-import Presentation at runtime to ensure it's available
        from pptx import Presentation as PPTXPresentation
        
        # Copy source file to temporary directory if temp_dir is provided
        if temp_dir:
            os.makedirs(temp_dir, exist_ok=True)
            temp_pptx_path = os.path.join(temp_dir, f"source_{os.urandom(8).hex()}.pptx")
            with open(temp_pptx_path, 'wb') as f:
                f.write(document.content)
            self.logger.debug(LogModule.TRANS, f"[PPTX] Copied source file to temporary location: {temp_pptx_path}")
            # Load presentation from temporary file
            prs = PPTXPresentation(temp_pptx_path)
        else:
            # Load presentation from bytes (original behavior)
            prs = PPTXPresentation(BytesIO(document.content))
        elements_to_translate = []
        original_texts = []

        # Extract text from slides
        for slide_idx, slide in enumerate(prs.slides):
            # Extract slide title
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    element_info = {
                        "type": "title",
                        "slide_idx": slide_idx,
                        "shape_idx": None,  # Title is a special shape
                        "paragraph_idx": 0,
                        "run_idx": 0,
                        "text_frame": slide.shapes.title.text_frame,
                        "paragraph": slide.shapes.title.text_frame.paragraphs[0] if slide.shapes.title.text_frame.paragraphs else None,
                        "run": slide.shapes.title.text_frame.paragraphs[0].runs[0] if slide.shapes.title.text_frame.paragraphs and slide.shapes.title.text_frame.paragraphs[0].runs else None,
                    }
                    elements_to_translate.append(element_info)
                    original_texts.append(title_text)

            # Extract text from shapes
            for shape_idx, shape in enumerate(slide.shapes):
                # Skip title shape (already processed)
                if shape == slide.shapes.title:
                    continue
                
                # Process text frames (text boxes, shapes with text)
                if self.translate_textboxes and shape.has_text_frame:
                    text_frame = shape.text_frame
                    for para_idx, paragraph in enumerate(text_frame.paragraphs):
                        # Collect all runs in paragraph
                        para_text_parts = []
                        runs_info = []
                        for run_idx, run in enumerate(paragraph.runs):
                            if run.text.strip():
                                para_text_parts.append(run.text)
                                runs_info.append({
                                    "run": run,
                                    "run_idx": run_idx,
                                })
                        
                        if para_text_parts:
                            para_text = "".join(para_text_parts)
                            if para_text.strip():
                                element_info = {
                                    "type": "text_frame",
                                    "slide_idx": slide_idx,
                                    "shape_idx": shape_idx,
                                    "paragraph_idx": para_idx,
                                    "runs_info": runs_info,
                                    "text_frame": text_frame,
                                    "paragraph": paragraph,
                                }
                                elements_to_translate.append(element_info)
                                original_texts.append(para_text)

                # Process tables
                elif self.translate_tables and shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            if cell_text:
                                element_info = {
                                    "type": "table_cell",
                                    "slide_idx": slide_idx,
                                    "shape_idx": shape_idx,
                                    "row_idx": row_idx,
                                    "col_idx": col_idx,
                                    "table": table,
                                    "cell": cell,
                                }
                                elements_to_translate.append(element_info)
                                original_texts.append(cell_text)

        # Extract notes pages (if enabled)
        if self.translate_notes:
            for slide_idx, slide in enumerate(prs.slides):
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        element_info = {
                            "type": "notes",
                            "slide_idx": slide_idx,
                            "notes_text_frame": slide.notes_slide.notes_text_frame,
                        }
                        elements_to_translate.append(element_info)
                        original_texts.append(notes_text)

        return prs, elements_to_translate, original_texts, temp_pptx_path

    def _get_excluded_segments(self, task_id: str | None) -> set[int]:
        """
        Get excluded segment indices using ExclusionManager (single source of truth).
        This ensures that manually excluded segments from Extract phase are correctly identified.
        """
        if not task_id:
            return set()
        try:
            from backend.app.services.task import task_manager
        except ImportError:
            return set()
        task_state = task_manager.get_task(task_id)
        if not task_state:
            return set()
        
        # CRITICAL: Use ExclusionManager.get_excluded_segments as the single source of truth
        # This ensures that manually excluded segments from Extract phase are correctly identified
        from exclusion.core import ExclusionManager
        excluded_segments_with_reasons = ExclusionManager.get_excluded_segments(task_state)
        excluded_set = set(excluded_segments_with_reasons.keys())
        
        if excluded_set:
            self.logger.info(LogModule.TRANS, f"[PPTX_TRANSLATOR] Task {task_id}: Retrieved {len(excluded_set)} excluded_segments from ExclusionManager "
                f"(single source of truth). Excluded indices: {sorted(excluded_set)[:20]}{'...' if len(excluded_set) > 20 else ''}"
            )
        else:
            self.logger.debug(
                f"[PPTX_TRANSLATOR] Task {task_id}: No excluded_segments found from ExclusionManager"
            )
        
        return excluded_set

    def _after_translate(self, prs: Any, elements_to_translate: List[Dict[str, Any]],
                        translated_texts: List[str], original_texts: List[str], temp_pptx_path: Optional[str] = None) -> bytes:
        """
        Write translated text back to PPTX, preserving all formatting.
        
        Args:
            prs: Presentation object
            elements_to_translate: List of element information
            translated_texts: List of translated texts (already mapped by index, same length as original_texts)
            original_texts: List of original texts
            
        Returns:
            Translated PPTX file as bytes
        """
        # CRITICAL: Use index-based access instead of text-based mapping to handle duplicate texts correctly
        # translated_texts is already mapped by index in translate/translate_async, so we can directly use index
        target_font = get_font_for_language_pptx(self.target_language)

        for i, element_info in enumerate(elements_to_translate):
            original_text = original_texts[i]
            # Use index-based access: translated_texts[i] corresponds to original_texts[i]
            translated_text = translated_texts[i] if i < len(translated_texts) else original_text

            # Determine final text based on insert mode
            if self.insert_mode == "replace":
                final_text = translated_text
            elif self.insert_mode == "append":
                final_text = original_text + self.separator + translated_text
            elif self.insert_mode == "prepend":
                final_text = translated_text + self.separator + original_text
            else:
                self.logger.error(LogModule.TRANS,"Invalid PptxTranslatorConfig parameter")
                final_text = translated_text

            element_type = element_info.get("type")

            if element_type == "title":
                # Update slide title
                text_frame = element_info["text_frame"]
                if text_frame.paragraphs:
                    paragraph = text_frame.paragraphs[0]
                    if paragraph.runs:
                        # Update first run, clear others
                        first_run = paragraph.runs[0]
                        first_run.text = final_text
                        # Set appropriate font
                        first_run.font.name = target_font
                        # Clear other runs
                        for run in paragraph.runs[1:]:
                            run.text = ""
                    else:
                        # No runs, add one
                        run = paragraph.add_run()
                        run.text = final_text
                        run.font.name = target_font

            elif element_type == "text_frame":
                # Update text frame paragraph
                paragraph = element_info["paragraph"]
                runs_info = element_info.get("runs_info", [])
                
                if runs_info:
                    # Update first run with translated text, preserve format
                    first_run = runs_info[0]["run"]
                    first_run.text = final_text
                    # Set appropriate font based on target language
                    first_run.font.name = target_font
                    
                    # Clear other runs
                    for run_info in runs_info[1:]:
                        run_info["run"].text = ""

            elif element_type == "table_cell":
                # Update table cell
                cell = element_info["cell"]
                cell.text = final_text
                # Set font for all paragraphs in cell
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = target_font

            elif element_type == "notes":
                # Update notes page
                notes_text_frame = element_info["notes_text_frame"]
                if notes_text_frame.paragraphs:
                    paragraph = notes_text_frame.paragraphs[0]
                    if paragraph.runs:
                        paragraph.runs[0].text = final_text
                        paragraph.runs[0].font.name = target_font
                        # Clear other runs
                        for run in paragraph.runs[1:]:
                            run.text = ""
                    else:
                        run = paragraph.add_run()
                        run.text = final_text
                        run.font.name = target_font

        # Save presentation
        if temp_pptx_path:
            # Save to temporary file (in-place update)
            prs.save(temp_pptx_path)
            self.logger.debug(LogModule.TRANS, f"[PPTX] Saved translated content to temporary file: {temp_pptx_path}")
            # Read the updated file and return as bytes
            with open(temp_pptx_path, 'rb') as f:
                return f.read()
        else:
            # Save to bytes buffer (original behavior)
            bio = BytesIO()
            prs.save(bio)
            return bio.getvalue()

    def translate(self, document: Document, temp_dir: Optional[str] = None) -> Self:
        """
        Synchronous translation method.
        
        Args:
            document: Document to translate
            temp_dir: Optional temporary directory path. If provided, source file will be copied there for safe modification.
        """
        prs, elements_to_translate, original_texts, temp_pptx_path = self._pre_translate(document, temp_dir)
        
        if not original_texts or self.skip_translate:
            # No text to translate or skip translation
            document.content = BytesIO(document.content).read()
            return self

        # Filter excluded segments if task_id is available
        excluded_set = self._get_excluded_segments(getattr(self, '_task_id', None))

        # Filter out excluded segments
        filtered_elements = []
        filtered_texts = []
        for i, (element_info, text) in enumerate(zip(elements_to_translate, original_texts)):
            if i not in excluded_set:
                filtered_elements.append(element_info)
                filtered_texts.append(text)

        if not filtered_texts:
            # All segments excluded
            document.content = BytesIO(document.content).read()
            return self

        # Translate using agent
        if self.translate_agent:
            translated_texts = self.translate_agent.send_segments(filtered_texts, self.chunk_size)
        else:
            translated_texts = filtered_texts

        # Map translated texts back to original positions using index-based mapping
        # CRITICAL: Use index-based mapping instead of text-based mapping to handle duplicate texts correctly
        # filtered_texts and translated_texts have the same length and order
        filtered_idx_to_translated = {}
        filtered_idx = 0
        for i, text in enumerate(original_texts):
            if i not in excluded_set:
                if filtered_idx < len(translated_texts):
                    filtered_idx_to_translated[i] = translated_texts[filtered_idx]
                    filtered_idx += 1
                else:
                    filtered_idx_to_translated[i] = text  # Fallback if index out of range
        
        all_translated_texts = [filtered_idx_to_translated.get(i, text) if i not in excluded_set else text 
                                for i, text in enumerate(original_texts)]

        # Write translated text back
        document.content = self._after_translate(prs, elements_to_translate, all_translated_texts, original_texts, temp_pptx_path)
        
        # Store temporary file path for export (don't delete it yet)
        # The file will be used for export and cleaned up later
        if temp_pptx_path and os.path.exists(temp_pptx_path):
            self.temp_pptx_path = temp_pptx_path
            self.logger.debug(LogModule.TRANS, f"[PPTX] Stored temporary file path for export: {temp_pptx_path}")
        
        return self

    async def translate_async(self, document: Document, progress_callback=None, temp_dir: Optional[str] = None) -> Self:
        """
        Asynchronous translation method.
        
        Args:
            document: Document to translate
            progress_callback: Optional progress callback function
            temp_dir: Optional temporary directory path. If provided, source file will be copied there for safe modification.
        """
        prs, elements_to_translate, original_texts, temp_pptx_path = await asyncio.to_thread(
            self._pre_translate, document, temp_dir
        )
        
        if not original_texts or self.skip_translate:
            # No text to translate or skip translation
            document.content = BytesIO(document.content).read()
            return self

        # Filter excluded segments if task_id is available
        excluded_set = self._get_excluded_segments(getattr(self, '_task_id', None))

        # Filter out excluded segments
        filtered_elements = []
        filtered_texts = []
        for i, (element_info, text) in enumerate(zip(elements_to_translate, original_texts)):
            if i not in excluded_set:
                filtered_elements.append(element_info)
                filtered_texts.append(text)

        if not filtered_texts:
            # All segments excluded
            document.content = BytesIO(document.content).read()
            return self

        # Translate using agent (async)
        if self.translate_agent:
            # Use chunk translation helper for segment recording
            if hasattr(self, '_task_id') and self._task_id:
                try:
                    from utils.chunk_translation_helper import translate_segments_with_agent_async
                    from backend.app.services.task import task_manager
                    task_state = task_manager.get_task(self._task_id) if self._task_id else None
                    
                    translated_texts, _ = await translate_segments_with_agent_async(
                        segments=filtered_texts,
                        chunk_size=self.chunk_size,
                        translate_agent=self.translate_agent,
                        task_id=self._task_id,
                        task_state=task_state,
                        original_filename=getattr(self, '_original_filename', None),
                        file_contents=document.content,
                        progress_callback=progress_callback,
                    )
                except Exception as e:
                    self.logger.warning(LogModule.TRANS, f"Failed to use chunk translation helper: {e}, falling back to direct translation")
                    translated_texts = await self.translate_agent.send_segments_async(filtered_texts, self.chunk_size, progress_callback)
            else:
                translated_texts = await self.translate_agent.send_segments_async(filtered_texts, self.chunk_size, progress_callback)
        else:
            translated_texts = filtered_texts

        # Map translated texts back to original positions using index-based mapping
        # CRITICAL: Use index-based mapping instead of text-based mapping to handle duplicate texts correctly
        # filtered_texts and translated_texts have the same length and order
        filtered_idx_to_translated = {}
        filtered_idx = 0
        for i, text in enumerate(original_texts):
            if i not in excluded_set:
                if filtered_idx < len(translated_texts):
                    filtered_idx_to_translated[i] = translated_texts[filtered_idx]
                    filtered_idx += 1
                else:
                    filtered_idx_to_translated[i] = text  # Fallback if index out of range
        
        all_translated_texts = [filtered_idx_to_translated.get(i, text) if i not in excluded_set else text 
                                for i, text in enumerate(original_texts)]

        # Save translated segments to task_state for segment recording
        # This ensures we use the actual translated segments (from AI) instead of re-extracting from PPTX
        if hasattr(self, '_task_id') and self._task_id:
            try:
                from backend.app.services.task import task_manager
                task_state = task_manager.get_task(self._task_id)
                if task_state:
                    # Store translated segments for later use in record_translation_segments
                    # all_translated_texts has the same length and order as original_texts
                    task_state["pptx_translated_segments"] = all_translated_texts
                    self.logger.debug(LogModule.TRANS, f"[PPTX] Saved {len(all_translated_texts)} translated segments to task_state "
                        f"for task {self._task_id}"
                    )
            except Exception as e:
                self.logger.warning(LogModule.TRANS, f"[PPTX] Failed to save translated segments to task_state: {e}")

        # Write translated text back
        document.content = await asyncio.to_thread(
            self._after_translate, prs, elements_to_translate, all_translated_texts, original_texts, temp_pptx_path
        )
        
        # Store temporary file path for export (don't delete it yet)
        # The file will be used for export and cleaned up later
        if temp_pptx_path and os.path.exists(temp_pptx_path):
            self.temp_pptx_path = temp_pptx_path
            self.logger.debug(LogModule.TRANS, f"[PPTX] Stored temporary file path for export: {temp_pptx_path}")
        
        return self

