# SPDX-FileCopyrightText: 2026 Zampher
# SPDX-License-Identifier: MPL-2.0

"""Unit tests for office compare-reading HTML preview helpers."""

from __future__ import annotations

import unittest
from io import BytesIO


class OfficePreviewUtilsTests(unittest.TestCase):
    def test_pptx_bytes_to_html_contains_slide_text(self) -> None:
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx not installed")

        from utils.office_preview_utils import pptx_bytes_to_html

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        box = slide.shapes.add_textbox(0, 0, 3000000, 500000)
        box.text_frame.text = "Hello Compare PPTX"
        buf = BytesIO()
        prs.save(buf)
        html = pptx_bytes_to_html(buf.getvalue())
        self.assertIn("Hello Compare PPTX", html)
        self.assertIn("Slide 1", html)

    def test_xlsx_bytes_to_html_contains_sheet_and_cell(self) -> None:
        try:
            import openpyxl
        except ImportError:
            self.skipTest("openpyxl not installed")

        from utils.office_preview_utils import xlsx_bytes_to_html

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "SheetA"
        ws["A1"] = "alpha-cell"
        buf = BytesIO()
        wb.save(buf)
        html = xlsx_bytes_to_html(buf.getvalue(), max_rows=50)
        self.assertIn("SheetA", html)
        self.assertIn("alpha-cell", html)

    def test_docx_bytes_to_html_when_mammoth_available(self) -> None:
        try:
            import mammoth  # noqa: F401
            from docx import Document
        except ImportError:
            self.skipTest("mammoth/python-docx not installed")

        from utils.office_preview_utils import docx_bytes_to_html

        doc = Document()
        doc.add_paragraph("Hello Compare DOCX")
        buf = BytesIO()
        doc.save(buf)
        html = docx_bytes_to_html(buf.getvalue())
        self.assertIn("Hello Compare DOCX", html)


if __name__ == "__main__":
    unittest.main()
